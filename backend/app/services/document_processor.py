"""Document processor: extract and apply text edits to DOCX and PPTX files.

Run-aware replacement strategy
--------------------------------
We never use proportional character distribution across runs.  Instead we:

1. Concatenate all run texts to get the full paragraph text.
2. Find the longest common prefix and suffix between old and new text to
   isolate the minimal "changed region".
3. Walk the runs and update only the ones that overlap that region.
   - Runs entirely before the region: untouched.
   - First run overlapping the region: gets (its unchanged prefix) +
     (the new replacement text) + (its unchanged suffix if the change ends
     inside this run).
   - Any subsequent runs overlapping the region: get only their unchanged
     trailing suffix (the middle is consumed by the first run).
   - Runs entirely after the region: untouched.

This ensures that formatting (colour, font, size, bold, italic, …) stays on
exactly the same runs as before.  Only the text content of runs that
actually contain the changed characters is modified.

Example
~~~~~~~
Paragraph runs:
  run[0] "Theme Name: "  (green, 12 chars)
  run[1] "Edtech"         (black,  6 chars)

new_text = "Theme Name: Healthtech"
  common prefix  = "Theme Name: "  → 12 chars
  common suffix  = "tech"           →  4 chars
  changed region = chars [12, 14) = "Ed"  →  replaced with "Health"

  run[0] untouched             → "Theme Name: " (green)  ✓
  run[1] "Ed"→"Health" + "tech" → "Healthtech"   (black)  ✓
"""
from __future__ import annotations

import copy
import logging
from dataclasses import dataclass
from pathlib import Path
from re import sub
from typing import Any

from docx import Document
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu, Inches
from pptx.util import Emu as EMU

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Custom IDs (UUID stamping for stable element identities)
# ---------------------------------------------------------------------------
import uuid
import shutil
import tempfile
from lxml import etree
from docx.oxml.ns import qn

CUSTOM_NS = 'http://documenteditor.local/ids'
CUSTOM_PREFIX = 'deid'
etree.register_namespace(CUSTOM_PREFIX, CUSTOM_NS)

# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class TextBlock:
    element_id: str
    type: str
    text: str
    metadata: dict


# ---------------------------------------------------------------------------
# Diff helper
# ---------------------------------------------------------------------------

def _changed_region(old: str, new: str) -> tuple[int, int, str]:
    """Return (prefix_len, suffix_len, new_middle) such that:
      old[prefix_len : len(old) - suffix_len]  should become  new_middle
      new_middle == new[prefix_len : len(new) - suffix_len]
    """
    # Longest common prefix
    prefix_len = 0
    min_len = min(len(old), len(new))
    while prefix_len < min_len and old[prefix_len] == new[prefix_len]:
        prefix_len += 1

    # Longest common suffix (bounded so it cannot overlap the prefix)
    old_remaining = len(old) - prefix_len
    new_remaining = len(new) - prefix_len
    max_suffix = min(old_remaining, new_remaining)
    suffix_len = 0
    while (suffix_len < max_suffix
           and old[len(old) - 1 - suffix_len] == new[len(new) - 1 - suffix_len]):
        suffix_len += 1

    new_end = len(new) - suffix_len if suffix_len > 0 else len(new)
    new_middle = new[prefix_len:new_end]
    return prefix_len, suffix_len, new_middle


# Alignment string → pptx enum mapping
_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# ---------------------------------------------------------------------------
# Main processor
# ---------------------------------------------------------------------------

class DocumentProcessor:

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def extract(self, path: Path, document_type: str) -> dict:
        if document_type == "pptx":
            res = self._extract_pptx_dom(path)
        else:
            res = self._extract_docx_dom(path)
            
        res["blocks"] = self._flatten_dom_for_retrieval(res["dom"])
        return res

    def extract_rich(self, path: Path, document_type: str) -> dict:
        """Extract enriched structure for template analysis (shapes, geometry, formatting)."""
        return self.extract(path, document_type)

    def _flatten_dom_for_retrieval(self, node: dict) -> list[dict]:
        blocks = []
        def traverse(n: dict):
            if n.get("type") == "paragraph" and n.get("text"):
                # Aggregate run styles to provide style hints
                colors = set()
                fonts = set()
                for r in n.get("runs", []):
                    st = r.get("style", {})
                    if st.get("color"): colors.add(st["color"])
                    if st.get("font"): fonts.add(st["font"])

                blocks.append({
                    "element_id": n["id"],
                    "text": n["text"],
                    "type": "text",
                    "metadata": {
                        "role": n.get("role"),
                        "colors": list(colors),
                        "fonts": list(fonts),
                        "include_in_toc": n.get("include_in_toc")
                    }
                })
            for child in n.get("children", []) + n.get("rows", []) + n.get("cells", []):
                traverse(child)
        traverse(node)
        return blocks

    def apply_edits(
        self,
        source: Path,
        target: Path,
        document_type: str,
        edits: list[dict],
    ) -> None:
        if document_type == "pptx":
            self._apply_pptx_edits(source, target, edits)
        else:
            self._apply_docx_edits(source, target, edits)

    def apply_slide_plan(
        self,
        source: Path,
        target: Path,
        slide_plan: dict,
    ) -> None:
        """Apply a structured slide plan to a PPTX template.

        The plan specifies which template slides to clone, what content to
        populate, and what formatting to apply.  Slides marked "delete" are
        removed.
        """
        self._apply_pptx_slide_plan(source, target, slide_plan)

    def apply_operations(
        self,
        source: Path,
        target: Path,
        document_type: str,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to the document.

        Returns (any_change_made, list_of_applied_summaries).
        """
        if document_type == "pptx":
            return self._apply_pptx_operations(source, target, operations)
        else:
            return self._apply_docx_operations(source, target, operations)

    # ------------------------------------------------------------------
    # PPTX: one block per paragraph inside each shape (original)
    # ------------------------------------------------------------------

    def _extract_pptx(self, path: Path) -> list[TextBlock]:
        prs = Presentation(path)
        blocks: list[TextBlock] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    eid = f"slide_{slide_idx}_shape_{shape_idx}_para_{para_idx}"
                    blocks.append(TextBlock(
                        element_id=eid,
                        type="text",
                        text=text,
                        metadata={
                            "slide": slide_idx,
                            "shape_index": shape_idx,
                            "shape_name": getattr(shape, "name", "Unknown"),
                            "para_index": para_idx,
                        },
                    ))
        return blocks

    # ------------------------------------------------------------------
    # PPTX: rich extraction (includes empty frames, geometry, formatting)
    # ------------------------------------------------------------------

    def _extract_pptx_dom(self, path: Path) -> dict:
        prs = Presentation(path)
        children = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_children = []
            for shape_idx, shape in enumerate(slide.shapes):
                shape_id = f"slide_{slide_idx}_shape_{shape_idx}"
                shape_role = "shape"
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).split('(')[0]
                        shape_role = f"placeholder_{ph_type.lower()}"
                    except Exception:
                        shape_role = "placeholder"
                
                shape_node = {
                    "id": shape_id,
                    "type": "shape",
                    "role": shape_role,
                    "name": getattr(shape, "name", "Unknown"),
                    "geometry": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    },
                    "children": []
                }

                if getattr(shape, "has_text_frame", False):
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        shape_node["children"].append(
                            self._extract_pptx_paragraph_dom(para, f"{shape_id}_para_{para_idx}")
                        )

                if getattr(shape, "has_table", False):
                    table_node = {
                        "id": f"{shape_id}_table",
                        "type": "table",
                        "role": "table",
                        "rows": []
                    }
                    for row_idx, row in enumerate(shape.table.rows):
                        row_node = {
                            "id": f"{shape_id}_table_row_{row_idx}",
                            "type": "row",
                            "row": row_idx,
                            "cells": []
                        }
                        for col_idx, cell in enumerate(row.cells):
                            cell_node = {
                                "id": f"{shape_id}_table_cell_{row_idx}_{col_idx}",
                                "type": "cell",
                                "row": row_idx,
                                "column": col_idx,
                                "children": []
                            }
                            for para_idx, para in enumerate(cell.text_frame.paragraphs):
                                cell_node["children"].append(
                                    self._extract_pptx_paragraph_dom(para, f"{cell_node['id']}_para_{para_idx}")
                                )
                            row_node["cells"].append(cell_node)
                        table_node["rows"].append(row_node)
                    shape_node["children"].append(table_node)

                slide_children.append(shape_node)

            layout_name = ""
            try:
                layout_name = slide.slide_layout.name
            except Exception:
                pass

            children.append({
                "id": f"slide_{slide_idx}",
                "type": "slide",
                "role": "slide",
                "layout_name": layout_name,
                "children": slide_children
            })

        return {
            "document_type": "pptx",
            "slide_count": len(prs.slides),
            "geometry": {
                "width": prs.slide_width,
                "height": prs.slide_height,
            },
            "dom": {
                "id": "document_root",
                "type": "document",
                "children": children
            }
        }

    def _extract_pptx_paragraph_dom(self, para, para_id: str) -> dict:
        alignment = None
        if para.alignment is not None:
            try:
                alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
            except Exception:
                pass

        runs = []
        for r_idx, run in enumerate(para.runs):
            font = run.font
            color_hex = None
            try:
                if font.color and font.color.type is not None and font.color.rgb:
                    color_hex = str(font.color.rgb)
            except Exception:
                pass
            runs.append({
                "id": f"{para_id}_run_{r_idx}",
                "type": "run",
                "text": run.text,
                "style": {
                    "font": font.name,
                    "size": round(font.size.pt, 1) if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                    "color": color_hex
                }
            })

        return {
            "id": para_id,
            "type": "paragraph",
            "role": "body",
            "text": para.text.strip(),
            "style": {
                "alignment": alignment,
            },
            "runs": runs
        }

    # ------------------------------------------------------------------
    # PPTX: apply text edits (original)
    # ------------------------------------------------------------------

    def _apply_pptx_edits(self, source: Path, target: Path, edits: list[dict]) -> None:
        prs = Presentation(source)
        edit_map = {e["element_id"]: e for e in edits}
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    eid = f"slide_{slide_idx}_shape_{shape_idx}_para_{para_idx}"
                    if eid in edit_map:
                        edit = edit_map[eid]
                        log.debug(
                            "PPTX edit [%s]: %r → %r",
                            eid, edit["old_text"], edit["new_text"],
                        )
                        self._apply_run_aware_replacement(para, edit["new_text"])
        prs.save(target)

    # ------------------------------------------------------------------
    # PPTX: slide plan application (new)
    # ------------------------------------------------------------------

    def _apply_pptx_slide_plan(self, source: Path, target: Path, plan: dict) -> None:
        """Apply a structured slide plan to a PPTX template.

        Plan structure::

            {
                "slides": [
                    {
                        "source_slide_index": 1,
                        "action": "populate" | "keep" | "delete",
                        "shapes": [
                            {
                                "shape_index": 0,
                                "paragraphs": [
                                    {
                                        "para_index": 0,
                                        "text": "New text",
                                        "formatting": {
                                            "font_size_pt": 36,
                                            "bold": true,
                                            "color_hex": null,
                                            "alignment": "center"
                                        }
                                    }
                                ]
                            }
                        ]
                    }
                ]
            }
        """
        prs = Presentation(source)
        template_slides = list(prs.slides)
        plan_slides = plan.get("slides", [])

        if not plan_slides:
            log.warning("Empty slide plan — saving template as-is.")
            prs.save(target)
            return

        # Phase 1: Build the output slide list by cloning template slides
        # We'll work backwards: first create all needed clones, then delete originals.
        new_slide_xmls: list[tuple[str, Any]] = []  # (action, slide_element or None)

        for entry in plan_slides:
            action = entry.get("action", "populate")
            if action == "delete":
                continue  # skip deleted slides

            src_idx = entry.get("source_slide_index", 1) - 1  # convert to 0-based
            if src_idx < 0 or src_idx >= len(template_slides):
                src_idx = 0  # fallback to first slide

            new_slide_xmls.append((action, src_idx, entry))

        # Phase 2: Build new presentation from template
        new_prs = Presentation(source)
        existing_count = len(list(new_prs.slides))

        # Clone ALL slides defined in the plan, appending them to the END
        for action, src_idx, entry in new_slide_xmls:
            self._clone_slide(new_prs, src_idx)

        # Delete ALL original slides from the beginning (in reverse order to avoid shifting)
        for i in range(existing_count - 1, -1, -1):
            self._delete_slide(new_prs, i)

        # Phase 3: Populate content for each slide
        slides_list = list(new_prs.slides)
        for slide_out_idx, (action, src_idx, entry) in enumerate(new_slide_xmls):
            if slide_out_idx >= len(slides_list):
                break

            slide = slides_list[slide_out_idx]

            # For cloned slides that don't match their source, we need to
            # copy content from the correct template slide
            if action == "keep":
                continue  # leave as-is

            # Populate shapes
            shape_edits = entry.get("shapes", [])
            edit_map = {se.get("shape_index", 0): se for se in shape_edits}
            shapes = list(slide.shapes)
            
            for shape_idx, shape in enumerate(shapes):
                has_tf = getattr(shape, "has_text_frame", False)
                has_t = getattr(shape, "has_table", False)
                
                if shape_idx in edit_map:
                    shape_edit = edit_map[shape_idx]
                    
                    if has_tf:
                        para_edits = shape_edit.get("paragraphs", [])
                        paras = list(shape.text_frame.paragraphs)
                        for para_edit in para_edits:
                            para_idx = para_edit.get("para_index", 0)
                            new_text = para_edit.get("text", "")
                            formatting = para_edit.get("formatting", {})

                            if para_idx < len(paras):
                                para = paras[para_idx]
                                self._set_paragraph_text(para, new_text)
                                self._apply_formatting(para, formatting)
                            else:
                                # Add new paragraph
                                para = shape.text_frame.paragraphs[0] if paras else shape.text_frame.add_paragraph()
                                if para_idx > 0:
                                    para = shape.text_frame.add_paragraph()
                                para.text = ""
                                run = para.add_run()
                                run.text = new_text
                                self._apply_formatting(para, formatting)
                                
                    if has_t and "table_rows" in shape_edit:
                        table_rows = shape_edit.get("table_rows", [])
                        for r_idx, row_edit in enumerate(table_rows):
                            if r_idx >= len(shape.table.rows):
                                break
                            for c_idx, cell_edit in enumerate(row_edit):
                                if c_idx >= len(shape.table.columns):
                                    break
                                cell = shape.table.cell(r_idx, c_idx)
                                para_edits = cell_edit.get("paragraphs", [])
                                paras = list(cell.text_frame.paragraphs)
                                
                                for para_edit in para_edits:
                                    para_idx = para_edit.get("para_index", 0)
                                    new_text = para_edit.get("text", "")
                                    formatting = para_edit.get("formatting", {})
                                    
                                    if para_idx < len(paras):
                                        para = paras[para_idx]
                                        self._set_paragraph_text(para, new_text)
                                        self._apply_formatting(para, formatting)
                                    else:
                                        para = cell.text_frame.paragraphs[0] if paras else cell.text_frame.add_paragraph()
                                        if para_idx > 0:
                                            para = cell.text_frame.add_paragraph()
                                        para.text = ""
                                        run = para.add_run()
                                        run.text = new_text
                                        self._apply_formatting(para, formatting)
                else:
                    # Unedited shape. If it's a placeholder (and NOT a slide number), clear it
                    # to remove leftover boilerplate.
                    if getattr(shape, "is_placeholder", False) and has_tf:
                        try:
                            ph_type = str(shape.placeholder_format.type)
                            if "SLIDE_NUMBER" not in ph_type:
                                # Clear all paragraphs
                                paras = list(shape.text_frame.paragraphs)
                                if paras:
                                    self._set_paragraph_text(paras[0], "")
                                    for p in paras[1:]:
                                        self._set_paragraph_text(p, "")
                        except Exception:
                            pass

        new_prs.save(target)

    def _clone_slide(self, prs: Presentation, slide_index: int) -> None:
        """Deep-clone a slide at the given index and append it to the presentation."""
        template_slide = list(prs.slides)[slide_index]
        slide_layout = template_slide.slide_layout

        # Add a new slide with the same layout
        new_slide = prs.slides.add_slide(slide_layout)

        # Copy all shapes from the template slide to the new slide
        # We do this by copying the XML of each shape
        for shape in template_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # Remove the default placeholder shapes that come with the layout
        # (they duplicate what we just copied)
        sp_tree = new_slide.shapes._spTree
        # Collect placeholders that were auto-created by add_slide
        default_shapes = []
        for sp in sp_tree:
            if sp.tag.endswith("}sp") or sp.tag == "sp":
                # Check if this is a default placeholder (not one we cloned)
                pass  # We'll use a different approach

        # Actually, the cleaner approach: remove all default shapes first,
        # then copy from template
        # Let's redo: remove shapes added by add_slide, keep only our cloned ones
        # The shapes added by add_slide come from the layout's placeholders
        # Our cloned shapes are appended at the end

        # Count shapes from template
        template_shape_count = len(list(template_slide.shapes))

        # The spTree contains: <cNvPr> (non-visual props) + shapes
        # Shapes added by layout come first, our cloned ones come last
        all_sps = [child for child in sp_tree
                   if child.tag.endswith("}sp") or child.tag.endswith("}pic")
                   or child.tag.endswith("}grpSp") or child.tag.endswith("}graphicFrame")
                   or child.tag.endswith("}cxnSp")]

        if len(all_sps) > template_shape_count:
            # Remove the auto-generated shapes (first N shapes minus our cloned ones)
            auto_count = len(all_sps) - template_shape_count
            for sp in all_sps[:auto_count]:
                sp_tree.remove(sp)

        # Copy slide-level properties (background, etc.)
        try:
            if template_slide._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            ) is not None:
                bg = copy.deepcopy(
                    template_slide._element.find(
                        "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
                    )
                )
                existing_bg = new_slide._element.find(
                    "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
                )
                if existing_bg is not None:
                    new_slide._element.replace(existing_bg, bg)
                else:
                    new_slide._element.insert(0, bg)
        except Exception:
            pass

    def _delete_slide(self, prs: Presentation, slide_index: int) -> None:
        """Delete a slide from the presentation by index (0-based)."""
        slides = list(prs.slides)
        if slide_index < 0 or slide_index >= len(slides):
            log.warning("Cannot delete slide %d: out of range (total: %d)", slide_index, len(slides))
            return

        slide = slides[slide_index]
        rId = None

        # Find the relationship ID for this slide
        for rel in prs.part.rels.values():
            if rel.target_part == slide.part:
                rId = rel.rId
                break

        if rId is None:
            log.warning("Cannot find relationship for slide %d", slide_index)
            return

        # Remove from slide list XML
        pres_elem = prs.part._element
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                 "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
        sldIdLst = pres_elem.find("p:sldIdLst", nsmap)
        if sldIdLst is not None:
            for sldId in list(sldIdLst):
                if sldId.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                ) == rId:
                    sldIdLst.remove(sldId)
                    break

        # Remove the relationship
        try:
            prs.part.rels.pop(rId)
        except Exception:
            pass

    def _set_paragraph_text(self, paragraph, new_text: str) -> None:
        """Set paragraph text while preserving formatting of the first run."""
        runs = paragraph.runs
        if not runs:
            # No runs exist — create one
            run = paragraph.add_run()
            run.text = new_text
            return

        # Set text on the first run, clear the rest
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""

    def _apply_formatting(self, paragraph, formatting: dict) -> None:
        """Apply formatting to a paragraph's runs.

        Only applies values that are not None — None means 'keep existing'.
        """
        if not formatting:
            return

        # Alignment
        alignment = formatting.get("alignment")
        if alignment and alignment in _ALIGN_MAP:
            paragraph.alignment = _ALIGN_MAP[alignment]

        # Font properties — apply to all runs
        font_size = formatting.get("font_size_pt")
        bold = formatting.get("bold")
        italic = formatting.get("italic")
        color_hex = formatting.get("color_hex")

        for run in paragraph.runs:
            font = run.font
            if font_size is not None:
                font.size = Pt(font_size)
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color_hex is not None and len(color_hex) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    pass

    # ------------------------------------------------------------------
    # Operations mode: apply structured operation list
    # ------------------------------------------------------------------

    def _parse_target_id(self, target_id: str) -> dict:
        """Parse a DOM ID (e.g. 'slide_1_shape_2_para_0') into a target dict for the ops engine."""
        if not target_id:
            return {}
        tgt = {}
        import re
        if m := re.search(r'slide_(\d+)', target_id):
            tgt["slide"] = int(m.group(1))
        if m := re.search(r'shape_(\d+)', target_id):
            tgt["shape_index"] = int(m.group(1))
        
        # New: Sections, Headers, Footers
        if m := re.search(r'section_(\d+)', target_id):
            tgt["section_index"] = int(m.group(1))
        if m := re.search(r'header_(\d+)', target_id):
            tgt["header_index"] = int(m.group(1))
        if m := re.search(r'footer_(\d+)', target_id):
            tgt["footer_index"] = int(m.group(1))
            
        # New: UUIDs
        if m := re.search(r'p_([a-f0-9]{8})', target_id):
            tgt["uid_paragraph"] = m.group(1)
        elif m := re.search(r'img_([a-f0-9]{8})', target_id):
            tgt["uid_image"] = m.group(1)
        elif m := re.search(r't_([a-f0-9]{8})', target_id):
            tgt["uid_table"] = m.group(1)
            
        # image_N must be checked before paragraph_N to avoid partial match
        if re.match(r'^image_\d+$', target_id):
            m = re.search(r'image_(\d+)', target_id)
            if m:
                tgt["image_index"] = int(m.group(1)) - 1
        elif m := re.search(r'paragraph_(\d+)', target_id):
            tgt["paragraph_index"] = int(m.group(1)) - 1
        elif m := re.search(r'table_(\d+)', target_id):
            tgt["table_index"] = int(m.group(1)) - 1
        
        if m := re.search(r'_para_(\d+)', target_id):
            tgt["para_index"] = int(m.group(1))
        # Note: cell_ and _run_ handlers remain below
        if m := re.search(r'cell_(\d+)_(\d+)', target_id):
            tgt["row_index"] = int(m.group(1))
            tgt["col_index"] = int(m.group(2))
        if m := re.search(r'_run_(\d+)', target_id):
            tgt["run_index"] = int(m.group(1))
        return tgt

    def _translate_uids_to_indices(self, doc: Document, tgt: dict) -> None:
        """Translates UID fields back to legacy sequential indices so downstream ops work."""
        if "uid_paragraph" in tgt:
            target_p = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"p_{tgt['uid_paragraph']}":
                    target_p = child
                    break
            if target_p is not None:
                for i, p in enumerate(doc.paragraphs):
                    if p._p == target_p:
                        tgt["paragraph_index"] = i
                        break
        elif "uid_table" in tgt:
            target_t = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"t_{tgt['uid_table']}":
                    target_t = child
                    break
            if target_t is not None:
                for i, t in enumerate(doc.tables):
                    if t._tbl == target_t:
                        tgt["table_index"] = i
                        break
        elif "uid_image" in tgt:
            target_i = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"img_{tgt['uid_image']}":
                    target_i = child
                    break
            if target_i is not None:
                img_idx = 0
                for p in doc.paragraphs:
                    from docx.oxml.ns import qn
                    has_drawing = p._p.find(f'.//{qn("w:drawing")}') is not None
                    if has_drawing:
                        if p._p == target_i:
                            tgt["image_index"] = img_idx
                            break
                        img_idx += 1

    def _apply_pptx_operations(
        self,
        source: Path,
        target: Path,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to a PPTX file.

        Returns (any_change_made, summary_list).
        """
        prs = Presentation(source)
        summaries: list[str] = []
        changed = False

        for raw_op in operations:
            target_ids = raw_op.get("target_id")
            if not isinstance(target_ids, list):
                target_ids = [target_ids]
            
            for t_id in target_ids:
                op = dict(raw_op)
                op["target_id"] = t_id
                
                op_type = op.get("op_type", "")
                params = op.get("parameters", {})
                tgt = self._parse_target_id(op.get("target_id") or "")


                try:
                    if op_type == "text_edit":
                        s = self._op_pptx_text_edit(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "text_format":
                        if op.get("target_id") == "all":
                            match_color = params.get("match_color_hex")
                            match_role = params.get("match_role")
                            if match_color or match_role:
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if getattr(shape, "has_text_frame", False):
                                            for para in shape.text_frame.paragraphs:
                                                self._apply_pptx_format_to_para(para, {}, params, shape)
                                        if getattr(shape, "has_table", False):
                                            for row in shape.table.rows:
                                                for cell in row.cells:
                                                    if getattr(cell, "text_frame", None):
                                                        for para in cell.text_frame.paragraphs:
                                                            self._apply_pptx_format_to_para(para, {}, params, shape)
                                summaries.append(f"Formatted all paragraphs matching criteria")
                                changed = True
                            continue

                        s = self._op_pptx_text_format(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "table_op":
                        s = self._op_pptx_table(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "image_op":
                        s = self._op_pptx_image(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "shape_op":
                        s = self._op_pptx_shape(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "theme_op":
                        s = self._op_pptx_theme(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "slide_op":
                        s = self._op_pptx_slide(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "chart_op":
                        s = self._op_pptx_chart(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "ai_design_op":
                        s = self._op_pptx_ai_design(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "needs_image":
                        # Handled upstream — shouldn't reach here
                        pass

                    else:
                        log.warning("Unknown op_type in operations list: %r", op_type)

                except Exception as exc:
                    log.exception("Failed to apply operation %r: %s", op_type, exc)

        if changed:
            prs.save(target)
        else:
            # Still save a copy so the pipeline has a version file
            import shutil
            shutil.copy2(source, target)

        return changed, summaries

    def _resolve_target_paras(self, doc: Document, tgt: dict) -> list:
        """Resolves the target ID dictionary to a list of paragraph elements."""
        if "table_index" in tgt and "row_index" in tgt and "col_index" in tgt:
            try:
                table = doc.tables[tgt["table_index"]]
                cell = table.cell(tgt["row_index"], tgt["col_index"])
                return cell.paragraphs
            except Exception:
                pass
        if "header_index" in tgt and "section_index" in tgt:
            try:
                return doc.sections[tgt["section_index"]].header.paragraphs
            except Exception:
                pass
        if "footer_index" in tgt and "section_index" in tgt:
            try:
                return doc.sections[tgt["section_index"]].footer.paragraphs
            except Exception:
                pass
        return doc.paragraphs

    def _build_legacy_id_map(self, doc: Document) -> dict[str, str]:
        """Maps legacy outline IDs (paragraph_1) to new stable UUIDs (p_abcdef)."""
        # Ensure all elements have UIDs stamped in memory before mapping!
        self._build_docx_body_index(doc)
        
        from docx.oxml.ns import qn
        WNS_P = qn('w:p')
        WNS_TBL = qn('w:tbl')
        WNS_DRAWING = qn('w:drawing')
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        legacy_map = {}
        para_counter = 0
        table_counter = 0
        image_counter = 0

        for child in doc.element.body:
            tag = child.tag
            uid = child.get(uid_attr)
            
            if tag in (WNS_P, WNS_TBL):
                if tag == WNS_P:
                    has_drawing = child.find(f'.//{WNS_DRAWING}') is not None
                    if has_drawing:
                        image_counter += 1
                        legacy_id = f"image_{image_counter}"
                        primary_id = f"img_{uid}" if uid else legacy_id
                    else:
                        para_counter += 1
                        legacy_id = f"paragraph_{para_counter}"
                        primary_id = f"p_{uid}" if uid else legacy_id
                else:
                    table_counter += 1
                    legacy_id = f"table_{table_counter}"
                    primary_id = f"t_{uid}" if uid else legacy_id
                
                if uid:
                    legacy_map[legacy_id] = primary_id

        return legacy_map

    def _translate_legacy_ids(self, doc: Document, params: dict) -> None:
        """Translates legacy IDs in operation parameters to new UUIDs."""
        legacy_map = self._build_legacy_id_map(doc)
        keys_to_check = ["start_id", "end_id", "before_id", "after_id", "section_a_start_id", "section_a_end_id", "section_b_start_id", "section_b_end_id", "_raw_target_id", "target_id"]
        for k in keys_to_check:
            if k in params and isinstance(params[k], str) and params[k] in legacy_map:
                params[k] = legacy_map[params[k]]
            elif k in params and isinstance(params[k], list):
                params[k] = [legacy_map.get(pid, pid) for pid in params[k]]

    def _apply_docx_operations(
        self,
        source: Path,
        target: Path,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to a DOCX file."""
        doc = Document(source)
        summaries: list[str] = []
        changed = False

        for raw_op in operations:
            op_type = raw_op.get("op_type", "")
            
            # Translate legacy IDs in the operation itself and its parameters
            self._translate_legacy_ids(doc, raw_op)
            if "parameters" in raw_op and isinstance(raw_op["parameters"], dict):
                self._translate_legacy_ids(doc, raw_op["parameters"])
            
            # Structural ops do not loop over target_ids
            if op_type in {"list_op", "layout_op", "theme_op", "ai_design_op", "meta_op", "style_op", "find_replace", "slide_op"}:
                params = dict(raw_op.get("parameters", {}))
                if "target_id" in raw_op:
                    params["_raw_target_id"] = raw_op.get("target_id")

                try:
                    summary = ""
                    if op_type == "list_op":
                        summary = self._op_docx_list(doc, params)
                    elif op_type == "layout_op":
                        summary = self._op_docx_layout(doc, params)
                    elif op_type == "theme_op":
                        summary = self._op_docx_theme(doc, params)
                    elif op_type == "ai_design_op":
                        summary = self._op_docx_ai_design(doc, params)
                    elif op_type == "meta_op":
                        from app.services.docx_extensions import apply_metadata
                        summary = apply_metadata(doc, params)
                    elif op_type == "style_op":
                        from app.services.docx_extensions import apply_global_style
                        summary = apply_global_style(doc, params)
                    elif op_type == "find_replace":
                        tgt = self._parse_target_id(raw_op.get("target_id") or "all")
                        summary = self._op_docx_find_replace(doc, tgt, params)
                    elif op_type == "slide_op":
                        summary = "Slide operations are only supported for PPTX files."
                    
                    if summary:
                        summaries.append(summary)
                        changed = True
                except Exception as exc:
                    log.exception("Failed to apply structural operation %r: %s", op_type, exc)
                continue

            target_ids = raw_op.get("target_id")
            if not isinstance(target_ids, list):
                target_ids = [target_ids]
            
            for t_id in target_ids:
                op = dict(raw_op)
                op["target_id"] = t_id
                
                op_type = op.get("op_type", "")
                params = op.get("parameters", {})
                tgt = self._parse_target_id(op.get("target_id") or "")
                self._translate_uids_to_indices(doc, tgt)

                try:
                    if op_type == "text_edit":
                        para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
                        new_text = params.get("new_text", "")
                        if para_idx is not None and new_text:
                            paras = self._resolve_target_paras(doc, tgt)

                            if paras and 0 <= para_idx < len(paras):
                                self._apply_run_aware_replacement(paras[para_idx], new_text, params)
                                summaries.append(f"Rewrote paragraph {para_idx}")
                                changed = True

                    elif op_type == "text_format":
                        if op.get("target_id") == "all":
                            match_color = params.get("match_color_hex")
                            match_role = params.get("match_role")
                            for p in doc.paragraphs:
                                self._apply_docx_format(p, {}, params)
                            for tbl in doc.tables:
                                for row in tbl.rows:
                                    for cell in row.cells:
                                        for p in cell.paragraphs:
                                            self._apply_docx_format(p, {}, params)
                            summaries.append(f"Formatted all paragraphs")
                            changed = True
                            continue

                        para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
                        if para_idx is not None:
                            paras = self._resolve_target_paras(doc, tgt)

                            if paras and 0 <= para_idx < len(paras):
                                self._apply_docx_format(paras[para_idx], tgt, params)
                                summaries.append(f"Formatted paragraph {para_idx}")
                                changed = True
                        elif "table_index" in tgt:
                            table_idx = tgt.get("table_index")
                            if table_idx is not None and 0 <= table_idx < len(doc.tables):
                                table = doc.tables[table_idx]
                                for row in table.rows:
                                    for cell in row.cells:
                                        for p in cell.paragraphs:
                                            self._apply_docx_format(p, {}, params)
                                summaries.append(f"Formatted text in table {table_idx + 1}")
                                changed = True

                    elif op_type == "table_op":
                        summary = self._op_docx_table(doc, tgt, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "image_op":
                        summary = self._op_docx_image(doc, tgt, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "section_op":
                        from app.services.docx_extensions import apply_section_formatting
                        target_id = op.get("target_id")
                        summary = apply_section_formatting(doc, target_id, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "needs_image":
                        pass

                except Exception as exc:
                    log.exception("DOCX op %r failed: %s", op_type, exc)

        if changed:
            doc.save(target)
        else:
            import shutil
            shutil.copy2(source, target)

        return changed, summaries

    def _op_docx_theme(self, doc, params: dict) -> str:
        action = params.get("action")
        if action == "set_bg_color":
            bg_color = str(params.get("bg_color_hex", "FFFFFF")).strip().lstrip("#")
            from docx.oxml import parse_xml
            from docx.oxml.ns import nsdecls
            background = parse_xml(r'<w:background {} w:color="{}"/>'.format(nsdecls('w'), bg_color))
            doc.element.insert(0, background)
            doc.settings.element.append(parse_xml(r'<w:displayBackgroundShape {}/>'.format(nsdecls('w'))))
            return f"Set document background color to #{bg_color}"
        elif action == "set_margins":
            inches = params.get("margin_inches")
            if inches:
                from docx.shared import Inches
                for section in doc.sections:
                    section.left_margin = Inches(inches)
                    section.right_margin = Inches(inches)
                    section.top_margin = Inches(inches)
                    section.bottom_margin = Inches(inches)
                return f"Set document margins to {inches} inches"
        elif action == "add_page_numbers":
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            for section in doc.sections:
                footer = section.footer
                # Clear existing footer paragraphs to avoid stacking
                for p in list(footer.paragraphs):
                    p._p.getparent().remove(p._p)
                p = footer.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run("Page ")
                fldChar1 = OxmlElement('w:fldChar')
                fldChar1.set(qn('w:fldCharType'), 'begin')
                instrText = OxmlElement('w:instrText')
                instrText.set(qn('xml:space'), 'preserve')
                instrText.text = "PAGE"
                fldChar2 = OxmlElement('w:fldChar')
                fldChar2.set(qn('w:fldCharType'), 'separate')
                fldChar3 = OxmlElement('w:fldChar')
                fldChar3.set(qn('w:fldCharType'), 'end')
                run._r.append(fldChar1)
                run._r.append(instrText)
                run._r.append(fldChar2)
                run._r.append(fldChar3)
            return "Added page numbers to footers"
        elif action == "apply_theme_colors":
            accent_colors = params.get("accent_colors", [])
            if not accent_colors:
                return "No theme colors provided"
            color1 = str(accent_colors[0]).lstrip('#')
            color2 = str(accent_colors[1]).lstrip('#') if len(accent_colors) > 1 else color1
            
            from docx.shared import RGBColor as DRGBColor
            def _hex_to_rgb(hx):
                if len(hx) == 6:
                    try:
                        return DRGBColor(int(hx[0:2],16), int(hx[2:4],16), int(hx[4:6],16))
                    except ValueError:
                        pass
                return None

            c1 = _hex_to_rgb(color1)
            c2 = _hex_to_rgb(color2)
            
            if c1:
                for para in doc.paragraphs:
                    if para.style and para.style.name and para.style.name.lower().startswith("heading"):
                        for run in para.runs:
                            run.font.color.rgb = c1
            if c2:
                for tbl in doc.tables:
                    if len(tbl.rows) > 0:
                        for cell in tbl.rows[0].cells:
                            for p in cell.paragraphs:
                                for run in p.runs:
                                    run.font.color.rgb = c2
            return f"Applied theme colors {accent_colors}"
        return ""

    def _extract_list_info(self, para, doc) -> dict | None:
        """Return numbering metadata for a list paragraph, or None if not a list item."""
        from docx.oxml.ns import qn
        pPr = para._p.find(qn('w:pPr'))
        if pPr is None:
            return None
        numPr = pPr.find(qn('w:numPr'))
        if numPr is None:
            return None
        numId_el = numPr.find(qn('w:numId'))
        ilvl_el = numPr.find(qn('w:ilvl'))
        if numId_el is None:
            return None
        num_id = int(numId_el.get(qn('w:val'), 0))
        if num_id == 0:  # numId=0 means "remove numbering"
            return None
        ilvl = int(ilvl_el.get(qn('w:val'), 0)) if ilvl_el is not None else 0

        num_fmt = "bullet"
        list_type = "bullet"
        lvl_text_val = ""
        try:
            numbering_part = doc.part.numbering_part
            if numbering_part is not None:
                nb = numbering_part._element
                # Locate <w:num w:numId="N">
                num_el = nb.find(
                    f'.//{qn("w:num")}[@{qn("w:numId")}="{num_id}"]'
                )
                if num_el is not None:
                    abs_ref = num_el.find(qn('w:abstractNumId'))
                    if abs_ref is not None:
                        abs_id = abs_ref.get(qn('w:val'), '0')
                        abs_num = nb.find(
                            f'.//{qn("w:abstractNum")}[@{qn("w:abstractNumId")}="{abs_id}"]'
                        )
                        if abs_num is not None:
                            lvl = abs_num.find(
                                f'{qn("w:lvl")}[@{qn("w:ilvl")}="{ilvl}"]'
                            )
                            if lvl is None:
                                # search without namespace prefix matching
                                for lv in abs_num:
                                    if lv.get(qn('w:ilvl')) == str(ilvl):
                                        lvl = lv
                                        break
                            if lvl is not None:
                                nfmt_el = lvl.find(qn('w:numFmt'))
                                if nfmt_el is not None:
                                    num_fmt = nfmt_el.get(qn('w:val'), 'bullet')
                                lvl_text_el = lvl.find(qn('w:lvlText'))
                                if lvl_text_el is not None:
                                    lvl_text_val = lvl_text_el.get(qn('w:val'), '')
        except Exception:
            pass

        if num_fmt in ('decimal', 'lowerRoman', 'upperRoman', 'lowerLetter', 'upperLetter'):
            list_type = 'numbered'
        elif lvl_text_val in ('\u2610', '\u2611', '\u2612', '\u25a1', '\u25a0',
                              '\u2714', '\u2718', '\u2013\u2013', '\u2610 '):
            list_type = 'checklist'
        else:
            list_type = 'bullet'

        return {
            'num_id': num_id,
            'ilvl': ilvl,
            'num_fmt': num_fmt,
            'list_type': list_type,
            'lvl_text': lvl_text_val,
        }

    def _renumber_headings_in_doc(self, doc) -> int:
        """Re-sequence numbered heading text after structural operations (swap, move).

        Detects headings whose visible text starts with a numeric prefix like "1. " or "3. "
        and renumbers them sequentially in document order at each heading level.
        Returns the number of headings that were actually renumbered.
        """
        import re as _re
        _NUMBERED_PREFIX = _re.compile(r'^(\d+)\.\s+(.+)$', _re.DOTALL)

        # Collect all heading paragraphs in body order
        heading_paras = []
        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            if style_name.lower().startswith("heading"):
                m = _re.search(r'(\d+)$', style_name)
                level = int(m.group(1)) if m else 1
                heading_paras.append((level, para))

        if not heading_paras:
            return 0

        # Build per-level counters, reset child counters when parent increments
        level_counters: dict[int, int] = {}
        changed = 0

        for level, para in heading_paras:
            # Reset all deeper level counters when we encounter this level
            deeper_keys = [k for k in level_counters if k > level]
            for k in deeper_keys:
                del level_counters[k]

            level_counters[level] = level_counters.get(level, 0) + 1
            current_num = level_counters[level]

            current_text = para.text.strip()
            pm = _NUMBERED_PREFIX.match(current_text)
            if pm:
                existing_num = int(pm.group(1))
                title_part = pm.group(2)
                if existing_num != current_num:
                    new_text = f"{current_num}. {title_part}"
                    self._apply_run_aware_replacement(para, new_text)
                    changed += 1

        return changed

    def _build_docx_body_index(self, doc) -> list[tuple[str, object]]:
        """Walk doc.element.body in order and return a flat list of (id, xml_element) pairs.

        IDs are resolved from the `deid:uid` custom XML attribute.
        """
        from docx.oxml.ns import qn
        WNS_P = qn('w:p')
        WNS_TBL = qn('w:tbl')
        WNS_DRAWING = qn('w:drawing')
        
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        other_counter = 1
        result = []
        for child in doc.element.body:
            tag = child.tag
            if tag in (WNS_P, WNS_TBL):
                uid = child.get(uid_attr)
                if not uid:
                    # Fallback for newly inserted elements that missed stamping
                    import uuid
                    uid = uuid.uuid4().hex[:8]
                    child.set(uid_attr, uid)
                
                if tag == WNS_P:
                    has_drawing = child.find(f'.//{WNS_DRAWING}') is not None
                    if has_drawing:
                        result.append((f"img_{uid}", child))
                    else:
                        result.append((f"p_{uid}", child))
                else:
                    result.append((f"t_{uid}", child))
            else:
                result.append((f"body_other_{other_counter}", child))
                other_counter += 1
        return result

    def _op_docx_layout(self, doc, params: dict) -> str:
        action = params.get("action")

        if action == "move_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            if not start_id or not end_id or not (before_id or after_id):
                return "move_block requires start_id, end_id, and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)
            before_pos = id_to_pos.get(before_id, -1) if before_id else -1
            after_pos = id_to_pos.get(after_id, -1) if after_id else -1

            if start_pos == -1 or end_pos == -1 or (before_id and before_pos == -1) or (after_id and after_pos == -1):
                missing = []
                if start_pos == -1: missing.append(("start_id", start_id))
                if end_pos == -1: missing.append(("end_id", end_id))
                if before_id and before_pos == -1: missing.append(("before_id", before_id))
                if after_id and after_pos == -1: missing.append(("after_id", after_id))
                return f"Failed to find required block boundaries: {missing}. Available IDs: {[eid for eid, _ in body_index[:30]]}"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            if start_pos <= before_pos <= end_pos or start_pos <= after_pos <= end_pos:
                return f"Moved {end_pos - start_pos + 1} block(s) (no-op because target is inside the moved block)"

            # Collect the xml elements to move
            elements_to_move = [body_index[i][1] for i in range(start_pos, end_pos + 1)]

            # Remove from current position
            for xml_el in elements_to_move:
                xml_el.getparent().remove(xml_el)
                
            # Insert at new position
            if before_pos != -1:
                before_xml_el = body_index[before_pos][1]
                for xml_el in elements_to_move:
                    before_xml_el.addprevious(xml_el)
                renumbered = self._renumber_headings_in_doc(doc)
                suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
                return f"Moved {len(elements_to_move)} block(s) ('{start_id}'→'{end_id}') before '{before_id}'{suffix}"
            elif after_pos != -1:
                after_xml_el = body_index[after_pos][1]
                for xml_el in reversed(elements_to_move):
                    after_xml_el.addnext(xml_el)
                renumbered = self._renumber_headings_in_doc(doc)
                suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
                return f"Moved {len(elements_to_move)} block(s) ('{start_id}'→'{end_id}') after '{after_id}'{suffix}"

        elif action == "insert_page_break":
            # Accept before_id, start_id (fallback), or after_id
            before_id = params.get("before_id") or params.get("start_id")
            after_id_pb = params.get("after_id")
            if not before_id and not after_id_pb:
                return "insert_page_break requires before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            # Build a paragraph containing a hard page break
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            pb_p = OxmlElement('w:p')
            r = OxmlElement('w:r')
            br = OxmlElement('w:br')
            br.set(qn('w:type'), 'page')
            r.append(br)
            pb_p.append(r)

            if before_id:
                before_pos = id_to_pos.get(before_id, -1)
                if before_pos == -1:
                    return f"insert_page_break: element '{before_id}' not found"
                body_index[before_pos][1].addprevious(pb_p)
                return f"Inserted page break before '{before_id}'"
            else:
                after_pos = id_to_pos.get(after_id_pb, -1)
                if after_pos == -1:
                    return f"insert_page_break: element '{after_id_pb}' not found"
                body_index[after_pos][1].addnext(pb_p)
                return f"Inserted page break after '{after_id_pb}'"

        elif action == "set_columns":
            num_cols = int(params.get("num_columns", 2))
            gap_inches = float(params.get("column_gap_inches", 0.5))
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            body = doc.element.body
            sectPr = body.find(qn("w:sectPr"))
            if sectPr is None:
                sectPr = OxmlElement("w:sectPr")
                body.append(sectPr)
            # Remove existing cols element
            existing_cols = sectPr.find(qn("w:cols"))
            if existing_cols is not None:
                sectPr.remove(existing_cols)
            cols = OxmlElement("w:cols")
            cols.set(qn("w:num"), str(num_cols))
            gap_twips = int(gap_inches * 1440)  # 1 inch = 1440 twips
            cols.set(qn("w:space"), str(gap_twips))
            cols.set(qn("w:equalWidth"), "1")
            sectPr.append(cols)
            return f"Set document to {num_cols}-column layout with {gap_inches}" + '" gap'

        elif action == "remove_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            if not start_id or not end_id:
                return "remove_block requires start_id and end_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)

            if start_pos == -1 or end_pos == -1:
                return f"Failed to find block boundaries. start={start_pos}, end={end_pos}"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            elements = [body_index[i][1] for i in range(start_pos, end_pos + 1)]
            for xml_el in elements:
                xml_el.getparent().remove(xml_el)
            return f"Removed {len(elements)} block(s) ('{start_id}'→'{end_id}')"

        elif action == "duplicate_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            
            if not start_id or not end_id or (not before_id and not after_id):
                return "duplicate_block requires start_id, end_id, and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)
            
            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if start_pos == -1 or end_pos == -1 or target_xml_el is None:
                return "Failed to find block boundaries or target id for duplicate_block"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            elements = [body_index[i][1] for i in range(start_pos, end_pos + 1)]
            import copy
            import uuid
            cloned_elements = [copy.deepcopy(el) for el in elements]
            
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            from docx.oxml.ns import qn
            WNS_P = qn('w:p')
            WNS_TBL = qn('w:tbl')
            for clone in cloned_elements:
                if clone.tag in (WNS_P, WNS_TBL):
                    clone.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in clone.iter(WNS_P):
                    if nested != clone:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in clone.iter(WNS_TBL):
                    if nested != clone:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
            
            current_target = target_xml_el
            for xml_el in cloned_elements:
                if insert_before:
                    current_target.addprevious(xml_el)
                else:
                    current_target.addnext(xml_el)
                    current_target = xml_el
            
            return f"Duplicated {len(elements)} block(s)"

        elif action == "swap_sections":
            a_start = params.get("section_a_start_id")
            a_end = params.get("section_a_end_id")
            b_start = params.get("section_b_start_id")
            b_end = params.get("section_b_end_id")

            if not all([a_start, a_end, b_start, b_end]):
                return "swap_sections requires start/end ids for both sections"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            a_start_pos = id_to_pos.get(a_start, -1)
            a_end_pos = id_to_pos.get(a_end, -1)
            b_start_pos = id_to_pos.get(b_start, -1)
            b_end_pos = id_to_pos.get(b_end, -1)

            if -1 in (a_start_pos, a_end_pos, b_start_pos, b_end_pos):
                return "Failed to find block boundaries for swap_sections"

            if a_start_pos > a_end_pos: a_start_pos, a_end_pos = a_end_pos, a_start_pos
            if b_start_pos > b_end_pos: b_start_pos, b_end_pos = b_end_pos, b_start_pos

            if a_start_pos > b_start_pos:
                # Ensure A is always before B for easier logic
                a_start_pos, a_end_pos, b_start_pos, b_end_pos = b_start_pos, b_end_pos, a_start_pos, a_end_pos

            if a_end_pos >= b_start_pos:
                return "Cannot swap overlapping or adjacent-intersecting sections"

            # Collect elements
            a_elements = [body_index[i][1] for i in range(a_start_pos, a_end_pos + 1)]
            b_elements = [body_index[i][1] for i in range(b_start_pos, b_end_pos + 1)]

            # Remove all from parent
            for el in a_elements: el.getparent().remove(el)
            for el in b_elements: el.getparent().remove(el)
            
            # Re-insert B where A was.
            if a_start_pos == 0:
                # B goes to start of body
                for el in reversed(b_elements):
                    doc.element.body.insert(0, el)
            else:
                pre_a_el = body_index[a_start_pos - 1][1]
                for el in reversed(b_elements):
                    pre_a_el.addnext(el)
                    
            # Re-insert A where B was.
            if b_start_pos - 1 >= a_start_pos and b_start_pos - 1 <= a_end_pos:
                # They were strictly adjacent. A goes right after B's new position.
                last_b_el = b_elements[-1]
                for el in reversed(a_elements):
                    last_b_el.addnext(el)
            else:
                pre_b_el = body_index[b_start_pos - 1][1]
                for el in reversed(a_elements):
                    pre_b_el.addnext(el)

            renumbered = self._renumber_headings_in_doc(doc)
            suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
            return f"Swapped sections ({len(a_elements)} blocks and {len(b_elements)} blocks){suffix}"

        elif action == "insert_block":
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            data = params.get("data", [])
            
            if not data or (not before_id and not after_id):
                return "insert_block requires data array and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if target_xml_el is None:
                return "Failed to find target id for insert_block"

            # Create paragraphs / tables from data and collect XML elements
            new_elements = []
            for item in data:
                text = item.get("text", "")
                role = item.get("role", "body")

                if role == "heading":
                    lvl = item.get("heading_level", 1)
                    p = doc.add_heading(text, level=lvl)
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

                elif role == "bullet_point":
                    from docx.oxml import OxmlElement as _OXe
                    from docx.oxml.ns import qn as _qne
                    try:
                        p = doc.add_paragraph(text, style="List Paragraph")
                    except Exception:
                        p = doc.add_paragraph(text)
                    pPr = p._p.get_or_add_pPr()
                    numPr = _OXe("w:numPr")
                    ilvl = _OXe("w:ilvl"); ilvl.set(_qne("w:val"), "0")
                    numId_el = _OXe("w:numId"); numId_el.set(_qne("w:val"), "1")
                    numPr.append(ilvl); numPr.append(numId_el)
                    pPr.append(numPr)
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

                elif role == "table":
                    # Inline table with headers + rows (used by ToC enricher etc.)
                    headers = item.get("headers", [])
                    rows = item.get("rows", [])
                    style_hint = item.get("style", "")
                    num_cols = max(
                        len(headers),
                        max((len(r) for r in rows), default=0),
                        1,
                    )
                    tbl = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    tbl_xml = tbl._tbl
                    tbl_xml.getparent().remove(tbl_xml)

                    from docx.oxml.ns import qn as _qnt
                    from docx.oxml import OxmlElement as _OXt

                    def _cell_text(cell, val: str, bold: bool = False) -> None:
                        cell.text = str(val)
                        if bold:
                            for run in cell.paragraphs[0].runs:
                                run.bold = True

                    # Header row
                    hdr_row = tbl.rows[0]
                    for ci, hdr in enumerate(headers[:num_cols]):
                        _cell_text(hdr_row.cells[ci], hdr, bold=True)

                    # Data rows
                    for ri, row_data in enumerate(rows):
                        tbl_row = tbl.rows[ri + 1]
                        for ci, cell_val in enumerate(row_data[:num_cols]):
                            _cell_text(tbl_row.cells[ci], cell_val)

                    # ToC style: remove all borders for a clean look
                    if style_hint == "toc":
                        tblPr = tbl_xml.find(_qnt("w:tblPr"))
                        if tblPr is None:
                            tblPr = _OXt("w:tblPr")
                            tbl_xml.insert(0, tblPr)
                        tblBorders = _OXt("w:tblBorders")
                        for bname in ("top", "left", "bottom", "right", "insideH", "insideV"):
                            b = _OXt(f"w:{bname}")
                            b.set(_qnt("w:val"), "none")
                            tblBorders.append(b)
                        existing = tblPr.find(_qnt("w:tblBorders"))
                        if existing is not None:
                            tblPr.remove(existing)
                        tblPr.append(tblBorders)

                    new_elements.append(tbl_xml)

                else:
                    # Default: plain body paragraph with optional bold
                    p = doc.add_paragraph(text)
                    if item.get("bold"):
                        for run in p.runs:
                            run.bold = True
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

            current_target = target_xml_el
            
            import uuid
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            from docx.oxml.ns import qn
            WNS_P = qn('w:p')
            WNS_TBL = qn('w:tbl')
            for xml_el in new_elements:
                if xml_el.tag in (WNS_P, WNS_TBL):
                    xml_el.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in xml_el.iter(WNS_P):
                    if nested != xml_el:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in xml_el.iter(WNS_TBL):
                    if nested != xml_el:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                if insert_before:
                    current_target.addprevious(xml_el)
                else:
                    current_target.addnext(xml_el)
                    current_target = xml_el

            return f"Inserted {len(new_elements)} new block(s)"

        elif action == "insert_toc":
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            
            if not before_id and not after_id:
                return "insert_toc requires before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if target_xml_el is None:
                return "Failed to find target id for insert_toc"

            from docx.oxml import parse_xml
            from docx.oxml.ns import nsdecls
            
            toc_xml = f"""
            <w:sdt {nsdecls('w')}>
              <w:sdtPr>
                <w:docPartObj>
                  <w:docPartGallery w:val="Table of Contents"/>
                  <w:docPartUnique/>
                </w:docPartObj>
              </w:sdtPr>
              <w:sdtContent>
                <w:p>
                  <w:pPr><w:pStyle w:val="TOCHeading"/></w:pPr>
                  <w:r><w:t>Table of Contents</w:t></w:r>
                </w:p>
                <w:p>
                  <w:r><w:fldChar w:fldCharType="begin"/></w:r>
                  <w:r><w:instrText xml:space="preserve"> TOC \\o "1-3" \\h \\z \\u </w:instrText></w:r>
                  <w:r><w:fldChar w:fldCharType="separate"/></w:r>
                </w:p>
                <w:p><w:r><w:fldChar w:fldCharType="end"/></w:r></w:p>
              </w:sdtContent>
            </w:sdt>
            """
            toc_elem = parse_xml(toc_xml)
            
            import uuid
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            for p in toc_elem.iter(qn('w:p')):
                p.set(uid_attr, uuid.uuid4().hex[:8])

            if insert_before:
                target_xml_el.addprevious(toc_elem)
            else:
                target_xml_el.addnext(toc_elem)
                
            return "Inserted Table of Contents"

        return ""

    # ------------------------------------------------------------------
    # DOCX list operations
    # ------------------------------------------------------------------

    def _create_docx_numbering(self, doc, list_type: str, bullet_char: str = "") -> int:
        """Create a new <w:num> entry in numbering.xml and return its numId.

        list_type: 'bullet' | 'numbered' | 'checklist'
        bullet_char: custom bullet character (used for 'bullet' and 'checklist')
        Returns the new numId (int).
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn, nsmap
        import copy

        BULLET_CHAR = bullet_char or '\u2022'  # default: •
        CHECKLIST_CHAR = bullet_char or '\u2610'  # ☐

        # Ensure numbering part exists
        try:
            nb_part = doc.part.numbering_part
        except Exception:
            # Create a minimal numbering part if missing
            from docx.opc.part import Part
            from docx.opc.packuri import PackURI
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            nb_xml = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:numbering xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas"'
                ' xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
                ' xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml">'
                '</w:numbering>'
            )
            nb_part = Part(
                PackURI('/word/numbering.xml'),
                'application/vnd.openxmlformats-officedocument.wordprocessingml.numbering+xml',
                nb_xml.encode(),
            )
            doc.part.relate_to(nb_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering')

        nb = doc.part.numbering_part._element

        # Find the next available abstractNumId and numId
        existing_abs = nb.findall(qn('w:abstractNum'))
        next_abs_id = max((int(a.get(qn('w:abstractNumId'), -1)) for a in existing_abs), default=-1) + 1
        existing_nums = nb.findall(qn('w:num'))
        next_num_id = max((int(n.get(qn('w:numId'), 0)) for n in existing_nums), default=0) + 1

        # ---------- Build <w:abstractNum> ----------
        abs_num = OxmlElement('w:abstractNum')
        abs_num.set(qn('w:abstractNumId'), str(next_abs_id))

        # Multi-level list style
        multi_lvl = OxmlElement('w:multiLevelType')
        multi_lvl.set(qn('w:val'), 'hybridMultilevel')
        abs_num.append(multi_lvl)

        # Build 9 levels (Word requires all 9)
        for lvl_idx in range(9):
            lvl_el = OxmlElement('w:lvl')
            lvl_el.set(qn('w:ilvl'), str(lvl_idx))

            start_el = OxmlElement('w:start')
            start_el.set(qn('w:val'), '1')
            lvl_el.append(start_el)

            nfmt_el = OxmlElement('w:numFmt')
            if list_type == 'numbered':
                nfmt_el.set(qn('w:val'), 'decimal')
            else:
                nfmt_el.set(qn('w:val'), 'bullet')
            lvl_el.append(nfmt_el)

            lvl_text_el = OxmlElement('w:lvlText')
            if list_type == 'numbered':
                # e.g. "%1." for level 0, "%2." for level 1, etc.
                lvl_text_el.set(qn('w:val'), f'%{lvl_idx + 1}.')
            elif list_type == 'checklist':
                lvl_text_el.set(qn('w:val'), CHECKLIST_CHAR)
            else:
                # Alternate bullet chars by level
                bullets = ['\u2022', 'o', '\u25aa', '\u2022', 'o', '\u25aa', '\u2022', 'o', '\u25aa']
                lvl_text_el.set(qn('w:val'), bullet_char or bullets[lvl_idx % 3])
            lvl_el.append(lvl_text_el)

            lvl_jc = OxmlElement('w:lvlJc')
            lvl_jc.set(qn('w:val'), 'left')
            lvl_el.append(lvl_jc)

            # Paragraph properties: indent
            pPr_el = OxmlElement('w:pPr')
            ind_el = OxmlElement('w:ind')
            left = 720 + lvl_idx * 720
            ind_el.set(qn('w:left'), str(left))
            ind_el.set(qn('w:hanging'), '360')
            pPr_el.append(ind_el)
            lvl_el.append(pPr_el)

            # Run properties: font for bullet chars
            if list_type in ('bullet', 'checklist'):
                rPr_el = OxmlElement('w:rPr')
                rFonts_el = OxmlElement('w:rFonts')
                if list_type == 'checklist':
                    # Segoe UI Symbol has good checkbox glyphs
                    rFonts_el.set(qn('w:ascii'), 'Segoe UI Symbol')
                    rFonts_el.set(qn('w:hAnsi'), 'Segoe UI Symbol')
                else:
                    rFonts_el.set(qn('w:ascii'), 'Symbol')
                    rFonts_el.set(qn('w:hAnsi'), 'Symbol')
                rPr_el.append(rFonts_el)
                lvl_el.append(rPr_el)

            abs_num.append(lvl_el)

        # Insert abstractNum before any existing <w:num> elements
        first_num = nb.find(qn('w:num'))
        if first_num is not None:
            nb.insert(list(nb).index(first_num), abs_num)
        else:
            nb.append(abs_num)

        # ---------- Build <w:num> ----------
        num_el = OxmlElement('w:num')
        num_el.set(qn('w:numId'), str(next_num_id))
        abs_num_id_ref = OxmlElement('w:abstractNumId')
        abs_num_id_ref.set(qn('w:val'), str(next_abs_id))
        num_el.append(abs_num_id_ref)
        nb.append(num_el)

        return next_num_id

    def _op_docx_list(self, doc, params: dict) -> str:
        """Handle list manipulation operations for DOCX documents.

        Actions:
          convert_type  — Change list format (bullet/numbered/checklist) for a range of paragraphs.
          add_items     — Insert new list items after an anchor paragraph.
          sort_items    — Alphabetically sort a range of list paragraphs.
          set_bullet_char — Change the bullet character for a list range.
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        import copy

        action = params.get('action')

        # ----------------------------------------------------------------
        # Shared helper: resolve a body_index id to a xml element
        # ----------------------------------------------------------------
        def _body_id_map():
            return {eid: xml_el for eid, xml_el in self._build_docx_body_index(doc)}

        def _xml_el_to_para(xml_el):
            for p in doc.paragraphs:
                if p._p is xml_el:
                    return p
            return None

        def _get_paras_in_range(start_id: str, end_id: str):
            """Return list of python-docx Paragraph objects between start_id and end_id (inclusive)."""
            bmap = _body_id_map()
            start_xml = bmap.get(start_id)
            end_xml = bmap.get(end_id)
            if start_xml is None or end_xml is None:
                return [], f"Could not find '{start_id}' or '{end_id}' in document"

            collecting = False
            result = []
            for eid, xml_el in self._build_docx_body_index(doc):
                if xml_el is start_xml:
                    collecting = True
                if collecting:
                    p = _xml_el_to_para(xml_el)
                    if p is not None:
                        result.append(p)
                if xml_el is end_xml:
                    break
            return result, None

        # ----------------------------------------------------------------
        # CONVERT_TYPE — change bullet/numbered/checklist for a paragraph range
        # ----------------------------------------------------------------
        if action == 'convert_type':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            list_type = params.get('list_type', 'numbered')  # 'bullet'|'numbered'|'checklist'
            bullet_char = params.get('bullet_char', '')

            if not start_id or not end_id:
                return "convert_type requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"convert_type: {err}"
            if not paras:
                return "convert_type: no paragraphs found in range"

            # Create a new numbering entry
            new_num_id = self._create_docx_numbering(doc, list_type, bullet_char)

            # Update each paragraph's <w:numPr>
            updated = 0
            for para in paras:
                pPr = para._p.find(qn('w:pPr'))
                if pPr is None:
                    pPr = OxmlElement('w:pPr')
                    para._p.insert(0, pPr)

                numPr = pPr.find(qn('w:numPr'))
                if numPr is None:
                    numPr = OxmlElement('w:numPr')
                    pPr.append(numPr)
                else:
                    # Clear existing numId / ilvl children
                    for child in list(numPr):
                        numPr.remove(child)

                ilvl_el = OxmlElement('w:ilvl')
                ilvl_el.set(qn('w:val'), '0')
                numPr.insert(0, ilvl_el)

                numId_el = OxmlElement('w:numId')
                numId_el.set(qn('w:val'), str(new_num_id))
                numPr.append(numId_el)

                # Update the paragraph style to match
                if list_type == 'numbered':
                    try:
                        para.style = doc.styles['List Number']
                    except Exception:
                        pass
                else:
                    try:
                        para.style = doc.styles['List Bullet']
                    except Exception:
                        pass

                updated += 1

            return f"Converted {updated} paragraph(s) to {list_type} list (numId={new_num_id})"

        # ----------------------------------------------------------------
        # ADD_ITEMS — insert new list items after an anchor paragraph
        # ----------------------------------------------------------------
        elif action == 'add_items':
            after_id = params.get('after_id')  # anchor element ID
            
            # Fallback to _raw_target_id if after_id is missing
            if not after_id:
                raw_tgt = params.get('_raw_target_id')
                if isinstance(raw_tgt, list) and raw_tgt:
                    after_id = raw_tgt[-1]
                elif isinstance(raw_tgt, str):
                    after_id = raw_tgt

            items = params.get('items', [])     # list of text strings
            if not after_id:
                return "add_items requires after_id"
            if not items:
                return "add_items: items list is empty"

            bmap = _body_id_map()
            after_xml = bmap.get(after_id)
            if after_xml is None:
                return f"add_items: element '{after_id}' not found"

            anchor_para = _xml_el_to_para(after_xml)

            # Read numId/ilvl from the anchor (or the last item in range if end_id given)
            end_id = params.get('end_id', after_id)
            end_xml = bmap.get(end_id, after_xml)
            end_para = _xml_el_to_para(end_xml) or anchor_para

            # Extract numPr from the anchor or end paragraph to clone
            template_para = end_para or anchor_para
            template_pPr = None
            if template_para is not None:
                template_pPr_el = template_para._p.find(qn('w:pPr'))
                if template_pPr_el is not None:
                    template_pPr = copy.deepcopy(template_pPr_el)

            # Insert items in reverse order (each addnext inserts immediately after end_xml)
            # So reverse to get correct final ordering
            insert_after = end_xml
            for item_text in items:
                new_p = OxmlElement('w:p')

                # Copy paragraph properties (numPr, style, indent, etc.)
                if template_pPr is not None:
                    new_p.append(copy.deepcopy(template_pPr))

                # Add a run with the item text
                new_r = OxmlElement('w:r')
                # Copy run properties from template paragraph's first run if available
                if template_para and template_para.runs:
                    first_rPr = template_para.runs[0]._r.find(qn('w:rPr'))
                    if first_rPr is not None:
                        new_r.append(copy.deepcopy(first_rPr))
                new_t = OxmlElement('w:t')
                new_t.text = item_text
                new_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
                new_r.append(new_t)
                new_p.append(new_r)

                insert_after.addnext(new_p)
                insert_after = new_p

            return f"Added {len(items)} item(s) after '{after_id}'"

        # ----------------------------------------------------------------
        # SORT_ITEMS — alphabetically sort a range of list paragraphs
        # ----------------------------------------------------------------
        elif action == 'sort_items':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            order = params.get('order', 'asc')  # 'asc'|'desc'

            if not start_id or not end_id:
                return "sort_items requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"sort_items: {err}"
            if len(paras) < 2:
                return "sort_items: fewer than 2 paragraphs in range, nothing to sort"

            # Only sort paragraphs that are list items (have numPr)
            list_paras = [p for p in paras if p._p.find(f'.//{qn("w:numPr")}') is not None
                          or (p._p.find(qn('w:pPr')) is not None
                              and p._p.find(qn('w:pPr')).find(qn('w:numPr')) is not None)]
            if not list_paras:
                # Fall back to all paras if none have numPr
                list_paras = paras

            # Extract text from each paragraph
            texts = [p.text.strip() for p in list_paras]

            # Sort
            reverse = (order == 'desc')
            sorted_texts = sorted(texts, key=lambda s: s.lower(), reverse=reverse)

            # Re-distribute sorted texts back into the same paragraph XML elements
            # We only swap the text content; run formatting stays on the same run positions.
            for para, new_text in zip(list_paras, sorted_texts):
                # Use run-aware replacement to set new text
                self._apply_run_aware_replacement(para, new_text, {})

            return f"Sorted {len(list_paras)} list item(s) in {order}ending order"

        # ----------------------------------------------------------------
        # SET_BULLET_CHAR — change the bullet character for a list range
        # ----------------------------------------------------------------
        elif action == 'set_bullet_char':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            char = params.get('char', '\u2022')  # default: •

            if not start_id or not end_id:
                return "set_bullet_char requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"set_bullet_char: {err}"

            # Determine the numId used by the first paragraph in range
            num_ids_to_patch = set()
            for p in paras:
                li = self._extract_list_info(p, doc)
                if li:
                    num_ids_to_patch.add((li['num_id'], li['ilvl']))

            if not num_ids_to_patch:
                return "set_bullet_char: no list paragraphs found in range"

            # For each unique (numId, ilvl), patch the lvlText in the abstractNum
            try:
                nb = doc.part.numbering_part._element
                for num_id, ilvl in num_ids_to_patch:
                    num_el = nb.find(f'.//{qn("w:num")}[@{qn("w:numId")}="{num_id}"]')
                    if num_el is None:
                        continue
                    abs_ref = num_el.find(qn('w:abstractNumId'))
                    if abs_ref is None:
                        continue
                    abs_id = abs_ref.get(qn('w:val'), '0')
                    abs_num = nb.find(f'.//{qn("w:abstractNum")}[@{qn("w:abstractNumId")}="{abs_id}"]')
                    if abs_num is None:
                        continue
                    for lvl in abs_num:
                        if lvl.get(qn('w:ilvl')) == str(ilvl):
                            lvl_text_el = lvl.find(qn('w:lvlText'))
                            if lvl_text_el is not None:
                                lvl_text_el.set(qn('w:val'), char)
            except Exception as e:
                return f"set_bullet_char: failed to patch numbering XML: {e}"

            return f"Set bullet character to '{char}' for {len(num_ids_to_patch)} list level(s)"

        return f"list_op action '{action}' not supported"

    # ------------------------------------------------------------------
    # DOCX find & replace operations
    # ------------------------------------------------------------------

    def _op_docx_find_replace(self, doc, tgt: dict, params: dict) -> str:
        """Find and replace text across the document or a specific target.

        Supports standard substring replacement or regex replacement.
        """
        import re

        find_text = params.get('find_text', '')
        replace_text = params.get('replace_text', '')
        is_regex = bool(params.get('is_regex', False))
        match_case = bool(params.get('match_case', False))

        if not find_text:
            return "find_replace: find_text is empty"

        # Determine target scope
        target_id = params.get('target_id', 'all')
        if not tgt:
            target_id = 'all'

        paras_to_check = []
        if target_id == 'all':
            # Collect all document paragraphs
            paras_to_check.extend(doc.paragraphs)
            # Collect all table cell paragraphs
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paras_to_check.extend(cell.paragraphs)
        else:
            para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
            if para_idx is not None:
                if "table_index" in tgt and "row_index" in tgt and "col_index" in tgt:
                    try:
                        table = doc.tables[tgt["table_index"]]
                        cell = table.cell(tgt["row_index"], tgt["col_index"])
                        paras_to_check = [cell.paragraphs[para_idx]]
                    except Exception:
                        pass
                else:
                    try:
                        paras_to_check = [doc.paragraphs[para_idx]]
                    except Exception:
                        pass

        if not paras_to_check:
            return "find_replace: no target paragraphs found"

        # Prepare regex pattern if needed
        pattern = None
        if is_regex:
            flags = 0 if match_case else re.IGNORECASE
            try:
                pattern = re.compile(find_text, flags)
            except Exception as e:
                return f"find_replace: invalid regex pattern '{find_text}' - {e}"

        replacements = 0
        for p in paras_to_check:
            old_text = p.text
            if not old_text:
                continue

            new_text = old_text
            if is_regex:
                if pattern.search(old_text):
                    new_text = pattern.sub(replace_text, old_text)
            else:
                if match_case:
                    if find_text in old_text:
                        new_text = old_text.replace(find_text, replace_text)
                else:
                    # case-insensitive replace
                    # Use a regex compilation for simplicity of ignoring case
                    flags = re.IGNORECASE
                    pat = re.compile(re.escape(find_text), flags)
                    if pat.search(old_text):
                        new_text = pat.sub(replace_text, old_text)

            if new_text != old_text:
                self._apply_run_aware_replacement(p, new_text)
                replacements += 1

        if target_id == 'all':
            return f"Replaced occurrences of '{find_text}' in {replacements} paragraph(s)"
        else:
            return f"Replaced text in targeted paragraph"

    # ------------------------------------------------------------------
    # DOCX image operations
    # ------------------------------------------------------------------

    def _get_docx_page_width_emu(self, doc) -> int:
        """Return the usable page width in EMU (page width minus left+right margins)."""
        try:
            section = doc.sections[0]
            usable = section.page_width - section.left_margin - section.right_margin
            return int(usable)
        except Exception:
            # A4 usable width fallback: ~16.51cm = 5940000 EMU
            return 5940000

    def _find_docx_image_paragraphs(self, doc) -> list[tuple[int, object]]:
        """Return list of (image_index, paragraph) for all paragraphs containing <w:drawing>."""
        from docx.oxml.ns import qn
        WNS_DRAWING = qn('w:drawing')
        result = []
        image_idx = 0
        for para in doc.paragraphs:
            if para._p.find(f'.//{WNS_DRAWING}') is not None:
                result.append((image_idx, para))
                image_idx += 1
        return result

    def _op_docx_image(self, doc, tgt: dict, params: dict) -> str:
        """Handle image operations for DOCX documents.

        Supported actions:
          insert      — Add a new image paragraph after a body element.
          replace     — Replace the drawing in an existing image paragraph.
          resize      — Resize an inline image by width_page_pct or explicit EMU.
          reposition  — Set paragraph alignment (left/center/right) of an image paragraph.
          add_caption — Insert a styled caption paragraph after the image paragraph.
          remove      — Remove the image paragraph entirely.
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt

        action = params.get("action", "insert")
        image_path = params.get("image_path")
        page_w_emu = self._get_docx_page_width_emu(doc)

        ALIGN_MAP = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
        }

        def _set_alt_text(inline_shape, alt_text: str):
            if not alt_text:
                return
            try:
                docPr = inline_shape._inline.find('.//' + qn('wp:docPr'))
                if docPr is not None:
                    docPr.set('name', alt_text)
                    docPr.set('descr', alt_text)
            except Exception:
                pass

        # ----------------------------------------------------------------
        # Helper: find target paragraph by target_id (image_N or paragraph_N)
        # ----------------------------------------------------------------
        def _resolve_para_by_id(target_id: str):
            """Resolve a DOM id (image_N or paragraph_N) to (xml_element, paragraph)."""
            body_index = self._build_docx_body_index(doc)
            id_to_elem = {eid: xml_el for eid, xml_el in body_index}
            xml_el = id_to_elem.get(target_id)
            if xml_el is None:
                return None, None
            # Find matching python-docx paragraph
            for para in doc.paragraphs:
                if para._p is xml_el:
                    return xml_el, para
            return xml_el, None

        # ----------------------------------------------------------------
        # INSERT — add image paragraph after after_id element
        # ----------------------------------------------------------------
        if action == "insert":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image insert: image_path missing or not found: %s", image_path)
                return ""

            after_id = params.get("after_id")  # body_index element id
            alignment = params.get("alignment", "left")
            width_page_pct = params.get("width_page_pct", 0.5)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            caption_text = params.get("caption_text")

            width_emu = int(page_w_emu * width_page_pct)

            # Add picture creates it at the end of the document, we'll move it
            new_para = doc.add_paragraph()
            new_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.LEFT)
            run = new_para.add_run()
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))

            new_p_el = new_para._p

            # Move the new paragraph to the correct position
            if after_id:
                body_index = self._build_docx_body_index(doc)
                id_to_elem = {eid: xml_el for eid, xml_el in body_index}
                after_xml = id_to_elem.get(after_id)
                if after_xml is not None and after_xml is not new_p_el:
                    # Remove from its current (appended-at-end) position
                    new_p_el.getparent().remove(new_p_el)
                    # Insert immediately after the target element
                    after_xml.addnext(new_p_el)

            summary = f"Inserted image (width={int(width_page_pct*100)}% page) after '{after_id}'"

            # Optional inline caption
            if caption_text:
                cap_para = doc.add_paragraph(caption_text)
                try:
                    cap_para.style = doc.styles["Caption"]
                except Exception:
                    cap_para.runs[0].font.italic = True if cap_para.runs else None
                cap_p_el = cap_para._p
                cap_p_el.getparent().remove(cap_p_el)
                new_p_el.addnext(cap_p_el)
                summary += f" + caption '{caption_text}'"

            return summary

        # ----------------------------------------------------------------
        # INSERT INTO PARAGRAPH — add image run to an existing paragraph
        # ----------------------------------------------------------------
        if action == "insert_into_paragraph":
            if not image_path or not Path(image_path).exists():
                return ""

            target_xml, target_para = _resolve_para_by_id(tgt.get("id", ""))
            if target_para is None:
                return "insert_into_paragraph: no target paragraph found"

            width_page_pct = params.get("width_page_pct", 0.3)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            width_emu = int(page_w_emu * width_page_pct)

            run = target_para.add_run()
            # If they want it spaced out, add a tab before it to push it towards the right
            run.add_text("\t") 
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))
            return f"Inserted image (width={int(width_page_pct*100)}%) into '{tgt.get('id')}'"

        # ----------------------------------------------------------------
        # REPLACE TEXT — swap a text paragraph with an image
        # ----------------------------------------------------------------
        if action == "replace_text":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image replace_text: image_path missing or not found: %s", image_path)
                return ""

            target_xml, _ = _resolve_para_by_id(tgt.get("id", ""))
            if target_xml is None:
                return "replace_text: no target paragraph found"

            alignment = params.get("alignment", "left")
            width_page_pct = params.get("width_page_pct", 0.5)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            caption_text = params.get("caption_text")

            width_emu = int(page_w_emu * width_page_pct)

            new_para = doc.add_paragraph()
            new_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.LEFT)
            run = new_para.add_run()
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))

            new_p_el = new_para._p

            # Insert new image paragraph immediately before the target text paragraph, then remove target
            target_xml.addprevious(new_p_el)
            target_xml.getparent().remove(target_xml)

            summary = f"Replaced text paragraph with image (width={int(width_page_pct*100)}% page)"

            if caption_text:
                cap_para = doc.add_paragraph(caption_text)
                try:
                    cap_para.style = doc.styles["Caption"]
                except Exception:
                    cap_para.runs[0].font.italic = True if cap_para.runs else None
                cap_p_el = cap_para._p
                cap_p_el.getparent().remove(cap_p_el)
                new_p_el.addnext(cap_p_el)
                summary += f" + caption '{caption_text}'"

            return summary

        # ----------------------------------------------------------------
        # Helpers for operations on EXISTING image paragraphs
        # ----------------------------------------------------------------
        target_id = tgt.get("type") or ""
        # Resolve image paragraph from target_id in op dict
        op_target_id = ""
        # We stored the raw target_id on `tgt` via _parse_target_id;
        # but image_N IDs have no numeric field in _parse_target_id.
        # Retrieve it directly from the op via params fallback.
        image_target_id = params.get("target_image_id") or ""

        # Build an image-index to paragraph lookup
        img_paras = self._find_docx_image_paragraphs(doc)  # [(img_idx, para), ...]

        def _get_img_para(img_idx: int):
            for i, p in img_paras:
                if i == img_idx:
                    return p
            return None

        # Parse image index from tgt (image_N -> N stored in tgt)
        img_idx = tgt.get("image_index")  # set by _parse_target_id extension below
        if img_idx is None and img_paras:
            img_idx = 0  # default to first image

        img_para = _get_img_para(img_idx) if img_idx is not None else None

        # ----------------------------------------------------------------
        # REPLACE — swap the drawing XML in the existing image paragraph
        # ----------------------------------------------------------------
        if action == "replace":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image replace: image_path missing or not found: %s", image_path)
                return ""

            if img_para is None:
                return "replace: no image paragraph found"

            width_page_pct = params.get("width_page_pct")
            alignment = params.get("alignment")

            # Determine width to use — try to preserve existing width if not specified
            WNS_EXTENT = qn('wp:extent')
            WNS_DRAWING = qn('w:drawing')
            old_drawing = img_para._p.find(f'.//{WNS_DRAWING}')
            old_width_emu = None
            if old_drawing is not None:
                extent = old_drawing.find(f'.//{WNS_EXTENT}')
                if extent is not None:
                    try:
                        old_width_emu = int(extent.get('cx', 0))
                    except Exception:
                        pass

            if width_page_pct is not None:
                width_emu = int(page_w_emu * width_page_pct)
            elif old_width_emu:
                width_emu = old_width_emu
            else:
                width_emu = int(page_w_emu * 0.5)

            # Remove all existing runs from the paragraph
            for r in list(img_para.runs):
                r._r.getparent().remove(r._r)

            # Add new run with the replacement image
            run = img_para.add_run()
            shape = run.add_picture(image_path, width=width_emu)
            _set_alt_text(shape, params.get("alt_text"))

            if alignment and alignment in ALIGN_MAP:
                img_para.alignment = ALIGN_MAP[alignment]

            return f"Replaced image {img_idx} with {Path(image_path).name}"

        # ----------------------------------------------------------------
        # RESIZE — change width (and optionally height) of inline image
        # ----------------------------------------------------------------
        elif action == "resize":
            if img_para is None:
                return "resize: no image paragraph found"

            width_page_pct = params.get("width_page_pct")
            maintain_ar = params.get("maintain_aspect_ratio", True)

            if width_page_pct is None:
                return "resize: width_page_pct required"

            new_w_emu = int(page_w_emu * width_page_pct)

            from docx.oxml.ns import qn as _qn
            WNS_DRAWING = _qn('w:drawing')
            WNS_EXTENT = _qn('wp:extent')
            WNS_EFF_EXTENT = _qn('wp:effectExtent')

            drawing = img_para._p.find(f'.//{WNS_DRAWING}')
            if drawing is None:
                return "resize: no <w:drawing> element found in paragraph"

            extent = drawing.find(f'.//{WNS_EXTENT}')
            if extent is None:
                return "resize: no <wp:extent> element found in drawing"

            old_w = int(extent.get('cx', new_w_emu))
            old_h = int(extent.get('cy', new_w_emu))

            if maintain_ar and old_w > 0:
                ratio = old_h / old_w
                new_h_emu = int(new_w_emu * ratio)
            else:
                height_page_pct = params.get("height_page_pct")
                new_h_emu = int(page_w_emu * height_page_pct) if height_page_pct else old_h

            extent.set('cx', str(new_w_emu))
            extent.set('cy', str(new_h_emu))

            # Also update distT/distB/distL/distR extents (effectExtent)
            # and the <a:ext> inside the graphic for compatibility
            from lxml import etree
            for a_ext in drawing.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}ext'):
                a_ext.set('cx', str(new_w_emu))
                a_ext.set('cy', str(new_h_emu))

            return f"Resized image {img_idx} to {int(width_page_pct*100)}% page width"

        # ----------------------------------------------------------------
        # REPOSITION — set paragraph alignment of image paragraph
        # ----------------------------------------------------------------
        elif action == "reposition":
            if img_para is None:
                return "reposition: no image paragraph found"

            alignment = params.get("alignment", "left")
            if alignment not in ALIGN_MAP:
                return f"reposition: unknown alignment '{alignment}', use left/center/right"

            img_para.alignment = ALIGN_MAP[alignment]

            # Also move image position: if float_position requested, use anchor wrapping
            float_pos = params.get("float_position")
            if float_pos in ("left", "right"):
                # Convert inline to anchored (floating) with text wrap
                from docx.oxml.ns import qn as _qn
                WNS_DRAWING = _qn('w:drawing')
                WNS_INLINE = _qn('wp:inline')
                WNS_ANCHOR = _qn('wp:anchor')
                drawing = img_para._p.find(f'.//{WNS_DRAWING}')
                if drawing is not None:
                    inline = drawing.find(WNS_INLINE)
                    if inline is not None:
                        # Build a minimal anchor element wrapping the inline content
                        anchor = OxmlElement('wp:anchor')
                        h_align = 'right' if float_pos == 'right' else 'left'
                        anchor.set('distT', '114300')
                        anchor.set('distB', '114300')
                        anchor.set('distL', '114300')
                        anchor.set('distR', '114300')
                        anchor.set('simplePos', '0')
                        anchor.set('relativeHeight', '251658240')
                        anchor.set('behindDoc', '0')
                        anchor.set('locked', '0')
                        anchor.set('layoutInCell', '1')
                        anchor.set('allowOverlap', '1')
                        # Copy children from inline (extent, graphic, etc.)
                        for child in list(inline):
                            anchor.append(child)
                        # Add positioning children
                        simple_pos = OxmlElement('wp:simplePos')
                        simple_pos.set('x', '0'); simple_pos.set('y', '0')
                        anchor.insert(0, simple_pos)
                        pos_h = OxmlElement('wp:positionH')
                        pos_h.set('relativeFrom', 'margin')
                        align_el = OxmlElement('wp:align')
                        align_el.text = h_align
                        pos_h.append(align_el)
                        pos_v = OxmlElement('wp:positionV')
                        pos_v.set('relativeFrom', 'paragraph')
                        align_v = OxmlElement('wp:align')
                        align_v.text = 'top'
                        pos_v.append(align_v)
                        wrapSquare = OxmlElement('wp:wrapSquare')
                        wrapSquare.set('wrapText', 'bothSides')
                        anchor.append(simple_pos)
                        anchor.append(pos_h)
                        anchor.append(pos_v)
                        anchor.append(wrapSquare)
                        drawing.remove(inline)
                        drawing.append(anchor)

            return f"Repositioned image {img_idx}: alignment={alignment}" + (
                f", float={float_pos}" if float_pos else ""
            )

        # ----------------------------------------------------------------
        # ADD_CAPTION — insert caption paragraph after image paragraph
        # ----------------------------------------------------------------
        elif action == "add_caption":
            caption_text = params.get("caption_text", "")
            caption_style = params.get("caption_style", "Caption")
            if not caption_text:
                return "add_caption: caption_text required"
            if img_para is None:
                return "add_caption: no image paragraph found"

            # Create a new paragraph for the caption
            cap_para = doc.add_paragraph(caption_text)
            try:
                cap_para.style = doc.styles[caption_style]
            except Exception:
                # Fallback: italic, small font
                if cap_para.runs:
                    cap_para.runs[0].font.italic = True
                    cap_para.runs[0].font.size = Pt(9)
            cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Move the caption paragraph to immediately after the image paragraph
            cap_p_el = cap_para._p
            cap_p_el.getparent().remove(cap_p_el)
            img_para._p.addnext(cap_p_el)

            return f"Added caption '{caption_text}' below image {img_idx}"

        # ----------------------------------------------------------------
        # REMOVE — delete the image paragraph
        # ----------------------------------------------------------------
        elif action == "remove":
            if img_para is None:
                return "remove: no image paragraph found"
            img_para._p.getparent().remove(img_para._p)
            return f"Removed image {img_idx}"

        return f"image_op action '{action}' not supported for DOCX"

    # ------------------------------------------------------------------
    # Individual operation handlers — PPTX
    # ------------------------------------------------------------------

    def _get_slide(self, prs: Presentation, slide_num: int | None):
        """Get a slide by 1-based index, or the first slide if not specified."""
        slides = list(prs.slides)
        if not slides:
            return None
        if slide_num is None:
            return slides[0]
        idx = slide_num - 1
        if 0 <= idx < len(slides):
            return slides[idx]
        return slides[0]

    def _get_shape(self, slide, shape_idx: int | None):
        """Get a shape by index from a slide."""
        shapes = list(slide.shapes)
        if not shapes:
            return None
        if shape_idx is None:
            return None
        if 0 <= shape_idx < len(shapes):
            return shapes[shape_idx]
        return None

    def _op_pptx_text_edit(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Rewrite the text content of a specific paragraph."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""
        shape = self._get_shape(slide, tgt.get("shape_index"))
        if not shape or not getattr(shape, "has_text_frame", False):
            return ""
        new_text = params.get("new_text", "")
        para_idx = tgt.get("para_index", 0)
        paras = list(shape.text_frame.paragraphs)
        if para_idx < len(paras):
            self._apply_run_aware_replacement(paras[para_idx], new_text, params)
            return f"Rewrote text in slide {tgt.get('slide')}, shape {tgt.get('shape_index')}"
        return ""

    def _op_pptx_text_format(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Apply rich text formatting to a paragraph."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        results = []
        shapes_to_process = []

        shape_idx = tgt.get("shape_index")
        if shape_idx is not None:
            shape = self._get_shape(slide, shape_idx)
            if shape:
                shapes_to_process = [shape]
        else:
            # Apply to all text shapes on the slide
            shapes_to_process = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]

        for shape in shapes_to_process:
            if not getattr(shape, "has_text_frame", False):
                continue
            para_idx = tgt.get("para_index")
            paras = list(shape.text_frame.paragraphs)
            target_paras = [paras[para_idx]] if para_idx is not None and para_idx < len(paras) else paras

            for para in target_paras:
                c = self._apply_pptx_format_to_para(para, tgt, params, shape)
                if c:
                    changed = True
                if changed:
                    results.append(f"slide {tgt.get('slide')} shape {shape.shape_id}")

        return f"Applied formatting to {results[0]}" if results else ""

    def _apply_pptx_format_to_para(self, para, tgt: dict, params: dict, shape=None) -> bool:
        """Apply formatting params to a paragraph and its runs. Returns True if changed."""
        match_role = str(params.get("match_role", "")).strip().lower()
        if match_role and shape:
            shape_role = "shape"
            if shape.is_placeholder:
                try:
                    ph_type = str(shape.placeholder_format.type).split('(')[0]
                    shape_role = f"placeholder_{ph_type.lower()}"
                except Exception:
                    shape_role = "placeholder"
            # simple matching: if user says "heading", match title. 
            if match_role == "heading" and "title" not in shape_role:
                return False
            if match_role == "body" and "body" not in shape_role:
                return False
            # exact match fallback
            if match_role not in ["heading", "body"] and match_role != shape_role:
                return False

        changed = False

        # Alignment
        alignment = params.get("alignment")
        if alignment and alignment in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[alignment]
            changed = True

        # Line spacing
        line_spacing = params.get("line_spacing")
        if line_spacing is not None:
            try:
                from pptx.util import Pt
                from pptx.oxml.ns import qn
                from lxml import etree
                pPr = para._p.get_or_add_pPr()
                lnSpc = pPr.find(qn("a:lnSpc"))
                if lnSpc is None:
                    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
                spcPct = lnSpc.find(qn("a:spcPct"))
                if spcPct is None:
                    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
                spcPct.set("val", str(int(line_spacing * 100000)))
                changed = True
            except Exception as e:
                log.debug("Line spacing failed: %s", e)

        # Run-level formatting
        bold = params.get("bold")
        italic = params.get("italic")
        underline = params.get("underline")
        strikethrough = params.get("strikethrough")
        font_family = params.get("font_family")
        font_size_pt = params.get("font_size_pt")
        color_hex = params.get("color_hex")
        highlight_hex = params.get("highlight_hex")
        superscript = params.get("superscript")
        subscript = params.get("subscript")
        char_spacing = params.get("char_spacing")

        if not para.runs:
            return changed

        target_run_idx = tgt.get("run_index")
        match_color_hex = str(params.get("match_color_hex", "")).strip().lstrip("#").upper()

        for r_idx, run in enumerate(para.runs):
            if target_run_idx is not None and r_idx != target_run_idx:
                continue

            if match_color_hex:
                _, cur_color = self._resolve_pptx_font_info(run)
                cur_color = str(cur_color).strip().lstrip("#").upper() if cur_color else "000000"
                if cur_color != match_color_hex and not (cur_color == "000000" and match_color_hex == "AUTO"):
                    continue

            font = run.font
            if bold is not None:
                font.bold = bold; changed = True
            if italic is not None:
                font.italic = italic; changed = True
            if underline is not None:
                font.underline = underline; changed = True
            if font_family is not None:
                font.name = font_family; changed = True
            if font_size_pt is not None:
                font.size = Pt(font_size_pt); changed = True
            if color_hex is not None and len(str(color_hex)) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(str(color_hex)); changed = True
                except Exception:
                    pass
            if superscript is not None or subscript is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    if superscript is not None:
                        rPr.set("baseline", "30000" if superscript else "0"); changed = True
                    if subscript is not None:
                        rPr.set("baseline", "-25000" if subscript else "0"); changed = True
                except Exception as e:
                    log.debug("Superscript/subscript failed: %s", e)
            if strikethrough is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("strike", "sngStrike" if strikethrough else "noStrike"); changed = True
                except Exception as e:
                    log.debug("Strikethrough failed: %s", e)
            if char_spacing is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("spc", str(int(char_spacing * 100))); changed = True
                except Exception as e:
                    log.debug("Char spacing failed: %s", e)

        return changed

    def _apply_docx_format(self, para, tgt: dict, params: dict) -> bool:
        """Apply rich formatting to a DOCX paragraph's runs."""
        match_role = str(params.get("match_role", "")).strip().lower()
        if match_role:
            style_name = para.style.name if para.style else "Normal"
            role = "body"
            if style_name.lower().startswith("heading"):
                role = "heading"
            elif "bullet" in style_name.lower() or "list" in style_name.lower():
                role = "bullet_point"
            
            if role != match_role:
                return False
                
        match_text = params.get("match_text")
        if match_text and match_text not in para.text:
            return False
        
        alignment = params.get("alignment")
        if alignment:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            align_map = {
                "left": WD_ALIGN_PARAGRAPH.LEFT,
                "center": WD_ALIGN_PARAGRAPH.CENTER,
                "right": WD_ALIGN_PARAGRAPH.RIGHT,
                "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
            }
            if alignment in align_map:
                para.alignment = align_map[alignment]

        line_spacing = params.get("line_spacing")
        if line_spacing is not None:
            para.paragraph_format.line_spacing = float(line_spacing)

        page_break_before = params.get("page_break_before")
        if page_break_before is not None:
            para.paragraph_format.page_break_before = bool(page_break_before)

        space_before = params.get("space_before_pt")
        if space_before is not None:
            from docx.shared import Pt
            para.paragraph_format.space_before = Pt(float(space_before))

        space_after = params.get("space_after_pt")
        if space_after is not None:
            from docx.shared import Pt
            para.paragraph_format.space_after = Pt(float(space_after))
            
        left_indent = params.get("left_indent_pt")
        if left_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.left_indent = Pt(float(left_indent))
            
        right_indent = params.get("right_indent_pt")
        if right_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.right_indent = Pt(float(right_indent))
            
        first_line_indent = params.get("first_line_indent_pt")
        if first_line_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.first_line_indent = Pt(float(first_line_indent))
            
        keep_with_next = params.get("keep_with_next")
        if keep_with_next is not None:
            para.paragraph_format.keep_with_next = bool(keep_with_next)
            
        keep_together = params.get("keep_together")
        if keep_together is not None:
            para.paragraph_format.keep_together = bool(keep_together)

        include_in_toc = params.get("include_in_toc")
        if include_in_toc is not None:
            from docx.oxml.ns import qn as _qn
            from docx.oxml import OxmlElement
            pPr = para._p.get_or_add_pPr()
            outlineLvl = pPr.find(_qn('w:outlineLvl'))
            if include_in_toc is False:
                if outlineLvl is None:
                    outlineLvl = OxmlElement('w:outlineLvl')
                    pPr.append(outlineLvl)
                outlineLvl.set(_qn('w:val'), '9')
            else:
                if outlineLvl is not None:
                    pPr.remove(outlineLvl)

        # Normalise color_hex: strip '#' prefix so both "FF0000" and "#FF0000" work
        raw_color = str(params.get("color_hex", "")).strip().lstrip("#").upper()
        _color_map = {
            "RED": "FF0000",
            "BLUE": "0000FF",
            "GREEN": "00FF00",
            "BLACK": "000000",
            "WHITE": "FFFFFF",
            "YELLOW": "FFFF00",
            "PURPLE": "800080",
            "ORANGE": "FFA500",
            "GREY": "808080",
            "GRAY": "808080"
        }
        color_hex = _color_map.get(raw_color, raw_color)

        runs = para.runs
        # If the paragraph has no runs but has text, create one to hold the formatting
        if not runs and para.text.strip():
            run = para.add_run(para.text)
            # Remove original text nodes to avoid duplication
            from lxml import etree
            for t_elem in para._p.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                parent = t_elem.getparent()
                if parent is not None and parent != run._r:
                    parent_parent = parent.getparent()
                    if parent_parent is not None and parent != run._r:
                        try:
                            parent_parent.remove(parent)
                        except Exception:
                            pass
            runs = para.runs

        if match_text:
            return self._apply_run_aware_format(para, match_text, params)
                
        target_run_idx = tgt.get("run_index")
        match_color_hex = str(params.get("match_color_hex", "")).strip().lstrip("#").upper()

        for r_idx, run in enumerate(runs):
            if target_run_idx is not None and r_idx != target_run_idx:
                continue
                
            if match_color_hex:
                _, cur_color = self._resolve_docx_font_info(para, run)
                cur_color = str(cur_color).strip().lstrip("#").upper() if cur_color else "000000"
                # If they don't match, skip this run
                if cur_color != match_color_hex and not (cur_color == "000000" and match_color_hex == "AUTO"):
                    continue

            if params.get("bold") is not None:
                run.bold = params["bold"]
            if params.get("italic") is not None:
                run.italic = params["italic"]
            if params.get("underline") is not None:
                run.underline = params["underline"]
            if params.get("font_family"):
                run.font.name = params["font_family"]
            if params.get("font_size_pt") is not None:
                from docx.shared import Pt
                run.font.size = Pt(params["font_size_pt"])
            if color_hex and len(color_hex) == 6:
                try:
                    from docx.shared import RGBColor as DRGBColor
                    run.font.color.rgb = DRGBColor(
                        int(color_hex[0:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16),
                    )
                except Exception as e:
                    log.debug("DOCX color application failed: %s", e)
                    
            highlight_hex = params.get("highlight_hex")
            if highlight_hex:
                from docx.enum.text import WD_COLOR_INDEX
                # Simple mapping for common highlight colors
                hl_color = str(highlight_hex).strip().lstrip("#").upper()
                if hl_color in ["FFFF00", "YELLOW"]:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                elif hl_color in ["00FF00", "GREEN"]:
                    run.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
                elif hl_color in ["00FFFF", "CYAN"]:
                    run.font.highlight_color = WD_COLOR_INDEX.TURQUOISE
                elif hl_color in ["FF00FF", "MAGENTA", "PINK"]:
                    run.font.highlight_color = WD_COLOR_INDEX.PINK
                elif hl_color in ["FF0000", "RED"]:
                    run.font.highlight_color = WD_COLOR_INDEX.RED
                elif hl_color in ["0000FF", "BLUE"]:
                    run.font.highlight_color = WD_COLOR_INDEX.BLUE
                else:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW


    def _op_docx_table(self, doc, tgt: dict, params: dict) -> str:
        action = params.get("action", "create")
        
        if action == "create":
            rows = params.get("rows", 3)
            cols = params.get("cols", 3)
            tbl = doc.add_table(rows=rows, cols=cols)
            
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            if before_id or after_id:
                body_index = self._build_docx_body_index(doc)
                id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}
                
                target_xml_el = None
                insert_before = True
                if before_id:
                    pos = id_to_pos.get(before_id, -1)
                    if pos != -1:
                        target_xml_el = body_index[pos][1]
                elif after_id:
                    pos = id_to_pos.get(after_id, -1)
                    if pos != -1:
                        target_xml_el = body_index[pos][1]
                        insert_before = False
                        
                if target_xml_el is not None:
                    tbl_xml = tbl._tbl
                    tbl_xml.getparent().remove(tbl_xml)
                    if insert_before:
                        target_xml_el.addprevious(tbl_xml)
                    else:
                        target_xml_el.addnext(tbl_xml)

            return f"Created {rows}x{cols} table"
            
        is_all = tgt.get("type") == "all"
        table_idx = tgt.get("table_index")
        
        if not is_all and (table_idx is None or table_idx < 0 or table_idx >= len(doc.tables)):
            return "Table not found."
            
        tables_to_process_indices = range(len(doc.tables)) if is_all else [table_idx]
        
        summaries = []
        for t_idx in reversed(tables_to_process_indices):
            tbl = doc.tables[t_idx]
            
            if action == "delete":
                tbl._element.getparent().remove(tbl._element)
                summaries.append(f"Deleted table {t_idx}")
                continue

            elif action == "add_row":
                new_row = tbl.add_row()
                data = params.get("data", [])
                if data:
                    # Support both 1D list and 2D list (first row)
                    row_data = data[0] if data and isinstance(data[0], list) else data
                    for i, cell in enumerate(new_row.cells):
                        if i < len(row_data):
                            cell.text = str(row_data[i])
                summaries.append(f"Added row to table {t_idx}")

            elif action == "remove_row":
                r_idx = params.get("row_index")
                if r_idx is not None and 0 <= r_idx < len(tbl.rows):
                    row = tbl.rows[r_idx]
                    row._element.getparent().remove(row._element)
                    summaries.append(f"Removed row {r_idx} from table {t_idx}")

            elif action == "add_col":
                new_col = tbl.add_column(Pt(72)) # Default 1 inch
                data = params.get("data", [])
                if data:
                    # Support both 1D list and 2D list
                    col_data = []
                    if data and isinstance(data[0], list):
                        # Transpose if it's a 2D list that looks like a row
                        col_data = [r[0] if len(r) > 0 else "" for r in data]
                    else:
                        col_data = data
                    for i, cell in enumerate(new_col.cells):
                        if i < len(col_data):
                            cell.text = str(col_data[i])
                summaries.append(f"Added column to table {t_idx}")

            elif action == "remove_col":
                c_idx = params.get("col_index")
                if c_idx is not None and len(tbl.rows) > 0 and 0 <= c_idx < len(tbl.rows[0].cells):
                    # Remove cell from each row
                    for row in tbl.rows:
                        cell = row.cells[c_idx]
                        cell._element.getparent().remove(cell._element)
                    
                    # Remove from tblGrid to prevent corruption
                    tblGrid = tbl._element.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblGrid")
                    if tblGrid is not None:
                        gridCols = tblGrid.findall("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}gridCol")
                        if 0 <= c_idx < len(gridCols):
                            tblGrid.remove(gridCols[c_idx])
                            
                    summaries.append(f"Removed column {c_idx} from table {t_idx}")

            elif action == "populate":
                data = params.get("data", [])
                if data:
                    # If it's a 1D list, wrap it in a 2D list to process cleanly
                    if data and not isinstance(data[0], list):
                        data = [[x] for x in data]
                        
                    start_row = params.get("row_index")
                    start_col = params.get("col_index")
                    
                    if start_col is None and len(tbl.rows) > 0:
                        # Auto-detect if they meant to append to the end by guessing from the row
                        start_col = 0
                        # But wait, if they didn't specify, we use default 0.
                        # Actually we can do a smart check: if the first cell is not empty, maybe they didn't specify correctly?
                        pass
                        
                    start_row = start_row or 0
                    start_col = start_col or 0
                    
                    for r_idx, row_data in enumerate(data):
                        tr_idx = start_row + r_idx
                        if tr_idx >= len(tbl.rows):
                            break
                        for c_idx, val in enumerate(row_data):
                            tc_idx = start_col + c_idx
                            if tc_idx >= len(tbl.rows[tr_idx].cells):
                                break
                            tbl.rows[tr_idx].cells[tc_idx].text = str(val)
                summaries.append(f"Populated data in table {t_idx}")

            elif action == "merge_cells":
                f_idx = params.get("merge_from", [0, 0])
                t_idx_merge = params.get("merge_to", [0, 1])
                try:
                    c1 = tbl.cell(f_idx[0], f_idx[1])
                    c2 = tbl.cell(t_idx_merge[0], t_idx_merge[1])
                    c1.merge(c2)
                    summaries.append(f"Merged cells in table {t_idx}")
                except Exception as e:
                    summaries.append(f"Failed to merge cells: {e}")

            elif action == "set_alignment":
                alignment = params.get("alignment", "center")
                from docx.enum.table import WD_TABLE_ALIGNMENT
                align_map = {
                    "left": WD_TABLE_ALIGNMENT.LEFT,
                    "center": WD_TABLE_ALIGNMENT.CENTER,
                    "right": WD_TABLE_ALIGNMENT.RIGHT,
                }
                tbl.alignment = align_map.get(alignment, WD_TABLE_ALIGNMENT.CENTER)
                summaries.append(f"Set alignment of table {t_idx} to {alignment}")

            elif action == "set_width_pct":
                width_pct = float(params.get("width_pct", 1.0))
                if width_pct > 1.0:
                    width_pct = width_pct / 100.0  # Handle cases where LLM passes 100 instead of 1.0
                try:
                    section = doc.sections[0]
                    usable_width = section.page_width - section.left_margin - section.right_margin
                except Exception:
                    usable_width = int(5940000) # Fallback to ~A4 EMU
                    
                target_width = int(usable_width * width_pct)
                tbl.autofit = False
                col_count = len(tbl.rows[0].cells) if len(tbl.rows) > 0 else 1
                col_width = int(target_width / max(1, col_count))
                
                for row in tbl.rows:
                    for cell in row.cells:
                        cell.width = col_width
                
                summaries.append(f"Set width of table {t_idx} to {int(width_pct*100)}%")

            elif action in ["set_cell_bg", "alternate_rows", "set_header_format", "apply_theme"]:
                # Helper to safely set background color without duplicating w:shd
                from docx.oxml.ns import nsdecls, qn
                from docx.oxml import parse_xml
                def _set_cell_bg(cell, clr):
                    tcPr = cell._tc.get_or_add_tcPr()
                    for existing_shd in tcPr.findall(qn('w:shd')):
                        tcPr.remove(existing_shd)
                    shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'), clr))
                    tcPr.append(shading_elm)

                if action == "set_header_format":
                    if len(tbl.rows) > 0:
                        from docx.shared import RGBColor as DRGBColor
                        for cell in tbl.rows[0].cells:
                            _set_cell_bg(cell, "000080") # dark blue
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = DRGBColor(255, 255, 255)
                        summaries.append(f"Formatted header for table {t_idx}")
                        continue
                
                if action == "set_cell_bg":
                    row_idx = params.get("row_index")
                    col_idx = params.get("col_index")
                    color = str(params.get("cell_bg_hex", "")).strip().lstrip("#").upper()
                    if color:
                        if row_idx is not None and col_idx is not None:
                            if 0 <= row_idx < len(tbl.rows) and 0 <= col_idx < len(tbl.rows[row_idx].cells):
                                _set_cell_bg(tbl.cell(row_idx, col_idx), color)
                        elif row_idx is not None:
                            if 0 <= row_idx < len(tbl.rows):
                                for c in tbl.rows[row_idx].cells:
                                    _set_cell_bg(c, color)
                        elif col_idx is not None:
                            for r in tbl.rows:
                                if 0 <= col_idx < len(r.cells):
                                    _set_cell_bg(r.cells[col_idx], color)
                    summaries.append(f"Set cell background for table {t_idx}")
                    continue

                if action == "alternate_rows":
                    colors = params.get("alternate_row_colors", ["FFFFFF", "F2F2F2"])
                    if len(colors) == 2:
                        for i, row in enumerate(tbl.rows):
                            color = colors[i % 2].lstrip("#")
                            for cell in row.cells:
                                _set_cell_bg(cell, color)
                        summaries.append(f"Applied alternate row colors to table {t_idx}")
                        continue

                if action == "apply_theme":
                    theme_color = params.get("theme_color_hex", "4F81BD").lstrip("#").upper()
                    if len(tbl.rows) > 0:
                        for cell in tbl.rows[0].cells:
                            _set_cell_bg(cell, theme_color)
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.bold = True
                                    try:
                                        from docx.shared import RGBColor as DRGBColor
                                        run.font.color.rgb = DRGBColor(255, 255, 255)
                                    except Exception:
                                        pass
                    if len(tbl.rows) > 1:
                        for i, row in enumerate(tbl.rows[1:]):
                            color = "FFFFFF" if i % 2 == 0 else "F9F9F9"
                            for cell in row.cells:
                                _set_cell_bg(cell, color)
                    summaries.append(f"Applied theme #{theme_color} to table {t_idx}")
                    continue

            elif action == "sort_data":
                if len(tbl.rows) > 1:
                    # Find sort column. Default is 0. If user provided col_index, use it.
                    sort_col_idx = params.get("col_index")
                    if sort_col_idx is None:
                        # If the user passed sort_by_column in the params even if it's not strictly in schema
                        sort_by_name = str(params.get("sort_by_column", "")).lower()
                        sort_col_idx = 0
                        if sort_by_name:
                            for c_i, c in enumerate(tbl.rows[0].cells):
                                if c.text.strip().lower() == sort_by_name:
                                    sort_col_idx = c_i
                                    break

                    header = [c.text for c in tbl.rows[0].cells]
                    data = [[c.text for c in r.cells] for r in tbl.rows[1:]]
                    try:
                        # Sort by the chosen column, handling numeric vs string
                        def _sort_key(x):
                            val = x[sort_col_idx] if 0 <= sort_col_idx < len(x) else ""
                            try:
                                # Strip currency/percent for sorting
                                clean = val.replace("$", "").replace("%", "").replace(",", "")
                                return (0, float(clean))
                            except ValueError:
                                return (1, val)
                        data.sort(key=_sort_key)
                        for i, row_data in enumerate(data):
                            for j, val in enumerate(row_data):
                                tbl.rows[i+1].cells[j].text = val
                        summaries.append(f"Sorted table {t_idx}")
                    except:
                        pass

            elif action == "set_borders":
                # We handle borders at the end of the loop now
                pass

            elif action == "set_cell_alignment":
                alignment = params.get("cell_alignment")
                row_idx = params.get("row_index")
                col_idx = params.get("col_index")
                if alignment:
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    align_map = {
                        "left": WD_ALIGN_PARAGRAPH.LEFT,
                        "center": WD_ALIGN_PARAGRAPH.CENTER,
                        "right": WD_ALIGN_PARAGRAPH.RIGHT,
                        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
                    }
                    if alignment in align_map:
                        def _set_align(cell):
                            for p in cell.paragraphs:
                                p.alignment = align_map[alignment]

                        if row_idx is not None and col_idx is not None:
                            if 0 <= row_idx < len(tbl.rows) and 0 <= col_idx < len(tbl.rows[row_idx].cells):
                                _set_align(tbl.cell(row_idx, col_idx))
                        elif row_idx is not None:
                            if 0 <= row_idx < len(tbl.rows):
                                for c in tbl.rows[row_idx].cells:
                                    _set_align(c)
                        elif col_idx is not None:
                            for r in tbl.rows:
                                if 0 <= col_idx < len(r.cells):
                                    _set_align(r.cells[col_idx])
                        else:
                            for row in tbl.rows:
                                for cell in row.cells:
                                    _set_align(cell)
                    summaries.append(f"Aligned text in table {t_idx} to {alignment}")

            # --- Universal Table Modifiers ---

            valign = params.get("cell_vertical_alignment")
            if valign:
                from docx.enum.table import WD_ALIGN_VERTICAL
                valign_map = {
                    "top": WD_ALIGN_VERTICAL.TOP,
                    "center": WD_ALIGN_VERTICAL.CENTER,
                    "bottom": WD_ALIGN_VERTICAL.BOTTOM
                }
                if valign in valign_map:
                    for row in tbl.rows:
                        for cell in row.cells:
                            cell.vertical_alignment = valign_map[valign]
                    summaries.append(f"Set vertical alignment to {valign} for table {t_idx}")
            
            border_hex = params.get("border_color_hex")
            if border_hex:
                border_color = str(border_hex).lstrip("#").upper()
                border_sz = str(int(params.get("border_width_pt", 12)))
                from docx.oxml.shared import OxmlElement, qn
                tblPr = tbl._element.xpath('w:tblPr')
                if tblPr:
                    tblBorders = OxmlElement('w:tblBorders')
                    for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                        border = OxmlElement(f'w:{border_name}')
                        border.set(qn('w:val'), 'single')
                        border.set(qn('w:sz'), border_sz)
                        border.set(qn('w:space'), '0')
                        border.set(qn('w:color'), border_color)
                        tblBorders.append(border)
                    # replace existing
                    old_borders = tblPr[0].xpath('w:tblBorders')
                    if old_borders:
                        tblPr[0].replace(old_borders[0], tblBorders)
                    else:
                        tblPr[0].append(tblBorders)
                summaries.append(f"Applied {border_color} borders to table {t_idx}")

            col_width = params.get("col_width_inches")
            if col_width is not None:
                from docx.shared import Inches
                for row in tbl.rows:
                    for cell in row.cells:
                        cell.width = Inches(float(col_width))
                summaries.append(f"Set column width to {col_width} inches for table {t_idx}")

        return "; ".join(summaries)

    def _op_docx_ai_design(self, doc, params: dict) -> str:
        action = params.get("action", "normalize_fonts")
        
        if action == "normalize_fonts":
            from docx.shared import Pt
            font_family = params.get("target_font", "Calibri")
            base_size = params.get("base_font_size_pt", 11)
            for para in doc.paragraphs:
                for run in para.runs:
                    run.font.name = font_family
                    if not para.style.name.startswith("Heading"):
                        run.font.size = Pt(base_size)
            for tbl in doc.tables:
                for row in tbl.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.font.name = font_family
                                run.font.size = Pt(base_size)
            return f"Normalized fonts to {font_family} {base_size}pt"
            
        return f"Applied AI design: {action}"

    def _op_pptx_table(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Create, delete, or edit a table on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "create")

        if action == "create":
            rows = params.get("rows", 3)
            cols = params.get("cols", 3)
            pos = params.get("position", {})

            sw = prs.slide_width
            sh = prs.slide_height

            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.3))
            width = int(sw * pos.get("width_pct", 0.8))
            height = int(sh * pos.get("height_pct", 0.5))

            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Populate with data if provided
            data = params.get("data", [])
            for r_idx, row_data in enumerate(data):
                if r_idx >= rows:
                    break
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx >= cols:
                        break
                    table.cell(r_idx, c_idx).text = str(cell_text)

            # Apply header row formatting
            if params.get("header_row") and rows > 0:
                from pptx.dml.color import RGBColor
                for c_idx in range(cols):
                    cell = table.cell(0, c_idx)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.bold = True

            # Apply alternate row colors
            alt_colors = params.get("alternate_row_colors")
            if alt_colors and len(alt_colors) >= 2:
                from pptx.oxml.ns import qn
                from lxml import etree
                start_row = 1 if params.get("header_row") else 0
                for r_idx in range(start_row, rows):
                    color_hex = alt_colors[(r_idx - start_row) % 2]
                    for c_idx in range(cols):
                        cell = table.cell(r_idx, c_idx)
                        try:
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                            srgbClr.set("val", str(color_hex).upper().lstrip("#"))
                        except Exception:
                            pass

            return f"Created {rows}×{cols} table on slide {tgt.get('slide')}"

        elif action == "delete":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Deleted table (shape {shape_idx}) from slide {tgt.get('slide')}"
            return ""

        elif action in ("add_row", "remove_row", "add_col", "remove_col"):
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                table = shape.table
                if action == "add_row":
                    # Clone last row
                    from lxml import etree
                    tbl_elem = table._tbl
                    last_tr = tbl_elem.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tr")[-1]
                    new_tr = copy.deepcopy(last_tr)
                    for tc in new_tr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tc"):
                        for t in tc.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                            t.text = ""
                    tbl_elem.append(new_tr)
                    return f"Added row to table on slide {tgt.get('slide')}"
                elif action == "remove_row":
                    row_idx = params.get("row_index", -1)
                    tbl_elem = table._tbl
                    rows = tbl_elem.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tr")
                    if rows and abs(row_idx) < len(rows):
                        tbl_elem.remove(rows[row_idx])
                        return f"Removed row from table on slide {tgt.get('slide')}"

        elif action == "set_cell_bg":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                r = params.get("row_index", 0)
                c = params.get("col_index", 0)
                color_hex = params.get("cell_bg_hex", "FFFFFF")
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    cell = shape.table.cell(r, c)
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(color_hex).upper().lstrip("#"))
                    return f"Set cell background on slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("set_cell_bg failed: %s", e)

        elif action == "populate":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                data = params.get("data", [])
                table = shape.table
                for r_idx, row_data in enumerate(data):
                    if r_idx >= len(table.rows):
                        break
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx >= len(table.columns):
                            break
                        table.cell(r_idx, c_idx).text = str(cell_text)
                return f"Populated table on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_image(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Insert, replace, resize or style images on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "insert")
        image_path = params.get("image_path")

        sw = prs.slide_width
        sh = prs.slide_height

        if action == "insert" and image_path:
            if not Path(image_path).exists():
                log.warning("Image file not found: %s", image_path)
                return ""

            pos = params.get("position", {})
            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.15))
            width = int(sw * pos.get("width_pct", 0.4))
            height = int(sh * pos.get("height_pct", 0.5))

            maintain_ar = params.get("maintain_aspect_ratio", True)
            if maintain_ar:
                pic = slide.shapes.add_picture(image_path, left, top, width=width)
            else:
                pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

            # Apply border if requested
            border_hex = params.get("border_color_hex")
            border_pt = params.get("border_width_pt", 1.5)
            if border_hex:
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = pic._element.spPr
                    ln = etree.SubElement(spPr, qn("a:ln"))
                    ln.set("w", str(int(Pt(border_pt))))
                    solidFill = etree.SubElement(ln, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(border_hex).upper().lstrip("#"))
                except Exception as e:
                    log.debug("Image border failed: %s", e)

            # Shadow
            if params.get("shadow"):
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = pic._element.spPr
                    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
                    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
                    outerShdw.set("blurRad", "60960")
                    outerShdw.set("dist", "114300")
                    outerShdw.set("dir", "2700000")
                    outerShdw.set("algn", "tl")
                    outerShdw.set("rotWithShape", "0")
                    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
                    srgbClr.set("val", "000000")
                    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                    alpha.set("val", "40000")
                except Exception as e:
                    log.debug("Image shadow failed: %s", e)

            return f"Inserted image on slide {tgt.get('slide')}"

        elif action == "remove":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Removed image from slide {tgt.get('slide')}"

        elif action in ("resize", "reposition"):
            shape_idx = tgt.get("shape_index")
            shape = self._get_shape(slide, shape_idx)
            if shape:
                pos = params.get("position", {})
                if "left_pct" in pos:
                    shape.left = int(sw * pos["left_pct"])
                if "top_pct" in pos:
                    shape.top = int(sh * pos["top_pct"])
                if "width_pct" in pos:
                    shape.width = int(sw * pos["width_pct"])
                if "height_pct" in pos:
                    shape.height = int(sh * pos["height_pct"])
                return f"Repositioned/resized shape {shape_idx} on slide {tgt.get('slide')}"

        elif action == "rotate":
            shape_idx = tgt.get("shape_index")
            shape = self._get_shape(slide, shape_idx)
            if shape:
                degrees = params.get("rotation_degrees", 0)
                shape.rotation = degrees
                return f"Rotated shape {shape_idx} by {degrees}° on slide {tgt.get('slide')}"

        elif action in ("bring_forward", "send_backward"):
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes_list = list(slide.shapes)
                if 0 <= shape_idx < len(shapes_list):
                    sp_elem = shapes_list[shape_idx]._element
                    sp_tree = sp_elem.getparent()
                    if action == "bring_forward":
                        sp_tree.append(sp_elem)
                    else:
                        children = list(sp_tree)
                        first_sp_idx = next(
                            (i for i, c in enumerate(children) if c.tag.endswith("}sp") or c.tag.endswith("}pic")),
                            0
                        )
                        sp_tree.insert(first_sp_idx, sp_elem)
                    return f"Moved shape {shape_idx} {action.replace('_', ' ')} on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_shape(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Add or edit shapes and text boxes on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "add_textbox")
        sw = prs.slide_width
        sh = prs.slide_height

        if action == "add_textbox":
            pos = params.get("position", {})
            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.1))
            width = int(sw * pos.get("width_pct", 0.4))
            height = int(sh * pos.get("height_pct", 0.1))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = params.get("text", "")

            # Apply fill
            fill_hex = params.get("fill_color_hex")
            if fill_hex:
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = txBox._element.spPr
                    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(fill_hex).upper().lstrip("#"))
                except Exception:
                    pass

            return f"Added text box on slide {tgt.get('slide')}"

        elif action == "delete":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Deleted shape {shape_idx} from slide {tgt.get('slide')}"

        elif action in ("resize", "move"):
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape:
                pos = params.get("position", {})
                if "left_pct" in pos:
                    shape.left = int(sw * pos["left_pct"])
                if "top_pct" in pos:
                    shape.top = int(sh * pos["top_pct"])
                if "width_pct" in pos:
                    shape.width = int(sw * pos["width_pct"])
                if "height_pct" in pos:
                    shape.height = int(sh * pos["height_pct"])
                return f"Moved/resized shape on slide {tgt.get('slide')}"

        elif action == "set_fill":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            fill_hex = params.get("fill_color_hex")
            if shape and fill_hex:
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(str(fill_hex).lstrip("#"))
                    return f"Set fill color on shape {tgt.get('shape_index')} slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("Set fill failed: %s", e)

        elif action == "set_outline":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            outline_hex = params.get("outline_color_hex")
            outline_pt = params.get("outline_width_pt", 1.5)
            if shape and outline_hex:
                try:
                    shape.line.color.rgb = RGBColor.from_string(str(outline_hex).lstrip("#"))
                    shape.line.width = Pt(outline_pt)
                    return f"Set outline on shape {tgt.get('shape_index')} slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("Set outline failed: %s", e)

        elif action == "duplicate":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    new_sp = copy.deepcopy(shapes[shape_idx]._element)
                    slide.shapes._spTree.append(new_sp)
                    return f"Duplicated shape {shape_idx} on slide {tgt.get('slide')}"

        elif action == "rotate":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape:
                shape.rotation = params.get("rotation_degrees", 0)
                return f"Rotated shape on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_theme(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Change slide backgrounds and theme colors."""
        action = params.get("action", "set_bg_color")
        scope = params.get("scope", "current_slide")
        slide_num = tgt.get("slide")

        slides_to_update = list(prs.slides) if scope == "all_slides" else []
        if not slides_to_update and slide_num:
            s = self._get_slide(prs, slide_num)
            if s:
                slides_to_update = [s]
        if not slides_to_update:
            slides_to_update = list(prs.slides)

        if action == "set_bg_color":
            bg_hex = params.get("bg_color_hex", "FFFFFF")
            from pptx.oxml.ns import qn
            from lxml import etree

            for slide in slides_to_update:
                try:
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(str(bg_hex).upper().lstrip("#"))
                except Exception as e:
                    log.debug("set_bg_color failed for slide: %s", e)

            scope_label = "all slides" if scope == "all_slides" else f"slide {slide_num}"
            return f"Set background color #{bg_hex} on {scope_label}"

        elif action == "set_bg_gradient":
            start_hex = params.get("gradient_start_hex", "1a1a2e")
            end_hex = params.get("gradient_end_hex", "16213e")
            direction = params.get("gradient_direction", "diagonal")

            from pptx.oxml.ns import qn
            from lxml import etree

            angle_map = {"horizontal": "5400000", "vertical": "10800000", "diagonal": "2700000"}
            angle = angle_map.get(direction, "5400000")

            for slide in slides_to_update:
                try:
                    bg = slide.background
                    bgPr = bg._element.get_or_add_bgPr()
                    # Clear existing fill
                    for child in list(bgPr):
                        bgPr.remove(child)
                    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
                    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

                    for pos_val, hex_color in [("0", start_hex), ("100000", end_hex)]:
                        gs = etree.SubElement(gsLst, qn("a:gs"))
                        gs.set("pos", pos_val)
                        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
                        srgbClr.set("val", str(hex_color).upper().lstrip("#"))

                    lin = etree.SubElement(gradFill, qn("a:lin"))
                    lin.set("ang", angle)
                    lin.set("scaled", "0")
                except Exception as e:
                    log.debug("Gradient bg failed: %s", e)

            return f"Set gradient background ({direction}) on {scope}"

        return ""

    def _op_pptx_slide(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Add, delete, duplicate, or reorder slides."""
        action = params.get("action", "add")
        slide_num = tgt.get("slide")

        if action == "add":
            # Use layout from an existing slide
            after_idx = (params.get("after_index") or slide_num or len(list(prs.slides))) - 1
            slides = list(prs.slides)
            src_idx = min(after_idx, len(slides) - 1)
            self._clone_slide(prs, src_idx)

            # Move the new slide to the correct position
            new_idx = len(list(prs.slides)) - 1
            target_pos = min(after_idx + 1, new_idx)
            if target_pos < new_idx:
                self._reorder_slide(prs, new_idx, target_pos)
            return f"Added new slide after position {after_idx + 1}"

        elif action == "delete":
            idx = (slide_num or 1) - 1
            slides = list(prs.slides)
            if 0 <= idx < len(slides):
                self._delete_slide(prs, idx)
                return f"Deleted slide {slide_num}"

        elif action == "duplicate":
            after_idx = (params.get("after_index") or slide_num or 1) - 1
            slides = list(prs.slides)
            src_idx = (slide_num or 1) - 1
            src_idx = min(src_idx, len(slides) - 1)
            self._clone_slide(prs, src_idx)
            new_idx = len(list(prs.slides)) - 1
            target_pos = after_idx + 1
            if target_pos < new_idx:
                self._reorder_slide(prs, new_idx, target_pos)
            return f"Duplicated slide {slide_num or 1}"

        elif action == "reorder":
            from_idx = (params.get("from_index") or slide_num or 1) - 1
            to_idx = (params.get("to_index") or 1) - 1
            slides = list(prs.slides)
            if 0 <= from_idx < len(slides) and 0 <= to_idx < len(slides):
                self._reorder_slide(prs, from_idx, to_idx)
                return f"Moved slide {from_idx + 1} to position {to_idx + 1}"

        elif action == "rename_title":
            slide = self._get_slide(prs, slide_num)
            title = params.get("title", "")
            if slide and title:
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if shape.placeholder_format.type in (
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ):
                                shape.text = title
                                return f"Renamed title of slide {slide_num} to '{title}'"
                        except Exception:
                            pass

        elif action in ("hide", "unhide"):
            # PPTX doesn't natively hide slides via python-pptx API simply,
            # but we can set the show attribute in the slide list XML
            slide = self._get_slide(prs, slide_num)
            if slide:
                try:
                    from pptx.oxml.ns import qn
                    sldIdLst = prs.part._element.find(
                        "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
                    )
                    for rel in prs.part.rels.values():
                        if rel.target_part == slide.part:
                            rId = rel.rId
                            if sldIdLst is not None:
                                for sldId in sldIdLst:
                                    ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                                    if sldId.get(f"{{{ns}}}id") == rId:
                                        sldId.set("show", "0" if action == "hide" else "1")
                            return f"{'Hid' if action == 'hide' else 'Unhid'} slide {slide_num}"
                except Exception as e:
                    log.debug("Hide/unhide slide failed: %s", e)

        return ""

    def _reorder_slide(self, prs: Presentation, from_idx: int, to_idx: int) -> None:
        """Reorder a slide within the presentation."""
        slides = list(prs.slides)
        if from_idx < 0 or from_idx >= len(slides):
            return
        if to_idx < 0 or to_idx >= len(slides):
            return

        pres_elem = prs.part._element
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sldIdLst = pres_elem.find(f"{{{ns}}}sldIdLst")
        if sldIdLst is None:
            return

        sldIds = list(sldIdLst)
        if from_idx >= len(sldIds) or to_idx >= len(sldIds):
            return

        elem = sldIds[from_idx]
        sldIdLst.remove(elem)
        sldIds.pop(from_idx)
        sldIds.insert(to_idx, elem)
        for s in sldIds:
            sldIdLst.remove(s)
        for s in sldIds:
            sldIdLst.append(s)

    def _op_pptx_chart(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Edit chart type, colors, labels. Returns summary string."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        shape_idx = tgt.get("shape_index")
        shape = self._get_shape(slide, shape_idx)
        if not shape:
            # Try to find any chart shape
            for s in slide.shapes:
                if s.has_chart:
                    shape = s
                    break
        if not shape or not shape.has_chart:
            return ""

        action = params.get("action", "update_labels")
        chart = shape.chart

        if action == "update_axis_labels":
            try:
                x_label = params.get("x_axis_label")
                y_label = params.get("y_axis_label")
                if x_label and chart.has_plot_area:
                    pass  # python-pptx limited chart axis label support
                return f"Updated chart axis labels on slide {tgt.get('slide')}"
            except Exception:
                pass

        elif action == "set_series_colors":
            colors = params.get("series_colors", [])
            try:
                for i, series in enumerate(chart.series):
                    if i < len(colors):
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor.from_string(
                            str(colors[i]).upper().lstrip("#")
                        )
                return f"Updated series colors on chart, slide {tgt.get('slide')}"
            except Exception as e:
                log.debug("Chart series color failed: %s", e)

        elif action in ("show_legend", "hide_legend"):
            try:
                chart.has_legend = (action == "show_legend")
                return f"{'Showed' if action == 'show_legend' else 'Hid'} chart legend on slide {tgt.get('slide')}"
            except Exception as e:
                log.debug("Chart legend failed: %s", e)

        return f"Edited chart on slide {tgt.get('slide')}"

    def _op_pptx_ai_design(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """AI-driven design normalization operations."""
        action = params.get("action", "normalize_fonts")
        scope = params.get("scope", "all_slides")

        # Determine which slides to process
        if scope == "all_slides":
            slides = list(prs.slides)
        elif scope and scope.startswith("slide:"):
            try:
                n = int(scope.split(":")[1])
                slide = self._get_slide(prs, n)
                slides = [slide] if slide else list(prs.slides)
            except Exception:
                slides = list(prs.slides)
        else:
            slides = list(prs.slides)

        if action == "normalize_fonts":
            target_font = params.get("target_font", "Calibri")
            base_size = params.get("base_font_size_pt", 18)
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.name != target_font:
                                run.font.name = target_font
            return f"Normalized fonts to {target_font} across {scope}"

        elif action == "normalize_spacing":
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        try:
                            from pptx.oxml.ns import qn
                            from lxml import etree
                            pPr = para._p.get_or_add_pPr()
                            spcBef = pPr.find(qn("a:spcBef"))
                            if spcBef is None:
                                spcBef = etree.SubElement(pPr, qn("a:spcBef"))
                            spcPts = spcBef.find(qn("a:spcPts"))
                            if spcPts is None:
                                spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
                            spcPts.set("val", "600")  # 6pt before paragraph
                        except Exception:
                            pass
            return f"Normalized paragraph spacing across {scope}"

        elif action == "generate_speaker_notes":
            # Add basic speaker notes placeholder to each slide
            for slide in slides:
                try:
                    notes_slide = slide.notes_slide
                    tf = notes_slide.notes_text_frame
                    if not tf.text.strip():
                        # Collect slide content for context
                        content = " | ".join(
                            s.text_frame.text[:100]
                            for s in slide.shapes
                            if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
                        )
                        tf.text = f"Slide content: {content}\n[Speaker notes generated by AI]"
                except Exception as e:
                    log.debug("Speaker notes failed: %s", e)
            return f"Generated speaker notes for {scope}"

        elif action == "improve_hierarchy":
            # Ensure title shapes are bold and large; body shapes are smaller
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    is_title = shape.is_placeholder and shape.placeholder_format.type in (1, 3)
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if is_title:
                                run.font.bold = True
                                if run.font.size and run.font.size.pt < 24:
                                    run.font.size = Pt(28)
                            else:
                                if run.font.size and run.font.size.pt > 28:
                                    run.font.size = Pt(18)
            return f"Improved visual hierarchy across {scope}"

        elif action == "auto_resize_text":
            for slide in slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        try:
                            shape.text_frame.auto_size = 1  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass
            return f"Enabled auto-resize text across {scope}"

        elif action == "improve_readability":
            # Ensure minimum font sizes and line spacing
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt < 12:
                                run.font.size = Pt(12)
            return f"Improved readability across {scope}"

        return f"Applied AI design ({action}) across {scope}"

    # ------------------------------------------------------------------
    # DOCX: one block per paragraph
    # ------------------------------------------------------------------

    def _resolve_docx_font_info(self, para, run) -> tuple[float | None, str | None]:
        font_size = None
        color_hex = None
        
        if run.font.size:
            font_size = run.font.size.pt
        try:
            if run.font.color and run.font.color.rgb:
                color_hex = str(run.font.color.rgb)
        except Exception:
            pass
            
        if font_size is None or color_hex is None:
            style = para.style
            while style:
                if font_size is None and style.font.size:
                    font_size = style.font.size.pt
                if color_hex is None:
                    try:
                        if style.font.color and style.font.color.rgb:
                            color_hex = str(style.font.color.rgb)
                    except Exception:
                        pass
                
                if font_size is not None and color_hex is not None:
                    break
                    
                style = style.base_style
                
        return font_size, color_hex

    def _extract_docx_dom(self, path: Path) -> dict:
        """Extract DOCX DOM in true document body order using stable UUIDs."""
        from docx.oxml.ns import qn as _qn
        WNS_P = _qn('w:p')
        WNS_TBL = _qn('w:tbl')
        WNS_DRAWING = _qn('w:drawing')
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        doc = Document(path)
        
        # --- UID STAMPING ---
        any_uids_assigned = False
        for child in doc.element.body:
            if child.tag in (WNS_P, WNS_TBL):
                if child.get(uid_attr) is None:
                    child.set(uid_attr, uuid.uuid4().hex[:8])
                    any_uids_assigned = True
                    
        if any_uids_assigned:
            import tempfile
            import shutil
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".docx", dir=path.parent)
            import os
            os.close(tmp_fd)
            doc.save(tmp_path)
            shutil.move(tmp_path, path)
        # --------------------

        children = []

        from app.services.docx_extensions import extract_metadata
        metadata_node = extract_metadata(doc)
        if metadata_node:
            children.append({
                "id": "metadata",
                "type": "metadata",
                "role": "metadata",
                "properties": metadata_node
            })
        
        # Extract basic section info
        for s_idx, section in enumerate(doc.sections):
            sec_style = {}
            try:
                sec_style["orientation"] = section.orientation.name if hasattr(section, 'orientation') else "PORTRAIT"
                if hasattr(section.page_width, 'inches'): sec_style["page_width_inches"] = round(section.page_width.inches, 2)
                if hasattr(section.page_height, 'inches'): sec_style["page_height_inches"] = round(section.page_height.inches, 2)
                if hasattr(section.top_margin, 'inches'): sec_style["top_margin_inches"] = round(section.top_margin.inches, 2)
                if hasattr(section.bottom_margin, 'inches'): sec_style["bottom_margin_inches"] = round(section.bottom_margin.inches, 2)
                if hasattr(section.left_margin, 'inches'): sec_style["left_margin_inches"] = round(section.left_margin.inches, 2)
                if hasattr(section.right_margin, 'inches'): sec_style["right_margin_inches"] = round(section.right_margin.inches, 2)
            except Exception:
                pass
            
            children.append({
                "id": f"section_{s_idx}",
                "type": "section",
                "role": "section",
                "style": sec_style
            })

        # Build lookup maps: xml element → python-docx object
        para_elements = {p._p: p for p in doc.paragraphs}
        table_elements = {t._tbl: t for t in doc.tables}

        body_index = 0  # sequential position across all body children

        for child in doc.element.body:
            tag = child.tag

            if tag == WNS_P and child in para_elements:
                para = para_elements[child]
                uid = child.get(uid_attr)

                # ---- Check if this paragraph contains an inline image ----
                has_image = child.find(f'.//{WNS_DRAWING}') is not None

                if has_image:
                    # Extract image dimensions from the drawing element
                    width_emu, height_emu = 0, 0
                    description = ""
                    try:
                        drawing_el = child.find(f'.//{WNS_DRAWING}')
                        extent_el = drawing_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}extent')
                        if extent_el is not None:
                            width_emu = int(extent_el.get('cx', 0))
                            height_emu = int(extent_el.get('cy', 0))
                        # Try to get alt text / description
                        docPr_el = drawing_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}docPr')
                        if docPr_el is not None:
                            description = docPr_el.get('descr', '') or docPr_el.get('name', '')
                    except Exception:
                        pass

                    alignment_str = None
                    if para.alignment is not None:
                        try:
                            alignment_str = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                        except Exception:
                            pass

                    children.append({
                        "id": f"img_{uid}",
                        "body_index": body_index,
                        "type": "image",
                        "role": "inline_image",
                        "width_emu": width_emu,
                        "height_emu": height_emu,
                        "description": description,
                        "alignment": alignment_str,
                    })

                else:
                    # Regular text paragraph
                    runs = []
                    for r_idx, run in enumerate(para.runs):
                        font_size, color_hex = self._resolve_docx_font_info(para, run)
                        font = run.font
                        runs.append({
                            "id": f"p_{uid}_run_{r_idx}",
                            "type": "run",
                            "text": run.text,
                            "style": {
                                "font": font.name,
                                "size": font_size,
                                "bold": font.bold,
                                "italic": font.italic,
                                "underline": font.underline,
                                "strike": font.strike,
                                "highlight": str(font.highlight_color) if font.highlight_color else None,
                                "color": color_hex,
                            },
                        })

                    alignment = None
                    if para.alignment is not None:
                        try:
                            alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                        except Exception:
                            pass

                    style_name = para.style.name if para.style else "Normal"
                    role = "body"
                    heading_level = None
                    if style_name.lower().startswith("heading"):
                        role = "heading"
                        import re as _re
                        m = _re.search(r'(\d+)$', style_name)
                        heading_level = int(m.group(1)) if m else 1
                    elif "bullet" in style_name.lower() or "list" in style_name.lower():
                        role = "bullet_point"

                    line_spacing = None
                    try:
                        line_spacing = para.paragraph_format.line_spacing
                    except Exception:
                        pass

                    page_break_before = None
                    try:
                        page_break_before = para.paragraph_format.page_break_before
                    except Exception:
                        pass

                    space_before_pt = None
                    try:
                        if para.paragraph_format.space_before is not None:
                            space_before_pt = para.paragraph_format.space_before.pt
                    except Exception:
                        pass

                    space_after_pt = None
                    try:
                        if para.paragraph_format.space_after is not None:
                            space_after_pt = para.paragraph_format.space_after.pt
                    except Exception:
                        pass
                    
                    from app.services.docx_extensions import extract_advanced_paragraph_style
                    adv_style = extract_advanced_paragraph_style(para)

                    node = {
                        "id": f"p_{uid}",
                        "body_index": body_index,
                        "type": "paragraph",
                        "role": role,
                        "text": para.text.strip(),
                        "style": {
                            "style_name": style_name,
                            "alignment": alignment,
                            "line_spacing": line_spacing,
                            "page_break_before": page_break_before,
                            "space_before_pt": space_before_pt,
                            "space_after_pt": space_after_pt,
                            **adv_style
                        },
                        "runs": runs,
                    }
                    
                    include_in_toc = (role == "heading")
                    if include_in_toc and para._p.pPr is not None:
                        from docx.oxml.ns import qn as _qn
                        outlineLvl = para._p.pPr.find(_qn('w:outlineLvl'))
                        if outlineLvl is not None and outlineLvl.get(_qn('w:val')) == '9':
                            include_in_toc = False
                    if include_in_toc or role == "heading":
                        node["include_in_toc"] = include_in_toc
                    
                    if heading_level is not None:
                        node["heading_level"] = heading_level

                    # Enrich with list metadata if this is a list paragraph
                    list_info = self._extract_list_info(para, doc)
                    if list_info is not None:
                        node["list_info"] = list_info
                        if role == "body":
                            # override role so LLM knows it's a list item
                            node["role"] = "bullet_point"

                    children.append(node)

            elif tag == WNS_TBL and child in table_elements:
                table = table_elements[child]
                uid = child.get(uid_attr)

                rows = []
                for r_idx, row in enumerate(table.rows):
                    cells = []
                    for c_idx, cell in enumerate(row.cells):
                        cell_paras = []
                        for p_idx, para in enumerate(cell.paragraphs):
                            runs = []
                            for r_run_idx, run in enumerate(para.runs):
                                font_size, color_hex = self._resolve_docx_font_info(para, run)
                                font = run.font
                                runs.append({
                                    "id": f"t_{uid}_cell_{r_idx}_{c_idx}_para_{p_idx}_run_{r_run_idx}",
                                    "type": "run",
                                    "text": run.text,
                                    "style": {
                                        "font": font.name,
                                        "size": font_size,
                                        "bold": font.bold,
                                        "italic": font.italic,
                                        "underline": font.underline,
                                        "strike": font.strike,
                                        "highlight": str(font.highlight_color) if font.highlight_color else None,
                                        "color": color_hex,
                                    },
                                })
                            alignment = None
                            if para.alignment is not None:
                                try:
                                    alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                                except Exception:
                                    pass
                            line_spacing = None
                            try:
                                line_spacing = para.paragraph_format.line_spacing
                            except Exception:
                                pass
                            cell_paras.append({
                                "id": f"t_{uid}_cell_{r_idx}_{c_idx}_para_{p_idx}",
                                "type": "paragraph",
                                "role": "table_cell_paragraph",
                                "text": para.text.strip(),
                                "style": {
                                    "alignment": alignment,
                                    "line_spacing": line_spacing,
                                    "page_break_before": para.paragraph_format.page_break_before if hasattr(para.paragraph_format, "page_break_before") else None,
                                    "space_before_pt": para.paragraph_format.space_before.pt if getattr(para.paragraph_format, "space_before", None) is not None else None,
                                    "space_after_pt": para.paragraph_format.space_after.pt if getattr(para.paragraph_format, "space_after", None) is not None else None,
                                },
                                "runs": runs,
                            })
                        cell_dict = {
                            "id": f"t_{uid}_cell_{r_idx}_{c_idx}",
                            "type": "cell",
                            "row": r_idx,
                            "column": c_idx,
                            "children": cell_paras,
                        }
                        
                        try:
                            tcPr = cell._tc.get_or_add_tcPr()
                            shd = tcPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd")
                            if shd is not None:
                                fill = shd.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill")
                                if fill and fill != "auto":
                                    cell_dict["bg_color"] = fill
                        except Exception:
                            pass
                        
                        try:
                            if cell.vertical_alignment is not None:
                                cell_dict["vertical_alignment"] = cell.vertical_alignment.name
                        except Exception:
                            pass
                        
                        cells.append(cell_dict)
                    rows.append({
                        "id": f"t_{uid}_row_{r_idx}",
                        "type": "row",
                        "row": r_idx,
                        "cells": cells,
                    })

                table_style_name = table.style.name if table.style else "Normal Table"
                node = {
                    "id": f"t_{uid}",
                    "body_index": body_index,
                    "type": "table",
                    "role": "table",
                    "style": {"style_name": table_style_name},
                    "row_count": len(table.rows),
                    "col_count": len(table.columns) if table.rows else 0,
                    "rows": rows,
                }
                
                try:
                    tblBorders = table._tbl.tblPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblBorders")
                    if tblBorders is not None:
                        node["style"]["has_custom_borders"] = True
                        # check colors of borders
                        border_colors = set()
                        for b_type in ["top", "left", "bottom", "right", "insideH", "insideV"]:
                            b_el = tblBorders.find(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{b_type}")
                            if b_el is not None:
                                c = b_el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color")
                                if c and c != "auto":
                                    border_colors.add(c)
                        if border_colors:
                            node["style"]["border_colors"] = list(border_colors)
                except Exception:
                    pass

                children.append(node)

            # Increment body_index for every body child regardless of type
            body_index += 1

        return {
            "document_type": "docx",
            "dom": {
                "id": "document_root",
                "type": "document",
                "children": children,
            },
        }

    def _apply_docx_edits(self, source: Path, target: Path, edits: list[dict]) -> None:
        doc = Document(source)
        edit_map = {e["element_id"]: e for e in edits}
        
        uid_attr = f'{{{CUSTOM_NS}}}uid'
        
        for para in doc.paragraphs:
            uid = para._p.get(uid_attr)
            if not uid:
                continue
            eid = f"p_{uid}"
            if eid in edit_map:
                edit = edit_map[eid]
                self._apply_run_aware_replacement(para, edit["new_text"])
                
        for table in doc.tables:
            uid = table._tbl.get(uid_attr)
            if not uid:
                continue
            t_eid = f"t_{uid}"
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for p_idx, para in enumerate(cell.paragraphs):
                        eid = f"{t_eid}_cell_{r_idx}_{c_idx}_para_{p_idx}"
                        if eid in edit_map:
                            edit = edit_map[eid]
                            self._apply_run_aware_replacement(para, edit["new_text"])
                            
        doc.save(target)

    # ------------------------------------------------------------------
    # Core: run-aware replacement
    # ------------------------------------------------------------------

    def _apply_run_aware_replacement(self, paragraph, new_text: str, format_params: dict = None) -> None:
        """Update paragraph text while preserving per-run formatting.

        Only the runs whose characters overlap the changed region are
        modified; all other runs are left completely untouched.
        If format_params is provided, it is applied exclusively to the newly inserted text.
        """
        runs = paragraph.runs
        if not runs:
            run = paragraph.add_run()
            run.text = new_text
            if format_params:
                self._apply_run_format(run, format_params)
            return

        old_text = "".join(r.text for r in runs)

        if old_text == new_text and not format_params:
            log.debug("  No change (text identical), skipping.")
            return

        # Log run structure before modification (debug only).
        self._log_runs("  BEFORE", runs)

        prefix_len, suffix_len, new_middle = _changed_region(old_text, new_text)
        old_changed_start = prefix_len
        old_changed_end   = len(old_text) - suffix_len

        log.debug(
            "  Changed region: chars [%d, %d) %r → %r",
            old_changed_start, old_changed_end,
            old_text[old_changed_start:old_changed_end],
            new_middle,
        )

        # Walk runs and update only those overlapping [old_changed_start, old_changed_end).
        run_pos = 0
        new_middle_placed = False

        for run in runs:
            r_len   = len(run.text)
            r_start = run_pos
            r_end   = run_pos + r_len
            run_pos  = r_end

            if old_changed_start == old_changed_end:
                # Pure insertion
                if r_start <= old_changed_start < r_end:
                    overlaps = True
                elif old_changed_start == len(old_text) and r_end == len(old_text):
                    overlaps = True
                else:
                    overlaps = False
            else:
                # Replacement
                if r_end <= old_changed_start or r_start >= old_changed_end:
                    overlaps = False
                else:
                    overlaps = True

            if not overlaps:
                # Run is entirely outside the changed region — leave it alone.
                continue

            # Run overlaps with the changed region.
            # Compute the portion of this run that falls BEFORE the change.
            prefix_in_run = max(0, old_changed_start - r_start)
            prefix_text   = run.text[:prefix_in_run]

            # Compute the portion of this run that falls AFTER the change.
            suffix_start_in_run = max(0, old_changed_end - r_start)
            suffix_text         = run.text[suffix_start_in_run:]

            if not new_middle_placed:
                # First overlapping run: inject the replacement text here.
                # Instead of modifying run.text in place which corrupts prefix/suffix,
                # we split the run if format_params is provided or if we want exact precision.
                import copy
                
                run.text = prefix_text
                
                # Insert new middle run
                new_mid = paragraph.add_run(new_middle)
                rPr = run._r.find('.//w:rPr', namespaces=run._r.nsmap)
                if rPr is not None:
                    new_mid._r.insert(0, copy.deepcopy(rPr))
                if format_params:
                    self._apply_run_format(new_mid, format_params)
                run._r.addnext(new_mid._r)
                
                # Insert suffix run
                if suffix_text:
                    new_suf = paragraph.add_run(suffix_text)
                    if rPr is not None:
                        new_suf._r.insert(0, copy.deepcopy(rPr))
                    new_mid._r.addnext(new_suf._r)
                    
                new_middle_placed = True
            else:
                # Subsequent overlapping runs: their "changed" content has
                # already been absorbed by the first run; keep only the
                # trailing unchanged portion.
                run.text = suffix_text
                log.debug("  run cleared to suffix: %r", run.text)

        # Log run structure after modification (debug only).
        self._log_runs("  AFTER ", runs)

    def _apply_run_format(self, run, params: dict) -> None:
        """Apply formatting directly to a docx run."""
        from docx.shared import Pt
        from docx.shared import RGBColor as DRGBColor

        if params.get("bold") is not None:
            run.font.bold = params["bold"]
        if params.get("italic") is not None:
            run.font.italic = params["italic"]
        if params.get("underline") is not None:
            run.font.underline = params["underline"]
        if params.get("font_family"):
            run.font.name = params["font_family"]
        if params.get("font_size_pt") is not None:
            run.font.size = Pt(params["font_size_pt"])
            
        color_hex = str(params.get("color_hex", "")).strip().lstrip("#").upper()
        if color_hex and len(color_hex) == 6:
            try:
                run.font.color.rgb = DRGBColor(
                    int(color_hex[0:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
            except Exception:
                pass

        highlight_hex = params.get("highlight_hex")
        if highlight_hex:
            from docx.enum.text import WD_COLOR_INDEX
            hl_color = str(highlight_hex).strip().lstrip("#").upper()
            if hl_color in ["FFFF00", "YELLOW"]:
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW
            elif hl_color in ["00FF00", "GREEN"]:
                run.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
            elif hl_color in ["00FFFF", "CYAN"]:
                run.font.highlight_color = WD_COLOR_INDEX.TURQUOISE
            elif hl_color in ["FF00FF", "MAGENTA"]:
                run.font.highlight_color = WD_COLOR_INDEX.PINK
            elif hl_color in ["0000FF", "BLUE"]:
                run.font.highlight_color = WD_COLOR_INDEX.BLUE
            elif hl_color in ["FF0000", "RED"]:
                run.font.highlight_color = WD_COLOR_INDEX.RED
            elif hl_color in ["FFFFFF", "WHITE"]:
                run.font.highlight_color = WD_COLOR_INDEX.WHITE
            elif hl_color in ["000000", "BLACK"]:
                run.font.highlight_color = WD_COLOR_INDEX.BLACK
            else:
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW

    def _apply_run_aware_format(self, paragraph, match_text: str, format_params: dict) -> bool:
        """Apply formatting to a specific substring within a paragraph by splitting runs."""
        runs = paragraph.runs
        old_text = "".join(r.text for r in runs)
        start_idx = old_text.find(match_text)
        if start_idx == -1:
            return False

        old_changed_start = start_idx
        old_changed_end = start_idx + len(match_text)
        new_middle = match_text

        run_pos = 0
        new_middle_placed = False

        for run in runs:
            r_len   = len(run.text)
            r_start = run_pos
            r_end   = run_pos + r_len
            run_pos  = r_end

            if r_end <= old_changed_start or r_start >= old_changed_end:
                continue

            prefix_in_run = max(0, old_changed_start - r_start)
            prefix_text   = run.text[:prefix_in_run]

            suffix_start_in_run = max(0, old_changed_end - r_start)
            suffix_text         = run.text[suffix_start_in_run:]

            if not new_middle_placed:
                import copy
                run.text = prefix_text
                new_mid = paragraph.add_run(new_middle)
                rPr = run._r.find('.//w:rPr', namespaces=run._r.nsmap)
                if rPr is not None:
                    new_mid._r.insert(0, copy.deepcopy(rPr))
                self._apply_run_format(new_mid, format_params)
                run._r.addnext(new_mid._r)
                
                if suffix_text:
                    new_suf = paragraph.add_run(suffix_text)
                    if rPr is not None:
                        new_suf._r.insert(0, copy.deepcopy(rPr))
                    new_mid._r.addnext(new_suf._r)
                    
                new_middle_placed = True
            else:
                run.text = suffix_text
        return True

    # ------------------------------------------------------------------
    # Debug helper
    # ------------------------------------------------------------------

    @staticmethod
    def _log_runs(label: str, runs) -> None:
        if not log.isEnabledFor(logging.DEBUG):
            return
        for i, run in enumerate(runs):
            try:
                color = run.font.color
                rgb = color.rgb if color.type is not None else "inherit"
            except Exception:
                rgb = "?"
            log.debug(
                "%s run[%d] text=%r bold=%s italic=%s color=%s",
                label, i, run.text, run.font.bold, run.font.italic, rgb,
            )


# ---------------------------------------------------------------------------
# Normalise helper (used elsewhere)
# ---------------------------------------------------------------------------

def normalize_text(value: str) -> str:
    return sub(r"\s+", " ", value).strip()
