from __future__ import annotations

# ===== BEGIN core/config.py =====
from pathlib import Path

from pydantic_settings import BaseSettings, SettingsConfigDict


class Settings(BaseSettings):
    model_config = SettingsConfigDict(env_file=".env", extra="ignore")

    database_url: str = "sqlite:///./document_agent.db"
    storage_root: Path = Path("../storage/workspaces")
    frontend_origin: str = "http://localhost:5173"
    converter_url: str | None = None
    # URL that the OnlyOffice Docker container can use to reach THIS backend.
    # On macOS/Windows Docker Desktop use http://host.docker.internal:<port>
    backend_url: str = "http://host.docker.internal:8000"

    # ---------------------------------------------------------------------------
    # LLM provider — "gemini" (default) or "openai" (legacy fallback)
    # ---------------------------------------------------------------------------
    llm_provider: str = "openrouter"

    # Gemini (primary)
    gemini_api_key: str | None = None
    llm_model: str = "google/gemini-2.5-flash-lite"
    embedding_model: str = "openai/text-embedding-3-small"

    # OpenAI (legacy — kept so the OpenAI provider path still works during
    # any transition period; unused when llm_provider="gemini")
    openai_api_key: str | None = None
    openai_base_url: str | None = None

    open_router_api: str | None = None

    # ---------------------------------------------------------------------------
    # Vector store (Qdrant)
    # ---------------------------------------------------------------------------
    qdrant_url: str | None = None
    qdrant_api_key: str | None = None

    max_review_iterations: int = 3


settings = Settings()

# ===== END core/config.py =====

# ===== BEGIN db/session.py =====
from collections.abc import Generator

from sqlalchemy import create_engine
from sqlalchemy.orm import DeclarativeBase, Session, sessionmaker



connect_args = {"check_same_thread": False} if settings.database_url.startswith("sqlite") else {}
engine = create_engine(settings.database_url, connect_args=connect_args)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)


class Base(DeclarativeBase):
    None


def get_db() -> Generator[Session, None, None]:
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# ===== END db/session.py =====

# ===== BEGIN models.py =====
from datetime import datetime
from uuid import uuid4

from sqlalchemy import DateTime, ForeignKey, Integer, String, Text
from sqlalchemy.orm import Mapped, mapped_column, relationship
from sqlalchemy.types import JSON



def new_id() -> str:
    return str(uuid4())


class User(Base):
    __tablename__ = "users"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    email: Mapped[str] = mapped_column(String, unique=True, index=True)
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)


class Workspace(Base):
    __tablename__ = "workspaces"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    user_id: Mapped[str] = mapped_column(String, ForeignKey("users.id"), nullable=True)
    document_type: Mapped[str] = mapped_column(String(10))
    original_filename: Mapped[str] = mapped_column(String)
    current_version: Mapped[int] = mapped_column(Integer, default=1)
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)
    updated_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    messages: Mapped[list["Message"]] = relationship(cascade="all, delete-orphan")
    versions: Mapped[list["DocumentVersion"]] = relationship(cascade="all, delete-orphan")
    knowledge_documents: Mapped[list["KnowledgeDocument"]] = relationship(cascade="all, delete-orphan")


class Message(Base):
    __tablename__ = "messages"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    role: Mapped[str] = mapped_column(String(20))
    content: Mapped[str] = mapped_column(Text)
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)


class DocumentVersion(Base):
    __tablename__ = "document_versions"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    version_number: Mapped[int] = mapped_column(Integer)
    document_path: Mapped[str] = mapped_column(String)
    pdf_path: Mapped[str] = mapped_column(String)
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)


class AgentRun(Base):
    __tablename__ = "agent_runs"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    status: Mapped[str] = mapped_column(String(20), default="running")
    started_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)
    completed_at: Mapped[datetime] = mapped_column(DateTime, nullable=True)


class DocumentStructure(Base):
    __tablename__ = "document_structures"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    version_number: Mapped[int] = mapped_column(Integer)
    structure_json: Mapped[dict] = mapped_column(JSON)


# ---------------------------------------------------------------------------
# Knowledge Base Models
# ---------------------------------------------------------------------------

KB_MAX_DOCUMENTS_PER_WORKSPACE = 50
KB_MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024   # 50 MB per file
KB_MAX_TOTAL_BYTES_PER_WORKSPACE = 100 * 1024 * 1024  # 100 MB total


class KnowledgeDocument(Base):
    """A document uploaded to a workspace's knowledge base."""
    __tablename__ = "knowledge_documents"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    filename: Mapped[str] = mapped_column(String)
    file_type: Mapped[str] = mapped_column(String(10))   # "pdf", "docx", "txt", "md"
    file_path: Mapped[str] = mapped_column(String)
    file_size_bytes: Mapped[int] = mapped_column(Integer, default=0)
    chunk_count: Mapped[int] = mapped_column(Integer, default=0)
    status: Mapped[str] = mapped_column(String(20), default="processing")  # processing|indexed|failed
    error_message: Mapped[str] = mapped_column(Text, nullable=True)
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)

    chunks: Mapped[list["KnowledgeChunk"]] = relationship(cascade="all, delete-orphan")


class KnowledgeChunk(Base):
    """A text chunk from a knowledge document, indexed for retrieval."""
    __tablename__ = "knowledge_chunks"

    id: Mapped[str] = mapped_column(String, primary_key=True, default=new_id)
    document_id: Mapped[str] = mapped_column(String, ForeignKey("knowledge_documents.id"), index=True)
    workspace_id: Mapped[str] = mapped_column(String, ForeignKey("workspaces.id"), index=True)
    chunk_index: Mapped[int] = mapped_column(Integer)
    text: Mapped[str] = mapped_column(Text)
    chunk_metadata: Mapped[dict] = mapped_column(JSON, default=dict)  # page, section, source info
    created_at: Mapped[datetime] = mapped_column(DateTime, default=datetime.utcnow)

# ===== END models.py =====

# ===== BEGIN schemas.py =====
from datetime import datetime
from typing import Any

from pydantic import BaseModel


class MessageOut(BaseModel):
    id: str
    role: str
    content: str
    content_parsed: dict[str, Any] | None = None
    created_at: datetime


class VersionOut(BaseModel):
    id: str
    version_number: int
    document_url: str
    pdf_url: str
    created_at: datetime


class KnowledgeDocumentOut(BaseModel):
    id: str
    filename: str
    file_type: str
    file_size_bytes: int
    chunk_count: int
    status: str
    error_message: str | None = None
    created_at: datetime


class WorkspaceOut(BaseModel):
    id: str
    document_type: str
    original_filename: str
    current_version: int
    created_at: datetime
    updated_at: datetime
    messages: list[MessageOut]
    versions: list[VersionOut]
    knowledge_documents: list[KnowledgeDocumentOut] = []


class ChatRequest(BaseModel):
    content: str


class ChatResponse(BaseModel):
    run_id: str
    status: str

# ===== END schemas.py =====

# ===== BEGIN repositories.py =====
from typing import List
from sqlalchemy import select, func
from sqlalchemy.orm import Session



class WorkspaceRepository:
    def __init__(self, db: Session):
        self.db = db

    def get(self, workspace_id: str) -> Workspace | None:
        return self.db.get(Workspace, workspace_id)

    def list(self) -> List[Workspace]:
        return list(self.db.scalars(select(Workspace).order_by(Workspace.updated_at.desc())))

    def latest_version(self, workspace_id: str) -> DocumentVersion | None:
        stmt = (
            select(DocumentVersion)
            .where(DocumentVersion.workspace_id == workspace_id)
            .order_by(DocumentVersion.version_number.desc())
        )
        return self.db.scalars(stmt).first()

    def version(self, workspace_id: str, version_number: int) -> DocumentVersion | None:
        stmt = select(DocumentVersion).where(
            DocumentVersion.workspace_id == workspace_id,
            DocumentVersion.version_number == version_number,
        )
        return self.db.scalars(stmt).first()

    def structure(self, workspace_id: str, version_number: int) -> DocumentStructure | None:
        stmt = select(DocumentStructure).where(
            DocumentStructure.workspace_id == workspace_id,
            DocumentStructure.version_number == version_number,
        )
        return self.db.scalars(stmt).first()

    def messages(self, workspace_id: str) -> List[Message]:
        stmt = select(Message).where(Message.workspace_id == workspace_id).order_by(Message.created_at.asc())
        return list(self.db.scalars(stmt))

    def versions(self, workspace_id: str) -> List[DocumentVersion]:
        stmt = (
            select(DocumentVersion)
            .where(DocumentVersion.workspace_id == workspace_id)
            .order_by(DocumentVersion.version_number.asc())
        )
        return list(self.db.scalars(stmt))

    # ── Knowledge Base ──────────────────────────────────────────────────────

    def list_knowledge_documents(self, workspace_id: str) -> List[KnowledgeDocument]:
        stmt = (
            select(KnowledgeDocument)
            .where(KnowledgeDocument.workspace_id == workspace_id)
            .order_by(KnowledgeDocument.created_at.asc())
        )
        return list(self.db.scalars(stmt))

    def get_knowledge_document(self, document_id: str) -> KnowledgeDocument | None:
        return self.db.get(KnowledgeDocument, document_id)

    def count_knowledge_documents(self, workspace_id: str) -> int:
        stmt = (
            select(func.count())
            .select_from(KnowledgeDocument)
            .where(KnowledgeDocument.workspace_id == workspace_id)
        )
        return self.db.scalar(stmt) or 0

    def total_knowledge_size_bytes(self, workspace_id: str) -> int:
        stmt = (
            select(func.sum(KnowledgeDocument.file_size_bytes))
            .where(KnowledgeDocument.workspace_id == workspace_id)
        )
        return self.db.scalar(stmt) or 0

    def list_knowledge_chunks(self, workspace_id: str) -> List[KnowledgeChunk]:
        stmt = (
            select(KnowledgeChunk)
            .where(KnowledgeChunk.workspace_id == workspace_id)
            .order_by(KnowledgeChunk.document_id, KnowledgeChunk.chunk_index)
        )
        return list(self.db.scalars(stmt))

# ===== END repositories.py =====

# ===== BEGIN services/llm_client.py =====
"""Provider-agnostic LLM client abstraction.

All pipeline stages (Task Planner, Operation Generator, Reference Resolver,
Content Enricher, Reviewer, etc.) call through this single interface.
Provider-specific code lives ONLY in this module.

Supported providers
-------------------
  "gemini"  — Google Gemini via the ``google-genai`` SDK (default)
  "openai"  — OpenAI-compatible API via the ``openai`` SDK (fallback / legacy)

Configuration (via Settings / .env)
------------------------------------
  LLM_PROVIDER       = "gemini" | "openai"       (default: "gemini")
  GEMINI_API_KEY     = <your Gemini API key>      (required when provider=gemini)
  LLM_MODEL          = "gemini-3.1-flash-lite"    (default)
  OPENAI_API_KEY     = <key>                      (required when provider=openai)
  OPENAI_BASE_URL    = <optional base url>        (openai only)

Usage
-----

    llm = LLMClient()
    response = llm.complete(LLMRequest(
        system_prompt="You are a ...",
        user_prompt="Do X and return JSON.",
        json_mode=True,
        temperature=0,
        max_tokens=1024,
    ))
    data = response.json   # already parsed dict/list, or None if parse failed
    raw  = response.text   # raw string always available
"""

import json
import logging
import time
from collections import deque
from dataclasses import dataclass, field
from typing import Any


log = logging.getLogger(__name__)

# Track the timestamps of the last 15 Gemini requests
_gemini_request_times = deque(maxlen=15)


# ---------------------------------------------------------------------------
# Public data structures
# ---------------------------------------------------------------------------

@dataclass
class LLMRequest:
    """Provider-agnostic request specification.

    Callers never import ``openai`` or ``google.genai`` — they build an
    ``LLMRequest`` and pass it to ``LLMClient.complete()``.
    """
    system_prompt: str
    user_prompt: str
    temperature: float = 0.0
    max_tokens: int = 2048
    json_mode: bool = False
    """If True, the provider is asked to return valid JSON.
    ``response.json`` will be the parsed result (dict or list).
    Use this for every structured-output call."""


@dataclass
class LLMResponse:
    """Normalised response from any LLM provider."""
    text: str
    """Raw text content — always set, even when json_mode=True."""
    json: dict | list | None = None
    """Parsed JSON when json_mode=True and parsing succeeded, otherwise None."""
    usage: dict = field(default_factory=dict)
    """Token usage info keyed by provider-specific names."""


# ---------------------------------------------------------------------------
# Main client
# ---------------------------------------------------------------------------

class LLMClient:
    """Singleton-style LLM client — instantiate once, inject into services.

    The instance is created in ``DocumentAgentGraph.__init__`` and passed
    to every service via constructor injection so there is a single
    connection-pool / retry budget for the whole pipeline run.
    """

    def __init__(self, provider: str | None = None, model: str | None = None) -> None:
        self._provider: str = provider or settings.llm_provider
        self._model: str = model or settings.llm_model
        self._client = self._build_client()

    # ------------------------------------------------------------------
    # Client factory
    # ------------------------------------------------------------------

    def _build_client(self) -> Any:
        if self._provider == "gemini":
            from google import genai  # type: ignore[import]
            if not settings.gemini_api_key:
                raise RuntimeError(
                    "GEMINI_API_KEY is not set. "
                    "Add it to .env or set LLM_PROVIDER=openai to use OpenAI."
                )
            return genai.Client(api_key=settings.gemini_api_key)

        if self._provider == "openrouter":
            from openai import OpenAI
            if not settings.open_router_api:
                raise RuntimeError("OPEN_ROUTER_API is not set in .env")
            return OpenAI(
                api_key=settings.open_router_api,
                base_url="https://openrouter.ai/api/v1",
            )

        # Fallback: OpenAI-compatible
        from openai import OpenAI
        if not settings.openai_api_key:
            raise RuntimeError(
                "OPENAI_API_KEY is not set. "
                "Set GEMINI_API_KEY and LLM_PROVIDER=gemini to use Gemini."
            )
        return OpenAI(
            api_key=settings.openai_api_key,
            base_url=settings.openai_base_url or None,
        )

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def complete(self, request: LLMRequest) -> LLMResponse:
        """Execute a completion request against the configured provider."""
        if self._provider == "gemini":
            return self._complete_gemini(request)
        elif self._provider == "openrouter":
            return self._complete_openai(request) # OpenRouter is OpenAI compatible
        return self._complete_openai(request)

    # ------------------------------------------------------------------
    # Gemini provider
    # ------------------------------------------------------------------

    def _complete_gemini(self, req: LLMRequest) -> LLMResponse:
        from google.genai import types  # type: ignore[import]

        config_kwargs: dict[str, Any] = {
            "temperature": req.temperature,
            "max_output_tokens": req.max_tokens,
            "system_instruction": req.system_prompt,
        }
        if req.json_mode:
            # Ask Gemini to respond with JSON.
            # NOTE: Gemini 3.5 Flash supports constrained JSON output via
            # response_mime_type. We do NOT pass response_schema here because
            # our operations involve open-ended arrays that don't map cleanly
            # to a single Pydantic model. Use plain json mode + parse manually.
            config_kwargs["response_mime_type"] = "application/json"

        # Rate limiter: 15 requests per minute
        if len(_gemini_request_times) == 15:
            elapsed = time.time() - _gemini_request_times[0]
            if elapsed < 60.0:
                sleep_time = 60.0 - elapsed
                log.info(f"Gemini RPM limit reached. Sleeping for {sleep_time:.2f} seconds.")
                time.sleep(sleep_time)
        
        _gemini_request_times.append(time.time())

        try:
            response = self._client.models.generate_content(
                model=self._model,
                contents=req.user_prompt,
                config=types.GenerateContentConfig(**config_kwargs),
            )
        except Exception as exc:
            log.error("Gemini API call failed: %s", exc)
            raise

        text = (response.text or "").strip()
        parsed = self._try_parse_json(text) if req.json_mode else None
        return LLMResponse(text=text, json=parsed)

    # ------------------------------------------------------------------
    # OpenAI provider
    # ------------------------------------------------------------------

    def _complete_openai(self, req: LLMRequest) -> LLMResponse:
        kwargs: dict[str, Any] = {
            "model": self._model,
            "messages": [
                {"role": "system", "content": req.system_prompt},
                {"role": "user", "content": req.user_prompt},
            ],
            "temperature": req.temperature,
            "max_tokens": req.max_tokens,
        }
        if req.json_mode:
            kwargs["response_format"] = {"type": "json_object"}

        try:
            response = self._client.chat.completions.create(**kwargs)
        except Exception as exc:
            log.error("OpenAI API call failed: %s", exc)
            raise

        text = (response.choices[0].message.content or "").strip()
        usage = {}
        if response.usage:
            usage = {
                "prompt_tokens": response.usage.prompt_tokens,
                "completion_tokens": response.usage.completion_tokens,
            }
        parsed = self._try_parse_json(text) if req.json_mode else None
        return LLMResponse(text=text, json=parsed, usage=usage)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _try_parse_json(text: str) -> dict | list | None:
        """Attempt to parse text as JSON. Returns None on failure."""
        if not text:
            return None
        # Strip accidental markdown fences that some models add
        stripped = text
        if stripped.startswith("```json"):
            stripped = stripped[7:]
        if stripped.startswith("```"):
            stripped = stripped[3:]
        if stripped.endswith("```"):
            stripped = stripped[:-3]
        stripped = stripped.strip()

        try:
            return json.loads(stripped)
        except json.JSONDecodeError:
            log.warning("LLM returned non-JSON text (first 200 chars): %s", text[:200])
            return None

# ===== END services/llm_client.py =====

# ===== BEGIN services/embedding_client.py =====
"""Embedding client for gemini-embedding-2 (multimodal, GA).

Key differences from gemini-embedding-001 this code accounts for:
1. No ``task_type`` parameter — task instructions are prepended into the
   input text itself as a prefix string.
2. Passing multiple strings in one ``contents`` call returns a SINGLE
   aggregated embedding, not one-per-string. To get one vector per
   input (which is what a document-chunk index needs), call
   ``embed_content`` once per string.

Usage
-----

    client = GeminiEmbeddingClient()

    # Indexing document chunks:
    vectors = client.embed_documents(["Revenue grew 15%...", "Q3 highlights..."])

    # Querying:
    query_vec = client.embed_query("revenue growth metrics")
"""

import logging


log = logging.getLogger(__name__)


class TaskType:
    """Task instruction prefixes for gemini-embedding-2.

    NOTE: verify these exact prefix strings against the current
    'Task types with Embeddings 2' section of ai.google.dev/gemini-api/docs/embeddings
    before relying on them in production — Google's docs give worked
    examples per task rather than one fixed template, and the exact
    wording matters for embedding quality.
    """
    RETRIEVAL_DOCUMENT = "search result"   # use when indexing documents/chunks
    RETRIEVAL_QUERY    = "search query"    # use when embedding a user query
    SEMANTIC_SIMILARITY = "similarity"
    CLASSIFICATION     = "classification"
    CLUSTERING         = "clustering"


class GeminiEmbeddingClient:
    """Embedding client for gemini-embedding-2."""

    MODEL = "gemini-embedding-2"

    def __init__(self, output_dimensionality: int = 768) -> None:
        from google import genai  # type: ignore[import]
        if not settings.gemini_api_key:
            raise RuntimeError(
                "GEMINI_API_KEY is not set. Required for GeminiEmbeddingClient."
            )
        self._client = genai.Client(api_key=settings.gemini_api_key)
        # 768 recommended for storage efficiency; bump to 1536/3072 if
        # retrieval quality benchmarks demand it.
        # IMPORTANT: this must match the dimensionality your Qdrant
        # collection was created with — changing it requires a re-index.
        self._output_dim = output_dimensionality

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        """Embed a batch of document chunks for indexing.

        One embed_content call per text — NOT one call with all texts —
        because gemini-embedding-2 aggregates multi-input calls into a
        single vector rather than returning one per input.
        """
        vectors: list[list[float]] = []
        for text in texts:
            formatted = self._format_input(text, TaskType.RETRIEVAL_DOCUMENT, role="document")
            try:
                response = self._client.models.embed_content(
                    model=self.MODEL,
                    contents=formatted,
                    config={"output_dimensionality": self._output_dim},
                )
                vectors.append(list(response.embeddings[0].values))
            except Exception as exc:
                log.error("gemini-embedding-2 document embed failed: %s", exc)
                raise
        return vectors

    def embed_query(self, text: str) -> list[float]:
        """Embed a single query string for similarity search against the index."""
        formatted = self._format_input(text, TaskType.RETRIEVAL_QUERY, role="query")
        try:
            response = self._client.models.embed_content(
                model=self.MODEL,
                contents=formatted,
                config={"output_dimensionality": self._output_dim},
            )
            return list(response.embeddings[0].values)
        except Exception as exc:
            log.error("gemini-embedding-2 query embed failed: %s", exc)
            raise

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _format_input(self, text: str, task: str, role: str = "document") -> str:
        """Prepend the task instruction prefix required by gemini-embedding-2.

        ``role`` distinguishes asymmetric retrieval: documents get indexed
        under one framing, queries embedded under another, so a query
        and its matching document land close together in vector space.
        """
        return f"task: {task} | {role}: {text}"

class OpenAIEmbeddingClient:
    """Embedding client for OpenAI."""

    def __init__(self) -> None:
        from openai import OpenAI
        if not settings.openai_api_key:
            raise RuntimeError(
                "OPENAI_API_KEY is not set. Required for OpenAIEmbeddingClient."
            )
        self._client = OpenAI(
            api_key=settings.openai_api_key,
            base_url=settings.openai_base_url or None,
        )
        self.model = settings.embedding_model or "text-embedding-3-small"

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        try:
            res = self._client.embeddings.create(input=texts, model=self.model)
            return [data.embedding for data in res.data]
        except Exception as exc:
            log.error("openai embed failed: %s", exc)
            raise

    def embed_query(self, text: str) -> list[float]:
        try:
            res = self._client.embeddings.create(input=[text], model=self.model)
            return res.data[0].embedding
        except Exception as exc:
            log.error("openai query embed failed: %s", exc)
            raise

class OpenRouterEmbeddingClient:
    """Embedding client for OpenRouter."""

    def __init__(self) -> None:
        from openai import OpenAI
        import httpx
        if not settings.open_router_api:
            raise RuntimeError(
                "OPEN_ROUTER_API is not set. Required for OpenRouterEmbeddingClient."
            )
        self._client = OpenAI(
            api_key=settings.open_router_api,
            base_url="https://openrouter.ai/api/v1",
            timeout=httpx.Timeout(60.0, read=60.0, write=60.0, connect=10.0),
        )
        self.model = settings.embedding_model or "openai/text-embedding-3-small"

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        
        # Batch texts to prevent write operation timeouts
        batch_size = 16
        all_embeddings = []
        
        try:
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                res = self._client.embeddings.create(input=batch, model=self.model)
                all_embeddings.extend([data.embedding for data in res.data])
            return all_embeddings
        except Exception as exc:
            log.error("openrouter embed failed: %s", exc)
            raise

    def embed_query(self, text: str) -> list[float]:
        try:
            res = self._client.embeddings.create(input=[text], model=self.model)
            return res.data[0].embedding
        except Exception as exc:
            log.error("openrouter query embed failed: %s", exc)
            raise

# ===== END services/embedding_client.py =====

# ===== BEGIN services/operations.py =====
"""Structured operation definitions for the document editing pipeline.

Every operation the LLM can produce is described here.  The document
processor dispatches on ``op_type`` to call the correct handler.

Op type taxonomy
----------------
text_edit      — rewrite text content of a paragraph (existing behaviour)
text_format    — change formatting of a paragraph/run (bold, font, color …)
table_op       — create/edit/delete tables and their cells
image_op       — insert/replace/resize/style images
shape_op       — add/edit text boxes and shapes
theme_op       — change slide/document background, colors, gradients
slide_op       — add/delete/duplicate/reorder/hide slides
chart_op       — change chart type, colors, labels, data
ai_design_op   — AI-directed normalizations (font, spacing, hierarchy)
"""

from typing import Any, Literal


# ---------------------------------------------------------------------------
# Base
# ---------------------------------------------------------------------------

class Operation:
    """Base class — not used directly; here for documentation only."""
    None


# ---------------------------------------------------------------------------
# Type aliases (used as JSON keys)
# ---------------------------------------------------------------------------

OpType = Literal[
    "text_edit",
    "text_format",
    "table_op",
    "image_op",
    "shape_op",
    "theme_op",
    "slide_op",
    "chart_op",
    "ai_design_op",
    "layout_op",
    "needs_image",
]


# ---------------------------------------------------------------------------
# Parameter schemas (all fields optional — None means "keep existing")
# ---------------------------------------------------------------------------

class TextEditParams:
    """Parameters for a text content rewrite."""
    new_text: str


class TextFormatParams:
    """Parameters for formatting a targeted paragraph."""
    bold: bool | None
    italic: bool | None
    underline: bool | None
    strikethrough: bool | None
    font_family: str | None          # e.g. "Arial", "Calibri"
    font_size_pt: float | None
    color_hex: str | None            # 6-char hex e.g. "FF0000"
    highlight_hex: str | None
    alignment: str | None            # "left"|"center"|"right"|"justify"
    line_spacing: float | None       # multiplier e.g. 1.5
    char_spacing: float | None       # pt
    superscript: bool | None
    subscript: bool | None
    shadow: bool | None


class TableOpParams:
    """Parameters for a table CRUD operation."""
    action: str  # "create"|"delete"|"add_row"|"remove_row"|"add_col"|"remove_col"
                 # "merge_cells"|"set_cell_bg"|"set_borders"|"alternate_rows"
                 # "populate"|"set_header_format"|"sort_data"
    rows: int | None
    cols: int | None
    header_row: bool | None
    alternate_row_colors: list[str] | None   # two 6-char hex strings
    data: list[list[str]] | None             # row-major cell content
    row_index: int | None                    # for add/remove row
    col_index: int | None                    # for add/remove col
    merge_from: tuple[int, int] | None       # (row, col)
    merge_to: tuple[int, int] | None         # (row, col)
    cell_bg_hex: str | None
    border_color_hex: str | None
    border_width_pt: float | None
    position: dict[str, float] | None        # {left_pct, top_pct, width_pct, height_pct}
    cell_padding_pt: float | None
    cell_alignment: str | None


class ImageOpParams:
    """Parameters for image operations."""
    action: str  # "insert"|"replace"|"remove"|"resize"|"reposition"|"move"|"place_inline"|"add_caption"|"reposition_caption"|"format_caption"|"rotate"
                 # "bring_forward"|"send_backward"|"set_transparency"|"set_border"
                 # "rounded_corners"|"shadow"
    image_path: str | None           # server-side absolute path to uploaded image
    position: dict[str, float] | None   # {left_pct, top_pct, width_pct, height_pct}
    after_id: str | None             # DOCX anchor element for insertions
    before_id: str | None            # DOCX anchor element for insertions
    width_page_pct: float | None     # DOCX image width as fraction of usable page width
    height_page_pct: float | None
    alignment: str | None            # "left"|"center"|"right"
    float_position: str | None       # "left"|"right"
    caption_text: str | None
    alt_text: str | None
    maintain_aspect_ratio: bool | None
    rotation_degrees: float | None
    transparency_pct: float | None      # 0–100
    border_color_hex: str | None
    border_width_pt: float | None
    rounded_corners: bool | None
    shadow: bool | None
    crop: dict[str, float] | None       # {top_pct, left_pct, right_pct, bottom_pct}


class ShapeOpParams:
    """Parameters for shape / text-box operations."""
    action: str  # "add_textbox"|"delete"|"resize"|"move"|"rotate"|"duplicate"
                 # "set_fill"|"set_outline"|"set_transparency"|"group"|"ungroup"
                 # "bring_forward"|"send_backward"|"align"|"distribute"
    text: str | None
    position: dict[str, float] | None
    fill_color_hex: str | None
    outline_color_hex: str | None
    outline_width_pt: float | None
    transparency_pct: float | None
    rotation_degrees: float | None
    corner_radius_pt: float | None
    group_shape_indices: list[int] | None


class ThemeOpParams:
    """Parameters for theme/color/background operations."""
    action: str  # "set_bg_color"|"set_bg_gradient"|"set_bg_pattern"
                 # "apply_theme_colors"|"corporate_branding"
    scope: str | None                # "all_slides"|"current_slide"
    bg_color_hex: str | None
    gradient_start_hex: str | None
    gradient_end_hex: str | None
    gradient_direction: str | None   # "horizontal"|"vertical"|"diagonal"
    accent_colors: list[str] | None  # list of 6-char hex strings


class SlideOpParams:
    """Parameters for slide-level operations."""
    action: str  # "add"|"delete"|"duplicate"|"reorder"|"hide"|"unhide"
                 # "rename_title"|"apply_layout"|"change_size"
    after_index: int | None          # for "add"/"duplicate": insert after this 1-based index
    from_index: int | None           # for "reorder": source position (1-based)
    to_index: int | None             # for "reorder": target position (1-based)
    layout_name: str | None
    title: str | None                # for "rename_title"


class ChartOpParams:
    """Parameters for chart editing operations."""
    action: str  # "change_type"|"update_data"|"set_series_colors"|"update_labels"
                 # "update_axis_labels"|"show_legend"|"hide_legend"|"apply_theme"
    chart_type: str | None           # "bar"|"line"|"pie"|"scatter"|"column"
    series_colors: list[str] | None  # hex colors per series
    data: list[list[Any]] | None     # new chart data (row-major)
    legend_position: str | None      # "top"|"bottom"|"left"|"right"|"none"
    x_axis_label: str | None
    y_axis_label: str | None
    data_labels_visible: bool | None


class AiDesignOpParams:
    """Parameters for AI-driven design normalization."""
    action: str  # "normalize_fonts"|"normalize_spacing"|"improve_hierarchy"
                 # "balance_whitespace"|"remove_overlaps"|"auto_resize_text"
                 # "make_consistent"|"generate_speaker_notes"|"convert_bullets_to_diagram"
                 # "improve_readability"|"detect_clutter"
    scope: str | None                # "all_slides"|"slide:{n}"
    target_font: str | None          # for normalize_fonts
    base_font_size_pt: float | None


# ---------------------------------------------------------------------------
# Unified operation dict (as produced by the LLM and consumed by processor)
# ---------------------------------------------------------------------------

def validate_operation(op: dict) -> dict:
    """Light validation / normalisation of a raw LLM-produced operation dict.
    
    Returns the op dict with defaults filled in, or raises ValueError if
    the op is structurally invalid.
    """
    if not isinstance(op, dict):
        raise ValueError(f"Operation must be a dict, got {type(op)}")
    
    op_type = op.get("op_type")
    valid_types = {
        "text_edit", "text_format", "table_op", "image_op",
        "shape_op", "theme_op", "slide_op", "chart_op", "layout_op",
        "ai_design_op", "needs_image", "list_op", "find_replace",
    }
    if op_type not in valid_types:
        raise ValueError(f"Unknown op_type: {op_type!r}")
    
    op.setdefault("target_id", None)
    op.setdefault("parameters", {})
    
    return op


def needs_image_response(reason: str = "") -> dict:
    """Build a special operation that signals the agent needs an image upload."""
    return {
        "op_type": "needs_image",
        "target_id": None,
        "parameters": {
            "message": (
                reason or
                "To insert an image, please attach it to your next message "
                "using the 📎 paperclip icon below the chat input."
            )
        }
    }

# ===== END services/operations.py =====

# ===== BEGIN services/docx_extensions.py =====
import logging
from docx.shared import Pt, Inches

log = logging.getLogger(__name__)

def extract_metadata(doc) -> dict:
    """Extracts core properties as a metadata dictionary."""
    try:
        cp = doc.core_properties
        return {
            "title": cp.title or "",
            "author": cp.author or "",
            "subject": cp.subject or "",
            "keywords": cp.keywords or "",
        }
    except Exception as e:
        log.error(f"Failed to extract metadata: {e}")
        return {}

def extract_advanced_paragraph_style(para) -> dict:
    """Extracts advanced paragraph styles (indents, pagination)."""
    style = {}
    try:
        pf = para.paragraph_format
        if pf.left_indent is not None:
            style["left_indent_pt"] = pf.left_indent.pt
        if pf.right_indent is not None:
            style["right_indent_pt"] = pf.right_indent.pt
        if pf.first_line_indent is not None:
            style["first_line_indent_pt"] = pf.first_line_indent.pt
        if pf.keep_with_next is not None:
            style["keep_with_next"] = pf.keep_with_next
        if pf.keep_together is not None:
            style["keep_together"] = pf.keep_together
    except Exception as e:
        None
    return style

def apply_metadata(doc, params: dict) -> str:
    """Updates the core properties of the document."""
    try:
        cp = doc.core_properties
        if "title" in params:
            cp.title = params["title"]
        if "author" in params:
            cp.author = params["author"]
        if "subject" in params:
            cp.subject = params["subject"]
        if "keywords" in params:
            cp.keywords = params["keywords"]
        return "Updated document metadata."
    except Exception as e:
        log.error(f"Failed to update metadata: {e}")
        return f"Failed to update metadata: {e}"

def apply_section_formatting(doc, target_id: str, params: dict) -> str:
    """Updates page size, orientation, and margins for a section."""
    try:
        # Default to first section if no target specified
        section = doc.sections[0]
        
        # In python-docx, sections are usually accessed by index. 
        # If target_id is something like "section_0", we can parse it.
        if target_id and target_id.startswith("section_"):
            try:
                idx = int(target_id.split("_")[1])
                if 0 <= idx < len(doc.sections):
                    section = doc.sections[idx]
            except ValueError:
                None
                
        action = params.get("action")
        if action == "set_margins":
            margins = params.get("margins", {})
            if "top_inches" in margins:
                section.top_margin = Inches(margins["top_inches"])
            if "bottom_inches" in margins:
                section.bottom_margin = Inches(margins["bottom_inches"])
            if "left_inches" in margins:
                section.left_margin = Inches(margins["left_inches"])
            if "right_inches" in margins:
                section.right_margin = Inches(margins["right_inches"])
            return f"Updated margins for section."
            
        elif action == "set_page_size":
            from docx.enum.section import WD_ORIENT
            orientation = params.get("orientation", "portrait").lower()
            if orientation == "landscape":
                section.orientation = WD_ORIENT.LANDSCAPE
                # Swap width and height if changing to landscape
                if section.page_width < section.page_height:
                    section.page_width, section.page_height = section.page_height, section.page_width
            else:
                section.orientation = WD_ORIENT.PORTRAIT
                if section.page_width > section.page_height:
                    section.page_width, section.page_height = section.page_height, section.page_width
            
            # Optional: set explicit width/height
            if "width_inches" in params:
                section.page_width = Inches(params["width_inches"])
            if "height_inches" in params:
                section.page_height = Inches(params["height_inches"])
                
            return f"Updated page setup for section."
            
        return "No section action performed."
    except Exception as e:
        log.error(f"Failed to apply section formatting: {e}")
        return f"Failed to apply section formatting: {e}"

def apply_global_style(doc, params: dict) -> str:
    """Modifies a global named style in the document."""
    try:
        style_name = params.get("style_name")
        if not style_name or style_name not in doc.styles:
            return f"Style '{style_name}' not found."
            
        style = doc.styles[style_name]
        
        if "font_name" in params:
            style.font.name = params["font_name"]
        if "font_size_pt" in params:
            style.font.size = Pt(params["font_size_pt"])
        if "bold" in params:
            style.font.bold = params["bold"]
        if "color_hex" in params:
            from docx.shared import RGBColor
            c = str(params["color_hex"]).strip().lstrip("#")
            if len(c) == 6:
                style.font.color.rgb = RGBColor(int(c[:2], 16), int(c[2:4], 16), int(c[4:], 16))
                
        return f"Updated global style '{style_name}'."
    except Exception as e:
        log.error(f"Failed to apply global style: {e}")
        return f"Failed to apply global style: {e}"

# ===== END services/docx_extensions.py =====

# ===== BEGIN services/document_processor.py =====
"""Document processor: extract and apply text edits to DOCX and PPTX files.

Run-aware replacement strategy
--------------------------------
We never use proportional character distribution across runs.  Instead we:

1. Concatenate all run texts to get the full paragraph text.
2. Find the longest common prefix and suffix between old and new text to
   isolate the minimal "changed region".
3. Walk the runs and update only the ones that overlap that region.
   - Runs entirely before the region: untouched.
   - First run overlapping the region: gets (its unchanged prefix) +
     (the new replacement text) + (its unchanged suffix if the change ends
     inside this run).
   - Any subsequent runs overlapping the region: get only their unchanged
     trailing suffix (the middle is consumed by the first run).
   - Runs entirely after the region: untouched.

This ensures that formatting (colour, font, size, bold, italic, …) stays on
exactly the same runs as before.  Only the text content of runs that
actually contain the changed characters is modified.

Example
~~~~~~~
Paragraph runs:
  run[0] "Theme Name: "  (green, 12 chars)
  run[1] "Edtech"         (black,  6 chars)

new_text = "Theme Name: Healthtech"
  common prefix  = "Theme Name: "  → 12 chars
  common suffix  = "tech"           →  4 chars
  changed region = chars [12, 14) = "Ed"  →  replaced with "Health"

  run[0] untouched             → "Theme Name: " (green)  ✓
  run[1] "Ed"→"Health" + "tech" → "Healthtech"   (black)  ✓
"""

import copy
import logging
from dataclasses import dataclass
from pathlib import Path
from re import sub
from typing import Any

from docx import Document
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu, Inches
from pptx.util import Emu as EMU

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Custom IDs (UUID stamping for stable element identities)
# ---------------------------------------------------------------------------
import uuid
import shutil
import tempfile
from lxml import etree
from docx.oxml.ns import qn

CUSTOM_NS = 'http://documenteditor.local/ids'
CUSTOM_PREFIX = 'deid'
etree.register_namespace(CUSTOM_PREFIX, CUSTOM_NS)

def _enable_update_fields(doc) -> None:
    """Ensure <w:updateFields w:val="true"/> exists in doc.settings in ECMA-376 schema-compliant order."""
    try:
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls, qn

        settings_el = doc.settings.element
        existing = settings_el.find(qn('w:updateFields'))
        if existing is not None:
            existing.set(qn('w:val'), 'true')
            return

        uf = parse_xml(r'<w:updateFields %s w:val="true"/>' % nsdecls('w'))

        # Local element names in CT_Settings schema that MUST appear after w:updateFields
        after_local_names = (
            'footnotePr', 'endnotePr', 'compat', 'docVars',
            'rsids', 'mathPr', 'attachedTemplate', 'linkStyles',
            'stylePaneFormatFilter', 'stylePaneSortMethod',
            'clrSchemeMapping', 'doNotIncludeSubdocsInStats',
            'doNotAutoCompressPictures', 'shapeDefaults',
            'decimalSymbol', 'listSeparator',
            'docId', 'defaultImageDpi'
        )

        target = None
        for child in settings_el:
            local_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if local_name in after_local_names:
                target = child
                break

        if target is not None:
            target.addprevious(uf)
        else:
            settings_el.append(uf)
    except Exception as exc:
        log.warning("Failed to set updateFields in document settings: %s", exc)

# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class TextBlock:
    element_id: str
    type: str
    text: str
    metadata: dict


# ---------------------------------------------------------------------------
# Diff helper
# ---------------------------------------------------------------------------

def _changed_region(old: str, new: str) -> tuple[int, int, str]:
    """Return (prefix_len, suffix_len, new_middle) such that:
      old[prefix_len : len(old) - suffix_len]  should become  new_middle
      new_middle == new[prefix_len : len(new) - suffix_len]
    """
    # Longest common prefix
    prefix_len = 0
    min_len = min(len(old), len(new))
    while prefix_len < min_len and old[prefix_len] == new[prefix_len]:
        prefix_len += 1

    # Longest common suffix (bounded so it cannot overlap the prefix)
    old_remaining = len(old) - prefix_len
    new_remaining = len(new) - prefix_len
    max_suffix = min(old_remaining, new_remaining)
    suffix_len = 0
    while (suffix_len < max_suffix
           and old[len(old) - 1 - suffix_len] == new[len(new) - 1 - suffix_len]):
        suffix_len += 1

    new_end = len(new) - suffix_len if suffix_len > 0 else len(new)
    new_middle = new[prefix_len:new_end]
    return prefix_len, suffix_len, new_middle


# Alignment string → pptx enum mapping
_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# ---------------------------------------------------------------------------
# Main processor
# ---------------------------------------------------------------------------

class DocumentProcessor:

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def extract(self, path: Path, document_type: str) -> dict:
        if document_type == "pptx":
            res = self._extract_pptx_dom(path)
        else:
            res = self._extract_docx_dom(path)
            
        res["blocks"] = self._flatten_dom_for_retrieval(res["dom"])
        return res

    def extract_rich(self, path: Path, document_type: str) -> dict:
        """Extract enriched structure for template analysis (shapes, geometry, formatting)."""
        return self.extract(path, document_type)

    def _flatten_dom_for_retrieval(self, node: dict) -> list[dict]:
        blocks = []
        def traverse(n: dict):
            if n.get("type") == "paragraph" and n.get("text"):
                # Aggregate run styles to provide style hints
                colors = set()
                fonts = set()
                for r in n.get("runs", []):
                    st = r.get("style", {})
                    if st.get("color"): colors.add(st["color"])
                    if st.get("font"): fonts.add(st["font"])

                blocks.append({
                    "element_id": n["id"],
                    "text": n["text"],
                    "type": "text",
                    "metadata": {
                        "role": n.get("role"),
                        "colors": list(colors),
                        "fonts": list(fonts),
                        "include_in_toc": n.get("include_in_toc")
                    }
                })
            for child in n.get("children", []) + n.get("rows", []) + n.get("cells", []):
                traverse(child)
        traverse(node)
        return blocks

    def apply_edits(
        self,
        source: Path,
        target: Path,
        document_type: str,
        edits: list[dict],
    ) -> None:
        if document_type == "pptx":
            self._apply_pptx_edits(source, target, edits)
        else:
            self._apply_docx_edits(source, target, edits)

    def apply_slide_plan(
        self,
        source: Path,
        target: Path,
        slide_plan: dict,
    ) -> None:
        """Apply a structured slide plan to a PPTX template.

        The plan specifies which template slides to clone, what content to
        populate, and what formatting to apply.  Slides marked "delete" are
        removed.
        """
        self._apply_pptx_slide_plan(source, target, slide_plan)

    def apply_operations(
        self,
        source: Path,
        target: Path,
        document_type: str,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to the document.

        Returns (any_change_made, list_of_applied_summaries).
        """
        if document_type == "pptx":
            return self._apply_pptx_operations(source, target, operations)
        else:
            return self._apply_docx_operations(source, target, operations)

    # ------------------------------------------------------------------
    # PPTX: one block per paragraph inside each shape (original)
    # ------------------------------------------------------------------

    def _extract_pptx(self, path: Path) -> list[TextBlock]:
        prs = Presentation(path)
        blocks: list[TextBlock] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    eid = f"slide_{slide_idx}_shape_{shape_idx}_para_{para_idx}"
                    blocks.append(TextBlock(
                        element_id=eid,
                        type="text",
                        text=text,
                        metadata={
                            "slide": slide_idx,
                            "shape_index": shape_idx,
                            "shape_name": getattr(shape, "name", "Unknown"),
                            "para_index": para_idx,
                        },
                    ))
        return blocks

    # ------------------------------------------------------------------
    # PPTX: rich extraction (includes empty frames, geometry, formatting)
    # ------------------------------------------------------------------

    def _extract_pptx_dom(self, path: Path) -> dict:
        prs = Presentation(path)
        children = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_children = []
            for shape_idx, shape in enumerate(slide.shapes):
                shape_id = f"slide_{slide_idx}_shape_{shape_idx}"
                shape_role = "shape"
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).split('(')[0]
                        shape_role = f"placeholder_{ph_type.lower()}"
                    except Exception:
                        shape_role = "placeholder"
                
                shape_node = {
                    "id": shape_id,
                    "type": "shape",
                    "role": shape_role,
                    "name": getattr(shape, "name", "Unknown"),
                    "geometry": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    },
                    "children": []
                }

                if getattr(shape, "has_text_frame", False):
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        shape_node["children"].append(
                            self._extract_pptx_paragraph_dom(para, f"{shape_id}_para_{para_idx}")
                        )

                if getattr(shape, "has_table", False):
                    table_node = {
                        "id": f"{shape_id}_table",
                        "type": "table",
                        "role": "table",
                        "rows": []
                    }
                    for row_idx, row in enumerate(shape.table.rows):
                        row_node = {
                            "id": f"{shape_id}_table_row_{row_idx}",
                            "type": "row",
                            "row": row_idx,
                            "cells": []
                        }
                        for col_idx, cell in enumerate(row.cells):
                            cell_node = {
                                "id": f"{shape_id}_table_cell_{row_idx}_{col_idx}",
                                "type": "cell",
                                "row": row_idx,
                                "column": col_idx,
                                "children": []
                            }
                            for para_idx, para in enumerate(cell.text_frame.paragraphs):
                                cell_node["children"].append(
                                    self._extract_pptx_paragraph_dom(para, f"{cell_node['id']}_para_{para_idx}")
                                )
                            row_node["cells"].append(cell_node)
                        table_node["rows"].append(row_node)
                    shape_node["children"].append(table_node)

                slide_children.append(shape_node)

            layout_name = ""
            try:
                layout_name = slide.slide_layout.name
            except Exception:
                None

            children.append({
                "id": f"slide_{slide_idx}",
                "type": "slide",
                "role": "slide",
                "layout_name": layout_name,
                "children": slide_children
            })

        return {
            "document_type": "pptx",
            "slide_count": len(prs.slides),
            "geometry": {
                "width": prs.slide_width,
                "height": prs.slide_height,
            },
            "dom": {
                "id": "document_root",
                "type": "document",
                "children": children
            }
        }

    def _extract_pptx_paragraph_dom(self, para, para_id: str) -> dict:
        alignment = None
        if para.alignment is not None:
            try:
                alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
            except Exception:
                None

        runs = []
        for r_idx, run in enumerate(para.runs):
            font = run.font
            color_hex = None
            try:
                if font.color and font.color.type is not None and font.color.rgb:
                    color_hex = str(font.color.rgb)
            except Exception:
                None
            runs.append({
                "id": f"{para_id}_run_{r_idx}",
                "type": "run",
                "text": run.text,
                "style": {
                    "font": font.name,
                    "size": round(font.size.pt, 1) if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                    "color": color_hex
                }
            })

        return {
            "id": para_id,
            "type": "paragraph",
            "role": "body",
            "text": para.text.strip(),
            "style": {
                "alignment": alignment,
            },
            "runs": runs
        }

    # ------------------------------------------------------------------
    # PPTX: apply text edits (original)
    # ------------------------------------------------------------------

    def _apply_pptx_edits(self, source: Path, target: Path, edits: list[dict]) -> None:
        prs = Presentation(source)
        edit_map = {e["element_id"]: e for e in edits}
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    eid = f"slide_{slide_idx}_shape_{shape_idx}_para_{para_idx}"
                    if eid in edit_map:
                        edit = edit_map[eid]
                        log.debug(
                            "PPTX edit [%s]: %r → %r",
                            eid, edit["old_text"], edit["new_text"],
                        )
                        self._apply_run_aware_replacement(para, edit["new_text"])
        prs.save(target)

    # ------------------------------------------------------------------
    # PPTX: slide plan application (new)
    # ------------------------------------------------------------------

    def _apply_pptx_slide_plan(self, source: Path, target: Path, plan: dict) -> None:
        """Apply a structured slide plan to a PPTX template.

        Plan structure::

            {
                "slides": [
                    {
                        "source_slide_index": 1,
                        "action": "populate" | "keep" | "delete",
                        "shapes": [
                            {
                                "shape_index": 0,
                                "paragraphs": [
                                    {
                                        "para_index": 0,
                                        "text": "New text",
                                        "formatting": {
                                            "font_size_pt": 36,
                                            "bold": true,
                                            "color_hex": null,
                                            "alignment": "center"
                                        }
                                    }
                                ]
                            }
                        ]
                    }
                ]
            }
        """
        prs = Presentation(source)
        template_slides = list(prs.slides)
        plan_slides = plan.get("slides", [])

        if not plan_slides:
            log.warning("Empty slide plan — saving template as-is.")
            prs.save(target)
            return

        # Phase 1: Build the output slide list by cloning template slides
        # We'll work backwards: first create all needed clones, then delete originals.
        new_slide_xmls: list[tuple[str, Any]] = []  # (action, slide_element or None)

        for entry in plan_slides:
            action = entry.get("action", "populate")
            if action == "delete":
                continue  # skip deleted slides

            src_idx = entry.get("source_slide_index", 1) - 1  # convert to 0-based
            if src_idx < 0 or src_idx >= len(template_slides):
                src_idx = 0  # fallback to first slide

            new_slide_xmls.append((action, src_idx, entry))

        # Phase 2: Build new presentation from template
        new_prs = Presentation(source)
        existing_count = len(list(new_prs.slides))

        # Clone ALL slides defined in the plan, appending them to the END
        for action, src_idx, entry in new_slide_xmls:
            self._clone_slide(new_prs, src_idx)

        # Delete ALL original slides from the beginning (in reverse order to avoid shifting)
        for i in range(existing_count - 1, -1, -1):
            self._delete_slide(new_prs, i)

        # Phase 3: Populate content for each slide
        slides_list = list(new_prs.slides)
        for slide_out_idx, (action, src_idx, entry) in enumerate(new_slide_xmls):
            if slide_out_idx >= len(slides_list):
                break

            slide = slides_list[slide_out_idx]

            # For cloned slides that don't match their source, we need to
            # copy content from the correct template slide
            if action == "keep":
                continue  # leave as-is

            # Populate shapes
            shape_edits = entry.get("shapes", [])
            edit_map = {se.get("shape_index", 0): se for se in shape_edits}
            shapes = list(slide.shapes)
            
            for shape_idx, shape in enumerate(shapes):
                has_tf = getattr(shape, "has_text_frame", False)
                has_t = getattr(shape, "has_table", False)
                
                if shape_idx in edit_map:
                    shape_edit = edit_map[shape_idx]
                    
                    if has_tf:
                        para_edits = shape_edit.get("paragraphs", [])
                        paras = list(shape.text_frame.paragraphs)
                        for para_edit in para_edits:
                            para_idx = para_edit.get("para_index", 0)
                            new_text = para_edit.get("text", "")
                            formatting = para_edit.get("formatting", {})

                            if para_idx < len(paras):
                                para = paras[para_idx]
                                self._set_paragraph_text(para, new_text)
                                self._apply_formatting(para, formatting)
                            else:
                                # Add new paragraph
                                para = shape.text_frame.paragraphs[0] if paras else shape.text_frame.add_paragraph()
                                if para_idx > 0:
                                    para = shape.text_frame.add_paragraph()
                                para.text = ""
                                run = para.add_run()
                                run.text = new_text
                                self._apply_formatting(para, formatting)
                                
                    if has_t and "table_rows" in shape_edit:
                        table_rows = shape_edit.get("table_rows", [])
                        for r_idx, row_edit in enumerate(table_rows):
                            if r_idx >= len(shape.table.rows):
                                break
                            for c_idx, cell_edit in enumerate(row_edit):
                                if c_idx >= len(shape.table.columns):
                                    break
                                cell = shape.table.cell(r_idx, c_idx)
                                para_edits = cell_edit.get("paragraphs", [])
                                paras = list(cell.text_frame.paragraphs)
                                
                                for para_edit in para_edits:
                                    para_idx = para_edit.get("para_index", 0)
                                    new_text = para_edit.get("text", "")
                                    formatting = para_edit.get("formatting", {})
                                    
                                    if para_idx < len(paras):
                                        para = paras[para_idx]
                                        self._set_paragraph_text(para, new_text)
                                        self._apply_formatting(para, formatting)
                                    else:
                                        para = cell.text_frame.paragraphs[0] if paras else cell.text_frame.add_paragraph()
                                        if para_idx > 0:
                                            para = cell.text_frame.add_paragraph()
                                        para.text = ""
                                        run = para.add_run()
                                        run.text = new_text
                                        self._apply_formatting(para, formatting)
                else:
                    # Unedited shape. If it's a placeholder (and NOT a slide number), clear it
                    # to remove leftover boilerplate.
                    if getattr(shape, "is_placeholder", False) and has_tf:
                        try:
                            ph_type = str(shape.placeholder_format.type)
                            if "SLIDE_NUMBER" not in ph_type:
                                # Clear all paragraphs
                                paras = list(shape.text_frame.paragraphs)
                                if paras:
                                    self._set_paragraph_text(paras[0], "")
                                    for p in paras[1:]:
                                        self._set_paragraph_text(p, "")
                        except Exception:
                            None

        new_prs.save(target)

    def _clone_slide(self, prs: Presentation, slide_index: int) -> None:
        """Deep-clone a slide at the given index and append it to the presentation."""
        template_slide = list(prs.slides)[slide_index]
        slide_layout = template_slide.slide_layout

        # Add a new slide with the same layout
        new_slide = prs.slides.add_slide(slide_layout)

        # Copy all shapes from the template slide to the new slide
        # We do this by copying the XML of each shape
        for shape in template_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # Remove the default placeholder shapes that come with the layout
        # (they duplicate what we just copied)
        sp_tree = new_slide.shapes._spTree
        # Collect placeholders that were auto-created by add_slide
        default_shapes = []
        for sp in sp_tree:
            if sp.tag.endswith("}sp") or sp.tag == "sp":
                # Check if this is a default placeholder (not one we cloned)
                None  # We'll use a different approach

        # Actually, the cleaner approach: remove all default shapes first,
        # then copy from template
        # Let's redo: remove shapes added by add_slide, keep only our cloned ones
        # The shapes added by add_slide come from the layout's placeholders
        # Our cloned shapes are appended at the end

        # Count shapes from template
        template_shape_count = len(list(template_slide.shapes))

        # The spTree contains: <cNvPr> (non-visual props) + shapes
        # Shapes added by layout come first, our cloned ones come last
        all_sps = [child for child in sp_tree
                   if child.tag.endswith("}sp") or child.tag.endswith("}pic")
                   or child.tag.endswith("}grpSp") or child.tag.endswith("}graphicFrame")
                   or child.tag.endswith("}cxnSp")]

        if len(all_sps) > template_shape_count:
            # Remove the auto-generated shapes (first N shapes minus our cloned ones)
            auto_count = len(all_sps) - template_shape_count
            for sp in all_sps[:auto_count]:
                sp_tree.remove(sp)

        # Copy slide-level properties (background, etc.)
        try:
            if template_slide._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            ) is not None:
                bg = copy.deepcopy(
                    template_slide._element.find(
                        "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
                    )
                )
                existing_bg = new_slide._element.find(
                    "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
                )
                if existing_bg is not None:
                    new_slide._element.replace(existing_bg, bg)
                else:
                    new_slide._element.insert(0, bg)
        except Exception:
            None

    def _delete_slide(self, prs: Presentation, slide_index: int) -> None:
        """Delete a slide from the presentation by index (0-based)."""
        slides = list(prs.slides)
        if slide_index < 0 or slide_index >= len(slides):
            log.warning("Cannot delete slide %d: out of range (total: %d)", slide_index, len(slides))
            return

        slide = slides[slide_index]
        rId = None

        # Find the relationship ID for this slide
        for rel in prs.part.rels.values():
            if rel.target_part == slide.part:
                rId = rel.rId
                break

        if rId is None:
            log.warning("Cannot find relationship for slide %d", slide_index)
            return

        # Remove from slide list XML
        pres_elem = prs.part._element
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                 "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
        sldIdLst = pres_elem.find("p:sldIdLst", nsmap)
        if sldIdLst is not None:
            for sldId in list(sldIdLst):
                if sldId.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                ) == rId:
                    sldIdLst.remove(sldId)
                    break

        # Remove the relationship
        try:
            prs.part.rels.pop(rId)
        except Exception:
            None

    def _set_paragraph_text(self, paragraph, new_text: str) -> None:
        """Set paragraph text while preserving formatting of the first run."""
        runs = paragraph.runs
        if not runs:
            # No runs exist — create one
            run = paragraph.add_run()
            run.text = new_text
            return

        # Set text on the first run, clear the rest
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""

    def _apply_formatting(self, paragraph, formatting: dict) -> None:
        """Apply formatting to a paragraph's runs.

        Only applies values that are not None — None means 'keep existing'.
        """
        if not formatting:
            return

        # Alignment
        alignment = formatting.get("alignment")
        if alignment and alignment in _ALIGN_MAP:
            paragraph.alignment = _ALIGN_MAP[alignment]

        # Font properties — apply to all runs
        font_size = formatting.get("font_size_pt")
        bold = formatting.get("bold")
        italic = formatting.get("italic")
        color_hex = formatting.get("color_hex")

        for run in paragraph.runs:
            font = run.font
            if font_size is not None:
                font.size = Pt(font_size)
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color_hex is not None and len(color_hex) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(color_hex)
                except Exception:
                    None

    # ------------------------------------------------------------------
    # Operations mode: apply structured operation list
    # ------------------------------------------------------------------

    def _parse_target_id(self, target_id: str) -> dict:
        """Parse a DOM ID (e.g. 'slide_1_shape_2_para_0') into a target dict for the ops engine."""
        if not target_id:
            return {}
        tgt = {"id": target_id}
        import re
        if m := re.search(r'slide_(\d+)', target_id):
            tgt["slide"] = int(m.group(1))
        if m := re.search(r'shape_(\d+)', target_id):
            tgt["shape_index"] = int(m.group(1))
        
        # New: Sections, Headers, Footers
        if m := re.search(r'section_(\d+)', target_id):
            tgt["section_index"] = int(m.group(1))
        if m := re.search(r'header_(\d+)', target_id):
            tgt["header_index"] = int(m.group(1))
        if m := re.search(r'footer_(\d+)', target_id):
            tgt["footer_index"] = int(m.group(1))
            
        # New: UUIDs
        if m := re.search(r'p_([a-f0-9]{8})', target_id):
            tgt["uid_paragraph"] = m.group(1)
        elif m := re.search(r'img_([a-f0-9]{8})', target_id):
            tgt["uid_image"] = m.group(1)
        elif m := re.search(r't_([a-f0-9]{8})', target_id):
            tgt["uid_table"] = m.group(1)
            
        # image_N must be checked before paragraph_N to avoid partial match
        if re.match(r'^image_\d+$', target_id):
            m = re.search(r'image_(\d+)', target_id)
            if m:
                tgt["image_index"] = int(m.group(1)) - 1
        elif m := re.search(r'paragraph_(\d+)', target_id):
            tgt["paragraph_index"] = int(m.group(1)) - 1
        elif m := re.search(r'table_(\d+)', target_id):
            tgt["table_index"] = int(m.group(1)) - 1
        
        if m := re.search(r'_para_(\d+)', target_id):
            tgt["para_index"] = int(m.group(1))
        # Note: cell_ and _run_ handlers remain below
        if m := re.search(r'cell_(\d+)_(\d+)', target_id):
            tgt["row_index"] = int(m.group(1))
            tgt["col_index"] = int(m.group(2))
        if m := re.search(r'_run_(\d+)', target_id):
            tgt["run_index"] = int(m.group(1))
        return tgt

    def _translate_uids_to_indices(self, doc: Document, tgt: dict) -> None:
        """Translates UID fields back to legacy sequential indices so downstream ops work."""
        if "uid_paragraph" in tgt:
            target_p = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"p_{tgt['uid_paragraph']}":
                    target_p = child
                    break
            if target_p is not None:
                for i, p in enumerate(doc.paragraphs):
                    if p._p == target_p:
                        tgt["paragraph_index"] = i
                        break
        elif "uid_table" in tgt:
            target_t = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"t_{tgt['uid_table']}":
                    target_t = child
                    break
            if target_t is not None:
                for i, t in enumerate(doc.tables):
                    if t._tbl == target_t:
                        tgt["table_index"] = i
                        break
        elif "uid_image" in tgt:
            target_i = None
            for eid, child in self._build_docx_body_index(doc):
                if eid == f"img_{tgt['uid_image']}":
                    target_i = child
                    break
            if target_i is not None:
                img_idx = 0
                for p in doc.paragraphs:
                    from docx.oxml.ns import qn
                    has_drawing = p._p.find(f'.//{qn("w:drawing")}') is not None
                    if has_drawing:
                        if p._p == target_i:
                            tgt["image_index"] = img_idx
                            break
                        img_idx += 1

    def _apply_pptx_operations(
        self,
        source: Path,
        target: Path,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to a PPTX file.

        Returns (any_change_made, summary_list).
        """
        prs = Presentation(source)
        summaries: list[str] = []
        changed = False

        for raw_op in operations:
            target_ids = raw_op.get("target_id")
            if not isinstance(target_ids, list):
                target_ids = [target_ids]
            
            for t_id in target_ids:
                op = dict(raw_op)
                op["target_id"] = t_id
                
                op_type = op.get("op_type", "")
                params = op.get("parameters", {})
                tgt = self._parse_target_id(op.get("target_id") or "")


                try:
                    if op_type == "text_edit":
                        s = self._op_pptx_text_edit(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "text_format":
                        if op.get("target_id") == "all":
                            match_color = params.get("match_color_hex")
                            match_role = params.get("match_role")
                            if match_color or match_role:
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if getattr(shape, "has_text_frame", False):
                                            for para in shape.text_frame.paragraphs:
                                                self._apply_pptx_format_to_para(para, {}, params, shape)
                                        if getattr(shape, "has_table", False):
                                            for row in shape.table.rows:
                                                for cell in row.cells:
                                                    if getattr(cell, "text_frame", None):
                                                        for para in cell.text_frame.paragraphs:
                                                            self._apply_pptx_format_to_para(para, {}, params, shape)
                                summaries.append(f"Formatted all paragraphs matching criteria")
                                changed = True
                            continue

                        s = self._op_pptx_text_format(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "table_op":
                        s = self._op_pptx_table(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "image_op":
                        s = self._op_pptx_image(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "shape_op":
                        s = self._op_pptx_shape(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "theme_op":
                        s = self._op_pptx_theme(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "slide_op":
                        s = self._op_pptx_slide(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "chart_op":
                        s = self._op_pptx_chart(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "ai_design_op":
                        s = self._op_pptx_ai_design(prs, tgt, params)
                        if s:
                            summaries.append(s); changed = True

                    elif op_type == "needs_image":
                        # Handled upstream — shouldn't reach here
                        None

                    else:
                        log.warning("Unknown op_type in operations list: %r", op_type)

                except Exception as exc:
                    log.exception("Failed to apply operation %r: %s", op_type, exc)

        if changed:
            prs.save(target)
        else:
            # Still save a copy so the pipeline has a version file
            import shutil
            shutil.copy2(source, target)

        return changed, summaries

    def _resolve_target_paras(self, doc: Document, tgt: dict) -> list:
        """Resolves the target ID dictionary to a list of paragraph elements."""
        if "table_index" in tgt and "row_index" in tgt and "col_index" in tgt:
            try:
                table = doc.tables[tgt["table_index"]]
                cell = table.cell(tgt["row_index"], tgt["col_index"])
                return cell.paragraphs
            except Exception:
                None
        if "header_index" in tgt and "section_index" in tgt:
            try:
                return doc.sections[tgt["section_index"]].header.paragraphs
            except Exception:
                None
        if "footer_index" in tgt and "section_index" in tgt:
            try:
                return doc.sections[tgt["section_index"]].footer.paragraphs
            except Exception:
                None
        return doc.paragraphs

    def _build_legacy_id_map(self, doc: Document) -> dict[str, str]:
        """Maps legacy outline IDs (paragraph_1) to new stable UUIDs (p_abcdef)."""
        # Ensure all elements have UIDs stamped in memory before mapping!
        self._build_docx_body_index(doc)
        
        from docx.oxml.ns import qn
        WNS_P = qn('w:p')
        WNS_TBL = qn('w:tbl')
        WNS_DRAWING = qn('w:drawing')
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        legacy_map = {}
        para_counter = 0
        table_counter = 0
        image_counter = 0

        for child in doc.element.body:
            tag = child.tag
            uid = child.get(uid_attr)
            
            if tag in (WNS_P, WNS_TBL):
                if tag == WNS_P:
                    has_drawing = child.find(f'.//{WNS_DRAWING}') is not None
                    if has_drawing:
                        image_counter += 1
                        legacy_id = f"image_{image_counter}"
                        primary_id = f"img_{uid}" if uid else legacy_id
                    else:
                        para_counter += 1
                        legacy_id = f"paragraph_{para_counter}"
                        primary_id = f"p_{uid}" if uid else legacy_id
                else:
                    table_counter += 1
                    legacy_id = f"table_{table_counter}"
                    primary_id = f"t_{uid}" if uid else legacy_id
                
                if uid:
                    legacy_map[legacy_id] = primary_id

        return legacy_map

    def _translate_legacy_ids(self, doc: Document, params: dict) -> None:
        """Translates legacy IDs in operation parameters to new UUIDs."""
        legacy_map = self._build_legacy_id_map(doc)
        keys_to_check = ["start_id", "end_id", "before_id", "after_id", "section_a_start_id", "section_a_end_id", "section_b_start_id", "section_b_end_id", "_raw_target_id", "target_id"]
        for k in keys_to_check:
            if k in params and isinstance(params[k], str) and params[k] in legacy_map:
                params[k] = legacy_map[params[k]]
            elif k in params and isinstance(params[k], list):
                params[k] = [legacy_map.get(pid, pid) for pid in params[k]]

    def _apply_docx_operations(
        self,
        source: Path,
        target: Path,
        operations: list[dict],
    ) -> tuple[bool, list[str]]:
        """Apply a list of structured operations to a DOCX file."""
        doc = Document(source)
        summaries: list[str] = []
        changed = False

        for raw_op in operations:
            op_type = raw_op.get("op_type", "")
            
            # Translate legacy IDs in the operation itself and its parameters
            self._translate_legacy_ids(doc, raw_op)
            if "parameters" in raw_op and isinstance(raw_op["parameters"], dict):
                self._translate_legacy_ids(doc, raw_op["parameters"])
            
            # Structural ops do not loop over target_ids
            if op_type in {"list_op", "layout_op", "theme_op", "ai_design_op", "meta_op", "style_op", "find_replace", "slide_op"}:
                params = dict(raw_op.get("parameters", {}))
                if "target_id" in raw_op:
                    params["_raw_target_id"] = raw_op.get("target_id")

                try:
                    summary = ""
                    if op_type == "list_op":
                        summary = self._op_docx_list(doc, params)
                    elif op_type == "layout_op":
                        summary = self._op_docx_layout(doc, params)
                    elif op_type == "theme_op":
                        summary = self._op_docx_theme(doc, params)
                    elif op_type == "ai_design_op":
                        summary = self._op_docx_ai_design(doc, params)
                    elif op_type == "meta_op":
                        summary = apply_metadata(doc, params)
                    elif op_type == "style_op":
                        summary = apply_global_style(doc, params)
                    elif op_type == "find_replace":
                        tgt = self._parse_target_id(raw_op.get("target_id") or "all")
                        summary = self._op_docx_find_replace(doc, tgt, params)
                    elif op_type == "slide_op":
                        summary = "Slide operations are only supported for PPTX files."
                    
                    if summary:
                        summaries.append(summary)
                        changed = True
                except Exception as exc:
                    log.exception("Failed to apply structural operation %r: %s", op_type, exc)
                continue

            target_ids = raw_op.get("target_id")
            if not isinstance(target_ids, list):
                target_ids = [target_ids]
            
            for t_id in target_ids:
                op = dict(raw_op)
                op["target_id"] = t_id
                
                op_type = op.get("op_type", "")
                params = op.get("parameters", {})
                tgt = self._parse_target_id(op.get("target_id") or "")
                self._translate_uids_to_indices(doc, tgt)

                try:
                    if op_type == "text_edit":
                        para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
                        new_text = params.get("new_text", "")
                        if para_idx is not None and new_text:
                            paras = self._resolve_target_paras(doc, tgt)

                            if paras and 0 <= para_idx < len(paras):
                                self._apply_run_aware_replacement(paras[para_idx], new_text, params)
                                summaries.append(f"Rewrote paragraph {para_idx}")
                                changed = True

                    elif op_type == "text_format":
                        if op.get("target_id") == "all":
                            match_color = params.get("match_color_hex")
                            match_role = params.get("match_role")
                            for p in doc.paragraphs:
                                self._apply_docx_format(p, {}, params)
                            for tbl in doc.tables:
                                for row in tbl.rows:
                                    for cell in row.cells:
                                        for p in cell.paragraphs:
                                            self._apply_docx_format(p, {}, params)
                            summaries.append(f"Formatted all paragraphs")
                            changed = True
                            continue

                        para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
                        if para_idx is not None:
                            paras = self._resolve_target_paras(doc, tgt)

                            if paras and 0 <= para_idx < len(paras):
                                self._apply_docx_format(paras[para_idx], tgt, params)
                                summaries.append(f"Formatted paragraph {para_idx}")
                                changed = True
                        elif "table_index" in tgt:
                            table_idx = tgt.get("table_index")
                            if table_idx is not None and 0 <= table_idx < len(doc.tables):
                                table = doc.tables[table_idx]
                                for row in table.rows:
                                    for cell in row.cells:
                                        for p in cell.paragraphs:
                                            self._apply_docx_format(p, {}, params)
                                summaries.append(f"Formatted text in table {table_idx + 1}")
                                changed = True

                    elif op_type == "table_op":
                        summary = self._op_docx_table(doc, tgt, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "image_op":
                        summary = self._op_docx_image(doc, tgt, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "section_op":
                        target_id = op.get("target_id")
                        summary = apply_section_formatting(doc, target_id, params)
                        if summary:
                            summaries.append(summary)
                            changed = True

                    elif op_type == "needs_image":
                        None

                except Exception as exc:
                    log.exception("DOCX op %r failed: %s", op_type, exc)

        if changed:
            doc.save(target)
        else:
            import shutil
            shutil.copy2(source, target)

        return changed, summaries

    def _op_docx_theme(self, doc, params: dict) -> str:
        action = params.get("action")
        if action == "set_bg_color":
            bg_color = str(params.get("bg_color_hex", "FFFFFF")).strip().lstrip("#")
            from docx.oxml import parse_xml
            from docx.oxml.ns import nsdecls
            background = parse_xml(r'<w:background {} w:color="{}"/>'.format(nsdecls('w'), bg_color))
            doc.element.insert(0, background)
            doc.settings.element.append(parse_xml(r'<w:displayBackgroundShape {}/>'.format(nsdecls('w'))))
            return f"Set document background color to #{bg_color}"
        elif action == "set_margins":
            inches = params.get("margin_inches")
            if inches:
                from docx.shared import Inches
                for section in doc.sections:
                    section.left_margin = Inches(inches)
                    section.right_margin = Inches(inches)
                    section.top_margin = Inches(inches)
                    section.bottom_margin = Inches(inches)
                return f"Set document margins to {inches} inches"
        elif action == "add_page_numbers":
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            for section in doc.sections:
                footer = section.footer
                # Clear existing footer paragraphs to avoid stacking
                for p in list(footer.paragraphs):
                    p._p.getparent().remove(p._p)
                p = footer.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run("Page ")
                fldChar1 = OxmlElement('w:fldChar')
                fldChar1.set(qn('w:fldCharType'), 'begin')
                instrText = OxmlElement('w:instrText')
                instrText.set(qn('xml:space'), 'preserve')
                instrText.text = "PAGE"
                fldChar2 = OxmlElement('w:fldChar')
                fldChar2.set(qn('w:fldCharType'), 'separate')
                fldChar3 = OxmlElement('w:fldChar')
                fldChar3.set(qn('w:fldCharType'), 'end')
                run._r.append(fldChar1)
                run._r.append(instrText)
                run._r.append(fldChar2)
                run._r.append(fldChar3)
            return "Added page numbers to footers"
        elif action == "apply_theme_colors":
            accent_colors = params.get("accent_colors", [])
            if not accent_colors:
                return "No theme colors provided"
            color1 = str(accent_colors[0]).lstrip('#')
            color2 = str(accent_colors[1]).lstrip('#') if len(accent_colors) > 1 else color1
            
            from docx.shared import RGBColor as DRGBColor
            def _hex_to_rgb(hx):
                if len(hx) == 6:
                    try:
                        return DRGBColor(int(hx[0:2],16), int(hx[2:4],16), int(hx[4:6],16))
                    except ValueError:
                        None
                return None

            c1 = _hex_to_rgb(color1)
            c2 = _hex_to_rgb(color2)
            
            if c1:
                for para in doc.paragraphs:
                    if para.style and para.style.name and para.style.name.lower().startswith("heading"):
                        for run in para.runs:
                            run.font.color.rgb = c1
            if c2:
                for tbl in doc.tables:
                    if len(tbl.rows) > 0:
                        for cell in tbl.rows[0].cells:
                            for p in cell.paragraphs:
                                for run in p.runs:
                                    run.font.color.rgb = c2
            return f"Applied theme colors {accent_colors}"
        return ""

    def _extract_list_info(self, para, doc) -> dict | None:
        """Return numbering metadata for a list paragraph, or None if not a list item."""
        from docx.oxml.ns import qn
        pPr = para._p.find(qn('w:pPr'))
        if pPr is None:
            return None
        numPr = pPr.find(qn('w:numPr'))
        if numPr is None:
            return None
        numId_el = numPr.find(qn('w:numId'))
        ilvl_el = numPr.find(qn('w:ilvl'))
        if numId_el is None:
            return None
        num_id = int(numId_el.get(qn('w:val'), 0))
        if num_id == 0:  # numId=0 means "remove numbering"
            return None
        ilvl = int(ilvl_el.get(qn('w:val'), 0)) if ilvl_el is not None else 0

        num_fmt = "bullet"
        list_type = "bullet"
        lvl_text_val = ""
        try:
            numbering_part = doc.part.numbering_part
            if numbering_part is not None:
                nb = numbering_part._element
                # Locate <w:num w:numId="N">
                num_el = nb.find(
                    f'.//{qn("w:num")}[@{qn("w:numId")}="{num_id}"]'
                )
                if num_el is not None:
                    abs_ref = num_el.find(qn('w:abstractNumId'))
                    if abs_ref is not None:
                        abs_id = abs_ref.get(qn('w:val'), '0')
                        abs_num = nb.find(
                            f'.//{qn("w:abstractNum")}[@{qn("w:abstractNumId")}="{abs_id}"]'
                        )
                        if abs_num is not None:
                            lvl = abs_num.find(
                                f'{qn("w:lvl")}[@{qn("w:ilvl")}="{ilvl}"]'
                            )
                            if lvl is None:
                                # search without namespace prefix matching
                                for lv in abs_num:
                                    if lv.get(qn('w:ilvl')) == str(ilvl):
                                        lvl = lv
                                        break
                            if lvl is not None:
                                nfmt_el = lvl.find(qn('w:numFmt'))
                                if nfmt_el is not None:
                                    num_fmt = nfmt_el.get(qn('w:val'), 'bullet')
                                lvl_text_el = lvl.find(qn('w:lvlText'))
                                if lvl_text_el is not None:
                                    lvl_text_val = lvl_text_el.get(qn('w:val'), '')
        except Exception:
            None

        if num_fmt in ('decimal', 'lowerRoman', 'upperRoman', 'lowerLetter', 'upperLetter'):
            list_type = 'numbered'
        elif lvl_text_val in ('\u2610', '\u2611', '\u2612', '\u25a1', '\u25a0',
                              '\u2714', '\u2718', '\u2013\u2013', '\u2610 '):
            list_type = 'checklist'
        else:
            list_type = 'bullet'

        return {
            'num_id': num_id,
            'ilvl': ilvl,
            'num_fmt': num_fmt,
            'list_type': list_type,
            'lvl_text': lvl_text_val,
        }

    def _renumber_headings_in_doc(self, doc) -> int:
        """Re-sequence numbered heading text after structural operations (swap, move).

        Detects headings whose visible text starts with a numeric prefix like "1. " or "3. "
        and renumbers them sequentially in document order at each heading level.
        Returns the number of headings that were actually renumbered.
        """
        import re as _re
        _NUMBERED_PREFIX = _re.compile(r'^(\d+)\.\s+(.+)$', _re.DOTALL)

        # Collect all heading paragraphs in body order
        heading_paras = []
        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            if style_name.lower().startswith("heading"):
                m = _re.search(r'(\d+)$', style_name)
                level = int(m.group(1)) if m else 1
                heading_paras.append((level, para))

        if not heading_paras:
            return 0

        # Build per-level counters, reset child counters when parent increments
        level_counters: dict[int, int] = {}
        changed = 0

        for level, para in heading_paras:
            # Reset all deeper level counters when we encounter this level
            deeper_keys = [k for k in level_counters if k > level]
            for k in deeper_keys:
                del level_counters[k]

            level_counters[level] = level_counters.get(level, 0) + 1
            current_num = level_counters[level]

            current_text = para.text.strip()
            pm = _NUMBERED_PREFIX.match(current_text)
            if pm:
                existing_num = int(pm.group(1))
                title_part = pm.group(2)
                if existing_num != current_num:
                    new_text = f"{current_num}. {title_part}"
                    self._apply_run_aware_replacement(para, new_text)
                    changed += 1

        return changed

    def _build_docx_body_index(self, doc) -> list[tuple[str, object]]:
        """Walk doc.element.body in order and return a flat list of (id, xml_element) pairs.

        IDs are resolved from the `deid:uid` custom XML attribute.
        """
        from docx.oxml.ns import qn
        WNS_P = qn('w:p')
        WNS_TBL = qn('w:tbl')
        WNS_DRAWING = qn('w:drawing')
        
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        other_counter = 1
        result = []
        for child in doc.element.body:
            tag = child.tag
            if tag in (WNS_P, WNS_TBL):
                uid = child.get(uid_attr)
                if not uid:
                    # Fallback for newly inserted elements that missed stamping
                    import uuid
                    uid = uuid.uuid4().hex[:8]
                    child.set(uid_attr, uid)
                
                if tag == WNS_P:
                    has_drawing = child.find(f'.//{WNS_DRAWING}') is not None
                    if has_drawing:
                        result.append((f"img_{uid}", child))
                    else:
                        result.append((f"p_{uid}", child))
                else:
                    result.append((f"t_{uid}", child))
            else:
                result.append((f"body_other_{other_counter}", child))
                other_counter += 1
        return result

    def _op_docx_layout(self, doc, params: dict) -> str:
        action = params.get("action")

        if action == "move_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            if not start_id or not end_id or not (before_id or after_id):
                return "move_block requires start_id, end_id, and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)
            before_pos = id_to_pos.get(before_id, -1) if before_id else -1
            after_pos = id_to_pos.get(after_id, -1) if after_id else -1

            if start_pos == -1 or end_pos == -1 or (before_id and before_pos == -1) or (after_id and after_pos == -1):
                missing = []
                if start_pos == -1: missing.append(("start_id", start_id))
                if end_pos == -1: missing.append(("end_id", end_id))
                if before_id and before_pos == -1: missing.append(("before_id", before_id))
                if after_id and after_pos == -1: missing.append(("after_id", after_id))
                return f"Failed to find required block boundaries: {missing}. Available IDs: {[eid for eid, _ in body_index[:30]]}"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            if start_pos <= before_pos <= end_pos or start_pos <= after_pos <= end_pos:
                return f"Moved {end_pos - start_pos + 1} block(s) (no-op because target is inside the moved block)"

            # Collect the xml elements to move
            elements_to_move = [body_index[i][1] for i in range(start_pos, end_pos + 1)]

            # Remove from current position
            for xml_el in elements_to_move:
                xml_el.getparent().remove(xml_el)
                
            # Insert at new position
            if before_pos != -1:
                before_xml_el = body_index[before_pos][1]
                for xml_el in elements_to_move:
                    before_xml_el.addprevious(xml_el)
                renumbered = self._renumber_headings_in_doc(doc)
                suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
                return f"Moved {len(elements_to_move)} block(s) ('{start_id}'→'{end_id}') before '{before_id}'{suffix}"
            elif after_pos != -1:
                after_xml_el = body_index[after_pos][1]
                for xml_el in reversed(elements_to_move):
                    after_xml_el.addnext(xml_el)
                renumbered = self._renumber_headings_in_doc(doc)
                suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
                return f"Moved {len(elements_to_move)} block(s) ('{start_id}'→'{end_id}') after '{after_id}'{suffix}"

        elif action == "insert_page_break":
            # Accept before_id, start_id (fallback), or after_id
            before_id = params.get("before_id") or params.get("start_id")
            after_id_pb = params.get("after_id")
            if not before_id and not after_id_pb:
                return "insert_page_break requires before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            # Build a paragraph containing a hard page break
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            pb_p = OxmlElement('w:p')
            r = OxmlElement('w:r')
            br = OxmlElement('w:br')
            br.set(qn('w:type'), 'page')
            r.append(br)
            pb_p.append(r)

            if before_id:
                before_pos = id_to_pos.get(before_id, -1)
                if before_pos == -1:
                    return f"insert_page_break: element '{before_id}' not found"
                body_index[before_pos][1].addprevious(pb_p)
                return f"Inserted page break before '{before_id}'"
            else:
                after_pos = id_to_pos.get(after_id_pb, -1)
                if after_pos == -1:
                    return f"insert_page_break: element '{after_id_pb}' not found"
                body_index[after_pos][1].addnext(pb_p)
                return f"Inserted page break after '{after_id_pb}'"

        elif action == "set_columns":
            num_cols = int(params.get("num_columns", 2))
            gap_inches = float(params.get("column_gap_inches", 0.5))
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            body = doc.element.body
            sectPr = body.find(qn("w:sectPr"))
            if sectPr is None:
                sectPr = OxmlElement("w:sectPr")
                body.append(sectPr)
            # Remove existing cols element
            existing_cols = sectPr.find(qn("w:cols"))
            if existing_cols is not None:
                sectPr.remove(existing_cols)
            cols = OxmlElement("w:cols")
            cols.set(qn("w:num"), str(num_cols))
            gap_twips = int(gap_inches * 1440)  # 1 inch = 1440 twips
            cols.set(qn("w:space"), str(gap_twips))
            cols.set(qn("w:equalWidth"), "1")
            sectPr.append(cols)
            return f"Set document to {num_cols}-column layout with {gap_inches}" + '" gap'

        elif action == "remove_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            if not start_id or not end_id:
                return "remove_block requires start_id and end_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)

            if start_pos == -1 or end_pos == -1:
                return f"Failed to find block boundaries. start={start_pos}, end={end_pos}"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            elements = [body_index[i][1] for i in range(start_pos, end_pos + 1)]
            for xml_el in elements:
                xml_el.getparent().remove(xml_el)
            return f"Removed {len(elements)} block(s) ('{start_id}'→'{end_id}')"

        elif action == "duplicate_block":
            start_id = params.get("start_id")
            end_id = params.get("end_id")
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            
            if not start_id or not end_id or (not before_id and not after_id):
                return "duplicate_block requires start_id, end_id, and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            start_pos = id_to_pos.get(start_id, -1)
            end_pos = id_to_pos.get(end_id, -1)
            
            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if start_pos == -1 or end_pos == -1 or target_xml_el is None:
                return "Failed to find block boundaries or target id for duplicate_block"

            if start_pos > end_pos:
                start_pos, end_pos = end_pos, start_pos

            elements = [body_index[i][1] for i in range(start_pos, end_pos + 1)]
            import copy
            import uuid
            cloned_elements = [copy.deepcopy(el) for el in elements]
            
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            from docx.oxml.ns import qn
            WNS_P = qn('w:p')
            WNS_TBL = qn('w:tbl')
            for clone in cloned_elements:
                if clone.tag in (WNS_P, WNS_TBL):
                    clone.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in clone.iter(WNS_P):
                    if nested != clone:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in clone.iter(WNS_TBL):
                    if nested != clone:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
            
            current_target = target_xml_el
            for xml_el in cloned_elements:
                if insert_before:
                    current_target.addprevious(xml_el)
                else:
                    current_target.addnext(xml_el)
                    current_target = xml_el
            
            return f"Duplicated {len(elements)} block(s)"

        elif action == "swap_sections":
            a_start = params.get("section_a_start_id")
            a_end = params.get("section_a_end_id")
            b_start = params.get("section_b_start_id")
            b_end = params.get("section_b_end_id")

            if not all([a_start, a_end, b_start, b_end]):
                return "swap_sections requires start/end ids for both sections"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            a_start_pos = id_to_pos.get(a_start, -1)
            a_end_pos = id_to_pos.get(a_end, -1)
            b_start_pos = id_to_pos.get(b_start, -1)
            b_end_pos = id_to_pos.get(b_end, -1)

            if -1 in (a_start_pos, a_end_pos, b_start_pos, b_end_pos):
                return "Failed to find block boundaries for swap_sections"

            if a_start_pos > a_end_pos: a_start_pos, a_end_pos = a_end_pos, a_start_pos
            if b_start_pos > b_end_pos: b_start_pos, b_end_pos = b_end_pos, b_start_pos

            if a_start_pos > b_start_pos:
                # Ensure A is always before B for easier logic
                a_start_pos, a_end_pos, b_start_pos, b_end_pos = b_start_pos, b_end_pos, a_start_pos, a_end_pos

            if a_end_pos >= b_start_pos:
                return "Cannot swap overlapping or adjacent-intersecting sections"

            # Collect elements
            a_elements = [body_index[i][1] for i in range(a_start_pos, a_end_pos + 1)]
            b_elements = [body_index[i][1] for i in range(b_start_pos, b_end_pos + 1)]

            # Remove all from parent
            for el in a_elements: el.getparent().remove(el)
            for el in b_elements: el.getparent().remove(el)
            
            # Re-insert B where A was.
            if a_start_pos == 0:
                # B goes to start of body
                for el in reversed(b_elements):
                    doc.element.body.insert(0, el)
            else:
                pre_a_el = body_index[a_start_pos - 1][1]
                for el in reversed(b_elements):
                    pre_a_el.addnext(el)
                    
            # Re-insert A where B was.
            if b_start_pos - 1 >= a_start_pos and b_start_pos - 1 <= a_end_pos:
                # They were strictly adjacent. A goes right after B's new position.
                last_b_el = b_elements[-1]
                for el in reversed(a_elements):
                    last_b_el.addnext(el)
            else:
                pre_b_el = body_index[b_start_pos - 1][1]
                for el in reversed(a_elements):
                    pre_b_el.addnext(el)

            renumbered = self._renumber_headings_in_doc(doc)
            suffix = f", renumbered {renumbered} heading(s)" if renumbered else ""
            return f"Swapped sections ({len(a_elements)} blocks and {len(b_elements)} blocks){suffix}"

        elif action == "insert_block":
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            data = params.get("data", [])
            
            if not data or (not before_id and not after_id):
                return "insert_block requires data array and either before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if target_xml_el is None:
                return "Failed to find target id for insert_block"

            # Create paragraphs / tables from data and collect XML elements
            new_elements = []
            for item in data:
                text = item.get("text", "")
                role = item.get("role", "body")

                if role == "heading":
                    lvl = item.get("heading_level", 1)
                    p = doc.add_heading(text, level=lvl)
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

                elif role == "bullet_point":
                    from docx.oxml import OxmlElement as _OXe
                    from docx.oxml.ns import qn as _qne
                    try:
                        p = doc.add_paragraph(text, style="List Paragraph")
                    except Exception:
                        p = doc.add_paragraph(text)
                    pPr = p._p.get_or_add_pPr()
                    numPr = _OXe("w:numPr")
                    ilvl = _OXe("w:ilvl"); ilvl.set(_qne("w:val"), "0")
                    numId_el = _OXe("w:numId"); numId_el.set(_qne("w:val"), "1")
                    numPr.append(ilvl); numPr.append(numId_el)
                    pPr.append(numPr)
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

                elif role == "table":
                    # Inline table with headers + rows (used by ToC enricher etc.)
                    headers = item.get("headers", [])
                    rows = item.get("rows", [])
                    style_hint = item.get("style", "")
                    num_cols = max(
                        len(headers),
                        max((len(r) for r in rows), default=0),
                        1,
                    )
                    tbl = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    tbl_xml = tbl._tbl
                    tbl_xml.getparent().remove(tbl_xml)

                    from docx.oxml.ns import qn as _qnt
                    from docx.oxml import OxmlElement as _OXt

                    def _cell_text(cell, val: str, bold: bool = False) -> None:
                        cell.text = str(val)
                        if bold:
                            for run in cell.paragraphs[0].runs:
                                run.bold = True

                    # Dynamic bookmark ID allocation & deduplication
                    from docx.oxml.ns import qn as _qn_bmk
                    existing_bmk_ids = []
                    for bmk in doc.element.body.iter(_qn_bmk('w:bookmarkStart')):
                        try:
                            val = int(bmk.get(_qn_bmk('w:id'), '0'))
                            existing_bmk_ids.append(val)
                        except ValueError:
                            None
                    next_bmk_id = max(existing_bmk_ids, default=0) + 1

                    def get_next_bmk_id():
                        nonlocal next_bmk_id
                        curr = next_bmk_id
                        next_bmk_id += 1
                        return curr

                    # Helper for bookmark insertion in target heading paragraphs
                    def _ensure_heading_bookmark(doc_ref, bmk_name: str, heading_id: str | None, ri: int, heading_text: str = ""):
                        from docx.oxml import parse_xml
                        from docx.oxml.ns import nsdecls, qn
                        uid_attr = f'{{{CUSTOM_NS}}}uid'

                        for bmk in doc_ref.element.body.iter(qn('w:bookmarkStart')):
                            if bmk.get(qn('w:name')) == bmk_name:
                                return

                        target_p = None
                        if heading_id:
                            for p in doc_ref.paragraphs:
                                if p._p.get(uid_attr) == heading_id:
                                    target_p = p
                                    break

                        if target_p is None:
                            headings = [p for p in doc_ref.paragraphs if p.style and p.style.name and "Heading" in p.style.name]
                            if 0 <= ri < len(headings):
                                target_p = headings[ri]

                        if target_p is None and heading_text.strip():
                            clean_heading = heading_text.strip()
                            for p in doc_ref.paragraphs:
                                if clean_heading in p.text.strip():
                                    target_p = p
                                    break

                        if target_p is not None:
                            bmk_id = get_next_bmk_id()
                            bmk_start = parse_xml(r'<w:bookmarkStart %s w:id="%d" w:name="%s"/>' % (nsdecls('w'), bmk_id, bmk_name))
                            bmk_end = parse_xml(r'<w:bookmarkEnd %s w:id="%d"/>' % (nsdecls('w'), bmk_id))
                            target_p._p.insert(0, bmk_start)
                            target_p._p.append(bmk_end)

                    # Header row
                    hdr_row = tbl.rows[0]
                    for ci, hdr in enumerate(headers[:num_cols]):
                        _cell_text(hdr_row.cells[ci], hdr, bold=True)

                    # Data rows
                    has_pageref = False
                    for ri, row_data in enumerate(rows):
                        tbl_row = tbl.rows[ri + 1]
                        for ci, cell_val in enumerate(row_data[:num_cols]):
                            cell = tbl_row.cells[ci]
                            if isinstance(cell_val, dict) and "pageref" in cell_val:
                                has_pageref = True
                                bmk_name = cell_val["pageref"]
                                page_str = str(cell_val.get("page", "1"))
                                heading_id = cell_val.get("heading_id")
                                heading_text = str(row_data[0]) if row_data else ""

                                _ensure_heading_bookmark(doc, bmk_name, heading_id, ri, heading_text)

                                from docx.oxml import parse_xml
                                from docx.oxml.ns import nsdecls
                                from xml.sax.saxutils import escape as xml_escape

                                page_esc = xml_escape(page_str)
                                bmk_esc = xml_escape(bmk_name)

                                fld_xml = (
                                    r'<w:p %s>'
                                    r'  <w:r>'
                                    r'    <w:fldChar w:fldCharType="begin"/>'
                                    r'  </w:r>'
                                    r'  <w:r>'
                                    r'    <w:instrText xml:space="preserve"> PAGEREF %s \h </w:instrText>'
                                    r'  </w:r>'
                                    r'  <w:r>'
                                    r'    <w:fldChar w:fldCharType="separate"/>'
                                    r'  </w:r>'
                                    r'  <w:r>'
                                    r'    <w:t>%s</w:t>'
                                    r'  </w:r>'
                                    r'  <w:r>'
                                    r'    <w:fldChar w:fldCharType="end"/>'
                                    r'  </w:r>'
                                    r'</w:p>' % (nsdecls('w'), bmk_esc, page_esc)
                                )
                                cell.paragraphs[0]._p.getparent().replace(cell.paragraphs[0]._p, parse_xml(fld_xml))
                            else:
                                _cell_text(cell, cell_val)

                    if has_pageref:
                        _enable_update_fields(doc)

                    # ToC style: remove all borders for a clean look
                    if style_hint == "toc":
                        tblPr = tbl_xml.find(_qnt("w:tblPr"))
                        if tblPr is None:
                            tblPr = _OXt("w:tblPr")
                            tbl_xml.insert(0, tblPr)
                        tblBorders = _OXt("w:tblBorders")
                        for bname in ("top", "left", "bottom", "right", "insideH", "insideV"):
                            b = _OXt(f"w:{bname}")
                            b.set(_qnt("w:val"), "none")
                            tblBorders.append(b)
                        existing = tblPr.find(_qnt("w:tblBorders"))
                        if existing is not None:
                            tblPr.remove(existing)
                        tblPr.append(tblBorders)

                    new_elements.append(tbl_xml)

                elif role == "toc_field":
                    from docx.oxml import parse_xml
                    from docx.oxml.ns import nsdecls
                    from xml.sax.saxutils import escape as xml_escape

                    field_text = xml_escape(str(item.get("field_text", r'TOC \o "1-3" \h \z \u')))
                    entries = item.get("entries", [])

                    # Compute usable page width for right tab stop position (default 9360 twips = 6.5 in)
                    try:
                        sec = doc.sections[0]
                        p_width = sec.page_width.twips if sec.page_width else 12240
                        l_margin = sec.left_margin.twips if sec.left_margin else 1440
                        r_margin = sec.right_margin.twips if sec.right_margin else 1440
                        tab_pos = p_width - l_margin - r_margin
                    except Exception:
                        tab_pos = 9360

                    cached_runs_xml = ""
                    if entries:
                        for idx, entry in enumerate(entries):
                            entry_text = xml_escape(str(entry.get("text", "")))
                            page_str = xml_escape(str(entry.get("page", idx + 1)))
                            cached_runs_xml += f'<w:r><w:t xml:space="preserve">{entry_text}</w:t><w:tab/><w:t>{page_str}</w:t></w:r>'
                            if idx < len(entries) - 1:
                                cached_runs_xml += '<w:r><w:br/></w:r>'
                    else:
                        cached_runs_xml = '<w:r><w:t xml:space="preserve">Table of Contents entries will update on open.</w:t></w:r>'

                    fld_xml = (
                        r'<w:p %s>'
                        r'  <w:pPr>'
                        r'    <w:tabs>'
                        r'      <w:tab w:val="right" w:leader="dot" w:pos="%d"/>'
                        r'    </w:tabs>'
                        r'  </w:pPr>'
                        r'  <w:r>'
                        r'    <w:fldChar w:fldCharType="begin"/>'
                        r'  </w:r>'
                        r'  <w:r>'
                        r'    <w:instrText xml:space="preserve"> %s </w:instrText>'
                        r'  </w:r>'
                        r'  <w:r>'
                        r'    <w:fldChar w:fldCharType="separate"/>'
                        r'  </w:r>'
                        r'  %s'
                        r'  <w:r>'
                        r'    <w:fldChar w:fldCharType="end"/>'
                        r'  </w:r>'
                        r'</w:p>' % (nsdecls('w'), tab_pos, field_text, cached_runs_xml)
                    )
                    p_toc = parse_xml(fld_xml)
                    new_elements.append(p_toc)
                    _enable_update_fields(doc)

                else:
                    # Default: plain body paragraph with optional bold
                    p = doc.add_paragraph(text)
                    if item.get("bold"):
                        for run in p.runs:
                            run.bold = True
                    xml_el = p._p
                    xml_el.getparent().remove(xml_el)
                    new_elements.append(xml_el)

            current_target = target_xml_el
            
            import uuid
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            from docx.oxml.ns import qn
            WNS_P = qn('w:p')
            WNS_TBL = qn('w:tbl')
            for xml_el in new_elements:
                if xml_el.tag in (WNS_P, WNS_TBL):
                    xml_el.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in xml_el.iter(WNS_P):
                    if nested != xml_el:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                for nested in xml_el.iter(WNS_TBL):
                    if nested != xml_el:
                        nested.set(uid_attr, uuid.uuid4().hex[:8])
                if insert_before:
                    current_target.addprevious(xml_el)
                else:
                    current_target.addnext(xml_el)
                    current_target = xml_el

            return f"Inserted {len(new_elements)} new block(s)"

        elif action == "insert_toc":
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            
            if not before_id and not after_id:
                return "insert_toc requires before_id or after_id"

            body_index = self._build_docx_body_index(doc)
            id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}

            target_xml_el = None
            insert_before = True
            if before_id:
                pos = id_to_pos.get(before_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
            elif after_id:
                pos = id_to_pos.get(after_id, -1)
                if pos != -1:
                    target_xml_el = body_index[pos][1]
                    insert_before = False

            if target_xml_el is None:
                return "Failed to find target id for insert_toc"

            from docx.oxml import parse_xml
            from docx.oxml.ns import nsdecls
            
            toc_xml = f"""
            <w:sdt {nsdecls('w')}>
              <w:sdtPr>
                <w:docPartObj>
                  <w:docPartGallery w:val="Table of Contents"/>
                  <w:docPartUnique/>
                </w:docPartObj>
              </w:sdtPr>
              <w:sdtContent>
                <w:p>
                  <w:pPr><w:pStyle w:val="TOCHeading"/></w:pPr>
                  <w:r><w:t>Table of Contents</w:t></w:r>
                </w:p>
                <w:p>
                  <w:pPr>
                    <w:tabs>
                      <w:tab w:val="right" w:leader="dot" w:pos="9360"/>
                    </w:tabs>
                  </w:pPr>
                  <w:r><w:fldChar w:fldCharType="begin"/></w:r>
                  <w:r><w:instrText xml:space="preserve"> TOC \\o "1-3" \\h \\z \\u </w:instrText></w:r>
                  <w:r><w:fldChar w:fldCharType="separate"/></w:r>
                  <w:r><w:fldChar w:fldCharType="end"/></w:r>
                </w:p>
              </w:sdtContent>
            </w:sdt>
            """
            toc_elem = parse_xml(toc_xml)
            
            import uuid
            uid_attr = f'{{{CUSTOM_NS}}}uid'
            for p in toc_elem.iter(qn('w:p')):
                p.set(uid_attr, uuid.uuid4().hex[:8])

            if insert_before:
                target_xml_el.addprevious(toc_elem)
            else:
                target_xml_el.addnext(toc_elem)
                
            return "Inserted Table of Contents"

        return ""

    # ------------------------------------------------------------------
    # DOCX list operations
    # ------------------------------------------------------------------

    def _create_docx_numbering(self, doc, list_type: str, bullet_char: str = "") -> int:
        """Create a new <w:num> entry in numbering.xml and return its numId.

        list_type: 'bullet' | 'numbered' | 'checklist'
        bullet_char: custom bullet character (used for 'bullet' and 'checklist')
        Returns the new numId (int).
        """
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn, nsmap
        import copy

        BULLET_CHAR = bullet_char or '\u2022'  # default: •
        CHECKLIST_CHAR = bullet_char or '\u2610'  # ☐

        # Ensure numbering part exists
        try:
            nb_part = doc.part.numbering_part
        except Exception:
            # Create a minimal numbering part if missing
            from docx.opc.part import Part
            from docx.opc.packuri import PackURI
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            nb_xml = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:numbering xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas"'
                ' xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
                ' xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml">'
                '</w:numbering>'
            )
            nb_part = Part(
                PackURI('/word/numbering.xml'),
                'application/vnd.openxmlformats-officedocument.wordprocessingml.numbering+xml',
                nb_xml.encode(),
            )
            doc.part.relate_to(nb_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering')

        nb = doc.part.numbering_part._element

        # Find the next available abstractNumId and numId
        existing_abs = nb.findall(qn('w:abstractNum'))
        next_abs_id = max((int(a.get(qn('w:abstractNumId'), -1)) for a in existing_abs), default=-1) + 1
        existing_nums = nb.findall(qn('w:num'))
        next_num_id = max((int(n.get(qn('w:numId'), 0)) for n in existing_nums), default=0) + 1

        # ---------- Build <w:abstractNum> ----------
        abs_num = OxmlElement('w:abstractNum')
        abs_num.set(qn('w:abstractNumId'), str(next_abs_id))

        # Multi-level list style
        multi_lvl = OxmlElement('w:multiLevelType')
        multi_lvl.set(qn('w:val'), 'hybridMultilevel')
        abs_num.append(multi_lvl)

        # Build 9 levels (Word requires all 9)
        for lvl_idx in range(9):
            lvl_el = OxmlElement('w:lvl')
            lvl_el.set(qn('w:ilvl'), str(lvl_idx))

            start_el = OxmlElement('w:start')
            start_el.set(qn('w:val'), '1')
            lvl_el.append(start_el)

            nfmt_el = OxmlElement('w:numFmt')
            if list_type == 'numbered':
                nfmt_el.set(qn('w:val'), 'decimal')
            else:
                nfmt_el.set(qn('w:val'), 'bullet')
            lvl_el.append(nfmt_el)

            lvl_text_el = OxmlElement('w:lvlText')
            if list_type == 'numbered':
                # e.g. "%1." for level 0, "%2." for level 1, etc.
                lvl_text_el.set(qn('w:val'), f'%{lvl_idx + 1}.')
            elif list_type == 'checklist':
                lvl_text_el.set(qn('w:val'), CHECKLIST_CHAR)
            else:
                # Alternate bullet chars by level
                bullets = ['\u2022', 'o', '\u25aa', '\u2022', 'o', '\u25aa', '\u2022', 'o', '\u25aa']
                lvl_text_el.set(qn('w:val'), bullet_char or bullets[lvl_idx % 3])
            lvl_el.append(lvl_text_el)

            lvl_jc = OxmlElement('w:lvlJc')
            lvl_jc.set(qn('w:val'), 'left')
            lvl_el.append(lvl_jc)

            # Paragraph properties: indent
            pPr_el = OxmlElement('w:pPr')
            ind_el = OxmlElement('w:ind')
            left = 720 + lvl_idx * 720
            ind_el.set(qn('w:left'), str(left))
            ind_el.set(qn('w:hanging'), '360')
            pPr_el.append(ind_el)
            lvl_el.append(pPr_el)

            # Run properties: font for bullet chars
            if list_type in ('bullet', 'checklist'):
                rPr_el = OxmlElement('w:rPr')
                rFonts_el = OxmlElement('w:rFonts')
                if list_type == 'checklist':
                    # Segoe UI Symbol has good checkbox glyphs
                    rFonts_el.set(qn('w:ascii'), 'Segoe UI Symbol')
                    rFonts_el.set(qn('w:hAnsi'), 'Segoe UI Symbol')
                else:
                    rFonts_el.set(qn('w:ascii'), 'Symbol')
                    rFonts_el.set(qn('w:hAnsi'), 'Symbol')
                rPr_el.append(rFonts_el)
                lvl_el.append(rPr_el)

            abs_num.append(lvl_el)

        # Insert abstractNum before any existing <w:num> elements
        first_num = nb.find(qn('w:num'))
        if first_num is not None:
            nb.insert(list(nb).index(first_num), abs_num)
        else:
            nb.append(abs_num)

        # ---------- Build <w:num> ----------
        num_el = OxmlElement('w:num')
        num_el.set(qn('w:numId'), str(next_num_id))
        abs_num_id_ref = OxmlElement('w:abstractNumId')
        abs_num_id_ref.set(qn('w:val'), str(next_abs_id))
        num_el.append(abs_num_id_ref)
        nb.append(num_el)

        return next_num_id

    def _op_docx_list(self, doc, params: dict) -> str:
        """Handle list manipulation operations for DOCX documents.

        Actions:
          convert_type  — Change list format (bullet/numbered/checklist) for a range of paragraphs.
          add_items     — Insert new list items after an anchor paragraph.
          sort_items    — Alphabetically sort a range of list paragraphs.
          set_bullet_char — Change the bullet character for a list range.
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        import copy

        action = params.get('action')

        # ----------------------------------------------------------------
        # Shared helper: resolve a body_index id to a xml element
        # ----------------------------------------------------------------
        def _body_id_map():
            return {eid: xml_el for eid, xml_el in self._build_docx_body_index(doc)}

        def _xml_el_to_para(xml_el):
            for p in doc.paragraphs:
                if p._p is xml_el:
                    return p
            return None

        def _get_paras_in_range(start_id: str, end_id: str):
            """Return list of python-docx Paragraph objects between start_id and end_id (inclusive)."""
            bmap = _body_id_map()
            start_xml = bmap.get(start_id)
            end_xml = bmap.get(end_id)
            if start_xml is None or end_xml is None:
                return [], f"Could not find '{start_id}' or '{end_id}' in document"

            collecting = False
            result = []
            for eid, xml_el in self._build_docx_body_index(doc):
                if xml_el is start_xml:
                    collecting = True
                if collecting:
                    p = _xml_el_to_para(xml_el)
                    if p is not None:
                        result.append(p)
                if xml_el is end_xml:
                    break
            return result, None

        # ----------------------------------------------------------------
        # CONVERT_TYPE — change bullet/numbered/checklist for a paragraph range
        # ----------------------------------------------------------------
        if action == 'convert_type':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            list_type = params.get('list_type', 'numbered')  # 'bullet'|'numbered'|'checklist'
            bullet_char = params.get('bullet_char', '')

            if not start_id or not end_id:
                return "convert_type requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"convert_type: {err}"
            if not paras:
                return "convert_type: no paragraphs found in range"

            # Create a new numbering entry
            new_num_id = self._create_docx_numbering(doc, list_type, bullet_char)

            # Update each paragraph's <w:numPr>
            updated = 0
            for para in paras:
                pPr = para._p.find(qn('w:pPr'))
                if pPr is None:
                    pPr = OxmlElement('w:pPr')
                    para._p.insert(0, pPr)

                numPr = pPr.find(qn('w:numPr'))
                if numPr is None:
                    numPr = OxmlElement('w:numPr')
                    pPr.append(numPr)
                else:
                    # Clear existing numId / ilvl children
                    for child in list(numPr):
                        numPr.remove(child)

                ilvl_el = OxmlElement('w:ilvl')
                ilvl_el.set(qn('w:val'), '0')
                numPr.insert(0, ilvl_el)

                numId_el = OxmlElement('w:numId')
                numId_el.set(qn('w:val'), str(new_num_id))
                numPr.append(numId_el)

                # Update the paragraph style to match
                if list_type == 'numbered':
                    try:
                        para.style = doc.styles['List Number']
                    except Exception:
                        None
                else:
                    try:
                        para.style = doc.styles['List Bullet']
                    except Exception:
                        None

                updated += 1

            return f"Converted {updated} paragraph(s) to {list_type} list (numId={new_num_id})"

        # ----------------------------------------------------------------
        # ADD_ITEMS — insert new list items after an anchor paragraph
        # ----------------------------------------------------------------
        elif action == 'add_items':
            after_id = params.get('after_id')  # anchor element ID
            
            # Fallback to _raw_target_id if after_id is missing
            if not after_id:
                raw_tgt = params.get('_raw_target_id')
                if isinstance(raw_tgt, list) and raw_tgt:
                    after_id = raw_tgt[-1]
                elif isinstance(raw_tgt, str):
                    after_id = raw_tgt

            items = params.get('items', [])     # list of text strings
            if not after_id:
                return "add_items requires after_id"
            if not items:
                return "add_items: items list is empty"

            bmap = _body_id_map()
            after_xml = bmap.get(after_id)
            if after_xml is None:
                return f"add_items: element '{after_id}' not found"

            anchor_para = _xml_el_to_para(after_xml)

            # Read numId/ilvl from the anchor (or the last item in range if end_id given)
            end_id = params.get('end_id', after_id)
            end_xml = bmap.get(end_id, after_xml)
            end_para = _xml_el_to_para(end_xml) or anchor_para

            # Extract numPr from the anchor or end paragraph to clone
            template_para = end_para or anchor_para
            template_pPr = None
            if template_para is not None:
                template_pPr_el = template_para._p.find(qn('w:pPr'))
                if template_pPr_el is not None:
                    template_pPr = copy.deepcopy(template_pPr_el)

            # Insert items in reverse order (each addnext inserts immediately after end_xml)
            # So reverse to get correct final ordering
            insert_after = end_xml
            for item_text in items:
                new_p = OxmlElement('w:p')

                # Copy paragraph properties (numPr, style, indent, etc.)
                if template_pPr is not None:
                    new_p.append(copy.deepcopy(template_pPr))

                # Add a run with the item text
                new_r = OxmlElement('w:r')
                # Copy run properties from template paragraph's first run if available
                if template_para and template_para.runs:
                    first_rPr = template_para.runs[0]._r.find(qn('w:rPr'))
                    if first_rPr is not None:
                        new_r.append(copy.deepcopy(first_rPr))
                new_t = OxmlElement('w:t')
                new_t.text = item_text
                new_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
                new_r.append(new_t)
                new_p.append(new_r)

                insert_after.addnext(new_p)
                insert_after = new_p

            return f"Added {len(items)} item(s) after '{after_id}'"

        # ----------------------------------------------------------------
        # SORT_ITEMS — alphabetically sort a range of list paragraphs
        # ----------------------------------------------------------------
        elif action == 'sort_items':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            order = params.get('order', 'asc')  # 'asc'|'desc'

            if not start_id or not end_id:
                return "sort_items requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"sort_items: {err}"
            if len(paras) < 2:
                return "sort_items: fewer than 2 paragraphs in range, nothing to sort"

            # Only sort paragraphs that are list items (have numPr)
            list_paras = [p for p in paras if p._p.find(f'.//{qn("w:numPr")}') is not None
                          or (p._p.find(qn('w:pPr')) is not None
                              and p._p.find(qn('w:pPr')).find(qn('w:numPr')) is not None)]
            if not list_paras:
                # Fall back to all paras if none have numPr
                list_paras = paras

            # Extract text from each paragraph
            texts = [p.text.strip() for p in list_paras]

            # Sort
            reverse = (order == 'desc')
            sorted_texts = sorted(texts, key=lambda s: s.lower(), reverse=reverse)

            # Re-distribute sorted texts back into the same paragraph XML elements
            # We only swap the text content; run formatting stays on the same run positions.
            for para, new_text in zip(list_paras, sorted_texts):
                # Use run-aware replacement to set new text
                self._apply_run_aware_replacement(para, new_text, {})

            return f"Sorted {len(list_paras)} list item(s) in {order}ending order"

        # ----------------------------------------------------------------
        # SET_BULLET_CHAR — change the bullet character for a list range
        # ----------------------------------------------------------------
        elif action == 'set_bullet_char':
            start_id = params.get('start_id')
            end_id = params.get('end_id')
            char = params.get('char', '\u2022')  # default: •

            if not start_id or not end_id:
                return "set_bullet_char requires start_id and end_id"

            paras, err = _get_paras_in_range(start_id, end_id)
            if err:
                return f"set_bullet_char: {err}"

            # Determine the numId used by the first paragraph in range
            num_ids_to_patch = set()
            for p in paras:
                li = self._extract_list_info(p, doc)
                if li:
                    num_ids_to_patch.add((li['num_id'], li['ilvl']))

            if not num_ids_to_patch:
                return "set_bullet_char: no list paragraphs found in range"

            # For each unique (numId, ilvl), patch the lvlText in the abstractNum
            try:
                nb = doc.part.numbering_part._element
                for num_id, ilvl in num_ids_to_patch:
                    num_el = nb.find(f'.//{qn("w:num")}[@{qn("w:numId")}="{num_id}"]')
                    if num_el is None:
                        continue
                    abs_ref = num_el.find(qn('w:abstractNumId'))
                    if abs_ref is None:
                        continue
                    abs_id = abs_ref.get(qn('w:val'), '0')
                    abs_num = nb.find(f'.//{qn("w:abstractNum")}[@{qn("w:abstractNumId")}="{abs_id}"]')
                    if abs_num is None:
                        continue
                    for lvl in abs_num:
                        if lvl.get(qn('w:ilvl')) == str(ilvl):
                            lvl_text_el = lvl.find(qn('w:lvlText'))
                            if lvl_text_el is not None:
                                lvl_text_el.set(qn('w:val'), char)
            except Exception as e:
                return f"set_bullet_char: failed to patch numbering XML: {e}"

            return f"Set bullet character to '{char}' for {len(num_ids_to_patch)} list level(s)"

        return f"list_op action '{action}' not supported"

    # ------------------------------------------------------------------
    # DOCX find & replace operations
    # ------------------------------------------------------------------

    def _op_docx_find_replace(self, doc, tgt: dict, params: dict) -> str:
        """Find and replace text across the document or a specific target.

        Supports standard substring replacement or regex replacement.
        """
        import re

        find_text = params.get('find_text', '')
        replace_text = params.get('replace_text', '')
        is_regex = bool(params.get('is_regex', False))
        match_case = bool(params.get('match_case', False))

        if not find_text:
            return "find_replace: find_text is empty"

        # Determine target scope
        target_id = params.get('target_id', 'all')
        if not tgt:
            target_id = 'all'

        paras_to_check = []
        if target_id == 'all':
            # Collect all document paragraphs
            paras_to_check.extend(doc.paragraphs)
            # Collect all table cell paragraphs
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paras_to_check.extend(cell.paragraphs)
        else:
            para_idx = tgt.get("paragraph_index", tgt.get("para_index"))
            if para_idx is not None:
                if "table_index" in tgt and "row_index" in tgt and "col_index" in tgt:
                    try:
                        table = doc.tables[tgt["table_index"]]
                        cell = table.cell(tgt["row_index"], tgt["col_index"])
                        paras_to_check = [cell.paragraphs[para_idx]]
                    except Exception:
                        None
                else:
                    try:
                        paras_to_check = [doc.paragraphs[para_idx]]
                    except Exception:
                        None

        if not paras_to_check:
            return "find_replace: no target paragraphs found"

        # Prepare regex pattern if needed
        pattern = None
        if is_regex:
            flags = 0 if match_case else re.IGNORECASE
            try:
                pattern = re.compile(find_text, flags)
            except Exception as e:
                return f"find_replace: invalid regex pattern '{find_text}' - {e}"

        replacements = 0
        for p in paras_to_check:
            old_text = p.text
            if not old_text:
                continue

            new_text = old_text
            if is_regex:
                if pattern.search(old_text):
                    new_text = pattern.sub(replace_text, old_text)
            else:
                if match_case:
                    if find_text in old_text:
                        new_text = old_text.replace(find_text, replace_text)
                else:
                    # case-insensitive replace
                    # Use a regex compilation for simplicity of ignoring case
                    flags = re.IGNORECASE
                    pat = re.compile(re.escape(find_text), flags)
                    if pat.search(old_text):
                        new_text = pat.sub(replace_text, old_text)

            if new_text != old_text:
                self._apply_run_aware_replacement(p, new_text)
                replacements += 1

        if target_id == 'all':
            return f"Replaced occurrences of '{find_text}' in {replacements} paragraph(s)"
        else:
            return f"Replaced text in targeted paragraph"

    # ------------------------------------------------------------------
    # DOCX image operations
    # ------------------------------------------------------------------

    def _get_docx_page_width_emu(self, doc) -> int:
        """Return the usable page width in EMU (page width minus left+right margins)."""
        try:
            section = doc.sections[0]
            usable = section.page_width - section.left_margin - section.right_margin
            return int(usable)
        except Exception:
            # A4 usable width fallback: ~16.51cm = 5940000 EMU
            return 5940000

    def _find_docx_image_paragraphs(self, doc) -> list[tuple[int, object]]:
        """Return list of (image_index, paragraph) for all paragraphs containing <w:drawing>."""
        from docx.oxml.ns import qn
        WNS_DRAWING = qn('w:drawing')
        result = []
        image_idx = 0
        for para in doc.paragraphs:
            if para._p.find(f'.//{WNS_DRAWING}') is not None:
                result.append((image_idx, para))
                image_idx += 1
        return result

    def _op_docx_image(self, doc, tgt: dict, params: dict) -> str:
        """Handle image operations for DOCX documents.

        Supported actions:
          insert      — Add a new image paragraph before/after a body element.
          replace     — Replace the drawing in an existing image paragraph.
          resize      — Resize an inline image by width_page_pct or explicit EMU.
          reposition  — Set paragraph alignment (left/center/right) of an image paragraph and its caption.
          move        — Move an existing image paragraph and its caption before/after another element.
          place_inline — Move an existing image into a target text paragraph (e.g. right side of title).
          add_caption — Insert a styled caption paragraph after the image paragraph.
          reposition_caption — Set only the caption paragraph alignment.
          format_caption — Apply text formatting to the adjacent caption.
          remove      — Remove the image paragraph entirely.
        """
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt

        action = params.get("action")
        if not action:
            action = "reposition" if params.get("alignment") and tgt.get("id") else "insert"
        image_path = params.get("image_path")
        page_w_emu = self._get_docx_page_width_emu(doc)

        if action in ("caption", "add caption"):
            action = "add_caption"
        elif action in ("align", "alignment"):
            action = "reposition"

        position = params.get("position") if isinstance(params.get("position"), dict) else {}
        if params.get("width_page_pct") is None and position.get("width_pct") is not None:
            try:
                params["width_page_pct"] = float(position["width_pct"])
            except (TypeError, ValueError):
                None
        if params.get("height_page_pct") is None and position.get("height_pct") is not None:
            try:
                params["height_page_pct"] = float(position["height_pct"])
            except (TypeError, ValueError):
                None
        if not params.get("alignment") and position.get("left_pct") is not None:
            try:
                left_pct = float(position["left_pct"])
                if left_pct >= 0.55:
                    params["alignment"] = "right"
                elif left_pct >= 0.25:
                    params["alignment"] = "center"
                else:
                    params["alignment"] = "left"
            except (TypeError, ValueError):
                None

        if action == "reposition" and not params.get("alignment"):
            float_pos = params.get("float_position")
            if float_pos in ("left", "right"):
                params["alignment"] = float_pos

        # For DOCX insertions, a paragraph/table target means "insert after this element".
        if action == "insert" and not params.get("after_id") and not params.get("before_id") and tgt.get("id"):
            params["after_id"] = tgt["id"]

        ALIGN_MAP = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
        }

        def _is_caption_para(para) -> bool:
            text = (para.text or "").strip().lower()
            style_name = ""
            try:
                style_name = (para.style.name or "").lower() if para.style else ""
            except Exception:
                style_name = ""
            return (
                "caption" in style_name
                or text.startswith("figure ")
                or text.startswith("fig. ")
                or text.startswith("image ")
                or text.startswith("photo ")
            )

        def _caption_after_image_para(image_para):
            paras = list(doc.paragraphs)
            for idx, para in enumerate(paras):
                if para._p is image_para._p and idx + 1 < len(paras):
                    candidate = paras[idx + 1]
                    if _is_caption_para(candidate):
                        return candidate
                    break
            return None

        def _set_para_alignment(para, alignment: str) -> bool:
            if alignment not in ALIGN_MAP or para is None:
                return False
            para.alignment = ALIGN_MAP[alignment]
            return True

        def _set_image_group_alignment(image_para, alignment: str, include_caption: bool = True) -> bool:
            changed_local = _set_para_alignment(image_para, alignment)
            if include_caption:
                caption_para = _caption_after_image_para(image_para)
                if caption_para is not None:
                    changed_local = _set_para_alignment(caption_para, alignment) or changed_local
            return changed_local

        def _first_body_element_id(exclude_xml: list | None = None) -> str | None:
            exclude_xml = exclude_xml or []
            for eid, xml_el in self._build_docx_body_index(doc):
                if eid.startswith("body_other_") or xml_el in exclude_xml:
                    continue
                return eid
            return None

        def _set_alt_text(inline_shape, alt_text: str):
            if not alt_text:
                return
            try:
                docPr = inline_shape._inline.find('.//' + qn('wp:docPr'))
                if docPr is not None:
                    docPr.set('name', alt_text)
                    docPr.set('descr', alt_text)
            except Exception:
                None

        # ----------------------------------------------------------------
        # Helper: find target paragraph by target_id (image_N or paragraph_N)
        # ----------------------------------------------------------------
        def _resolve_para_by_id(target_id: str):
            """Resolve a DOM id (image_N or paragraph_N) to (xml_element, paragraph)."""
            body_index = self._build_docx_body_index(doc)
            id_to_elem = {eid: xml_el for eid, xml_el in body_index}
            xml_el = id_to_elem.get(target_id)
            if xml_el is None:
                return None, None
            # Find matching python-docx paragraph
            for para in doc.paragraphs:
                if para._p is xml_el:
                    return xml_el, para
            return xml_el, None

        # ----------------------------------------------------------------
        # INSERT — add image paragraph after after_id element
        # ----------------------------------------------------------------
        if action == "insert":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image insert: image_path missing or not found: %s", image_path)
                return ""

            after_id = params.get("after_id")  # body_index element id
            before_id = params.get("before_id")
            alignment = params.get("alignment", "center")
            width_page_pct = params.get("width_page_pct", 0.5)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            caption_text = params.get("caption_text")

            width_emu = int(page_w_emu * width_page_pct)

            # Add picture creates it at the end of the document, we'll move it
            new_para = doc.add_paragraph()
            new_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.LEFT)
            run = new_para.add_run()
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))

            new_p_el = new_para._p

            # Move the new paragraph to the correct position
            if before_id or after_id:
                body_index = self._build_docx_body_index(doc)
                id_to_elem = {eid: xml_el for eid, xml_el in body_index}
                anchor_xml = id_to_elem.get(before_id or after_id)
                if anchor_xml is not None and anchor_xml is not new_p_el:
                    # Remove from its current (appended-at-end) position
                    new_p_el.getparent().remove(new_p_el)
                    if before_id:
                        anchor_xml.addprevious(new_p_el)
                    else:
                        anchor_xml.addnext(new_p_el)

            anchor_id = before_id or after_id
            anchor_side = "before" if before_id else "after"
            summary = f"Inserted image (width={int(width_page_pct*100)}% page) {anchor_side} '{anchor_id}'"

            # Optional inline caption
            if caption_text:
                cap_para = doc.add_paragraph(caption_text)
                try:
                    cap_para.style = doc.styles["Caption"]
                except Exception:
                    cap_para.runs[0].font.italic = True if cap_para.runs else None
                cap_p_el = cap_para._p
                cap_p_el.getparent().remove(cap_p_el)
                new_p_el.addnext(cap_p_el)
                cap_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.CENTER)
                summary += f" + caption '{caption_text}'"

            return summary

        # ----------------------------------------------------------------
        # INSERT INTO PARAGRAPH — add image run to an existing paragraph
        # ----------------------------------------------------------------
        if action == "insert_into_paragraph":
            if not image_path or not Path(image_path).exists():
                return ""

            target_xml, target_para = _resolve_para_by_id(tgt.get("id", ""))
            if target_para is None:
                return "insert_into_paragraph: no target paragraph found"

            width_page_pct = params.get("width_page_pct", 0.3)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            width_emu = int(page_w_emu * width_page_pct)

            run = target_para.add_run()
            # If they want it spaced out, add a tab before it to push it towards the right
            run.add_text("\t") 
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))
            return f"Inserted image (width={int(width_page_pct*100)}%) into '{tgt.get('id')}'"

        # ----------------------------------------------------------------
        # REPLACE TEXT — swap a text paragraph with an image
        # ----------------------------------------------------------------
        if action == "replace_text":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image replace_text: image_path missing or not found: %s", image_path)
                return ""

            target_xml, _ = _resolve_para_by_id(tgt.get("id", ""))
            if target_xml is None:
                return "replace_text: no target paragraph found"

            alignment = params.get("alignment", "center")
            width_page_pct = params.get("width_page_pct", 0.5)
            maintain_ar = params.get("maintain_aspect_ratio", True)
            caption_text = params.get("caption_text")

            width_emu = int(page_w_emu * width_page_pct)

            new_para = doc.add_paragraph()
            new_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.LEFT)
            run = new_para.add_run()
            if maintain_ar:
                shape = run.add_picture(image_path, width=width_emu)
            else:
                height_pct = params.get("height_page_pct", width_page_pct * 0.75)
                height_emu = int(page_w_emu * height_pct)
                shape = run.add_picture(image_path, width=width_emu, height=height_emu)
            
            _set_alt_text(shape, params.get("alt_text"))

            new_p_el = new_para._p

            # Insert new image paragraph immediately before the target text paragraph, then remove target
            target_xml.addprevious(new_p_el)
            target_xml.getparent().remove(target_xml)

            summary = f"Replaced text paragraph with image (width={int(width_page_pct*100)}% page)"

            if caption_text:
                cap_para = doc.add_paragraph(caption_text)
                try:
                    cap_para.style = doc.styles["Caption"]
                except Exception:
                    cap_para.runs[0].font.italic = True if cap_para.runs else None
                cap_p_el = cap_para._p
                cap_p_el.getparent().remove(cap_p_el)
                new_p_el.addnext(cap_p_el)
                cap_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.CENTER)
                summary += f" + caption '{caption_text}'"

            return summary

        # ----------------------------------------------------------------
        # Helpers for operations on EXISTING image paragraphs
        # ----------------------------------------------------------------
        target_id = tgt.get("type") or ""
        # Resolve image paragraph from target_id in op dict
        op_target_id = ""
        # We stored the raw target_id on `tgt` via _parse_target_id;
        # but image_N IDs have no numeric field in _parse_target_id.
        # Retrieve it directly from the op via params fallback.
        image_target_id = params.get("target_image_id") or ""

        # Build an image-index to paragraph lookup
        img_paras = self._find_docx_image_paragraphs(doc)  # [(img_idx, para), ...]

        def _get_img_para(img_idx: int):
            for i, p in img_paras:
                if i == img_idx:
                    return p
            return None

        # Parse image index from tgt (image_N -> N stored in tgt)
        img_idx = tgt.get("image_index")  # set by _parse_target_id extension below
        if img_idx is None and img_paras:
            img_idx = 0  # default to first image

        img_para = _get_img_para(img_idx) if img_idx is not None else None

        # ----------------------------------------------------------------
        # REPLACE — swap the drawing XML in the existing image paragraph
        # ----------------------------------------------------------------
        if action == "replace":
            if not image_path or not Path(image_path).exists():
                log.warning("DOCX image replace: image_path missing or not found: %s", image_path)
                return ""

            if img_para is None:
                return "replace: no image paragraph found"

            width_page_pct = params.get("width_page_pct")
            alignment = params.get("alignment")

            # Determine width to use — try to preserve existing width if not specified
            WNS_EXTENT = qn('wp:extent')
            WNS_DRAWING = qn('w:drawing')
            old_drawing = img_para._p.find(f'.//{WNS_DRAWING}')
            old_width_emu = None
            if old_drawing is not None:
                extent = old_drawing.find(f'.//{WNS_EXTENT}')
                if extent is not None:
                    try:
                        old_width_emu = int(extent.get('cx', 0))
                    except Exception:
                        None

            if width_page_pct is not None:
                width_emu = int(page_w_emu * width_page_pct)
            elif old_width_emu:
                width_emu = old_width_emu
            else:
                width_emu = int(page_w_emu * 0.5)

            # Remove all existing runs from the paragraph
            for r in list(img_para.runs):
                r._r.getparent().remove(r._r)

            # Add new run with the replacement image
            run = img_para.add_run()
            shape = run.add_picture(image_path, width=width_emu)
            _set_alt_text(shape, params.get("alt_text"))

            if alignment and alignment in ALIGN_MAP:
                _set_image_group_alignment(img_para, alignment, include_caption=True)

            return f"Replaced image {img_idx} with {Path(image_path).name}"

        # ----------------------------------------------------------------
        # RESIZE — change width (and optionally height) of inline image
        # ----------------------------------------------------------------
        elif action == "resize":
            if img_para is None:
                return "resize: no image paragraph found"

            width_page_pct = params.get("width_page_pct")
            maintain_ar = params.get("maintain_aspect_ratio", True)

            if width_page_pct is None:
                return "resize: width_page_pct required"

            new_w_emu = int(page_w_emu * width_page_pct)

            from docx.oxml.ns import qn as _qn
            WNS_DRAWING = _qn('w:drawing')
            WNS_EXTENT = _qn('wp:extent')
            WNS_EFF_EXTENT = _qn('wp:effectExtent')

            drawing = img_para._p.find(f'.//{WNS_DRAWING}')
            if drawing is None:
                return "resize: no <w:drawing> element found in paragraph"

            extent = drawing.find(f'.//{WNS_EXTENT}')
            if extent is None:
                return "resize: no <wp:extent> element found in drawing"

            old_w = int(extent.get('cx', new_w_emu))
            old_h = int(extent.get('cy', new_w_emu))

            if maintain_ar and old_w > 0:
                ratio = old_h / old_w
                new_h_emu = int(new_w_emu * ratio)
            else:
                height_page_pct = params.get("height_page_pct")
                new_h_emu = int(page_w_emu * height_page_pct) if height_page_pct else old_h

            extent.set('cx', str(new_w_emu))
            extent.set('cy', str(new_h_emu))

            # Also update distT/distB/distL/distR extents (effectExtent)
            # and the <a:ext> inside the graphic for compatibility
            from lxml import etree
            for a_ext in drawing.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}ext'):
                a_ext.set('cx', str(new_w_emu))
                a_ext.set('cy', str(new_h_emu))

            return f"Resized image {img_idx} to {int(width_page_pct*100)}% page width"

        # ----------------------------------------------------------------
        # MOVE — relocate image paragraph and adjacent caption together
        # ----------------------------------------------------------------
        elif action == "move":
            if img_para is None:
                return "move: no image paragraph found"

            before_id = params.get("before_id")
            after_id = params.get("after_id")
            placement = str(params.get("placement") or "").lower()
            alignment = params.get("alignment")

            if not before_id and not after_id and placement in ("top", "top_page", "top_of_page", "top_right", "top_right_page"):
                caption_para_for_top = _caption_after_image_para(img_para)
                moving_xml = [img_para._p]
                if caption_para_for_top is not None:
                    moving_xml.append(caption_para_for_top._p)
                before_id = _first_body_element_id(exclude_xml=moving_xml)
                if not alignment and "right" in placement:
                    alignment = "right"

            if not before_id and not after_id:
                alignment = alignment or params.get("float_position")
                if alignment in ALIGN_MAP:
                    _set_image_group_alignment(img_para, alignment, include_caption=True)
                    return f"Moved image {img_idx}: alignment={alignment} with caption"
                return "move: before_id, after_id, or alignment required"

            body_index = self._build_docx_body_index(doc)
            id_to_elem = {eid: xml_el for eid, xml_el in body_index}
            anchor_xml = id_to_elem.get(before_id or after_id)
            if anchor_xml is None:
                return f"move: anchor '{before_id or after_id}' not found"

            caption_para = _caption_after_image_para(img_para)
            moving = [img_para._p]
            if caption_para is not None:
                moving.append(caption_para._p)

            if anchor_xml in moving:
                return "move: target anchor is inside the image/caption group"

            for xml_el in moving:
                xml_el.getparent().remove(xml_el)

            if before_id:
                for xml_el in moving:
                    anchor_xml.addprevious(xml_el)
            else:
                current = anchor_xml
                for xml_el in moving:
                    current.addnext(xml_el)
                    current = xml_el

            if alignment in ALIGN_MAP:
                _set_image_group_alignment(img_para, alignment, include_caption=True)

            side = "before" if before_id else "after"
            suffix = " with caption" if caption_para is not None else ""
            return f"Moved image {img_idx}{suffix} {side} '{before_id or after_id}'"

        # ----------------------------------------------------------------
        # PLACE_INLINE — move image into target paragraph, keeping caption after it
        # ----------------------------------------------------------------
        elif action == "place_inline":
            if img_para is None:
                return "place_inline: no image paragraph found"

            anchor_id = params.get("anchor_id") or params.get("paragraph_id") or params.get("after_id") or params.get("before_id")
            if not anchor_id:
                return "place_inline: anchor_id required"

            _, target_para = _resolve_para_by_id(anchor_id)
            if target_para is None:
                return f"place_inline: target paragraph '{anchor_id}' not found"

            caption_para = _caption_after_image_para(img_para)
            moving = [img_para._p]
            if caption_para is not None:
                moving.append(caption_para._p)

            if target_para._p in moving:
                return "place_inline: target paragraph is inside the image/caption group"

            for xml_el in moving:
                xml_el.getparent().remove(xml_el)

            current = target_para._p
            for xml_el in moving:
                current.addnext(xml_el)
                current = xml_el

            alignment = params.get("alignment", "right")
            if alignment in ALIGN_MAP:
                _set_image_group_alignment(img_para, alignment, include_caption=True)

            return f"Placed image {img_idx} after paragraph '{anchor_id}' aligned {alignment}"

        # ----------------------------------------------------------------
        # REPOSITION — set paragraph alignment of image paragraph
        # ----------------------------------------------------------------
        elif action == "reposition":
            if img_para is None:
                return "reposition: no image paragraph found"

            alignment = params.get("alignment", "left")
            if alignment not in ALIGN_MAP:
                return f"reposition: unknown alignment '{alignment}', use left/center/right"

            _set_image_group_alignment(img_para, alignment, include_caption=True)

            return f"Repositioned image {img_idx}: alignment={alignment}" + (
                " with caption" if _caption_after_image_para(img_para) is not None else ""
            )

        # ----------------------------------------------------------------
        # REPOSITION_CAPTION — set only the caption paragraph alignment
        # ----------------------------------------------------------------
        elif action == "reposition_caption":
            if img_para is None:
                return "reposition_caption: no image paragraph found"

            alignment = params.get("alignment", "center")
            if alignment not in ALIGN_MAP:
                return f"reposition_caption: unknown alignment '{alignment}', use left/center/right"

            caption_para = _caption_after_image_para(img_para)
            if caption_para is None:
                return "reposition_caption: no caption paragraph found"

            caption_para.alignment = ALIGN_MAP[alignment]
            return f"Repositioned caption for image {img_idx}: alignment={alignment}"

        # ----------------------------------------------------------------
        # FORMAT_CAPTION — apply text formatting to the adjacent caption
        # ----------------------------------------------------------------
        elif action == "format_caption":
            if img_para is None:
                return "format_caption: no image paragraph found"

            caption_para = _caption_after_image_para(img_para)
            if caption_para is None:
                return "format_caption: no caption paragraph found"

            self._apply_docx_format(caption_para, {}, params)
            return "Formatted caption for image {}".format(img_idx)

        # ----------------------------------------------------------------
        # ADD_CAPTION — insert caption paragraph after image paragraph
        # ----------------------------------------------------------------
        elif action == "add_caption":
            caption_text = params.get("caption_text", "")
            caption_style = params.get("caption_style", "Caption")
            alignment = params.get("alignment", "center")
            if not caption_text:
                return "add_caption: caption_text required"
            if img_para is None:
                return "add_caption: no image paragraph found"

            # Create a new paragraph for the caption
            cap_para = doc.add_paragraph(caption_text)
            try:
                cap_para.style = doc.styles[caption_style]
            except Exception:
                # Fallback: italic, small font
                if cap_para.runs:
                    cap_para.runs[0].font.italic = True
                    cap_para.runs[0].font.size = Pt(9)
            cap_para.alignment = ALIGN_MAP.get(alignment, WD_ALIGN_PARAGRAPH.CENTER)
            if alignment in ALIGN_MAP:
                img_para.alignment = ALIGN_MAP[alignment]

            # Move the caption paragraph to immediately after the image paragraph
            cap_p_el = cap_para._p
            cap_p_el.getparent().remove(cap_p_el)
            img_para._p.addnext(cap_p_el)

            return f"Added caption '{caption_text}' below image {img_idx}"

        # ----------------------------------------------------------------
        # REMOVE — delete the image paragraph
        # ----------------------------------------------------------------
        elif action == "remove":
            if img_para is None:
                return "remove: no image paragraph found"
            img_para._p.getparent().remove(img_para._p)
            return f"Removed image {img_idx}"

        return f"image_op action '{action}' not supported for DOCX"

    # ------------------------------------------------------------------
    # Individual operation handlers — PPTX
    # ------------------------------------------------------------------

    def _get_slide(self, prs: Presentation, slide_num: int | None):
        """Get a slide by 1-based index, or the first slide if not specified."""
        slides = list(prs.slides)
        if not slides:
            return None
        if slide_num is None:
            return slides[0]
        idx = slide_num - 1
        if 0 <= idx < len(slides):
            return slides[idx]
        return slides[0]

    def _get_shape(self, slide, shape_idx: int | None):
        """Get a shape by index from a slide."""
        shapes = list(slide.shapes)
        if not shapes:
            return None
        if shape_idx is None:
            return None
        if 0 <= shape_idx < len(shapes):
            return shapes[shape_idx]
        return None

    def _op_pptx_text_edit(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Rewrite the text content of a specific paragraph."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""
        shape = self._get_shape(slide, tgt.get("shape_index"))
        if not shape or not getattr(shape, "has_text_frame", False):
            return ""
        new_text = params.get("new_text", "")
        para_idx = tgt.get("para_index", 0)
        paras = list(shape.text_frame.paragraphs)
        if para_idx < len(paras):
            self._apply_run_aware_replacement(paras[para_idx], new_text, params)
            return f"Rewrote text in slide {tgt.get('slide')}, shape {tgt.get('shape_index')}"
        return ""

    def _op_pptx_text_format(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Apply rich text formatting to a paragraph."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        results = []
        shapes_to_process = []

        shape_idx = tgt.get("shape_index")
        if shape_idx is not None:
            shape = self._get_shape(slide, shape_idx)
            if shape:
                shapes_to_process = [shape]
        else:
            # Apply to all text shapes on the slide
            shapes_to_process = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]

        for shape in shapes_to_process:
            if not getattr(shape, "has_text_frame", False):
                continue
            para_idx = tgt.get("para_index")
            paras = list(shape.text_frame.paragraphs)
            target_paras = [paras[para_idx]] if para_idx is not None and para_idx < len(paras) else paras

            for para in target_paras:
                c = self._apply_pptx_format_to_para(para, tgt, params, shape)
                if c:
                    changed = True
                if changed:
                    results.append(f"slide {tgt.get('slide')} shape {shape.shape_id}")

        return f"Applied formatting to {results[0]}" if results else ""

    def _apply_pptx_format_to_para(self, para, tgt: dict, params: dict, shape=None) -> bool:
        """Apply formatting params to a paragraph and its runs. Returns True if changed."""
        match_role = str(params.get("match_role", "")).strip().lower()
        if match_role and shape:
            shape_role = "shape"
            if shape.is_placeholder:
                try:
                    ph_type = str(shape.placeholder_format.type).split('(')[0]
                    shape_role = f"placeholder_{ph_type.lower()}"
                except Exception:
                    shape_role = "placeholder"
            # simple matching: if user says "heading", match title. 
            if match_role == "heading" and "title" not in shape_role:
                return False
            if match_role == "body" and "body" not in shape_role:
                return False
            # exact match fallback
            if match_role not in ["heading", "body"] and match_role != shape_role:
                return False

        changed = False

        # Alignment
        alignment = params.get("alignment")
        if alignment and alignment in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[alignment]
            changed = True

        # Line spacing
        line_spacing = params.get("line_spacing")
        if line_spacing is not None:
            try:
                from pptx.util import Pt
                from pptx.oxml.ns import qn
                from lxml import etree
                pPr = para._p.get_or_add_pPr()
                lnSpc = pPr.find(qn("a:lnSpc"))
                if lnSpc is None:
                    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
                spcPct = lnSpc.find(qn("a:spcPct"))
                if spcPct is None:
                    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
                spcPct.set("val", str(int(line_spacing * 100000)))
                changed = True
            except Exception as e:
                log.debug("Line spacing failed: %s", e)

        # Run-level formatting
        bold = params.get("bold")
        italic = params.get("italic")
        underline = params.get("underline")
        strikethrough = params.get("strikethrough")
        font_family = params.get("font_family")
        font_size_pt = params.get("font_size_pt")
        color_hex = params.get("color_hex")
        highlight_hex = params.get("highlight_hex")
        superscript = params.get("superscript")
        subscript = params.get("subscript")
        char_spacing = params.get("char_spacing")

        if not para.runs:
            return changed

        target_run_idx = tgt.get("run_index")
        match_color_hex = str(params.get("match_color_hex", "")).strip().lstrip("#").upper()

        for r_idx, run in enumerate(para.runs):
            if target_run_idx is not None and r_idx != target_run_idx:
                continue

            if match_color_hex:
                _, cur_color = self._resolve_pptx_font_info(run)
                cur_color = str(cur_color).strip().lstrip("#").upper() if cur_color else "000000"
                if cur_color != match_color_hex and not (cur_color == "000000" and match_color_hex == "AUTO"):
                    continue

            font = run.font
            if bold is not None:
                font.bold = bold; changed = True
            if italic is not None:
                font.italic = italic; changed = True
            if underline is not None:
                font.underline = underline; changed = True
            if font_family is not None:
                font.name = font_family; changed = True
            if font_size_pt is not None:
                font.size = Pt(font_size_pt); changed = True
            if color_hex is not None and len(str(color_hex)) == 6:
                try:
                    font.color.rgb = RGBColor.from_string(str(color_hex)); changed = True
                except Exception:
                    None
            if superscript is not None or subscript is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    if superscript is not None:
                        rPr.set("baseline", "30000" if superscript else "0"); changed = True
                    if subscript is not None:
                        rPr.set("baseline", "-25000" if subscript else "0"); changed = True
                except Exception as e:
                    log.debug("Superscript/subscript failed: %s", e)
            if strikethrough is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("strike", "sngStrike" if strikethrough else "noStrike"); changed = True
                except Exception as e:
                    log.debug("Strikethrough failed: %s", e)
            if char_spacing is not None:
                try:
                    from pptx.oxml.ns import qn
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("spc", str(int(char_spacing * 100))); changed = True
                except Exception as e:
                    log.debug("Char spacing failed: %s", e)

        return changed

    def _apply_docx_format(self, para, tgt: dict, params: dict) -> bool:
        """Apply rich formatting to a DOCX paragraph's runs."""
        match_role = str(params.get("match_role", "")).strip().lower()
        if match_role:
            style_name = para.style.name if para.style else "Normal"
            role = "body"
            if style_name.lower().startswith("heading"):
                role = "heading"
            elif "bullet" in style_name.lower() or "list" in style_name.lower():
                role = "bullet_point"
            
            if role != match_role:
                return False
                
        match_text = params.get("match_text")
        if match_text and match_text not in para.text:
            return False
        
        alignment = params.get("alignment")
        if alignment:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            align_map = {
                "left": WD_ALIGN_PARAGRAPH.LEFT,
                "center": WD_ALIGN_PARAGRAPH.CENTER,
                "right": WD_ALIGN_PARAGRAPH.RIGHT,
                "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
            }
            if alignment in align_map:
                para.alignment = align_map[alignment]

        line_spacing = params.get("line_spacing")
        if line_spacing is not None:
            para.paragraph_format.line_spacing = float(line_spacing)

        page_break_before = params.get("page_break_before")
        if page_break_before is not None:
            para.paragraph_format.page_break_before = bool(page_break_before)

        space_before = params.get("space_before_pt")
        if space_before is not None:
            from docx.shared import Pt
            para.paragraph_format.space_before = Pt(float(space_before))

        space_after = params.get("space_after_pt")
        if space_after is not None:
            from docx.shared import Pt
            para.paragraph_format.space_after = Pt(float(space_after))
            
        left_indent = params.get("left_indent_pt")
        if left_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.left_indent = Pt(float(left_indent))
            
        right_indent = params.get("right_indent_pt")
        if right_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.right_indent = Pt(float(right_indent))
            
        first_line_indent = params.get("first_line_indent_pt")
        if first_line_indent is not None:
            from docx.shared import Pt
            para.paragraph_format.first_line_indent = Pt(float(first_line_indent))
            
        keep_with_next = params.get("keep_with_next")
        if keep_with_next is not None:
            para.paragraph_format.keep_with_next = bool(keep_with_next)
            
        keep_together = params.get("keep_together")
        if keep_together is not None:
            para.paragraph_format.keep_together = bool(keep_together)

        include_in_toc = params.get("include_in_toc")
        if include_in_toc is not None:
            from docx.oxml.ns import qn as _qn
            from docx.oxml import OxmlElement
            pPr = para._p.get_or_add_pPr()
            outlineLvl = pPr.find(_qn('w:outlineLvl'))
            if include_in_toc is False:
                if outlineLvl is None:
                    outlineLvl = OxmlElement('w:outlineLvl')
                    pPr.append(outlineLvl)
                outlineLvl.set(_qn('w:val'), '9')
            else:
                if outlineLvl is not None:
                    pPr.remove(outlineLvl)

        # Normalise color_hex: strip '#' prefix so both "FF0000" and "#FF0000" work
        raw_color = str(params.get("color_hex", "")).strip().lstrip("#").upper()
        _color_map = {
            "RED": "FF0000",
            "BLUE": "0000FF",
            "GREEN": "00FF00",
            "BLACK": "000000",
            "WHITE": "FFFFFF",
            "YELLOW": "FFFF00",
            "PURPLE": "800080",
            "ORANGE": "FFA500",
            "GREY": "808080",
            "GRAY": "808080"
        }
        color_hex = _color_map.get(raw_color, raw_color)

        runs = para.runs
        # If the paragraph has no runs but has text, create one to hold the formatting
        if not runs and para.text.strip():
            run = para.add_run(para.text)
            # Remove original text nodes to avoid duplication
            from lxml import etree
            for t_elem in para._p.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                parent = t_elem.getparent()
                if parent is not None and parent != run._r:
                    parent_parent = parent.getparent()
                    if parent_parent is not None and parent != run._r:
                        try:
                            parent_parent.remove(parent)
                        except Exception:
                            None
            runs = para.runs

        if match_text:
            return self._apply_run_aware_format(para, match_text, params)
                
        target_run_idx = tgt.get("run_index")
        match_color_hex = str(params.get("match_color_hex", "")).strip().lstrip("#").upper()

        for r_idx, run in enumerate(runs):
            if target_run_idx is not None and r_idx != target_run_idx:
                continue
                
            if match_color_hex:
                _, cur_color = self._resolve_docx_font_info(para, run)
                cur_color = str(cur_color).strip().lstrip("#").upper() if cur_color else "000000"
                # If they don't match, skip this run
                if cur_color != match_color_hex and not (cur_color == "000000" and match_color_hex == "AUTO"):
                    continue

            if params.get("bold") is not None:
                run.bold = params["bold"]
            if params.get("italic") is not None:
                run.italic = params["italic"]
            if params.get("underline") is not None:
                run.underline = params["underline"]
            if params.get("font_family"):
                run.font.name = params["font_family"]
            if params.get("font_size_pt") is not None:
                from docx.shared import Pt
                run.font.size = Pt(params["font_size_pt"])
            if color_hex and len(color_hex) == 6:
                try:
                    from docx.shared import RGBColor as DRGBColor
                    run.font.color.rgb = DRGBColor(
                        int(color_hex[0:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16),
                    )
                except Exception as e:
                    log.debug("DOCX color application failed: %s", e)
                    
            highlight_hex = params.get("highlight_hex")
            if highlight_hex:
                from docx.enum.text import WD_COLOR_INDEX
                # Simple mapping for common highlight colors
                hl_color = str(highlight_hex).strip().lstrip("#").upper()
                if hl_color in ["FFFF00", "YELLOW"]:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                elif hl_color in ["00FF00", "GREEN"]:
                    run.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
                elif hl_color in ["00FFFF", "CYAN"]:
                    run.font.highlight_color = WD_COLOR_INDEX.TURQUOISE
                elif hl_color in ["FF00FF", "MAGENTA", "PINK"]:
                    run.font.highlight_color = WD_COLOR_INDEX.PINK
                elif hl_color in ["FF0000", "RED"]:
                    run.font.highlight_color = WD_COLOR_INDEX.RED
                elif hl_color in ["0000FF", "BLUE"]:
                    run.font.highlight_color = WD_COLOR_INDEX.BLUE
                else:
                    run.font.highlight_color = WD_COLOR_INDEX.YELLOW


    def _op_docx_table(self, doc, tgt: dict, params: dict) -> str:
        action = params.get("action", "create")
        
        if action == "create":
            data = params.get("data") or []
            if data and not isinstance(data[0], list):
                data = [[x] for x in data]
            rows = params.get("rows") or (len(data) if data else 3)
            cols = params.get("cols") or (max((len(r) for r in data), default=0) if data else 3)
            rows = max(int(rows), 1)
            cols = max(int(cols), 1)
            tbl = doc.add_table(rows=rows, cols=cols)
            if data:
                for r_idx, row_data in enumerate(data[:rows]):
                    for c_idx, val in enumerate(row_data[:cols]):
                        tbl.rows[r_idx].cells[c_idx].text = str(val)
            
            before_id = params.get("before_id")
            after_id = params.get("after_id")
            if before_id or after_id:
                body_index = self._build_docx_body_index(doc)
                id_to_pos = {eid: i for i, (eid, _) in enumerate(body_index)}
                
                target_xml_el = None
                insert_before = True
                if before_id:
                    pos = id_to_pos.get(before_id, -1)
                    if pos != -1:
                        target_xml_el = body_index[pos][1]
                elif after_id:
                    pos = id_to_pos.get(after_id, -1)
                    if pos != -1:
                        target_xml_el = body_index[pos][1]
                        insert_before = False
                        
                if target_xml_el is not None:
                    tbl_xml = tbl._tbl
                    tbl_xml.getparent().remove(tbl_xml)
                    if insert_before:
                        target_xml_el.addprevious(tbl_xml)
                    else:
                        target_xml_el.addnext(tbl_xml)

            return f"Created {rows}x{cols} table" + (" with data" if data else "")
            
        is_all = tgt.get("type") == "all"
        table_idx = tgt.get("table_index")
        
        if not is_all and (table_idx is None or table_idx < 0 or table_idx >= len(doc.tables)):
            return "Table not found."
            
        tables_to_process_indices = range(len(doc.tables)) if is_all else [table_idx]
        
        summaries = []
        for t_idx in reversed(tables_to_process_indices):
            tbl = doc.tables[t_idx]
            
            if action == "delete":
                tbl._element.getparent().remove(tbl._element)
                summaries.append(f"Deleted table {t_idx}")
                continue

            elif action == "add_row":
                new_row = tbl.add_row()
                data = params.get("data", [])
                if data:
                    # Support both 1D list and 2D list (first row)
                    row_data = data[0] if data and isinstance(data[0], list) else data
                    for i, cell in enumerate(new_row.cells):
                        if i < len(row_data):
                            cell.text = str(row_data[i])
                summaries.append(f"Added row to table {t_idx}")

            elif action == "remove_row":
                r_idx = params.get("row_index")
                if r_idx is not None and 0 <= r_idx < len(tbl.rows):
                    row = tbl.rows[r_idx]
                    row._element.getparent().remove(row._element)
                    summaries.append(f"Removed row {r_idx} from table {t_idx}")

            elif action == "add_col":
                new_col = tbl.add_column(Pt(72)) # Default 1 inch
                data = params.get("data", [])
                if data:
                    # Support both 1D list and 2D list
                    col_data = []
                    if data and isinstance(data[0], list):
                        # Transpose if it's a 2D list that looks like a row
                        col_data = [r[0] if len(r) > 0 else "" for r in data]
                    else:
                        col_data = data
                    for i, cell in enumerate(new_col.cells):
                        if i < len(col_data):
                            cell.text = str(col_data[i])
                summaries.append(f"Added column to table {t_idx}")

            elif action == "remove_col":
                c_idx = params.get("col_index")
                if c_idx is not None and len(tbl.rows) > 0 and 0 <= c_idx < len(tbl.rows[0].cells):
                    # Remove cell from each row
                    for row in tbl.rows:
                        cell = row.cells[c_idx]
                        cell._element.getparent().remove(cell._element)
                    
                    # Remove from tblGrid to prevent corruption
                    tblGrid = tbl._element.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblGrid")
                    if tblGrid is not None:
                        gridCols = tblGrid.findall("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}gridCol")
                        if 0 <= c_idx < len(gridCols):
                            tblGrid.remove(gridCols[c_idx])
                            
                    summaries.append(f"Removed column {c_idx} from table {t_idx}")

            elif action == "populate":
                data = params.get("data", [])
                if data:
                    # If it's a 1D list, wrap it in a 2D list to process cleanly
                    if data and not isinstance(data[0], list):
                        data = [[x] for x in data]
                        
                    start_row = params.get("row_index")
                    start_col = params.get("col_index")
                    
                    if start_col is None and len(tbl.rows) > 0:
                        # Auto-detect if they meant to append to the end by guessing from the row
                        start_col = 0
                        # But wait, if they didn't specify, we use default 0.
                        # Actually we can do a smart check: if the first cell is not empty, maybe they didn't specify correctly?
                        None
                        
                    start_row = start_row or 0
                    start_col = start_col or 0
                    
                    for r_idx, row_data in enumerate(data):
                        tr_idx = start_row + r_idx
                        if tr_idx >= len(tbl.rows):
                            break
                        for c_idx, val in enumerate(row_data):
                            tc_idx = start_col + c_idx
                            if tc_idx >= len(tbl.rows[tr_idx].cells):
                                break
                            tbl.rows[tr_idx].cells[tc_idx].text = str(val)
                summaries.append(f"Populated data in table {t_idx}")

            elif action == "merge_cells":
                f_idx = params.get("merge_from", [0, 0])
                t_idx_merge = params.get("merge_to", [0, 1])
                try:
                    c1 = tbl.cell(f_idx[0], f_idx[1])
                    c2 = tbl.cell(t_idx_merge[0], t_idx_merge[1])
                    c1.merge(c2)
                    summaries.append(f"Merged cells in table {t_idx}")
                except Exception as e:
                    summaries.append(f"Failed to merge cells: {e}")

            elif action == "set_alignment":
                alignment = params.get("alignment", "center")
                from docx.enum.table import WD_TABLE_ALIGNMENT
                align_map = {
                    "left": WD_TABLE_ALIGNMENT.LEFT,
                    "center": WD_TABLE_ALIGNMENT.CENTER,
                    "right": WD_TABLE_ALIGNMENT.RIGHT,
                }
                tbl.alignment = align_map.get(alignment, WD_TABLE_ALIGNMENT.CENTER)
                summaries.append(f"Set alignment of table {t_idx} to {alignment}")

            elif action == "set_width_pct":
                width_pct = float(params.get("width_pct", 1.0))
                if width_pct > 1.0:
                    width_pct = width_pct / 100.0  # Handle cases where LLM passes 100 instead of 1.0
                try:
                    section = doc.sections[0]
                    usable_width = section.page_width - section.left_margin - section.right_margin
                except Exception:
                    usable_width = int(5940000) # Fallback to ~A4 EMU
                    
                target_width = int(usable_width * width_pct)
                tbl.autofit = False
                col_count = len(tbl.rows[0].cells) if len(tbl.rows) > 0 else 1
                col_width = int(target_width / max(1, col_count))
                
                for row in tbl.rows:
                    for cell in row.cells:
                        cell.width = col_width
                
                summaries.append(f"Set width of table {t_idx} to {int(width_pct*100)}%")

            elif action in ["set_cell_bg", "alternate_rows", "set_header_format", "apply_theme"]:
                # Helper to safely set background color without duplicating w:shd
                from docx.oxml.ns import nsdecls, qn
                from docx.oxml import parse_xml
                def _set_cell_bg(cell, clr):
                    tcPr = cell._tc.get_or_add_tcPr()
                    for existing_shd in tcPr.findall(qn('w:shd')):
                        tcPr.remove(existing_shd)
                    shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'), clr))
                    tcPr.append(shading_elm)

                if action == "set_header_format":
                    if len(tbl.rows) > 0:
                        from docx.shared import RGBColor as DRGBColor
                        for cell in tbl.rows[0].cells:
                            _set_cell_bg(cell, "000080") # dark blue
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = DRGBColor(255, 255, 255)
                        summaries.append(f"Formatted header for table {t_idx}")
                        continue
                
                if action == "set_cell_bg":
                    row_idx = params.get("row_index")
                    col_idx = params.get("col_index")
                    color = str(params.get("cell_bg_hex", "")).strip().lstrip("#").upper()
                    if color:
                        if row_idx is not None and col_idx is not None:
                            if 0 <= row_idx < len(tbl.rows) and 0 <= col_idx < len(tbl.rows[row_idx].cells):
                                _set_cell_bg(tbl.cell(row_idx, col_idx), color)
                        elif row_idx is not None:
                            if 0 <= row_idx < len(tbl.rows):
                                for c in tbl.rows[row_idx].cells:
                                    _set_cell_bg(c, color)
                        elif col_idx is not None:
                            for r in tbl.rows:
                                if 0 <= col_idx < len(r.cells):
                                    _set_cell_bg(r.cells[col_idx], color)
                    summaries.append(f"Set cell background for table {t_idx}")
                    continue

                if action == "alternate_rows":
                    colors = params.get("alternate_row_colors", ["FFFFFF", "F2F2F2"])
                    if len(colors) == 2:
                        for i, row in enumerate(tbl.rows):
                            color = colors[i % 2].lstrip("#")
                            for cell in row.cells:
                                _set_cell_bg(cell, color)
                        summaries.append(f"Applied alternate row colors to table {t_idx}")
                        continue

                if action == "apply_theme":
                    theme_color = params.get("theme_color_hex", "4F81BD").lstrip("#").upper()
                    if len(tbl.rows) > 0:
                        for cell in tbl.rows[0].cells:
                            _set_cell_bg(cell, theme_color)
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.bold = True
                                    try:
                                        from docx.shared import RGBColor as DRGBColor
                                        run.font.color.rgb = DRGBColor(255, 255, 255)
                                    except Exception:
                                        None
                    if len(tbl.rows) > 1:
                        for i, row in enumerate(tbl.rows[1:]):
                            color = "FFFFFF" if i % 2 == 0 else "F9F9F9"
                            for cell in row.cells:
                                _set_cell_bg(cell, color)
                    summaries.append(f"Applied theme #{theme_color} to table {t_idx}")
                    continue

            elif action == "sort_data":
                if len(tbl.rows) > 1:
                    # Find sort column. Default is 0. If user provided col_index, use it.
                    sort_col_idx = params.get("col_index")
                    if sort_col_idx is None:
                        # If the user passed sort_by_column in the params even if it's not strictly in schema
                        sort_by_name = str(params.get("sort_by_column", "")).lower()
                        sort_col_idx = 0
                        if sort_by_name:
                            for c_i, c in enumerate(tbl.rows[0].cells):
                                if c.text.strip().lower() == sort_by_name:
                                    sort_col_idx = c_i
                                    break

                    header = [c.text for c in tbl.rows[0].cells]
                    data = [[c.text for c in r.cells] for r in tbl.rows[1:]]
                    try:
                        # Sort by the chosen column, handling numeric vs string
                        def _sort_key(x):
                            val = x[sort_col_idx] if 0 <= sort_col_idx < len(x) else ""
                            try:
                                # Strip currency/percent for sorting
                                clean = val.replace("$", "").replace("%", "").replace(",", "")
                                return (0, float(clean))
                            except ValueError:
                                return (1, val)
                        data.sort(key=_sort_key)
                        for i, row_data in enumerate(data):
                            for j, val in enumerate(row_data):
                                tbl.rows[i+1].cells[j].text = val
                        summaries.append(f"Sorted table {t_idx}")
                    except:
                        None

            elif action == "set_borders":
                # We handle borders at the end of the loop now
                None

            elif action == "set_cell_alignment":
                alignment = params.get("cell_alignment")
                row_idx = params.get("row_index")
                col_idx = params.get("col_index")
                if alignment:
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    align_map = {
                        "left": WD_ALIGN_PARAGRAPH.LEFT,
                        "center": WD_ALIGN_PARAGRAPH.CENTER,
                        "right": WD_ALIGN_PARAGRAPH.RIGHT,
                        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
                    }
                    if alignment in align_map:
                        def _set_align(cell):
                            for p in cell.paragraphs:
                                p.alignment = align_map[alignment]

                        if row_idx is not None and col_idx is not None:
                            if 0 <= row_idx < len(tbl.rows) and 0 <= col_idx < len(tbl.rows[row_idx].cells):
                                _set_align(tbl.cell(row_idx, col_idx))
                        elif row_idx is not None:
                            if 0 <= row_idx < len(tbl.rows):
                                for c in tbl.rows[row_idx].cells:
                                    _set_align(c)
                        elif col_idx is not None:
                            for r in tbl.rows:
                                if 0 <= col_idx < len(r.cells):
                                    _set_align(r.cells[col_idx])
                        else:
                            for row in tbl.rows:
                                for cell in row.cells:
                                    _set_align(cell)
                    summaries.append(f"Aligned text in table {t_idx} to {alignment}")

            # --- Universal Table Modifiers ---

            valign = params.get("cell_vertical_alignment")
            if valign:
                from docx.enum.table import WD_ALIGN_VERTICAL
                valign_map = {
                    "top": WD_ALIGN_VERTICAL.TOP,
                    "center": WD_ALIGN_VERTICAL.CENTER,
                    "bottom": WD_ALIGN_VERTICAL.BOTTOM
                }
                if valign in valign_map:
                    for row in tbl.rows:
                        for cell in row.cells:
                            cell.vertical_alignment = valign_map[valign]
                    summaries.append(f"Set vertical alignment to {valign} for table {t_idx}")
            
            border_hex = params.get("border_color_hex")
            if border_hex:
                border_color = str(border_hex).lstrip("#").upper()
                border_sz = str(int(params.get("border_width_pt", 12)))
                from docx.oxml.shared import OxmlElement, qn
                tblPr = tbl._element.xpath('w:tblPr')
                if tblPr:
                    tblBorders = OxmlElement('w:tblBorders')
                    for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                        border = OxmlElement(f'w:{border_name}')
                        border.set(qn('w:val'), 'single')
                        border.set(qn('w:sz'), border_sz)
                        border.set(qn('w:space'), '0')
                        border.set(qn('w:color'), border_color)
                        tblBorders.append(border)
                    # replace existing
                    old_borders = tblPr[0].xpath('w:tblBorders')
                    if old_borders:
                        tblPr[0].replace(old_borders[0], tblBorders)
                    else:
                        tblPr[0].append(tblBorders)
                summaries.append(f"Applied {border_color} borders to table {t_idx}")

            col_width = params.get("col_width_inches")
            if col_width is not None:
                from docx.shared import Inches
                for row in tbl.rows:
                    for cell in row.cells:
                        cell.width = Inches(float(col_width))
                summaries.append(f"Set column width to {col_width} inches for table {t_idx}")

        return "; ".join(summaries)

    def _op_docx_ai_design(self, doc, params: dict) -> str:
        action = params.get("action", "normalize_fonts")
        
        if action == "normalize_fonts":
            from docx.shared import Pt
            font_family = params.get("target_font", "Calibri")
            base_size = params.get("base_font_size_pt", 11)
            for para in doc.paragraphs:
                for run in para.runs:
                    run.font.name = font_family
                    if not para.style.name.startswith("Heading"):
                        run.font.size = Pt(base_size)
            for tbl in doc.tables:
                for row in tbl.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.font.name = font_family
                                run.font.size = Pt(base_size)
            return f"Normalized fonts to {font_family} {base_size}pt"
            
        return f"Applied AI design: {action}"

    def _op_pptx_table(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Create, delete, or edit a table on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "create")

        if action == "create":
            rows = params.get("rows", 3)
            cols = params.get("cols", 3)
            pos = params.get("position", {})

            sw = prs.slide_width
            sh = prs.slide_height

            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.3))
            width = int(sw * pos.get("width_pct", 0.8))
            height = int(sh * pos.get("height_pct", 0.5))

            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Populate with data if provided
            data = params.get("data", [])
            for r_idx, row_data in enumerate(data):
                if r_idx >= rows:
                    break
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx >= cols:
                        break
                    table.cell(r_idx, c_idx).text = str(cell_text)

            # Apply header row formatting
            if params.get("header_row") and rows > 0:
                from pptx.dml.color import RGBColor
                for c_idx in range(cols):
                    cell = table.cell(0, c_idx)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.bold = True

            # Apply alternate row colors
            alt_colors = params.get("alternate_row_colors")
            if alt_colors and len(alt_colors) >= 2:
                from pptx.oxml.ns import qn
                from lxml import etree
                start_row = 1 if params.get("header_row") else 0
                for r_idx in range(start_row, rows):
                    color_hex = alt_colors[(r_idx - start_row) % 2]
                    for c_idx in range(cols):
                        cell = table.cell(r_idx, c_idx)
                        try:
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                            srgbClr.set("val", str(color_hex).upper().lstrip("#"))
                        except Exception:
                            None

            return f"Created {rows}×{cols} table on slide {tgt.get('slide')}"

        elif action == "delete":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Deleted table (shape {shape_idx}) from slide {tgt.get('slide')}"
            return ""

        elif action in ("add_row", "remove_row", "add_col", "remove_col"):
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                table = shape.table
                if action == "add_row":
                    # Clone last row
                    from lxml import etree
                    tbl_elem = table._tbl
                    last_tr = tbl_elem.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tr")[-1]
                    new_tr = copy.deepcopy(last_tr)
                    for tc in new_tr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tc"):
                        for t in tc.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                            t.text = ""
                    tbl_elem.append(new_tr)
                    return f"Added row to table on slide {tgt.get('slide')}"
                elif action == "remove_row":
                    row_idx = params.get("row_index", -1)
                    tbl_elem = table._tbl
                    rows = tbl_elem.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tr")
                    if rows and abs(row_idx) < len(rows):
                        tbl_elem.remove(rows[row_idx])
                        return f"Removed row from table on slide {tgt.get('slide')}"

        elif action == "set_cell_bg":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                r = params.get("row_index", 0)
                c = params.get("col_index", 0)
                color_hex = params.get("cell_bg_hex", "FFFFFF")
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    cell = shape.table.cell(r, c)
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(color_hex).upper().lstrip("#"))
                    return f"Set cell background on slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("set_cell_bg failed: %s", e)

        elif action == "populate":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape and getattr(shape, "has_table", False):
                data = params.get("data", [])
                table = shape.table
                for r_idx, row_data in enumerate(data):
                    if r_idx >= len(table.rows):
                        break
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx >= len(table.columns):
                            break
                        table.cell(r_idx, c_idx).text = str(cell_text)
                return f"Populated table on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_image(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Insert, replace, resize or style images on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "insert")
        image_path = params.get("image_path")

        sw = prs.slide_width
        sh = prs.slide_height

        if action == "insert" and image_path:
            if not Path(image_path).exists():
                log.warning("Image file not found: %s", image_path)
                return ""

            pos = params.get("position", {})
            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.15))
            width = int(sw * pos.get("width_pct", 0.4))
            height = int(sh * pos.get("height_pct", 0.5))

            maintain_ar = params.get("maintain_aspect_ratio", True)
            if maintain_ar:
                pic = slide.shapes.add_picture(image_path, left, top, width=width)
            else:
                pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

            # Apply border if requested
            border_hex = params.get("border_color_hex")
            border_pt = params.get("border_width_pt", 1.5)
            if border_hex:
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = pic._element.spPr
                    ln = etree.SubElement(spPr, qn("a:ln"))
                    ln.set("w", str(int(Pt(border_pt))))
                    solidFill = etree.SubElement(ln, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(border_hex).upper().lstrip("#"))
                except Exception as e:
                    log.debug("Image border failed: %s", e)

            # Shadow
            if params.get("shadow"):
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = pic._element.spPr
                    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
                    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
                    outerShdw.set("blurRad", "60960")
                    outerShdw.set("dist", "114300")
                    outerShdw.set("dir", "2700000")
                    outerShdw.set("algn", "tl")
                    outerShdw.set("rotWithShape", "0")
                    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
                    srgbClr.set("val", "000000")
                    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                    alpha.set("val", "40000")
                except Exception as e:
                    log.debug("Image shadow failed: %s", e)

            return f"Inserted image on slide {tgt.get('slide')}"

        elif action == "remove":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Removed image from slide {tgt.get('slide')}"

        elif action in ("resize", "reposition"):
            shape_idx = tgt.get("shape_index")
            shape = self._get_shape(slide, shape_idx)
            if shape:
                pos = params.get("position", {})
                if "left_pct" in pos:
                    shape.left = int(sw * pos["left_pct"])
                if "top_pct" in pos:
                    shape.top = int(sh * pos["top_pct"])
                if "width_pct" in pos:
                    shape.width = int(sw * pos["width_pct"])
                if "height_pct" in pos:
                    shape.height = int(sh * pos["height_pct"])
                return f"Repositioned/resized shape {shape_idx} on slide {tgt.get('slide')}"

        elif action == "rotate":
            shape_idx = tgt.get("shape_index")
            shape = self._get_shape(slide, shape_idx)
            if shape:
                degrees = params.get("rotation_degrees", 0)
                shape.rotation = degrees
                return f"Rotated shape {shape_idx} by {degrees}° on slide {tgt.get('slide')}"

        elif action in ("bring_forward", "send_backward"):
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes_list = list(slide.shapes)
                if 0 <= shape_idx < len(shapes_list):
                    sp_elem = shapes_list[shape_idx]._element
                    sp_tree = sp_elem.getparent()
                    if action == "bring_forward":
                        sp_tree.append(sp_elem)
                    else:
                        children = list(sp_tree)
                        first_sp_idx = next(
                            (i for i, c in enumerate(children) if c.tag.endswith("}sp") or c.tag.endswith("}pic")),
                            0
                        )
                        sp_tree.insert(first_sp_idx, sp_elem)
                    return f"Moved shape {shape_idx} {action.replace('_', ' ')} on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_shape(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Add or edit shapes and text boxes on a slide."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        action = params.get("action", "add_textbox")
        sw = prs.slide_width
        sh = prs.slide_height

        if action == "add_textbox":
            pos = params.get("position", {})
            left = int(sw * pos.get("left_pct", 0.1))
            top = int(sh * pos.get("top_pct", 0.1))
            width = int(sw * pos.get("width_pct", 0.4))
            height = int(sh * pos.get("height_pct", 0.1))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = params.get("text", "")

            # Apply fill
            fill_hex = params.get("fill_color_hex")
            if fill_hex:
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    spPr = txBox._element.spPr
                    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
                    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgbClr.set("val", str(fill_hex).upper().lstrip("#"))
                except Exception:
                    None

            return f"Added text box on slide {tgt.get('slide')}"

        elif action == "delete":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    sp = shapes[shape_idx]._element
                    sp.getparent().remove(sp)
                    return f"Deleted shape {shape_idx} from slide {tgt.get('slide')}"

        elif action in ("resize", "move"):
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape:
                pos = params.get("position", {})
                if "left_pct" in pos:
                    shape.left = int(sw * pos["left_pct"])
                if "top_pct" in pos:
                    shape.top = int(sh * pos["top_pct"])
                if "width_pct" in pos:
                    shape.width = int(sw * pos["width_pct"])
                if "height_pct" in pos:
                    shape.height = int(sh * pos["height_pct"])
                return f"Moved/resized shape on slide {tgt.get('slide')}"

        elif action == "set_fill":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            fill_hex = params.get("fill_color_hex")
            if shape and fill_hex:
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(str(fill_hex).lstrip("#"))
                    return f"Set fill color on shape {tgt.get('shape_index')} slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("Set fill failed: %s", e)

        elif action == "set_outline":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            outline_hex = params.get("outline_color_hex")
            outline_pt = params.get("outline_width_pt", 1.5)
            if shape and outline_hex:
                try:
                    shape.line.color.rgb = RGBColor.from_string(str(outline_hex).lstrip("#"))
                    shape.line.width = Pt(outline_pt)
                    return f"Set outline on shape {tgt.get('shape_index')} slide {tgt.get('slide')}"
                except Exception as e:
                    log.debug("Set outline failed: %s", e)

        elif action == "duplicate":
            shape_idx = tgt.get("shape_index")
            if shape_idx is not None:
                shapes = list(slide.shapes)
                if 0 <= shape_idx < len(shapes):
                    new_sp = copy.deepcopy(shapes[shape_idx]._element)
                    slide.shapes._spTree.append(new_sp)
                    return f"Duplicated shape {shape_idx} on slide {tgt.get('slide')}"

        elif action == "rotate":
            shape = self._get_shape(slide, tgt.get("shape_index"))
            if shape:
                shape.rotation = params.get("rotation_degrees", 0)
                return f"Rotated shape on slide {tgt.get('slide')}"

        return ""

    def _op_pptx_theme(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Change slide backgrounds and theme colors."""
        action = params.get("action", "set_bg_color")
        scope = params.get("scope", "current_slide")
        slide_num = tgt.get("slide")

        slides_to_update = list(prs.slides) if scope == "all_slides" else []
        if not slides_to_update and slide_num:
            s = self._get_slide(prs, slide_num)
            if s:
                slides_to_update = [s]
        if not slides_to_update:
            slides_to_update = list(prs.slides)

        if action == "set_bg_color":
            bg_hex = params.get("bg_color_hex", "FFFFFF")
            from pptx.oxml.ns import qn
            from lxml import etree

            for slide in slides_to_update:
                try:
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(str(bg_hex).upper().lstrip("#"))
                except Exception as e:
                    log.debug("set_bg_color failed for slide: %s", e)

            scope_label = "all slides" if scope == "all_slides" else f"slide {slide_num}"
            return f"Set background color #{bg_hex} on {scope_label}"

        elif action == "set_bg_gradient":
            start_hex = params.get("gradient_start_hex", "1a1a2e")
            end_hex = params.get("gradient_end_hex", "16213e")
            direction = params.get("gradient_direction", "diagonal")

            from pptx.oxml.ns import qn
            from lxml import etree

            angle_map = {"horizontal": "5400000", "vertical": "10800000", "diagonal": "2700000"}
            angle = angle_map.get(direction, "5400000")

            for slide in slides_to_update:
                try:
                    bg = slide.background
                    bgPr = bg._element.get_or_add_bgPr()
                    # Clear existing fill
                    for child in list(bgPr):
                        bgPr.remove(child)
                    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
                    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

                    for pos_val, hex_color in [("0", start_hex), ("100000", end_hex)]:
                        gs = etree.SubElement(gsLst, qn("a:gs"))
                        gs.set("pos", pos_val)
                        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
                        srgbClr.set("val", str(hex_color).upper().lstrip("#"))

                    lin = etree.SubElement(gradFill, qn("a:lin"))
                    lin.set("ang", angle)
                    lin.set("scaled", "0")
                except Exception as e:
                    log.debug("Gradient bg failed: %s", e)

            return f"Set gradient background ({direction}) on {scope}"

        return ""

    def _op_pptx_slide(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Add, delete, duplicate, or reorder slides."""
        action = params.get("action", "add")
        slide_num = tgt.get("slide")

        if action == "add":
            # Use layout from an existing slide
            after_idx = (params.get("after_index") or slide_num or len(list(prs.slides))) - 1
            slides = list(prs.slides)
            src_idx = min(after_idx, len(slides) - 1)
            self._clone_slide(prs, src_idx)

            # Move the new slide to the correct position
            new_idx = len(list(prs.slides)) - 1
            target_pos = min(after_idx + 1, new_idx)
            if target_pos < new_idx:
                self._reorder_slide(prs, new_idx, target_pos)
            return f"Added new slide after position {after_idx + 1}"

        elif action == "delete":
            idx = (slide_num or 1) - 1
            slides = list(prs.slides)
            if 0 <= idx < len(slides):
                self._delete_slide(prs, idx)
                return f"Deleted slide {slide_num}"

        elif action == "duplicate":
            after_idx = (params.get("after_index") or slide_num or 1) - 1
            slides = list(prs.slides)
            src_idx = (slide_num or 1) - 1
            src_idx = min(src_idx, len(slides) - 1)
            self._clone_slide(prs, src_idx)
            new_idx = len(list(prs.slides)) - 1
            target_pos = after_idx + 1
            if target_pos < new_idx:
                self._reorder_slide(prs, new_idx, target_pos)
            return f"Duplicated slide {slide_num or 1}"

        elif action == "reorder":
            from_idx = (params.get("from_index") or slide_num or 1) - 1
            to_idx = (params.get("to_index") or 1) - 1
            slides = list(prs.slides)
            if 0 <= from_idx < len(slides) and 0 <= to_idx < len(slides):
                self._reorder_slide(prs, from_idx, to_idx)
                return f"Moved slide {from_idx + 1} to position {to_idx + 1}"

        elif action == "rename_title":
            slide = self._get_slide(prs, slide_num)
            title = params.get("title", "")
            if slide and title:
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if shape.placeholder_format.type in (
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ):
                                shape.text = title
                                return f"Renamed title of slide {slide_num} to '{title}'"
                        except Exception:
                            None

        elif action in ("hide", "unhide"):
            # PPTX doesn't natively hide slides via python-pptx API simply,
            # but we can set the show attribute in the slide list XML
            slide = self._get_slide(prs, slide_num)
            if slide:
                try:
                    from pptx.oxml.ns import qn
                    sldIdLst = prs.part._element.find(
                        "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
                    )
                    for rel in prs.part.rels.values():
                        if rel.target_part == slide.part:
                            rId = rel.rId
                            if sldIdLst is not None:
                                for sldId in sldIdLst:
                                    ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                                    if sldId.get(f"{{{ns}}}id") == rId:
                                        sldId.set("show", "0" if action == "hide" else "1")
                            return f"{'Hid' if action == 'hide' else 'Unhid'} slide {slide_num}"
                except Exception as e:
                    log.debug("Hide/unhide slide failed: %s", e)

        return ""

    def _reorder_slide(self, prs: Presentation, from_idx: int, to_idx: int) -> None:
        """Reorder a slide within the presentation."""
        slides = list(prs.slides)
        if from_idx < 0 or from_idx >= len(slides):
            return
        if to_idx < 0 or to_idx >= len(slides):
            return

        pres_elem = prs.part._element
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sldIdLst = pres_elem.find(f"{{{ns}}}sldIdLst")
        if sldIdLst is None:
            return

        sldIds = list(sldIdLst)
        if from_idx >= len(sldIds) or to_idx >= len(sldIds):
            return

        elem = sldIds[from_idx]
        sldIdLst.remove(elem)
        sldIds.pop(from_idx)
        sldIds.insert(to_idx, elem)
        for s in sldIds:
            sldIdLst.remove(s)
        for s in sldIds:
            sldIdLst.append(s)

    def _op_pptx_chart(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """Edit chart type, colors, labels. Returns summary string."""
        slide = self._get_slide(prs, tgt.get("slide"))
        if not slide:
            return ""

        shape_idx = tgt.get("shape_index")
        shape = self._get_shape(slide, shape_idx)
        if not shape:
            # Try to find any chart shape
            for s in slide.shapes:
                if s.has_chart:
                    shape = s
                    break
        if not shape or not shape.has_chart:
            return ""

        action = params.get("action", "update_labels")
        chart = shape.chart

        if action == "update_axis_labels":
            try:
                x_label = params.get("x_axis_label")
                y_label = params.get("y_axis_label")
                if x_label and chart.has_plot_area:
                    None  # python-pptx limited chart axis label support
                return f"Updated chart axis labels on slide {tgt.get('slide')}"
            except Exception:
                None

        elif action == "set_series_colors":
            colors = params.get("series_colors", [])
            try:
                for i, series in enumerate(chart.series):
                    if i < len(colors):
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor.from_string(
                            str(colors[i]).upper().lstrip("#")
                        )
                return f"Updated series colors on chart, slide {tgt.get('slide')}"
            except Exception as e:
                log.debug("Chart series color failed: %s", e)

        elif action in ("show_legend", "hide_legend"):
            try:
                chart.has_legend = (action == "show_legend")
                return f"{'Showed' if action == 'show_legend' else 'Hid'} chart legend on slide {tgt.get('slide')}"
            except Exception as e:
                log.debug("Chart legend failed: %s", e)

        return f"Edited chart on slide {tgt.get('slide')}"

    def _op_pptx_ai_design(self, prs: Presentation, tgt: dict, params: dict) -> str:
        """AI-driven design normalization operations."""
        action = params.get("action", "normalize_fonts")
        scope = params.get("scope", "all_slides")

        # Determine which slides to process
        if scope == "all_slides":
            slides = list(prs.slides)
        elif scope and scope.startswith("slide:"):
            try:
                n = int(scope.split(":")[1])
                slide = self._get_slide(prs, n)
                slides = [slide] if slide else list(prs.slides)
            except Exception:
                slides = list(prs.slides)
        else:
            slides = list(prs.slides)

        if action == "normalize_fonts":
            target_font = params.get("target_font", "Calibri")
            base_size = params.get("base_font_size_pt", 18)
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.name != target_font:
                                run.font.name = target_font
            return f"Normalized fonts to {target_font} across {scope}"

        elif action == "normalize_spacing":
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        try:
                            from pptx.oxml.ns import qn
                            from lxml import etree
                            pPr = para._p.get_or_add_pPr()
                            spcBef = pPr.find(qn("a:spcBef"))
                            if spcBef is None:
                                spcBef = etree.SubElement(pPr, qn("a:spcBef"))
                            spcPts = spcBef.find(qn("a:spcPts"))
                            if spcPts is None:
                                spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
                            spcPts.set("val", "600")  # 6pt before paragraph
                        except Exception:
                            None
            return f"Normalized paragraph spacing across {scope}"

        elif action == "generate_speaker_notes":
            # Add basic speaker notes placeholder to each slide
            for slide in slides:
                try:
                    notes_slide = slide.notes_slide
                    tf = notes_slide.notes_text_frame
                    if not tf.text.strip():
                        # Collect slide content for context
                        content = " | ".join(
                            s.text_frame.text[:100]
                            for s in slide.shapes
                            if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
                        )
                        tf.text = f"Slide content: {content}\n[Speaker notes generated by AI]"
                except Exception as e:
                    log.debug("Speaker notes failed: %s", e)
            return f"Generated speaker notes for {scope}"

        elif action == "improve_hierarchy":
            # Ensure title shapes are bold and large; body shapes are smaller
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    is_title = shape.is_placeholder and shape.placeholder_format.type in (1, 3)
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if is_title:
                                run.font.bold = True
                                if run.font.size and run.font.size.pt < 24:
                                    run.font.size = Pt(28)
                            else:
                                if run.font.size and run.font.size.pt > 28:
                                    run.font.size = Pt(18)
            return f"Improved visual hierarchy across {scope}"

        elif action == "auto_resize_text":
            for slide in slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        try:
                            shape.text_frame.auto_size = 1  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            None
            return f"Enabled auto-resize text across {scope}"

        elif action == "improve_readability":
            # Ensure minimum font sizes and line spacing
            for slide in slides:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt < 12:
                                run.font.size = Pt(12)
            return f"Improved readability across {scope}"

        return f"Applied AI design ({action}) across {scope}"

    # ------------------------------------------------------------------
    # DOCX: one block per paragraph
    # ------------------------------------------------------------------

    def _resolve_docx_font_info(self, para, run) -> tuple[float | None, str | None]:
        font_size = None
        color_hex = None
        
        if run.font.size:
            font_size = run.font.size.pt
        try:
            if run.font.color and run.font.color.rgb:
                color_hex = str(run.font.color.rgb)
        except Exception:
            None
            
        if font_size is None or color_hex is None:
            style = para.style
            while style:
                if font_size is None and style.font.size:
                    font_size = style.font.size.pt
                if color_hex is None:
                    try:
                        if style.font.color and style.font.color.rgb:
                            color_hex = str(style.font.color.rgb)
                    except Exception:
                        None
                
                if font_size is not None and color_hex is not None:
                    break
                    
                style = style.base_style
                
        return font_size, color_hex

    def _extract_docx_dom(self, path: Path) -> dict:
        """Extract DOCX DOM in true document body order using stable UUIDs."""
        from docx.oxml.ns import qn as _qn
        WNS_P = _qn('w:p')
        WNS_TBL = _qn('w:tbl')
        WNS_DRAWING = _qn('w:drawing')
        uid_attr = f'{{{CUSTOM_NS}}}uid'

        doc = Document(path)
        
        # --- UID STAMPING ---
        any_uids_assigned = False
        for child in doc.element.body:
            if child.tag in (WNS_P, WNS_TBL):
                if child.get(uid_attr) is None:
                    child.set(uid_attr, uuid.uuid4().hex[:8])
                    any_uids_assigned = True
                    
        if any_uids_assigned:
            import tempfile
            import shutil
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".docx", dir=path.parent)
            import os
            os.close(tmp_fd)
            doc.save(tmp_path)
            shutil.move(tmp_path, path)
        # --------------------

        children = []

        metadata_node = extract_metadata(doc)
        if metadata_node:
            children.append({
                "id": "metadata",
                "type": "metadata",
                "role": "metadata",
                "properties": metadata_node
            })
        
        # Extract basic section info
        for s_idx, section in enumerate(doc.sections):
            sec_style = {}
            try:
                sec_style["orientation"] = section.orientation.name if hasattr(section, 'orientation') else "PORTRAIT"
                if hasattr(section.page_width, 'inches'): sec_style["page_width_inches"] = round(section.page_width.inches, 2)
                if hasattr(section.page_height, 'inches'): sec_style["page_height_inches"] = round(section.page_height.inches, 2)
                if hasattr(section.top_margin, 'inches'): sec_style["top_margin_inches"] = round(section.top_margin.inches, 2)
                if hasattr(section.bottom_margin, 'inches'): sec_style["bottom_margin_inches"] = round(section.bottom_margin.inches, 2)
                if hasattr(section.left_margin, 'inches'): sec_style["left_margin_inches"] = round(section.left_margin.inches, 2)
                if hasattr(section.right_margin, 'inches'): sec_style["right_margin_inches"] = round(section.right_margin.inches, 2)
            except Exception:
                None
            
            children.append({
                "id": f"section_{s_idx}",
                "type": "section",
                "role": "section",
                "style": sec_style
            })

        # Build lookup maps: xml element → python-docx object
        para_elements = {p._p: p for p in doc.paragraphs}
        table_elements = {t._tbl: t for t in doc.tables}

        body_index = 0  # sequential position across all body children

        for child in doc.element.body:
            tag = child.tag

            if tag == WNS_P and child in para_elements:
                para = para_elements[child]
                uid = child.get(uid_attr)

                # ---- Check if this paragraph contains an inline image ----
                has_image = child.find(f'.//{WNS_DRAWING}') is not None

                if has_image:
                    # Extract image dimensions from the drawing element
                    width_emu, height_emu = 0, 0
                    description = ""
                    image_name = ""
                    title = ""
                    try:
                        drawing_el = child.find(f'.//{WNS_DRAWING}')
                        extent_el = drawing_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}extent')
                        if extent_el is not None:
                            width_emu = int(extent_el.get('cx', 0))
                            height_emu = int(extent_el.get('cy', 0))
                        # Try to get alt text / description
                        docPr_el = drawing_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}docPr')
                        if docPr_el is not None:
                            description = docPr_el.get('descr', '') or ""
                            image_name = docPr_el.get('name', '') or ""
                            title = docPr_el.get('title', '') or ""
                    except Exception:
                        None

                    alignment_str = None
                    if para.alignment is not None:
                        try:
                            alignment_str = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                        except Exception:
                            None

                    children.append({
                        "id": f"img_{uid}",
                        "body_index": body_index,
                        "type": "image",
                        "role": "inline_image",
                        "width_emu": width_emu,
                        "height_emu": height_emu,
                        "description": description,
                        "name": image_name,
                        "title": title,
                        "text": " ".join(
                            x for x in (description, title, image_name) if x
                        ).strip(),
                        "alignment": alignment_str,
                    })

                else:
                    # Regular text paragraph
                    runs = []
                    for r_idx, run in enumerate(para.runs):
                        font_size, color_hex = self._resolve_docx_font_info(para, run)
                        font = run.font
                        runs.append({
                            "id": f"p_{uid}_run_{r_idx}",
                            "type": "run",
                            "text": run.text,
                            "style": {
                                "font": font.name,
                                "size": font_size,
                                "bold": font.bold,
                                "italic": font.italic,
                                "underline": font.underline,
                                "strike": font.strike,
                                "highlight": str(font.highlight_color) if font.highlight_color else None,
                                "color": color_hex,
                            },
                        })

                    alignment = None
                    if para.alignment is not None:
                        try:
                            alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                        except Exception:
                            None

                    style_name = para.style.name if para.style else "Normal"
                    role = "body"
                    heading_level = None
                    if style_name.lower().startswith("heading"):
                        role = "heading"
                        import re as _re
                        m = _re.search(r'(\d+)$', style_name)
                        heading_level = int(m.group(1)) if m else 1
                    elif "bullet" in style_name.lower() or "list" in style_name.lower():
                        role = "bullet_point"

                    line_spacing = None
                    try:
                        line_spacing = para.paragraph_format.line_spacing
                    except Exception:
                        None

                    page_break_before = None
                    try:
                        page_break_before = para.paragraph_format.page_break_before
                    except Exception:
                        None

                    space_before_pt = None
                    try:
                        if para.paragraph_format.space_before is not None:
                            space_before_pt = para.paragraph_format.space_before.pt
                    except Exception:
                        None

                    space_after_pt = None
                    try:
                        if para.paragraph_format.space_after is not None:
                            space_after_pt = para.paragraph_format.space_after.pt
                    except Exception:
                        None
                    
                    adv_style = extract_advanced_paragraph_style(para)

                    node = {
                        "id": f"p_{uid}",
                        "body_index": body_index,
                        "type": "paragraph",
                        "role": role,
                        "text": para.text.strip(),
                        "style": {
                            "style_name": style_name,
                            "alignment": alignment,
                            "line_spacing": line_spacing,
                            "page_break_before": page_break_before,
                            "space_before_pt": space_before_pt,
                            "space_after_pt": space_after_pt,
                            **adv_style
                        },
                        "runs": runs,
                    }
                    
                    include_in_toc = (role == "heading")
                    if include_in_toc and para._p.pPr is not None:
                        from docx.oxml.ns import qn as _qn
                        outlineLvl = para._p.pPr.find(_qn('w:outlineLvl'))
                        if outlineLvl is not None and outlineLvl.get(_qn('w:val')) == '9':
                            include_in_toc = False
                    if include_in_toc or role == "heading":
                        node["include_in_toc"] = include_in_toc
                    
                    if heading_level is not None:
                        node["heading_level"] = heading_level

                    # Enrich with list metadata if this is a list paragraph
                    list_info = self._extract_list_info(para, doc)
                    if list_info is not None:
                        node["list_info"] = list_info
                        if role == "body":
                            # override role so LLM knows it's a list item
                            node["role"] = "bullet_point"

                    children.append(node)

            elif tag == WNS_TBL and child in table_elements:
                table = table_elements[child]
                uid = child.get(uid_attr)

                rows = []
                for r_idx, row in enumerate(table.rows):
                    cells = []
                    for c_idx, cell in enumerate(row.cells):
                        cell_paras = []
                        for p_idx, para in enumerate(cell.paragraphs):
                            runs = []
                            for r_run_idx, run in enumerate(para.runs):
                                font_size, color_hex = self._resolve_docx_font_info(para, run)
                                font = run.font
                                runs.append({
                                    "id": f"t_{uid}_cell_{r_idx}_{c_idx}_para_{p_idx}_run_{r_run_idx}",
                                    "type": "run",
                                    "text": run.text,
                                    "style": {
                                        "font": font.name,
                                        "size": font_size,
                                        "bold": font.bold,
                                        "italic": font.italic,
                                        "underline": font.underline,
                                        "strike": font.strike,
                                        "highlight": str(font.highlight_color) if font.highlight_color else None,
                                        "color": color_hex,
                                    },
                                })
                            alignment = None
                            if para.alignment is not None:
                                try:
                                    alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
                                except Exception:
                                    None
                            line_spacing = None
                            try:
                                line_spacing = para.paragraph_format.line_spacing
                            except Exception:
                                None
                            cell_paras.append({
                                "id": f"t_{uid}_cell_{r_idx}_{c_idx}_para_{p_idx}",
                                "type": "paragraph",
                                "role": "table_cell_paragraph",
                                "text": para.text.strip(),
                                "style": {
                                    "alignment": alignment,
                                    "line_spacing": line_spacing,
                                    "page_break_before": para.paragraph_format.page_break_before if hasattr(para.paragraph_format, "page_break_before") else None,
                                    "space_before_pt": para.paragraph_format.space_before.pt if getattr(para.paragraph_format, "space_before", None) is not None else None,
                                    "space_after_pt": para.paragraph_format.space_after.pt if getattr(para.paragraph_format, "space_after", None) is not None else None,
                                },
                                "runs": runs,
                            })
                        cell_dict = {
                            "id": f"t_{uid}_cell_{r_idx}_{c_idx}",
                            "type": "cell",
                            "row": r_idx,
                            "column": c_idx,
                            "children": cell_paras,
                        }
                        
                        try:
                            tcPr = cell._tc.get_or_add_tcPr()
                            shd = tcPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd")
                            if shd is not None:
                                fill = shd.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill")
                                if fill and fill != "auto":
                                    cell_dict["bg_color"] = fill
                        except Exception:
                            None
                        
                        try:
                            if cell.vertical_alignment is not None:
                                cell_dict["vertical_alignment"] = cell.vertical_alignment.name
                        except Exception:
                            None
                        
                        cells.append(cell_dict)
                    rows.append({
                        "id": f"t_{uid}_row_{r_idx}",
                        "type": "row",
                        "row": r_idx,
                        "cells": cells,
                    })

                table_style_name = table.style.name if table.style else "Normal Table"
                node = {
                    "id": f"t_{uid}",
                    "body_index": body_index,
                    "type": "table",
                    "role": "table",
                    "style": {"style_name": table_style_name},
                    "row_count": len(table.rows),
                    "col_count": len(table.columns) if table.rows else 0,
                    "rows": rows,
                }
                
                try:
                    tblBorders = table._tbl.tblPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tblBorders")
                    if tblBorders is not None:
                        node["style"]["has_custom_borders"] = True
                        # check colors of borders
                        border_colors = set()
                        for b_type in ["top", "left", "bottom", "right", "insideH", "insideV"]:
                            b_el = tblBorders.find(f"{{http://schemas.openxmlformats.org/wordprocessingml/2006/main}}{b_type}")
                            if b_el is not None:
                                c = b_el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color")
                                if c and c != "auto":
                                    border_colors.add(c)
                        if border_colors:
                            node["style"]["border_colors"] = list(border_colors)
                except Exception:
                    None

                children.append(node)

            # Increment body_index for every body child regardless of type
            body_index += 1

        return {
            "document_type": "docx",
            "dom": {
                "id": "document_root",
                "type": "document",
                "children": children,
            },
        }

    def _apply_docx_edits(self, source: Path, target: Path, edits: list[dict]) -> None:
        doc = Document(source)
        edit_map = {e["element_id"]: e for e in edits}
        
        uid_attr = f'{{{CUSTOM_NS}}}uid'
        
        for para in doc.paragraphs:
            uid = para._p.get(uid_attr)
            if not uid:
                continue
            eid = f"p_{uid}"
            if eid in edit_map:
                edit = edit_map[eid]
                self._apply_run_aware_replacement(para, edit["new_text"])
                
        for table in doc.tables:
            uid = table._tbl.get(uid_attr)
            if not uid:
                continue
            t_eid = f"t_{uid}"
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for p_idx, para in enumerate(cell.paragraphs):
                        eid = f"{t_eid}_cell_{r_idx}_{c_idx}_para_{p_idx}"
                        if eid in edit_map:
                            edit = edit_map[eid]
                            self._apply_run_aware_replacement(para, edit["new_text"])
                            
        doc.save(target)

    # ------------------------------------------------------------------
    # Core: run-aware replacement
    # ------------------------------------------------------------------

    def _apply_run_aware_replacement(self, paragraph, new_text: str, format_params: dict = None) -> None:
        """Update paragraph text while preserving per-run formatting.

        Only the runs whose characters overlap the changed region are
        modified; all other runs are left completely untouched.
        If format_params is provided, it is applied exclusively to the newly inserted text.
        """
        runs = paragraph.runs
        if not runs:
            run = paragraph.add_run()
            run.text = new_text
            if format_params:
                self._apply_run_format(run, format_params)
            return

        old_text = "".join(r.text for r in runs)

        if old_text == new_text and not format_params:
            log.debug("  No change (text identical), skipping.")
            return

        # Log run structure before modification (debug only).
        self._log_runs("  BEFORE", runs)

        prefix_len, suffix_len, new_middle = _changed_region(old_text, new_text)
        old_changed_start = prefix_len
        old_changed_end   = len(old_text) - suffix_len

        log.debug(
            "  Changed region: chars [%d, %d) %r → %r",
            old_changed_start, old_changed_end,
            old_text[old_changed_start:old_changed_end],
            new_middle,
        )

        # Walk runs and update only those overlapping [old_changed_start, old_changed_end).
        run_pos = 0
        new_middle_placed = False

        for run in runs:
            r_len   = len(run.text)
            r_start = run_pos
            r_end   = run_pos + r_len
            run_pos  = r_end

            if old_changed_start == old_changed_end:
                # Pure insertion
                if r_start <= old_changed_start < r_end:
                    overlaps = True
                elif old_changed_start == len(old_text) and r_end == len(old_text):
                    overlaps = True
                else:
                    overlaps = False
            else:
                # Replacement
                if r_end <= old_changed_start or r_start >= old_changed_end:
                    overlaps = False
                else:
                    overlaps = True

            if not overlaps:
                # Run is entirely outside the changed region — leave it alone.
                continue

            # Run overlaps with the changed region.
            # Compute the portion of this run that falls BEFORE the change.
            prefix_in_run = max(0, old_changed_start - r_start)
            prefix_text   = run.text[:prefix_in_run]

            # Compute the portion of this run that falls AFTER the change.
            suffix_start_in_run = max(0, old_changed_end - r_start)
            suffix_text         = run.text[suffix_start_in_run:]

            if not new_middle_placed:
                # First overlapping run: inject the replacement text here.
                # Instead of modifying run.text in place which corrupts prefix/suffix,
                # we split the run if format_params is provided or if we want exact precision.
                import copy
                
                run.text = prefix_text
                
                # Insert new middle run
                new_mid = paragraph.add_run(new_middle)
                rPr = run._r.find('.//w:rPr', namespaces=run._r.nsmap)
                if rPr is not None:
                    new_mid._r.insert(0, copy.deepcopy(rPr))
                if format_params:
                    self._apply_run_format(new_mid, format_params)
                run._r.addnext(new_mid._r)
                
                # Insert suffix run
                if suffix_text:
                    new_suf = paragraph.add_run(suffix_text)
                    if rPr is not None:
                        new_suf._r.insert(0, copy.deepcopy(rPr))
                    new_mid._r.addnext(new_suf._r)
                    
                new_middle_placed = True
            else:
                # Subsequent overlapping runs: their "changed" content has
                # already been absorbed by the first run; keep only the
                # trailing unchanged portion.
                run.text = suffix_text
                log.debug("  run cleared to suffix: %r", run.text)

        # Log run structure after modification (debug only).
        self._log_runs("  AFTER ", runs)

    def _apply_run_format(self, run, params: dict) -> None:
        """Apply formatting directly to a docx run."""
        from docx.shared import Pt
        from docx.shared import RGBColor as DRGBColor

        if params.get("bold") is not None:
            run.font.bold = params["bold"]
        if params.get("italic") is not None:
            run.font.italic = params["italic"]
        if params.get("underline") is not None:
            run.font.underline = params["underline"]
        if params.get("font_family"):
            run.font.name = params["font_family"]
        if params.get("font_size_pt") is not None:
            run.font.size = Pt(params["font_size_pt"])
            
        color_hex = str(params.get("color_hex", "")).strip().lstrip("#").upper()
        if color_hex and len(color_hex) == 6:
            try:
                run.font.color.rgb = DRGBColor(
                    int(color_hex[0:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16),
                )
            except Exception:
                None

        highlight_hex = params.get("highlight_hex")
        if highlight_hex:
            from docx.enum.text import WD_COLOR_INDEX
            hl_color = str(highlight_hex).strip().lstrip("#").upper()
            if hl_color in ["FFFF00", "YELLOW"]:
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW
            elif hl_color in ["00FF00", "GREEN"]:
                run.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
            elif hl_color in ["00FFFF", "CYAN"]:
                run.font.highlight_color = WD_COLOR_INDEX.TURQUOISE
            elif hl_color in ["FF00FF", "MAGENTA"]:
                run.font.highlight_color = WD_COLOR_INDEX.PINK
            elif hl_color in ["0000FF", "BLUE"]:
                run.font.highlight_color = WD_COLOR_INDEX.BLUE
            elif hl_color in ["FF0000", "RED"]:
                run.font.highlight_color = WD_COLOR_INDEX.RED
            elif hl_color in ["FFFFFF", "WHITE"]:
                run.font.highlight_color = WD_COLOR_INDEX.WHITE
            elif hl_color in ["000000", "BLACK"]:
                run.font.highlight_color = WD_COLOR_INDEX.BLACK
            else:
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW

    def _apply_run_aware_format(self, paragraph, match_text: str, format_params: dict) -> bool:
        """Apply formatting to a specific substring within a paragraph by splitting runs."""
        runs = paragraph.runs
        old_text = "".join(r.text for r in runs)
        start_idx = old_text.find(match_text)
        if start_idx == -1:
            return False

        old_changed_start = start_idx
        old_changed_end = start_idx + len(match_text)
        new_middle = match_text

        run_pos = 0
        new_middle_placed = False

        for run in runs:
            r_len   = len(run.text)
            r_start = run_pos
            r_end   = run_pos + r_len
            run_pos  = r_end

            if r_end <= old_changed_start or r_start >= old_changed_end:
                continue

            prefix_in_run = max(0, old_changed_start - r_start)
            prefix_text   = run.text[:prefix_in_run]

            suffix_start_in_run = max(0, old_changed_end - r_start)
            suffix_text         = run.text[suffix_start_in_run:]

            if not new_middle_placed:
                import copy
                run.text = prefix_text
                new_mid = paragraph.add_run(new_middle)
                rPr = run._r.find('.//w:rPr', namespaces=run._r.nsmap)
                if rPr is not None:
                    new_mid._r.insert(0, copy.deepcopy(rPr))
                self._apply_run_format(new_mid, format_params)
                run._r.addnext(new_mid._r)
                
                if suffix_text:
                    new_suf = paragraph.add_run(suffix_text)
                    if rPr is not None:
                        new_suf._r.insert(0, copy.deepcopy(rPr))
                    new_mid._r.addnext(new_suf._r)
                    
                new_middle_placed = True
            else:
                run.text = suffix_text
        return True

    # ------------------------------------------------------------------
    # Debug helper
    # ------------------------------------------------------------------

    @staticmethod
    def _log_runs(label: str, runs) -> None:
        if not log.isEnabledFor(logging.DEBUG):
            return
        for i, run in enumerate(runs):
            try:
                color = run.font.color
                rgb = color.rgb if color.type is not None else "inherit"
            except Exception:
                rgb = "?"
            log.debug(
                "%s run[%d] text=%r bold=%s italic=%s color=%s",
                label, i, run.text, run.font.bold, run.font.italic, rgb,
            )


# ---------------------------------------------------------------------------
# Normalise helper (used elsewhere)
# ---------------------------------------------------------------------------

def normalize_text(value: str) -> str:
    return sub(r"\s+", " ", value).strip()

# ===== END services/document_processor.py =====

# ===== BEGIN services/storage.py =====
import shutil
from pathlib import Path



class StorageService:
    def workspace_dir(self, workspace_id: str) -> Path:
        path = settings.storage_root / workspace_id
        path.mkdir(parents=True, exist_ok=True)
        return path

    def version_document_path(self, workspace_id: str, version: int, document_type: str) -> Path:
        return self.workspace_dir(workspace_id) / f"v{version}.{document_type}"

    def version_pdf_path(self, workspace_id: str, version: int) -> Path:
        return self.workspace_dir(workspace_id) / f"v{version}.pdf"

    def copy_version(self, source: str, workspace_id: str, version: int, document_type: str) -> Path:
        target = self.version_document_path(workspace_id, version, document_type)
        shutil.copyfile(source, target)
        return target

# ===== END services/storage.py =====

# ===== BEGIN services/preview.py =====
import asyncio
import logging
import uuid
from pathlib import Path

import httpx
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


logger = logging.getLogger(__name__)


class PreviewService:
    async def convert_to_pdf(self, document_path: Path, pdf_path: Path) -> Path:
        """Convert a .docx or .pptx to PDF using OnlyOffice Document Server."""
        if settings.converter_url:
            try:
                return await self._convert_with_onlyoffice(document_path, pdf_path)
            except Exception as exc:
                logger.warning("OnlyOffice conversion failed (%s), falling back to placeholder.", exc)

        # Run blocking reportlab code in a thread so the event loop stays free
        await asyncio.to_thread(self._placeholder, document_path, pdf_path)
        return pdf_path

    async def _convert_with_onlyoffice(self, document_path: Path, pdf_path: Path) -> Path:
        """
        Use the OnlyOffice Document Server Conversion API (async).

        Flow:
          1. Build a URL that OnlyOffice (running in Docker) can use to fetch
             the source file from this backend's /api/source-files/ endpoint.
          2. POST JSON to /ConvertService.ashx — OnlyOffice pulls the file,
             converts it, and returns a JSON payload with a fileUrl.
          3. Download the PDF from that fileUrl and write it to pdf_path.

        IMPORTANT: Must be async — if blocking httpx is used here, the event loop
        freezes and OnlyOffice's file-fetch request to this same server can never
        be handled, causing a deadlock and timeout.
        """
        rel_path = document_path.resolve().relative_to(settings.storage_root.resolve())
        source_url = settings.backend_url.rstrip("/") + "/api/source-files/" + rel_path.as_posix()

        filetype = document_path.suffix.lstrip(".").lower()  # "docx" or "pptx"
        key = uuid.uuid4().hex

        endpoint = settings.converter_url.rstrip("/") + "/ConvertService.ashx"

        payload = {
            "async": False,
            "filetype": filetype,
            "key": key,
            "outputtype": "pdf",
            "title": pdf_path.name,
            "url": source_url,
        }

        logger.debug("OnlyOffice request → %s  payload=%s", endpoint, payload)

        # Step 2 — POST to OnlyOffice; it downloads the source file and converts
        async with httpx.AsyncClient(timeout=120) as client:
            response = await client.post(
                endpoint,
                json=payload,
                headers={"Accept": "application/json"},
            )

        logger.debug("OnlyOffice HTTP %s: %s", response.status_code, response.text)
        response.raise_for_status()

        data = response.json()
        error_code = data.get("error")
        if error_code:
            raise RuntimeError(f"OnlyOffice conversion error code: {error_code}")

        file_url: str | None = data.get("fileUrl")
        if not file_url:
            raise RuntimeError(f"OnlyOffice did not return a fileUrl: {data}")

        # Step 3 — download the converted PDF (fileUrl points to OnlyOffice's own cache)
        # OnlyOffice returns URLs referencing itself, so no hostname rewriting needed
        async with httpx.AsyncClient(timeout=60) as client:
            pdf_response = await client.get(file_url)
        pdf_response.raise_for_status()

        pdf_path.write_bytes(pdf_response.content)
        logger.info("OnlyOffice converted %s → %s", document_path.name, pdf_path.name)
        return pdf_path

    def _placeholder(self, document_path: Path, pdf_path: Path) -> None:
        pdf_path.parent.mkdir(parents=True, exist_ok=True)
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter
        c.setFont("Helvetica-Bold", 18)
        c.drawString(72, height - 96, "Preview placeholder")
        c.setFont("Helvetica", 11)
        c.drawString(72, height - 125, "Start the OnlyOffice Document Server to render live DOCX/PPTX previews.")
        c.drawString(72, height - 145, f"Document: {document_path.name}")
        c.showPage()
        c.save()

# ===== END services/preview.py =====

# ===== BEGIN services/retrieval.py =====
import re
import uuid
import logging
import hashlib


log = logging.getLogger(__name__)

DOCUMENT_RETRIEVAL_STOPWORDS = {
    "a", "all", "and", "are", "for", "find", "in",
    "make", "of", "on", "remove", "rewrite", "shorter",
    "the", "them", "to",
}

DOCUMENT_BLOCKS_COLLECTION_NAME = "document_blocks"


class RetrievalService:
    """Retrieves relevant document blocks for a given query.

    Uses Qdrant vector search when configured, otherwise falls back to
    local lexical scoring. Supports both OpenAI embeddings and Gemini embeddings
    via the GeminiEmbeddingClient.
    """

    def __init__(self, embed_client=None) -> None:
        self._embed_client = embed_client

    def _get_embedding_client(self):
        if self._embed_client:
            return self._embed_client
        
        if settings.llm_provider == "gemini":
            return GeminiEmbeddingClient()
        elif settings.llm_provider == "openrouter":
            return OpenRouterEmbeddingClient()
        else:
            # Fallback to OpenAI embedding method
            return OpenAIEmbeddingClient()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def retrieve(self, query: str, structure: dict, limit: int = 4, workspace_id: str | None = None) -> list[dict]:
        if settings.qdrant_url and (settings.gemini_api_key or settings.openai_api_key):
            try:
                return self._retrieve_semantic(query, structure, limit, workspace_id)
            except Exception as exc:
                log.warning("Semantic retrieval failed: %s. Falling back to local.", exc)
                None
        return self._retrieve_local(query, structure, limit, workspace_id)

    def index_workspace(self, workspace_id: str, structure: dict) -> None:
        """Upsert all document blocks for a workspace into Qdrant."""
        if not (settings.qdrant_url and (settings.gemini_api_key or settings.openai_api_key)):
            return
        try:
            self._upsert_blocks(workspace_id, structure)
        except Exception as exc:
            log.warning("Workspace indexing failed: %s", exc)
            None

    def delete_workspace(self, workspace_id: str) -> None:
        """Delete all document blocks for a workspace from Qdrant."""
        if not (settings.qdrant_url and (settings.gemini_api_key or settings.openai_api_key)):
            return
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue
            qdrant = self._get_qdrant_client()
            qdrant.delete(
                collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                points_selector=Filter(
                    must=[
                        FieldCondition(
                            key="workspace_id",
                            match=MatchValue(value=workspace_id),
                        ),
                    ],
                ),
            )
        except Exception as exc:
            log.warning("Workspace deletion failed: %s", exc)
            None

    # ------------------------------------------------------------------
    # Qdrant path
    # ------------------------------------------------------------------

    def _get_qdrant_client(self):
        from qdrant_client import QdrantClient

        return QdrantClient(
            url=settings.qdrant_url,
            api_key=settings.qdrant_api_key or None,
            timeout=60.0,
        )

    def _ensure_collection(self, qdrant, vector_size: int) -> None:
        from qdrant_client.models import Distance, VectorParams, PayloadSchemaType

        existing = {col.name for col in qdrant.get_collections().collections}
        if DOCUMENT_BLOCKS_COLLECTION_NAME not in existing:
            qdrant.create_collection(
                collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
            )
            qdrant.create_payload_index(
                collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                field_name="workspace_id",
                field_schema=PayloadSchemaType.KEYWORD,
            )

    def _upsert_blocks(self, workspace_id: str, structure: dict) -> None:
        from qdrant_client.models import PointStruct

        blocks = structure.get("blocks", [])
        if not blocks:
            return

        texts = [block["text"] for block in blocks]
        
        embedder = self._get_embedding_client()
        # Use appropriate method name depending on type
        if hasattr(embedder, "embed_documents"):
            embeddings = embedder.embed_documents(texts)
        else:
            embeddings = embedder.embed(texts)

        qdrant = self._get_qdrant_client()
        self._ensure_collection(qdrant, len(embeddings[0]))

        points = [
            PointStruct(
                id=str(uuid.uuid4()),
                vector=embedding,
                payload={
                    "workspace_id": workspace_id,
                    "element_id": block["element_id"],
                    "text": block["text"],
                    "text_hash": hashlib.md5(block["text"].encode("utf-8")).hexdigest(),
                    "metadata": block.get("metadata", {}),
                    "type": block.get("type", "text"),
                },
            )
            for block, embedding in zip(blocks, embeddings)
        ]
        qdrant.upsert(collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME, points=points)

    def sync_workspace(self, workspace_id: str, structure: dict) -> None:
        """Delta-sync the workspace blocks to Qdrant without unnecessary re-embeddings."""
        if not (settings.qdrant_url and (settings.gemini_api_key or settings.openai_api_key)):
            return
            
        blocks = structure.get("blocks", [])
        if not blocks:
            # If document is empty, just clear the workspace from Qdrant
            self.delete_workspace(workspace_id)
            return
            
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue, PointStruct, PointIdsList
            qdrant = self._get_qdrant_client()
            
            # Step 1: Scroll existing metadata with pagination
            old_points = []
            next_page_offset = None
            scroll_filter = Filter(must=[FieldCondition(key="workspace_id", match=MatchValue(value=workspace_id))])
            
            while True:
                pts, next_page_offset = qdrant.scroll(
                    collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                    scroll_filter=scroll_filter,
                    with_payload=True,
                    with_vectors=False,
                    limit=1000,
                    offset=next_page_offset
                )
                old_points.extend(pts)
                if next_page_offset is None:
                    break
            
            from collections import defaultdict
            hash_to_available_qids = defaultdict(list)
            for pt in old_points:
                thash = pt.payload.get("text_hash")
                if thash:
                    hash_to_available_qids[thash].append((pt.id, pt.payload.get("element_id")))
                    
            # Step 2: Categorize live blocks
            points_to_upsert = []
            texts_to_embed = []
            blocks_waiting_for_embed = []
            
            payload_updates = [] # list of (qid, new_payload)
            live_qids = set()
            
            for block in blocks:
                text = block.get("text", "")
                thash = hashlib.md5(text.encode("utf-8")).hexdigest()
                current_element_id = block["element_id"]
                
                if hash_to_available_qids[thash]:
                    # Prefer popping a QID that perfectly matches the element_id to avoid unnecessary updates
                    match_idx = -1
                    for idx, (qid, old_element_id) in enumerate(hash_to_available_qids[thash]):
                        if old_element_id == current_element_id:
                            match_idx = idx
                            break
                            
                    if match_idx != -1:
                        qid, _ = hash_to_available_qids[thash].pop(match_idx)
                        live_qids.add(qid)
                        continue # Unchanged
                    else:
                        qid, _ = hash_to_available_qids[thash].pop(0)
                        live_qids.add(qid)
                        # Shifted
                        payload_updates.append((qid, {
                            "workspace_id": workspace_id,
                            "element_id": current_element_id,
                            "text": text,
                            "text_hash": thash,
                            "metadata": block.get("metadata", {}),
                            "type": block.get("type", "text")
                        }))
                        continue
                        
                # New content
                qid = str(uuid.uuid4())
                live_qids.add(qid)
                texts_to_embed.append(text)
                blocks_waiting_for_embed.append((qid, block, thash))
                        
            # Step 3: Update payloads for Shifted Vectors
            for qid, payload in payload_updates:
                qdrant.set_payload(
                    collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                    payload=payload,
                    points=[qid]
                )
                        
            # Step 4: Embed New Content
            if texts_to_embed:
                embedder = self._get_embedding_client()
                
                # Deduplicate embeddings based on text_hash
                unique_texts = {}
                for text in texts_to_embed:
                    h = hashlib.md5(text.encode("utf-8")).hexdigest()
                    if h not in unique_texts:
                        unique_texts[h] = text
                        
                unique_hashes = list(unique_texts.keys())
                unique_strings = list(unique_texts.values())
                
                if hasattr(embedder, "embed_documents"):
                    unique_embeddings = embedder.embed_documents(unique_strings)
                else:
                    unique_embeddings = embedder.embed(unique_strings)
                    
                hash_to_embedding = dict(zip(unique_hashes, unique_embeddings))
                
                for qid, block, thash in blocks_waiting_for_embed:
                    if thash in hash_to_embedding:
                        points_to_upsert.append(PointStruct(
                            id=qid,
                            vector=hash_to_embedding[thash],
                            payload={
                                "workspace_id": workspace_id,
                                "element_id": block["element_id"],
                                "text": block["text"],
                                "text_hash": thash,
                                "metadata": block.get("metadata", {}),
                                "type": block.get("type", "text")
                            }
                        ))

            # Step 5: Upsert First (for New Content)
            if points_to_upsert:
                self._ensure_collection(qdrant, len(points_to_upsert[0].vector))
                qdrant.upsert(collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME, points=points_to_upsert)
                
            # Step 6: Delete Stale Points
            stale_qids = [pt.id for pt in old_points if pt.id not in live_qids]
            if stale_qids:
                qdrant.delete(
                    collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
                    points_selector=PointIdsList(points=stale_qids)
                )
                
        except Exception as exc:
            log.warning("Workspace sync failed: %s", exc)

    def _retrieve_semantic(self, query: str, structure: dict, limit: int, workspace_id: str | None = None) -> list[dict]:
        from qdrant_client.models import Filter, FieldCondition, MatchValue

        embedder = self._get_embedding_client()
        if hasattr(embedder, "embed_query"):
            query_vector = embedder.embed_query(query)
        else:
            query_vector = embedder.embed([query])[0]

        qdrant = self._get_qdrant_client()
        
        q_filter = None
        if workspace_id:
            q_filter = Filter(
                must=[
                    FieldCondition(
                        key="workspace_id",
                        match=MatchValue(value=workspace_id),
                    )
                ]
            )

        results = qdrant.query_points(
            collection_name=DOCUMENT_BLOCKS_COLLECTION_NAME,
            query=query_vector,
            limit=limit,
            score_threshold=0.45,
            query_filter=q_filter,
        ).points

        found: list[dict] = []
        for hit in results:
            if hit.payload:
                found.append({
                    "element_id": hit.payload.get("element_id"),
                    "text": hit.payload.get("text", ""),
                    "metadata": hit.payload.get("metadata", {}),
                    "type": hit.payload.get("type", "text")
                })
        return found

    # ------------------------------------------------------------------
    # Local lexical fallback
    # ------------------------------------------------------------------

    def _retrieve_local(self, query: str, structure: dict, limit: int, workspace_id: str | None = None) -> list[dict]:
        terms = {term for term in re.findall(r"[a-zA-Z0-9]+", query.lower()) if term not in DOCUMENT_RETRIEVAL_STOPWORDS}
        scored: list[tuple[int, dict]] = []
        for block in structure.get("blocks", []):
            haystack = block["text"].lower()
            score = sum(1 for term in terms if term in haystack)
            if score:
                scored.append((score, block))
        scored.sort(key=lambda item: item[0], reverse=True)
        return [block for _, block in scored[:limit]]

# ===== END services/retrieval.py =====

# ===== BEGIN services/kb_processor.py =====
"""Knowledge Base Document Processor.

Parses uploaded documents (PDF, DOCX, TXT, MD) into text chunks suitable
for embedding and semantic retrieval.

Chunking strategy:
  - Target ~500 tokens per chunk (~350 words) with 50-token overlap.
  - Preserve section headings in chunk metadata for citation context.
  - PDF: page-by-page extraction, then paragraph merging.
  - DOCX: heading-aware extraction preserving document hierarchy.
  - TXT/MD: Paragraph-based splitting on double newlines.
"""

import logging
import re
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)

# Rough heuristic: ~1.3 chars per token for English prose
CHARS_PER_TOKEN = 1.3
TARGET_CHUNK_TOKENS = 500
OVERLAP_TOKENS = 50
TARGET_CHUNK_CHARS = int(TARGET_CHUNK_TOKENS * CHARS_PER_TOKEN)
OVERLAP_CHARS = int(OVERLAP_TOKENS * CHARS_PER_TOKEN)


class KBProcessor:
    """Parses and chunks documents for the knowledge base."""

    SUPPORTED_TYPES = {"pdf", "docx", "txt", "md"}

    def process_document(self, file_path: Path, file_type: str) -> list[dict]:
        """Parse a document and return a list of text chunks with metadata.

        Each chunk dict has:
          - text: the chunk content
          - chunk_index: 0-based position
          - metadata: source info (page, section, filename, etc.)
        """
        file_type = file_type.lower().lstrip(".")
        if file_type not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}")

        filename = file_path.name

        if file_type == "pdf":
            sections = self._parse_pdf(file_path, filename)
        elif file_type == "docx":
            sections = self._parse_docx(file_path, filename)
        else:  # txt or md
            sections = self._parse_text(file_path, filename)

        chunks = self._chunk_sections(sections)
        return chunks

    # ------------------------------------------------------------------
    # Parsers
    # ------------------------------------------------------------------

    def _parse_pdf(self, path: Path, filename: str) -> list[dict]:
        """Extract text page-by-page from PDF."""
        try:
            import PyPDF2
        except ImportError:
            # Fallback: attempt raw text read
            log.warning("PyPDF2 not installed; trying plain text fallback for PDF.")
            return self._parse_text(path, filename)

        sections: list[dict] = []
        try:
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, start=1):
                    text = page.extract_text() or ""
                    text = text.strip()
                    if not text:
                        continue
                    # Split page text into paragraphs
                    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
                    for para in paragraphs:
                        sections.append({
                            "text": para,
                            "metadata": {
                                "source": filename,
                                "page": page_num,
                                "type": "pdf_paragraph",
                            },
                        })
        except Exception as exc:
            log.error("PDF parsing failed for %s: %s", path, exc)
            raise

        return sections

    def _parse_docx(self, path: Path, filename: str) -> list[dict]:
        """Extract paragraphs from DOCX, preserving heading hierarchy."""
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError("python-docx is required but not installed.")

        sections: list[dict] = []
        current_heading: str | None = None
        current_heading_level: int | None = None

        try:
            doc = Document(path)
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                style_name = para.style.name if para.style else "Normal"
                is_heading = style_name.startswith("Heading")
                heading_level = None

                if is_heading:
                    try:
                        heading_level = int(style_name.split()[-1])
                    except (ValueError, IndexError):
                        heading_level = 1
                    current_heading = text
                    current_heading_level = heading_level
                    # Add heading itself as a section marker
                    sections.append({
                        "text": text,
                        "metadata": {
                            "source": filename,
                            "section": text,
                            "heading_level": heading_level,
                            "type": "heading",
                        },
                    })
                else:
                    sections.append({
                        "text": text,
                        "metadata": {
                            "source": filename,
                            "section": current_heading,
                            "heading_level": current_heading_level,
                            "type": "paragraph",
                        },
                    })

            # Also extract table content
            for table_idx, table in enumerate(doc.tables, start=1):
                rows_text: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        rows_text.append(row_text)
                if rows_text:
                    table_text = "\n".join(rows_text)
                    sections.append({
                        "text": table_text,
                        "metadata": {
                            "source": filename,
                            "section": current_heading,
                            "type": "table",
                            "table_index": table_idx,
                        },
                    })
        except Exception as exc:
            log.error("DOCX KB parsing failed for %s: %s", path, exc)
            raise

        return sections

    def _parse_text(self, path: Path, filename: str) -> list[dict]:
        """Parse plain text or markdown files."""
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            log.error("Text parsing failed for %s: %s", path, exc)
            raise

        # Split on double newlines (markdown paragraphs) or single newlines for long lines
        paragraphs = re.split(r"\n{2,}", text)
        sections: list[dict] = []
        current_heading: str | None = None

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # Detect markdown headings
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", para, re.MULTILINE)
            if heading_match:
                current_heading = heading_match.group(2).strip()
                level = len(heading_match.group(1))
                sections.append({
                    "text": current_heading,
                    "metadata": {
                        "source": filename,
                        "section": current_heading,
                        "heading_level": level,
                        "type": "heading",
                    },
                })
            else:
                sections.append({
                    "text": para,
                    "metadata": {
                        "source": filename,
                        "section": current_heading,
                        "type": "paragraph",
                    },
                })

        return sections

    # ------------------------------------------------------------------
    # Chunking
    # ------------------------------------------------------------------

    def _chunk_sections(self, sections: list[dict]) -> list[dict]:
        """Combine small sections into chunks and split large sections.

        Attempts to keep section context intact. Each chunk includes
        its source metadata from the first constituent section.
        """
        chunks: list[dict] = []
        current_text = ""
        current_meta: dict[str, Any] = {}

        for section in sections:
            text = section["text"]
            meta = section["metadata"]

            # If current text is empty, start a new chunk
            if not current_text:
                current_text = text
                current_meta = meta
                continue

            # If adding this section would exceed the target, flush first
            if len(current_text) + len(text) + 1 > TARGET_CHUNK_CHARS:
                # Flush current
                if current_text.strip():
                    chunks.append({
                        "text": current_text.strip(),
                        "chunk_index": len(chunks),
                        "metadata": current_meta,
                    })
                    # Carry overlap from end of current chunk
                    overlap_start = max(0, len(current_text) - OVERLAP_CHARS)
                    current_text = current_text[overlap_start:].strip() + "\n" + text
                    current_meta = meta
                else:
                    current_text = text
                    current_meta = meta
            else:
                current_text += "\n" + text

        # Flush remaining
        if current_text.strip():
            chunks.append({
                "text": current_text.strip(),
                "chunk_index": len(chunks),
                "metadata": current_meta,
            })

        # Handle individual oversized sections by splitting them
        final_chunks: list[dict] = []
        for chunk in chunks:
            if len(chunk["text"]) > TARGET_CHUNK_CHARS * 2:
                sub = self._split_large_text(chunk["text"], chunk["metadata"])
                for i, s in enumerate(sub):
                    s["chunk_index"] = len(final_chunks)
                    final_chunks.append(s)
            else:
                chunk["chunk_index"] = len(final_chunks)
                final_chunks.append(chunk)

        return final_chunks

    def _split_large_text(self, text: str, meta: dict) -> list[dict]:
        """Split a large text block into overlapping chunks at sentence boundaries."""
        # Try to split at sentence boundaries
        sentences = re.split(r"(?<=[.!?])\s+", text)
        chunks: list[dict] = []
        current = ""

        for sentence in sentences:
            if len(current) + len(sentence) > TARGET_CHUNK_CHARS:
                if current.strip():
                    chunks.append({
                        "text": current.strip(),
                        "chunk_index": 0,
                        "metadata": meta,
                    })
                    # Start new chunk with overlap
                    overlap_start = max(0, len(current) - OVERLAP_CHARS)
                    current = current[overlap_start:].strip() + " " + sentence
                else:
                    current = sentence
            else:
                current += " " + sentence

        if current.strip():
            chunks.append({
                "text": current.strip(),
                "chunk_index": 0,
                "metadata": meta,
            })

        return chunks

# ===== END services/kb_processor.py =====

# ===== BEGIN services/kb_retrieval.py =====
"""Knowledge Base Retrieval Service.

Embeds KB chunks into a dedicated Qdrant collection ('knowledge_chunks')
and retrieves relevant passages for the document generation pipeline.

Uses a separate Qdrant collection from the existing 'document_blocks'
so document structure index and KB content index remain isolated.

Falls back to SQLite keyword search when Qdrant is not configured.
"""

import logging
import re
import uuid
from typing import Any


log = logging.getLogger(__name__)

KB_COLLECTION = "knowledge_chunks"
KB_RETRIEVAL_STOPWORDS = {
    "a", "all", "an", "and", "are", "as", "at", "be", "by", "for",
    "from", "has", "he", "in", "is", "it", "its", "of", "on", "or",
    "that", "the", "this", "to", "was", "were", "will", "with",
}


class KBRetrievalService:
    """Manages embedding and retrieval of knowledge base chunks."""

    def __init__(self, embed_client=None) -> None:
        self._embed_client = embed_client

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def index_document(
        self,
        workspace_id: str,
        doc_id: str,
        chunks: list[dict],
    ) -> None:
        """Embed and upsert chunks into Qdrant (or no-op if not configured)."""
        if not self._can_use_qdrant():
            log.info("Qdrant not configured — KB chunks stored in SQLite only.")
            return
        try:
            self._upsert_chunks(workspace_id, doc_id, chunks)
        except Exception as exc:
            log.warning("KB indexing failed (Qdrant): %s", exc)

    def retrieve(
        self,
        workspace_id: str,
        query: str,
        chunks_from_db: list[dict],
        limit: int = 15,
    ) -> list[dict]:
        """Retrieve the most relevant KB chunks for a generation query.

        Args:
            workspace_id: The workspace whose KB to search.
            query: User's generation request.
            chunks_from_db: All chunks for the workspace from SQLite
                            (used as fallback and for re-ranking).
            limit: Number of top chunks to return.
        """
        if self._can_use_qdrant():
            try:
                return self._retrieve_semantic(workspace_id, query, chunks_from_db, limit)
            except Exception as exc:
                log.warning("Qdrant KB retrieval failed: %s. Falling back.", exc)

        return self._retrieve_local(query, chunks_from_db, limit)

    def delete_document(self, workspace_id: str, doc_id: str) -> None:
        """Remove all Qdrant vectors for a specific KB document."""
        if not self._can_use_qdrant():
            return
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue
            qdrant = self._get_qdrant_client()
            qdrant.delete(
                collection_name=KB_COLLECTION,
                points_selector=Filter(
                    must=[
                        FieldCondition(key="workspace_id", match=MatchValue(value=workspace_id)),
                        FieldCondition(key="doc_id", match=MatchValue(value=doc_id)),
                    ]
                ),
            )
        except Exception as exc:
            if "Not found" not in str(exc) and "404" not in str(exc):
                log.warning("KB document delete from Qdrant failed: %s", exc)

    def delete_workspace_kb(self, workspace_id: str) -> None:
        """Remove all KB vectors for a workspace from Qdrant."""
        if not self._can_use_qdrant():
            return
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue
            qdrant = self._get_qdrant_client()
            qdrant.delete(
                collection_name=KB_COLLECTION,
                points_selector=Filter(
                    must=[
                        FieldCondition(key="workspace_id", match=MatchValue(value=workspace_id)),
                    ]
                ),
            )
        except Exception as exc:
            if "Not found" not in str(exc) and "404" not in str(exc):
                log.warning("KB workspace delete from Qdrant failed: %s", exc)

    def retrieve_for_section(
        self,
        workspace_id: str,
        section_query: str,
        fallback_chunks: list[dict],
        limit: int = 15,
    ) -> tuple[list[dict], bool]:
        """Retrieve the most relevant chunks for a specific section via Qdrant.
        
        Returns (chunks, used_semantic). used_semantic=False means keyword fallback was used.
        """
        if self._can_use_qdrant():
            try:
                results = self._retrieve_semantic(workspace_id, section_query, fallback_chunks, limit)
                if results:
                    scores = [r.get("score", 0) for r in results]
                    log.info(
                        "Section '%s': retrieved %d chunks, scores min=%.3f median=%.3f max=%.3f",
                        section_query[:60], len(results), min(scores), sorted(scores)[len(scores)//2], max(scores)
                    )
                    return results, True
            except Exception as exc:
                log.warning("Semantic section retrieval failed: %s. Falling back to keyword.", exc)
        
        # Keyword fallback — but flag it
        return self._retrieve_local(section_query, fallback_chunks, limit), False

    # ------------------------------------------------------------------
    # Qdrant internals
    # ------------------------------------------------------------------

    def _can_use_qdrant(self) -> bool:
        return bool(settings.qdrant_url and (settings.gemini_api_key or settings.openai_api_key))

    def _get_qdrant_client(self):
        from qdrant_client import QdrantClient
        return QdrantClient(url=settings.qdrant_url, api_key=settings.qdrant_api_key or None, timeout=60.0)

    def _get_embedding_client(self):
        if self._embed_client:
            return self._embed_client
        if settings.llm_provider == "gemini":
            return GeminiEmbeddingClient()
        if settings.llm_provider == "openrouter":
            return OpenRouterEmbeddingClient()
        return OpenAIEmbeddingClient()

    def _ensure_collection(self, qdrant, vector_size: int) -> None:
        from qdrant_client.models import Distance, VectorParams, PayloadSchemaType
        existing = {col.name for col in qdrant.get_collections().collections}
        if KB_COLLECTION not in existing:
            try:
                qdrant.create_collection(
                    collection_name=KB_COLLECTION,
                    vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
                )
                qdrant.create_payload_index(
                    collection_name=KB_COLLECTION,
                    field_name="workspace_id",
                    field_schema=PayloadSchemaType.KEYWORD,
                )
                qdrant.create_payload_index(
                    collection_name=KB_COLLECTION,
                    field_name="doc_id",
                    field_schema=PayloadSchemaType.KEYWORD,
                )
            except Exception as exc:
                if "already exists" not in str(exc) and "409" not in str(exc):
                    raise

    def _upsert_chunks(self, workspace_id: str, doc_id: str, chunks: list[dict]) -> None:
        from qdrant_client.models import PointStruct
        if not chunks:
            return

        texts = [c["text"] for c in chunks]
        embedder = self._get_embedding_client()
        embeddings = embedder.embed_documents(texts)

        qdrant = self._get_qdrant_client()
        self._ensure_collection(qdrant, len(embeddings[0]))

        points = [
            PointStruct(
                id=str(uuid.uuid5(uuid.NAMESPACE_URL, f"{workspace_id}/{doc_id}/{c['chunk_index']}")),
                vector=embedding,
                payload={
                    "workspace_id": workspace_id,
                    "doc_id": doc_id,
                    "chunk_index": c["chunk_index"],
                    "text": c["text"],
                    "metadata": c.get("metadata", {}),
                },
            )
            for c, embedding in zip(chunks, embeddings)
        ]
        qdrant.upsert(collection_name=KB_COLLECTION, points=points)

    def _retrieve_semantic(
        self,
        workspace_id: str,
        query: str,
        chunks_from_db: list[dict],
        limit: int,
    ) -> list[dict]:
        from qdrant_client.models import Filter, FieldCondition, MatchValue

        embedder = self._get_embedding_client()
        if hasattr(embedder, "embed_query"):
            query_vector = embedder.embed_query(query)
        else:
            query_vector = embedder.embed([query])[0]

        qdrant = self._get_qdrant_client()

        # Check collection exists
        existing = {col.name for col in qdrant.get_collections().collections}
        if KB_COLLECTION not in existing:
            return self._retrieve_local(query, chunks_from_db, limit)

        results = qdrant.query_points(
            collection_name=KB_COLLECTION,
            query=query_vector,
            query_filter=Filter(
                must=[FieldCondition(key="workspace_id", match=MatchValue(value=workspace_id))]
            ),
            limit=limit,
            score_threshold=0.25,
        ).points

        # Map qdrant results back to full chunk records
        idx_map = {(c.get("metadata", {}).get("doc_id"), c.get("chunk_index")): c for c in chunks_from_db}

        found: list[dict] = []
        for hit in results:
            payload = hit.payload or {}
            chunk = {
                "text": payload.get("text", ""),
                "chunk_index": payload.get("chunk_index"),
                "metadata": payload.get("metadata", {}),
                "score": hit.score,
            }
            found.append(chunk)

        return found

    # ------------------------------------------------------------------
    # Local keyword fallback
    # ------------------------------------------------------------------

    def _retrieve_local(self, query: str, chunks: list[dict], limit: int) -> list[dict]:
        terms = {t for t in re.findall(r"[a-zA-Z0-9]+", query.lower()) if t not in KB_RETRIEVAL_STOPWORDS}
        if not terms:
            return chunks[:limit]

        scored: list[tuple[int, dict]] = []
        for chunk in chunks:
            haystack = chunk.get("text", "").lower()
            score = sum(1 for term in terms if term in haystack)
            if score:
                scored.append((score, chunk))

        scored.sort(key=lambda x: x[0], reverse=True)
        return [c for _, c in scored[:limit]]

# ===== END services/kb_retrieval.py =====

# ===== BEGIN services/analyzer.py =====
import logging
import json

log = logging.getLogger(__name__)

class DocumentAnalyzer:
    """Analyzes document structure to extract high-level context (theme, purpose)."""

    def __init__(self, llm: LLMClient | None = None):
        self._llm = llm or LLMClient()

    def analyze(self, structure: dict) -> dict:
        """Analyze the document and return a summary of its style and purpose."""
        if not structure:
            return {}

        # Just take a highly simplified version of the structure to save tokens
        simplified = []
        for s in structure.get("sections", []):
            simplified.append(f"Section: {s.get('heading', 'Untitled')}")
            for el in s.get("elements", [])[:5]:  # Just first few elements
                simplified.append(f"- {el.get('type')}: {el.get('text_preview', '')}")

        doc_preview = "\n".join(simplified)
        
        system_prompt = (
            "You are an expert document analyst. "
            "Given a structural preview of a document, provide a high-level analysis.\n"
            "Identify the likely purpose of the document (e.g., Marketing Report, Business Proposal, Recipe).\n"
            "Identify the current stylistic theme if any (e.g., Formal, Modern, Minimalist).\n"
            "Provide a short 2-3 sentence summary of the document's overall context.\n"
            "Return JSON matching this schema:\n"
            "{\n"
            '  "purpose": "...",\n'
            '  "theme": "...",\n'
            '  "summary": "..."\n'
            "}"
        )

        try:
            response = self._llm.complete(LLMRequest(
                system_prompt=system_prompt,
                user_prompt=f"Document Preview:\n{doc_preview}",
                temperature=0,
                max_tokens=256,
                json_mode=True,
            ))
            return response.json or {}
        except Exception as e:
            log.error(f"Failed to analyze document: {e}")
            return {}

# ===== END services/analyzer.py =====

# ===== BEGIN services/task_planner.py =====
"""Task Planner — decomposes user requests into ordered lists of atomic tasks.

Replaces the regex-based intent classification and routes any request
to a sequence of structured editing/layout/formatting steps.
"""

import logging
from typing import Any


log = logging.getLogger(__name__)


PLANNER_SYSTEM_PROMPT = """You are a document editing task planner.
Given a user request, chat history, and a document outline, decompose the request into
an ordered list of atomic tasks. Each task represents one discrete change.

Available task types:
- text_edit: Rewrite text content of a specific element (editing text, sentences, paragraphs). NEVER use this for formatting (fonts, colors, alignment, spacing) - use text_format instead!
- text_format: Change formatting (font family/name, font style, bold, color, size, alignment, margins, spacing, bullets, etc.)
- table_op: Create, modify, or delete tables, columns, rows, or cell contents/styling
- image_op: Insert, replace, resize, style, reposition, or remove images
- layout_op: Move sections, swap sections, insert page breaks, remove sections, add/remove Table of Contents (TOC). Do NOT use for adding new content sections/paragraphs — use content_generation instead!
- list_op: Convert list formats (bullets, numbered, checklist), add list items, sort lists
- find_replace: Global text find and replace across the document
- theme_op: Slide/page background, margin settings, corporate themes, color palettes
- meta_op: Modify document metadata (Title, Author, Subject, Keywords)
- section_op: Modify document section properties (Margins, orientation, page size)
- style_op: Modify built-in global DOCX styles (Heading 1, Normal, etc.)
- slide_op: Add, delete, duplicate, hide, or reorder slides (for presentations only)
- generate: Create full slide presentation content from scratch (for presentations only)
- docx_generate: Generate a comprehensive, full document from scratch using the DOCX template structure and the workspace knowledge base. Use this ONLY for DOCX documents when the user asks to "create", "generate", "write", "draft", "produce", or "build" a full document or major new sections (e.g., "create a financial report", "generate meeting minutes", "write a Q3 summary report"). Do NOT use for targeted edits to existing content.
- content_generation: Generate new substantive section/paragraph content grounded in workspace Knowledge Base evidence. Use this for requests that add new factual sections, topics, or multi-paragraph content (e.g., "add a sustainability section", "write an ESG summary", "add a section on Q2 financial metrics", "add an environmental impact section"). Do NOT use for simple text rewrites, expanding/shortening an existing section, or formatting changes.

For each task, provide:
- task_type: one of the types above
- description: what to do, in plain language (e.g., "Change heading font color to dark green")
- target_hint: which element(s) to target, using names, ordinals, or roles from the outline
  (e.g., "Table 3", "the Conclusion section", "all headings", "paragraph 5", "the bulleted list")
- dependencies: list of 0-based task indices this task depends on (usually empty, unless task B must run after task A, e.g. add section then add content to it)

CRITICAL RULES:
1. Decompose EVERY distinct action in the request. If the user says "change all headings to green, set margins to 1 inch, and move Table 1 to the end", you MUST output 3 separate tasks.
2. Group repetitive actions: If adding multiple items to the same list, or applying the same format to a group of elements, combine them into a SINGLE task (e.g., "Add 3 new bullet points to the Highlights list"). Do NOT split them into one task per item. Note: This grouping rule does NOT apply to content_generation tasks for multiple distinct sections.
3. Ordering matters: tasks must be ordered logically so they can be executed sequentially.
4. Be precise with target_hint so the resolver can map them accurately. Use ordinal indicators from the outline (like "Table 1", "Section 2") if present. If the request applies to all instances of a type (e.g., "all headings", "all tables", "all images"), use exactly that phrase for target_hint (e.g., "all headings"). Do NOT use "entire document" for these.
5. DO NOT create image_op tasks (like adding or replacing images/logos) unless the user EXPLICITLY asks you to add, replace, or modify an image. Do not invent image tasks to "improve" the document.
6. FONT & STYLING INTENT: If the user says "change font to X", "make text Arial", or "use Helvetica", this is ALWAYS a `text_format` task (changing the font family). Do NOT use `text_edit` for these requests!
6A. SELECTIVE IMPORTANCE HIGHLIGHTING: If the user asks to "highlight important content" or "highlight key/critical content", emit a `text_format` task with target_hint "important content". NEVER use target_hint "all" or "all paragraphs" unless the user explicitly says to highlight everything.
7. DOCX GENERATION INTENT: If the user asks to "create", "generate", "write", "draft", or "produce" a complete document, report, or major new content (e.g., "create a Q3 financial report", "generate minutes of meeting", "write a compliance report"), and the document type is DOCX, use a single `docx_generate` task. The target_hint should describe the type of document to generate.
8. KNOW YOUR LIMITATIONS: The document engine natively supports:
   - Text: bold, italic, underline, strikethrough, highlight_color, font name, font size (pt), and font color (RGB hex).
   - Paragraphs: left/center/right/justify align, space before/after (pt), line spacing (e.g., 1.5), page breaks, left/right/first-line indents (pt), and keep with next / keep together.
   - Tables: modify columns/rows, cell text formatting, cell backgrounds, cell vertical alignments, column widths, row alternate colors, header formatting, and borders.
   - Layout & Sections: page orientation (landscape/portrait), margins, and exact page dimensions.
   - Global Styles: modify built-in DOCX styles globally.
   - Metadata: modify document properties (Title, Author, Subject, Keywords).
   - Headers/Footers: edit contents within headers and footers just like normal body text.
9. DO NOT invent unsupported tasks (e.g., floating images, rounded corners, drop shadows). For aesthetic requests (e.g., "make it modern"), creatively combine the SUPPORTED properties (like changing heading fonts to sans-serif, using elegant dark gray colors, adjusting page layout, and adding paragraph spacing).
10. ADDING NEW CONTENT SECTIONS / PARAGRAPHS: When the user asks to add, insert, or write new content sections, paragraphs, or factual topics (e.g., "add an environmental impact section", "add a sustainability section", "insert content below Executive Summary"), you MUST emit task_type: content_generation. NEVER use layout_op for adding new factual text content.
10A. EXPANDING EXISTING CONTENT: When the user asks to "increase", "expand", "lengthen", "elaborate", "make longer", "shorten", "reduce", or set a word count for an EXISTING section or paragraph, emit task_type: text_edit. NEVER use content_generation or layout_op for this. Preserve the existing heading and rewrite the existing body text in place; do not add a duplicate heading.
11. MULTI-SECTION CONTENT GENERATION: When a request asks to add multiple distinct new factual sections or topics (e.g. "add sustainability and ESG sections"), you MUST emit separate content_generation tasks — one per distinct section/topic (e.g. Task 1: "Generate sustainability section", Task 2: "Generate ESG metrics section") — so each section can independently retrieve its own Knowledge Base evidence. Do NOT combine multiple distinct new sections into a single task.
11A. CONCLUSION WITH SUMMARY TABLE: If the user asks to add a conclusion/final summary section with content and a table of ESG/other metrics/findings, emit exactly TWO tasks in order: (1) `content_generation` to add the Conclusion section at the end with ending/summary description text before the table, target_hint "after end of document"; (2) `table_op` to create and populate a detailed findings/metrics table under the Conclusion section, target_hint "after Conclusion section", dependent on task 0. Do NOT create separate ESG Metrics or Other Metrics sections unless explicitly requested.
11. SWAPPING SECTIONS: When the user asks to "swap", "exchange", or "switch" two sections, emit ONE layout_op task. The description should be "Swap [Section A] and [Section B] sections". The target_hint should name both sections, e.g., "Action Items and Key Metrics sections".
12. SECTION INSERTION HINT FORMAT: For target_hint when inserting content after a section, always use the pattern "after [Section Name] section" (e.g., "after Executive Summary section"). This tells the resolver to use the end of the section — after all existing content — as the insertion anchor.
13. MOVING SECTIONS (DISTINCT FROM SWAPPING): When the user says "move [Section X] above [Section Y]", "move [Section X] before [Section Y]", "move [Section X] below [Section Y]", or "move [Section X] after [Section Y]", this is a MOVE, not a swap. Emit ONE layout_op task with description "Move [Section X] section above/below [Section Y] section". The target_hint MUST use the pattern "[Section X] above [Section Y]" or "[Section X] below [Section Y]" (e.g., "Action Items above Highlights", "Company Overview below Business Objectives"). NEVER use "swap" for move requests — swap exchanges both sections, move only relocates one.
14. PAGE BREAKS: When the user asks to "insert a page break before/after [Section X]", emit ONE layout_op task. Description: "Insert a page break before [Section X] section". Target hint: "before [Section X] section" (e.g., "before Action Items section"). The resolver will provide the heading ID of that section as the anchor.
15. TABLE OF CONTENTS (TOC): When the user asks to "add a table of contents", "insert a TOC", or "create a table of contents", emit ONE layout_op task with description "Add Table of Contents" and target_hint "before top of document". NEVER use text_edit, text_format, or docx_generate for TOC creation.
16. RENUMBERING VS DUPLICATION: If the user says "make heading X say 3", "renumber section Y to 3", "change the number to 3", or similar — this is a text edit to existing heading text, emit a text_edit task targeting that specific heading element. Do NOT emit a duplicate_block layout task unless the user explicitly asks to create additional copies of a section.
17. FIXING TOC PAGE NUMBERS: If the user asks to "fix TOC page numbers", "update page numbers in TOC", or complains that TOC page numbers are wrong/missing, DO NOT emit text_edit, layout_op, or duplicate_block tasks! Page numbers in a Table of Contents are calculated automatically by Word/rendering engine on open/print via native TOC fields.
18. TOC DOT LEADERS & FORMATTING: Requests to adjust TOC dot leaders, extend dotted lines, or format TOC entry alignment must NEVER emit text_format, set_alignment, or generic paragraph operations! Dot leaders and entry alignments are handled natively in Word via TOC tab stops (w:leader='dot' w:val='right').
19. IMAGE INSERTION ANCHORS: For image insertion near text, set target_hint to the anchor element, not to "end of document". Examples: "Insert the company logo below the title" -> task_type image_op, target_hint "the title"; "insert image below Executive Summary" -> target_hint "Executive Summary".
20. IMAGE IDENTITY: For existing image edits, preserve the user's semantic label in target_hint. Examples: "company logo", "placeholder image", "sales chart", "Image 1". Do NOT use "all images" unless the user explicitly says all/every image.
21. IMAGE CAPTIONS: "Add a caption below the image..." is image_op targeting the referenced existing image. It does NOT require adding a new image. "Center/left/right align just the caption" and "increase caption font size" are also image_op tasks targeting the image/caption pair, not text_edit.
22. IMAGE PAGE/TITLE PLACEMENT: "top right corner of the page" means move/insert the image in a new paragraph at the top of the document and right-align it. "right side/corner of the title" means place the image inline in the title paragraph on the right side.

FEW-SHOT CLASSIFICATION BOUNDARIES:
- "add an environmental impact section" -> task_type: "content_generation"
- "insert a sustainability section after Executive Summary" -> task_type: "content_generation"
- "add an ESG metrics summary" -> task_type: "content_generation"
- "add a page break before Action Items" -> task_type: "layout_op"
- "add a table of contents" -> task_type: "layout_op"
- "move Section 2 above Section 1" -> task_type: "layout_op"
- "swap Highlights and Financials" -> task_type: "layout_op"
- "change heading font color to dark green" -> task_type: "text_format"
- "highlight all the important content in the document" -> task_type: "text_format", target_hint: "important content"
- "rewrite conclusion to be shorter" -> task_type: "text_edit"
- "increase the Risk Assessment section to 100 words" -> task_type: "text_edit", target_hint: "Risk Assessment section"
- "expand the Executive Summary section" -> task_type: "text_edit", target_hint: "Executive Summary section"
- "add 3 new items to the Key Highlights list" -> task_type: "list_op"
- "add a row to Table 1" -> task_type: "table_op"
- "insert the company logo below the title" -> task_type: "image_op", target_hint: "the title"
- "replace the placeholder image with a sales chart" -> task_type: "image_op", target_hint: "placeholder image"
- "resize the image to 40% of the page width" -> task_type: "image_op", target_hint: "the image"
- "move the company logo to the right" -> task_type: "image_op", target_hint: "company logo"
- "add a caption below the image saying Figure 1: Quarterly Sales." -> task_type: "image_op", target_hint: "the image"
- "move just the caption to the center" -> task_type: "image_op", target_hint: "the image"
- "increase the caption font size to 12" -> task_type: "image_op", target_hint: "the image"
- "place the company logo on the top right corner of the page" -> task_type: "image_op", target_hint: "company logo"
- "place the company logo on the right side of the title" -> task_type: "image_op", target_hint: "company logo"
- "add a final summary table in the end with both ESG and other metrics under a conclusion section with content before the table" -> tasks: content_generation target_hint "after end of document", then table_op target_hint "after Conclusion section"

Return ONLY a JSON object:
{
  "tasks": [
    {
      "task_type": "...",
      "description": "...",
      "target_hint": "...",
      "dependencies": []
    }
  ]
}
"""


class TaskPlanner:
    """Decomposes complex requests into atomic task lists."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    def plan(
        self,
        request: str,
        outline: dict,
        chat_history: list[dict] = None,
        analysis: dict = None,
        relevant_blocks: dict = None,
    ) -> list[dict]:
        """Generate a task list from a user request and outline."""
        llm = self._llm or LLMClient()

        history_str = ""
        if chat_history:
            history_str = "Recent conversation history:\n"
            for msg in chat_history[-5:]:
                role = "User" if msg["role"] == "user" else "Agent"
                history_str += f"{role}: {msg['content']}\n"
            history_str += "\n"

        import json
        import re
        outline_summary = json.dumps({
            "document_type": outline.get("document_type"),
            "title": outline.get("title"),
            "element_count": outline.get("element_count"),
            "sections": [
                {
                    "heading": s.get("heading"),
                    "heading_id": s.get("heading_id"),
                    "semantic_type": s.get("semantic_type"),
                    "ordinal": s.get("ordinal"),
                    "elements": [
                        {
                            "type": el.get("type"),
                            "role": el.get("role"),
                            "id": el.get("id"),
                            "ordinal_label": el.get("ordinal_label"),
                            "text_preview": el.get("text_preview"),
                            "alt_text": el.get("alt_text"),
                            "name": el.get("name"),
                            "title": el.get("title"),
                            "size": el.get("size"),
                            "alignment": el.get("alignment"),
                        } for el in s.get("elements", [])
                    ]
                } for s in outline.get("sections", [])
            ]
        }, indent=2)

        context_str = ""
        if analysis:
            context_str += f"Global Document Analysis:\n{json.dumps(analysis, indent=2)}\n\n"
        if relevant_blocks:
            context_str += f"Relevant Full-Text Blocks (from Semantic Search):\n{json.dumps(relevant_blocks, indent=2)}\n\n"

        user_prompt = (
            f"{history_str}"
            f"Document Outline (Truncated Preview):\n{outline_summary}\n\n"
            f"{context_str}"
            f"User Request: {request}"
        )

        response = llm.complete(LLMRequest(
            system_prompt=PLANNER_SYSTEM_PROMPT,
            user_prompt=user_prompt,
            temperature=0.0,
            max_tokens=1024,
            json_mode=True,
        ))

        parsed = response.json or {}
        tasks = parsed.get("tasks", [])

        request_l = str(request or "").lower()
        if (
            "conclusion" in request_l
            and "table" in request_l
            and any(word in request_l for word in ("metric", "metrics", "esg", "finding", "findings", "summary"))
            and any(word in request_l for word in ("end", "final"))
        ):
            return [
                {
                    "task_type": "content_generation",
                    "description": "Add a Conclusion section at the end with concise ending summary content before the table.",
                    "target_hint": "after end of document",
                    "dependencies": [],
                },
                {
                    "task_type": "table_op",
                    "description": "Create a detailed final findings table under the Conclusion section covering ESG metrics and other relevant metrics.",
                    "target_hint": "after Conclusion section",
                    "dependencies": [0],
                },
            ]
        
        # Simple validation
        validated = []
        for task in tasks:
            if isinstance(task, dict) and task.get("task_type") and task.get("description"):
                task_type = task["task_type"]
                description_l = str(task.get("description", "")).lower()
                target_hint_l = str(task.get("target_hint", "")).lower()
                existing_rewrite_intent = (
                    any(word in f"{request_l} {description_l}" for word in (
                        "increase", "expand", "lengthen", "elaborate", "make longer",
                        "shorten", "reduce", "condense", "summarize", "summarise",
                    ))
                    or bool(re.search(r"\b\d+\s*(?:words?|word)\b", f"{request_l} {description_l}"))
                )
                section_or_paragraph_target = (
                    "section" in target_hint_l
                    or "paragraph" in target_hint_l
                    or any(
                        str(s.get("heading", "")).lower() in target_hint_l
                        for s in outline.get("sections", [])
                        if s.get("heading")
                    )
                )
                if task_type == "content_generation" and existing_rewrite_intent and section_or_paragraph_target:
                    task_type = "text_edit"
                if (
                    task_type == "text_format"
                    and "highlight" in f"{request_l} {description_l}"
                    and any(word in f"{request_l} {description_l}" for word in ("important", "key", "critical", "main"))
                    and target_hint_l in ("all", "all paragraphs", "entire document", "document", "all content")
                ):
                    target_hint_l = "important content"
                    task["target_hint"] = "important content"
                validated.append({
                    "task_type": task_type,
                    "description": task["description"],
                    "target_hint": task.get("target_hint") or "all",
                    "dependencies": [d for d in task.get("dependencies", []) if isinstance(d, int)],
                })
        return validated

# ===== END services/task_planner.py =====

# ===== BEGIN services/outline_builder.py =====
"""Document Outline Builder — builds a hierarchical semantic outline of a document.

Replaces the flat, lossy 80-character truncated structure summary.
Supports both DOCX and PPTX formats.
Provides pre-built lookup maps for deterministic ordinal/named/semantic resolution.
"""

import logging
import re
from typing import Any

log = logging.getLogger(__name__)


class OutlineBuilder:
    """Builds a hierarchical semantic outline of a document.

    Outputs a structured outline representation that can be fed to the
    Task Planner and Reference Resolver.
    """

    @classmethod
    def build(cls, structure: dict, document_type: str) -> dict:
        """Build the outline.

        Returns a dictionary containing:
        - document_type (str)
        - title (str)
        - element_count (int)
        - sections (list[dict])
        - indices (dict): lookup tables for structural/ordinal resolution
        """
        if document_type == "pptx":
            return cls._build_pptx(structure)
        return cls._build_docx(structure)

    @classmethod
    def _build_docx(cls, structure: dict) -> dict:
        dom_children = structure.get("dom", {}).get("children", [])
        blocks = structure.get("blocks", [])

        indices: dict[str, Any] = {
            "tables_by_ordinal": {},
            "images_by_ordinal": {},
            "headings_by_name": {},
            "last_element_id": None,
            "first_content_id": None,
            "all_elements": [],
        }

        # If dom is empty, construct outline from flat blocks list
        if not dom_children:
            return cls._build_docx_fallback(blocks)

        table_counter = 0
        image_counter = 0
        last_id = None
        first_id = None

        # Build flat list of all elements first for easy indexing
        flat_elements = []
        for el in dom_children:
            el_type = el.get("type", "paragraph")
            el_id = el.get("id")
            
            if not el_id:
                continue
            
            if el_type not in ("metadata", "section"):
                last_id = el_id
                if first_id is None:
                    first_id = el_id

            el_summary = {
                "id": el_id,
                "type": el_type,
            }

            if el_type == "paragraph":
                role = el.get("role", "body")
                text = el.get("text", "").strip()
                el_summary["role"] = role
                el_summary["text_preview"] = text[:120]
                el_summary["text_length"] = len(text)
                
                # Extract style summary
                styles = set()
                for run in el.get("runs", []):
                    s = run.get("style", {})
                    if s.get("bold"): styles.add("bold")
                    if s.get("italic"): styles.add("italic")
                    if s.get("underline"): styles.add("underline")
                    if s.get("strike"): styles.add("strikethrough")
                    if s.get("highlight"): styles.add(f"highlight:{s['highlight']}")
                    if s.get("color"): styles.add(f"color:{s['color']}")
                    if s.get("size"):
                        try:
                            # Round to nearest integer to avoid trivial diff failures
                            styles.add(f"size:{round(float(s['size']))}")
                        except (ValueError, TypeError):
                            styles.add(f"size:{s['size']}")
                    if s.get("font"): styles.add(f"font:{s['font']}")
                
                # Also extract paragraph-level styles that verifier might care about
                p_style = el.get("style", {})
                if p_style.get("line_spacing"):
                    try:
                        styles.add(f"spacing:{round(float(p_style['line_spacing']), 1)}")
                    except (ValueError, TypeError):
                        None
                if p_style.get("page_break_before"):
                    styles.add("page_break_before")
                if p_style.get("space_before_pt") is not None:
                    styles.add(f"space_before:{round(float(p_style['space_before_pt']))}pt")
                if p_style.get("space_after_pt") is not None:
                    styles.add(f"space_after:{round(float(p_style['space_after_pt']))}pt")
                if p_style.get("left_indent_pt") is not None:
                    styles.add(f"left_indent:{round(float(p_style['left_indent_pt']))}pt")
                if p_style.get("right_indent_pt") is not None:
                    styles.add(f"right_indent:{round(float(p_style['right_indent_pt']))}pt")
                if p_style.get("first_line_indent_pt") is not None:
                    styles.add(f"first_line_indent:{round(float(p_style['first_line_indent_pt']))}pt")
                if p_style.get("keep_with_next"):
                    styles.add("keep_with_next")
                if p_style.get("keep_together"):
                    styles.add("keep_together")

                list_info = el.get("list_info")
                if list_info:
                    if list_info.get("list_type"):
                        styles.add(f"list_type:{list_info.get('list_type')}")
                    if list_info.get("num_fmt"):
                        styles.add(f"list_fmt:{list_info.get('num_fmt')}")
                    if list_info.get("lvl_text"):
                        styles.add(f"list_lvl_text:{list_info.get('lvl_text')}")

                if styles:
                    el_summary["style_summary"] = sorted(list(styles))
                if p_style.get("alignment"):
                    el_summary["alignment"] = p_style["alignment"]

                if role == "heading":
                    hlvl = el.get("heading_level", 1)
                    el_summary["heading_level"] = hlvl
                    # Add to headings index
                    normalized_heading = text.lower().strip()
                    if normalized_heading:
                        indices["headings_by_name"][normalized_heading] = el_id
            elif el_type == "table":
                table_counter += 1
                rows = el.get("row_count", len(el.get("rows", [])))
                cols = el.get("col_count", 0)
                el_summary["role"] = "table"
                el_summary["ordinal_label"] = f"Table {table_counter}"
                el_summary["rows"] = rows
                el_summary["cols"] = cols
                t_style = el.get("style", {})
                t_styles = []
                if t_style.get("style_name"):
                    t_styles.append(f"table_style:{t_style['style_name']}")
                if t_style.get("has_custom_borders"):
                    t_styles.append("custom_borders")
                if t_style.get("border_colors"):
                    t_styles.append(f"border_colors:{','.join(t_style['border_colors'])}")
                
                # Check for custom cell background colors
                cell_bgs = set()
                for row in el.get("rows", []):
                    for cell in row.get("cells", []):
                        if cell.get("bg_color"):
                            cell_bgs.add(cell["bg_color"])
                if cell_bgs:
                    t_styles.append(f"cell_bgs:{','.join(list(cell_bgs))}")
                
                # Check for cell vertical alignments and paragraph alignments
                valigns = set()
                aligns = set()
                table_text = []
                for row in el.get("rows", []):
                    for cell in row.get("cells", []):
                        if cell.get("vertical_alignment"):
                            valigns.add(cell["vertical_alignment"])
                        for child in cell.get("children", []):
                            if child.get("text"):
                                table_text.append(child["text"])
                            c_style = child.get("style", {})
                            if c_style and c_style.get("alignment"):
                                aligns.add(c_style["alignment"])
                if valigns:
                    t_styles.append(f"valign:{','.join(sorted(valigns))}")
                if aligns:
                    t_styles.append(f"align:{','.join(sorted(aligns))}")

                if table_text:
                    full_text = " | ".join(table_text)
                    el_summary["text_preview"] = full_text[:120]
                    el_summary["text_length"] = len(full_text)

                if t_styles:
                    el_summary["style_summary"] = t_styles
                
                indices["tables_by_ordinal"][str(table_counter)] = el_id
            elif el_type == "image":
                image_counter += 1
                w_emu = el.get("width_emu", 0)
                h_emu = el.get("height_emu", 0)
                w_cm = round(w_emu / 360000, 1) if w_emu else "?"
                h_cm = round(h_emu / 360000, 1) if h_emu else "?"
                image_text = " ".join(
                    str(el.get(k, "")).strip()
                    for k in ("description", "title", "name", "text")
                    if str(el.get(k, "")).strip()
                )
                el_summary["role"] = "image"
                el_summary["ordinal_label"] = f"Image {image_counter}"
                el_summary["size"] = f"{w_cm}cm x {h_cm}cm"
                el_summary["alt_text"] = el.get("description", "")
                el_summary["name"] = el.get("name", "")
                el_summary["title"] = el.get("title", "")
                el_summary["text_preview"] = image_text[:120]
                if el.get("alignment"):
                    el_summary["alignment"] = el.get("alignment")
                indices["images_by_ordinal"][str(image_counter)] = el_id
            elif el_type == "section":
                el_summary["role"] = "section"
                if el.get("style"):
                    el_summary["style_summary"] = [f"{k}:{v}" for k, v in el["style"].items()]
            elif el_type == "metadata":
                el_summary["role"] = "metadata"
                if el.get("properties"):
                    el_summary["style_summary"] = [f"{k}:{v}" for k, v in el["properties"].items()]
            
            flat_elements.append(el_summary)

        indices["last_element_id"] = last_id
        indices["first_content_id"] = first_id
        indices["all_elements"] = flat_elements

        # Group flat elements hierarchically into sections
        headings_and_start = []
        headings_and_start.append({
            "idx": -1, 
            "el": {"id": "start", "heading_level": 0, "text_preview": "Document Start", "role": "heading"}
        })
        for i, el in enumerate(flat_elements):
            if el.get("role") == "heading":
                headings_and_start.append({"idx": i, "el": el})
                
        flat_sections = []
        for i, h_info in enumerate(headings_and_start):
            h_idx = h_info["idx"]
            h_el = h_info["el"]
            
            # end_idx for flat (non-overlapping) elements
            end_idx_flat = len(flat_elements) - 1
            if i + 1 < len(headings_and_start):
                end_idx_flat = headings_and_start[i + 1]["idx"] - 1
            
            start_elem = h_idx if h_idx >= 0 else 0
            section_elements = flat_elements[start_elem:end_idx_flat + 1] if start_elem <= end_idx_flat else []
            
            # end_idx for the HIERARCHICAL section (same or higher level)
            end_idx_hier = len(flat_elements) - 1
            for j in range(i + 1, len(headings_and_start)):
                if headings_and_start[j]["el"].get("heading_level", 1) <= h_el.get("heading_level", 0):
                    end_idx_hier = headings_and_start[j]["idx"] - 1
                    break
            
            section_elements_hier = flat_elements[start_elem:end_idx_hier + 1] if start_elem <= end_idx_hier else []
            
            # Boundary conditions:
            # - If section_elements_hier is empty, start_id == end_id == heading_id
            # - If it's the last section, end_id will automatically be the last element's id
            hier_start_id = section_elements_hier[0]["id"] if section_elements_hier else h_el["id"]
            hier_end_id = section_elements_hier[-1]["id"] if section_elements_hier else h_el["id"]
            
            is_toc = False
            text = h_el.get("text_preview", "")
            if h_el["id"] != "start":
                is_toc = "table of contents" in text.lower() or "contents" in text.lower() or "toc" == text.lower().strip()
                
            sec = {
                "heading": text,
                "heading_id": h_el["id"],
                "semantic_type": "toc" if is_toc else "section",
                "heading_level": h_el.get("heading_level", 0),
                "ordinal": i,
                "elements": section_elements,
                "section_start_id": hier_start_id,
                "section_end_id": hier_end_id,
                "child_sections": []
            }
            if section_elements or h_el["id"] != "start":
                flat_sections.append(sec)

        # Build child_sections tree
        stack = []
        tree_sections = []
        for sec in flat_sections:
            while stack and stack[-1]["heading_level"] >= sec["heading_level"]:
                stack.pop()
            
            if stack:
                stack[-1]["child_sections"].append(sec)
            else:
                tree_sections.append(sec)
            stack.append(sec)

        # Document title is inferred from first section heading
        title = "Untitled Document"
        if flat_sections and len(flat_sections) > 0:
            first_sect = flat_sections[0]
            if first_sect["heading_id"] != "start":
                title = first_sect["heading"]
            elif len(flat_sections) > 1:
                title = flat_sections[1]["heading"]

        return {
            "document_type": "docx",
            "title": title,
            "element_count": len(flat_elements),
            "sections": flat_sections,
            "tree_sections": tree_sections,
            "indices": {
                "tables_by_ordinal": indices["tables_by_ordinal"],
                "images_by_ordinal": indices["images_by_ordinal"],
                "headings_by_name": indices["headings_by_name"],
                "last_element_id": indices["last_element_id"],
                "first_content_id": indices["first_content_id"],
            }
        }

    @classmethod
    def _build_docx_fallback(cls, blocks: list) -> dict:
        """Fallback to build outline from flat blocks array."""
        flat_elements = []
        indices: dict[str, Any] = {
            "tables_by_ordinal": {},
            "images_by_ordinal": {},
            "headings_by_name": {},
            "last_element_id": None,
            "first_content_id": None,
        }

        table_counter = 0
        image_counter = 0
        last_id = None
        first_id = None

        for idx, b in enumerate(blocks):
            el_id = b.get("element_id")
            if not el_id:
                continue
            
            last_id = el_id
            if first_id is None:
                first_id = el_id

            meta = b.get("metadata", {})
            role = meta.get("role", "body")
            text = b.get("text", "").strip()

            el_summary = {
                "id": el_id,
                "type": "paragraph",
                "body_index": idx,
                "role": role,
                "text_preview": text[:120],
                "text_length": len(text),
            }

            if role == "heading":
                el_summary["heading_level"] = meta.get("heading_level", 1)
                normalized_heading = text.lower().strip()
                if normalized_heading:
                    indices["headings_by_name"][normalized_heading] = el_id
            
            flat_elements.append(el_summary)

        indices["last_element_id"] = last_id
        indices["first_content_id"] = first_id

        # Since it's fallback flat data, return a single big section
        sections = [{
            "heading": "Document Body",
            "heading_id": "body",
            "semantic_type": "section",
            "heading_level": 0,
            "ordinal": 1,
            "elements": flat_elements,
            "section_start_id": first_id if first_id else "body",
            "section_end_id": last_id if last_id else "body",
            "child_sections": [],
        }]

        return {
            "document_type": "docx",
            "title": "Document (Flat fallback)",
            "element_count": len(flat_elements),
            "sections": sections,
            "indices": indices,
        }

    @classmethod
    def _build_pptx(cls, structure: dict) -> dict:
        slides = structure.get("slides", [])
        sections = []
        indices: dict[str, Any] = {
            "slides_by_index": {},
            "images_by_ordinal": {},
            "tables_by_ordinal": {},
        }

        image_counter = 0
        table_counter = 0

        for slide in slides:
            slide_idx = slide.get("slide_index", 1)
            layout = slide.get("layout_name", "unknown")
            elements = []

            for shape in slide.get("shapes", []):
                shape_idx = shape.get("shape_index", 0)
                shape_type = "shape"
                
                el_summary = {
                    "id": f"slide_{slide_idx}_shape_{shape_idx}",
                    "shape_index": shape_idx,
                    "shape_name": shape.get("shape_name", ""),
                }

                if shape.get("has_text_frame"):
                    shape_type = "text_frame"
                    text = " ".join(p.get("text", "") for p in shape.get("paragraphs", []))
                    el_summary["type"] = "text_frame"
                    el_summary["text_preview"] = text[:120]
                    el_summary["text_length"] = len(text)
                elif shape.get("has_table"):
                    shape_type = "table"
                    table_counter += 1
                    el_summary["type"] = "table"
                    el_summary["ordinal_label"] = f"Table {table_counter}"
                    el_summary["rows"] = len(shape.get("table_rows", []))
                    indices["tables_by_ordinal"][str(table_counter)] = el_summary["id"]
                elif shape.get("has_image") or "image" in shape.get("shape_name", "").lower():
                    shape_type = "image"
                    image_counter += 1
                    el_summary["type"] = "image"
                    el_summary["ordinal_label"] = f"Image {image_counter}"
                    indices["images_by_ordinal"][str(image_counter)] = el_summary["id"]

                elements.append(el_summary)

            slide_title = next((el.get("text_preview", "") for el in elements if el.get("type") == "text_frame" and el.get("text_preview")), f"Slide {slide_idx}")
            
            section = {
                "heading": slide_title,
                "heading_id": f"slide_{slide_idx}",
                "semantic_type": "slide",
                "ordinal": slide_idx,
                "layout": layout,
                "elements": elements,
            }
            sections.append(section)
            indices["slides_by_index"][str(slide_idx)] = f"slide_{slide_idx}"

        return {
            "document_type": "pptx",
            "title": "Presentation Structure",
            "element_count": len(slides),
            "sections": sections,
            "indices": indices,
        }

# ===== END services/outline_builder.py =====

# ===== BEGIN services/context_fetcher.py =====
"""Context Fetcher — Stage 3 helper of the task graph pipeline.

Pulls full text/structural content for targeted element IDs from the DOM.
Avoids truncation of text fields so the LLM gets full context for editing.
"""

import logging

log = logging.getLogger(__name__)


class ContextFetcher:
    """Retrieves full content of elements by ID from the document structure."""

    @classmethod
    def fetch(cls, element_ids: list[str], structure: dict) -> dict[str, Any]:
        """Return a mapping of {element_id: content_details} for requested IDs."""
        if not element_ids:
            return {}

        dom = structure.get("dom", {})
        id_to_element = {}
        cls._walk_dom(dom, id_to_element)

        # For slide presentations, structure might be flat blocks
        blocks = structure.get("blocks", [])
        for b in blocks:
            eid = b.get("element_id")
            if eid and eid not in id_to_element:
                id_to_element[eid] = {
                    "id": eid,
                    "type": "paragraph",
                    "text": b.get("text", ""),
                    "metadata": b.get("metadata", {}),
                }

        results = {}
        for eid in element_ids:
            el = id_to_element.get(eid)
            if el:
                results[eid] = el
            else:
                results[eid] = {"id": eid, "text": "", "error": "Not found in document structure"}

        return results

    @classmethod
    def _walk_dom(cls, node: dict, acc: dict[str, dict]):
        nid = node.get("id")
        if nid:
            acc[nid] = {
                "id": nid,
                "type": node.get("type", "paragraph"),
                "text": node.get("text", ""),
                "role": node.get("role", "body"),
                "heading_level": node.get("heading_level"),
                "style": node.get("style", {}),
                "list_info": node.get("list_info", {}),
                "rows": node.get("rows", []),
                "row_count": node.get("row_count"),
                "col_count": node.get("col_count"),
                "width_emu": node.get("width_emu"),
                "height_emu": node.get("height_emu"),
                "description": node.get("description", ""),
                "alignment": node.get("alignment", ""),
            }
        
        children = node.get("children", [])
        rows = node.get("rows", [])
        cells = node.get("cells", [])

        for child in children + rows + cells:
            if isinstance(child, dict):
                cls._walk_dom(child, acc)
        
        # PPTX shapes have slide_index and shape_index inside slides
        slides = node.get("slides", [])
        for slide in slides:
            slide_idx = slide.get("slide_index")
            for shape in slide.get("shapes", []):
                shape_idx = shape.get("shape_index")
                sid = f"slide_{slide_idx}_shape_{shape_idx}"
                acc[sid] = {
                    "id": sid,
                    "type": "shape",
                    "slide_index": slide_idx,
                    "shape_index": shape_idx,
                    "shape_name": shape.get("shape_name", ""),
                    "has_text_frame": shape.get("has_text_frame"),
                    "has_table": shape.get("has_table"),
                    "has_image": shape.get("has_image"),
                    "paragraphs": shape.get("paragraphs", []),
                    "table_rows": shape.get("table_rows", []),
                }

# ===== END services/context_fetcher.py =====

# ===== BEGIN services/reference_resolver.py =====
"""Reference Resolver — resolves target hints to stable DOM element IDs.

Two-pass resolution Strategy:
1. Deterministic Structural Matching (No LLM): Handles ordinal lookups
   ("table 2", "third paragraph"), exact named sections ("Conclusion"),
   and structural roles ("all headings", "bulleted lists").
2. LLM-assisted Semantic Matching: Falls back to Gemini to resolve
   ambiguous descriptions ("the section discussing growth metrics").
"""

import logging
import re
from typing import Any


log = logging.getLogger(__name__)


class ReferenceResolver:
    """Resolves target_hint to concrete DOM element IDs from the outline."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    def resolve(self, target_hint: str, outline: dict, task_description: str = "") -> dict:
        """Resolve a target hint string to a dictionary with matching element IDs and optional section range."""
        if not target_hint:
            return {"ids": []}

        hint_lower = target_hint.lower().strip()

        # Pass 1: Deterministic matching
        resolved = self._structural_resolve(hint_lower, outline)
        if resolved and resolved.get("ids"):
            log.info("Pass 1: Deterministic resolve matched target '%s' to: %s", target_hint, resolved)
            return resolved

        # Pass 2: Semantic fallback (LLM-assisted)
        resolved_ids = self._semantic_resolve(target_hint, outline, task_description)
        log.info("Pass 2: Semantic fallback resolved target '%s' to: %s", target_hint, resolved_ids)
        return {"ids": resolved_ids}

    def _section_last_content_id(self, section: dict) -> str | None:
        """Return the ID of the last non-heading element in a section."""
        elements = section.get("elements", [])
        for el in reversed(elements):
            if el.get("role") != "heading" and el.get("id"):
                return el["id"]
        return section.get("section_end_id") or section.get("heading_id")

    def _find_move_sections(self, hint: str, outline: dict) -> dict | None:
        """Detect 'X above Y', 'X below Y', 'X before Y', 'X after Y' move patterns.

        Returns a dict with:
          - section_to_move: section dict
          - anchor_section: section dict
          - direction: 'before' | 'after'
        Or None if not matched.
        """
        sections = outline.get("sections", [])

        def _match_section(query: str) -> dict | None:
            query = query.strip()
            query_normalized = re.sub(r'^\d+\.\s*', '', query).strip()
            query_no_section = query_normalized.replace("the ", "").replace(" section", "").strip()
            for s in sections:
                if s.get("heading_id") == "start":
                    continue
                name = s.get("heading", "").lower()
                name_normalized = re.sub(r'^\d+\.\s*', '', name).strip()
                if (query == name or query_normalized == name_normalized or
                        query_no_section == name_normalized or
                        query_normalized in name_normalized or
                        name_normalized in query_normalized):
                    return s
            return None

        for pattern, direction in [
            (r'^(.+?)\s+(?:above|before)\s+(.+)$', 'before'),
            (r'^(.+?)\s+(?:below|after)\s+(.+)$', 'after'),
        ]:
            m = re.match(pattern, hint.strip())
            if m:
                a_query = m.group(1).strip()
                b_query = m.group(2).strip()
                # Strip leading verbs from section A
                for verb in ("move ", "place ", "put ", "swap ", "exchange "):
                    if a_query.startswith(verb):
                        a_query = a_query[len(verb):]
                # Strip trailing qualifiers
                a_query = a_query.replace(" sections", "").replace(" section", "").strip()
                b_query = b_query.replace(" sections", "").replace(" section", "").strip()

                sec_a = _match_section(a_query)
                sec_b = _match_section(b_query)
                if sec_a and sec_b and sec_a != sec_b:
                    return {
                        "section_to_move": sec_a,
                        "anchor_section": sec_b,
                        "direction": direction,
                    }

        return None

    def _find_two_sections(self, hint: str, outline: dict) -> list[dict] | None:
        """Try to identify exactly two sections from a hint containing 'and'.

        Returns a list of two section dicts, or None if fewer than 2 were found.
        """
        sections = outline.get("sections", [])
        
        def _match_section(query: str) -> dict | None:
            query = query.strip()
            query_normalized = re.sub(r'^\d+\.\s*', '', query).strip()
            query_no_section = query_normalized.replace("the ", "").replace(" section", "").strip()
            for s in sections:
                if s.get("heading_id") == "start":
                    continue
                name = s.get("heading", "").lower()
                name_normalized = re.sub(r'^\d+\.\s*', '', name).strip()
                if (query == name or query_normalized == name_normalized or
                        query_no_section == name_normalized or
                        query_normalized in name_normalized or
                        name_normalized in query_normalized):
                    return s
            return None
        
        # Split on " and " — first try splitting evenly
        parts = hint.split(" and ")
        if len(parts) < 2:
            return None
        
        # Handle "swap A and B" — strip leading verbs
        clean_parts = []
        for p in parts:
            p = p.strip()
            for verb in ("swap ", "exchange ", "switch ", "move "):
                if p.startswith(verb):
                    p = p[len(verb):]
            # Strip trailing qualifiers
            p = p.replace(" sections", "").replace(" section", "").strip()
            clean_parts.append(p)
        
        matched = []
        for part in clean_parts:
            s = _match_section(part)
            if s and s not in matched:
                matched.append(s)
        
        return matched if len(matched) == 2 else None

    def _structural_resolve(self, hint: str, outline: dict) -> dict | None:
        indices = outline.get("indices", {})
        sections = outline.get("sections", [])

        # Detect move intent: "X above Y", "X below Y", "X before Y", "X after Y"
        # Must check BEFORE the swap check because "X above Y" doesn't contain " and "
        if any(kw in hint for kw in (" above ", " below ")):
            move_result = self._find_move_sections(hint, outline)
            if move_result:
                s_move = move_result["section_to_move"]
                s_anchor = move_result["anchor_section"]
                direction = move_result["direction"]
                result: dict = {
                    "ids": [s_move["heading_id"]],
                    "section_range": {
                        "start_id": s_move.get("section_start_id"),
                        "end_id": s_move.get("section_end_id"),
                    },
                }
                if direction == "before":
                    result["before_anchor_id"] = s_anchor.get("section_start_id")
                else:
                    result["after_anchor_id"] = s_anchor.get("section_end_id")
                return result

        # Detect swap intent: "action items and key metrics sections" or similar
        # Pattern: "[Name A] and [Name B] sections" or "swap [Name A] and [Name B]"
        if " and " in hint and ("section" in hint or "swap" in hint or "exchange" in hint or "switch" in hint):
            matched_sections = self._find_two_sections(hint, outline)
            if matched_sections and len(matched_sections) == 2:
                s_a, s_b = matched_sections
                return {
                    "ids": [s_a["heading_id"], s_b["heading_id"]],
                    "section_a_range": {
                        "start_id": s_a.get("section_start_id"),
                        "end_id": s_a.get("section_end_id"),
                    },
                    "section_b_range": {
                        "start_id": s_b.get("section_start_id"),
                        "end_id": s_b.get("section_end_id"),
                    },
                }

        # Match "table of contents" or "toc"
        if "table of contents" in hint or hint == "toc":
            for s in sections:
                if s.get("semantic_type") == "toc":
                    return {"ids": [s["heading_id"]]}
            # Return first section if it matches TOC heading
            if sections and "contents" in sections[0]["heading"].lower():
                return {"ids": [sections[0]["heading_id"]]}

        # Match ordinals: "table N" or "table number N"
        table_match = re.search(r"\btable\s*(?:number\s*)?(\d+)\b", hint)
        if table_match:
            table_num = table_match.group(1)
            tbl_id = indices.get("tables_by_ordinal", {}).get(table_num)
            if tbl_id:
                return {"ids": [tbl_id]}

        # Match ordinals: "image N" or "image number N" or "figure N" or "figure number N"
        image_match = re.search(r"\b(?:image|figure|photo)\s*(?:number\s*)?(\d+)\b", hint)
        if image_match:
            img_num = image_match.group(1)
            img_id = indices.get("images_by_ordinal", {}).get(img_num)
            if img_id:
                return {"ids": [img_id]}

        image_words = {
            "image", "picture", "photo", "figure", "logo", "placeholder", "chart",
            "graph", "company", "sales",
        }
        if any(word in hint for word in image_words):
            image_matches = []
            for sec in sections:
                for el in sec.get("elements", []):
                    if el.get("type") != "image":
                        continue
                    haystack = " ".join(
                        str(el.get(k, "") or "")
                        for k in ("text_preview", "alt_text", "name", "title", "ordinal_label")
                    ).lower()
                    if haystack and (
                        hint in haystack
                        or any(tok and tok in haystack for tok in re.findall(r"[a-z0-9]+", hint) if tok not in {"the", "a", "an", "image", "picture", "photo", "figure"})
                    ):
                        image_matches.append(el["id"])
            if len(image_matches) == 1:
                return {"ids": image_matches}

        # Match "all headings" or "headings"
        if hint in ("all headings", "headings", "section headings", "heading"):
            heading_ids = []
            for s in sections:
                if s.get("heading_id") and s["heading_id"] != "start":
                    heading_ids.append(s["heading_id"])
            if heading_ids:
                return {"ids": heading_ids}

        if hint in ("important content", "key content", "critical content", "main content", "important points", "key points"):
            scored: list[tuple[int, str]] = []
            signal_words = {
                "key", "important", "critical", "significant", "risk", "risks",
                "revenue", "growth", "decline", "increase", "decrease", "metric",
                "metrics", "esg", "finding", "findings", "recommendation",
                "recommendations", "conclusion", "objective", "objectives",
                "cybersecurity", "compliance", "supply", "opportunity", "impact",
            }
            for sec in sections:
                for el in sec.get("elements", []):
                    if el.get("type") != "paragraph" or el.get("role") == "heading":
                        continue
                    text = str(el.get("text_preview", "") or "")
                    words = set(re.findall(r"[a-zA-Z]+", text.lower()))
                    score = len(words & signal_words)
                    if re.search(r"[$€£]?\d+(?:\.\d+)?\s*(?:%|percent|million|billion|k|m|bn)?", text, re.I):
                        score += 3
                    if el.get("style_summary"):
                        score += 1
                    if score > 0 and el.get("id"):
                        scored.append((score, el["id"]))
            scored.sort(key=lambda item: item[0], reverse=True)
            selected = [eid for _, eid in scored[:8]]
            if selected:
                return {"ids": selected}

        # Match "all tables" or "tables"
        if hint in ("all tables", "tables", "both tables"):
            return {"ids": list(indices.get("tables_by_ordinal", {}).values())}

        # Match only explicit plural/all-image requests deterministically. Singular
        # semantic labels like "logo" or "placeholder image" must go through the
        # semantic resolver so we do not edit the wrong image.
        if hint in ("all images", "images", "both images"):
            return {"ids": list(indices.get("images_by_ordinal", {}).values())}

        if hint in ("the image", "image", "the picture", "picture", "the photo", "photo", "the figure", "figure"):
            image_ids = list(indices.get("images_by_ordinal", {}).values())
            if len(image_ids) == 1:
                return {"ids": image_ids}

        # Match "last paragraph" or "last element"
        if hint in ("last paragraph", "last element", "end of document", "the end", "the end of the document", "after end of document", "after the end of the document"):
            last_id = indices.get("last_element_id")
            if last_id:
                result = {"ids": [last_id]}
                if "after" in hint:
                    result["after_anchor_id"] = last_id
                return result

        # Match "first paragraph" or "first element"
        if hint in ("first paragraph", "first element", "beginning of document", "the beginning", "the beginning of the document", "start of document"):
            first_id = indices.get("first_content_id")
            if first_id:
                return {"ids": [first_id]}

        # Try heading name matching with numbered prefix stripping
        # e.g. hint "executive summary" matches stored name "1. executive summary"
        headings_by_name = indices.get("headings_by_name", {})

        matched_h_id = None

        # Normalize hint: strip leading number prefix like "1." or "2."
        hint_normalized = re.sub(r'^\d+\.\s*', '', hint).strip()
        hint_no_section = hint_normalized.replace("the ", "").replace(" section", "").strip()

        # First: exact match (covers plain and numbered-stripped cases)
        if hint in headings_by_name:
            matched_h_id = headings_by_name[hint]
        elif hint_normalized in headings_by_name:
            matched_h_id = headings_by_name[hint_normalized]
        else:
            # Substring match
            for name, h_id in headings_by_name.items():
                name_normalized = re.sub(r'^\d+\.\s*', '', name).strip()
                
                is_match = (
                    hint == name or
                    hint_normalized == name_normalized or
                    hint_normalized in name_normalized or
                    name_normalized in hint_normalized or
                    hint_no_section == name_normalized or
                    hint_no_section in name_normalized
                )
                
                if is_match:
                    matched_h_id = h_id
                    break

        if matched_h_id:
            # Determine intent: section-level or just the heading element?
            # Keywords that indicate section-level intent
            is_section_intent = (
                "section" in hint or
                "after" in hint or
                "below" in hint or
                "end of" in hint or
                "bottom of" in hint or
                "inside" in hint
            )

            # Find the matching section object to get boundary IDs
            matched_section = None
            for s in outline.get("sections", []):
                if s.get("heading_id") == matched_h_id:
                    matched_section = s
                    break

            if is_section_intent and matched_section:
                section_ids = [matched_h_id]
                section_range = {
                    "start_id": matched_section.get("section_start_id"),
                    "end_id": matched_section.get("section_end_id")
                }
                for el in matched_section.get("elements", []):
                    if el.get("id") and el["id"] != matched_h_id:
                        section_ids.append(el["id"])

                # after_anchor_id = last CONTENT element in section (not the heading)
                after_anchor_id = self._section_last_content_id(matched_section)

                result = {"ids": section_ids, "section_range": section_range}
                if after_anchor_id:
                    result["after_anchor_id"] = after_anchor_id
                return result

            # Heading-only intent (e.g. editing heading text/format)
            return {"ids": [matched_h_id]}

        # Match "slide N" or "slide number N" (for PPTX)
        slide_match = re.search(r"\bslide\s*(?:number\s*)?(\d+)\b", hint)
        if slide_match:
            slide_num = slide_match.group(1)
            slide_id = indices.get("slides_by_index", {}).get(slide_num)
            if slide_id:
                return {"ids": [slide_id]}

        return None

    def _semantic_resolve(self, target_hint: str, outline: dict, task_description: str = "") -> list[str]:
        llm = self._llm or LLMClient()

        # Build a list of candidate elements for the LLM to choose from
        candidates = []
        for s in outline.get("sections", []):
            if s.get("heading_id") and s["heading_id"] != "start":
                candidates.append({
                    "id": s["heading_id"],
                    "type": "heading",
                    "text": s["heading"],
                })
            for el in s.get("elements", []):
                if el.get("type") in ("paragraph", "table", "image"):
                    candidates.append({
                        "id": el["id"],
                        "type": el["type"],
                        "text": el.get("text_preview", ""),
                        "alt_text": el.get("alt_text", ""),
                        "name": el.get("name", ""),
                        "title": el.get("title", ""),
                        "size": el.get("size", ""),
                        "alignment": el.get("alignment", ""),
                        "ordinal_label": el.get("ordinal_label", ""),
                    })

        # Cut down candidate list to prevent token overload
        candidates = candidates[:150]

        import json
        system_prompt = (
            "You are a document target resolver.\n"
            "Given a target description (how the user referred to some part of the document) "
            "and a list of all elements in the document, identify the element ID(s) that "
            "the user is targeting.\n\n"
            "CRITICAL RULES:\n"
            "1. Choose ONLY from the list of provided elements. Do NOT invent IDs.\n"
            "2. If the user targets a section (e.g. 'the executive summary'), return ALL element IDs contained within that section (e.g. the heading ID, followed by all paragraphs, tables, images, etc. in that section).\n"
            "3. If the user targets a specific paragraph (e.g. 'the paragraph about revenue growth'), return the ID of that paragraph.\n"
            "4. Return multiple IDs ONLY if the reference clearly targets multiple elements (e.g. 'all paragraphs in section X', or 'the entire section').\n"
            "5. For 'important/key/critical content', return only a focused subset of high-signal body paragraphs containing metrics, risks, findings, recommendations, conclusions, dates, percentages, financial values, or other decision-critical claims. NEVER return all paragraphs for this request.\n"
            "6. The provided elements may lack explicit types (e.g., all elements might be marked as 'paragraph'). You MUST infer semantic roles (headings, list items, etc.) based on the actual 'text' content and structural patterns (like short phrases followed by multiple sentences).\n\n"
            "Return a JSON object with a single key 'ids' containing an array of matched element IDs:\n"
            "{\n"
            '  "ids": ["id_1", "id_2"]\n'
            "}"
        )

        user_prompt = (
            f"Task context: {task_description}\n"
            f"Target hint: {target_hint}\n\n"
            f"Candidate document elements:\n{json.dumps(candidates, indent=2)}"
        )

        import logging
        log = logging.getLogger(__name__)
        
        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0,
            max_tokens=4096,
            json_mode=True,
        ))

        parsed = response.json or {}
        ids = parsed.get("ids", [])
        
        log.warning(f"TARGET_HINT: {target_hint}")
        log.warning(f"RAW LLM RESPONSE: {response.text}")
        log.warning(f"PARSED JSON: {parsed}")
        print(f"TARGET_HINT: {target_hint}")
        print(f"RAW LLM RESPONSE: {response.text}")
        print(f"PARSED JSON: {parsed}")
        
        return [str(i) for i in ids if isinstance(i, str)]

# ===== END services/reference_resolver.py =====

# ===== BEGIN services/content_enricher.py =====
"""Content Enricher — post-processes generated operations to fill in substantive content.

Two responsibilities:
1. **Section content generation**: Any `insert_block` operation whose `data` array is
   empty or contains only placeholder text gets replaced with real, contextually relevant
   content via a single batched LLM call.

2. **Visible ToC generation**: `insert_toc` operations are converted into visible
   2-column tables (Section | Page) built from the document's actual heading structure,
   so the ToC renders correctly in the PDF preview (Word's native TOC field requires
   a field update, which never happens during server-side conversion).

Uses LLMClient for provider-agnostic LLM calls (Gemini by default).
"""

import json
import logging
import re
from datetime import datetime

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Placeholder detection
# ---------------------------------------------------------------------------

_PLACEHOLDER_MARKERS = (
    "[placeholder]", "[content]", "[text]", "[add", "[insert",
    "lorem ipsum", "placeholder", "content goes here", "...", "tbd", "n/a",
)


def _is_placeholder(text: str) -> bool:
    t = text.strip().lower()
    if not t:
        return True
    return any(m in t for m in _PLACEHOLDER_MARKERS)


def _block_needs_enrichment(op: dict) -> bool:
    """Return True if this insert_block op needs content generation."""
    if op.get("op_type") != "layout_op":
        return False
    params = op.get("parameters", {})
    if params.get("action") != "insert_block":
        return False
    data = params.get("data", [])
    if not data:
        return True
    body_items = [d for d in data if d.get("role") != "heading"]
    if not body_items:
        return True

    for d in body_items:
        if d.get("role") in ("table", "toc_field") and (d.get("rows") or d.get("headers") or d.get("field_text")):
            return False
        if d.get("text", "").strip() and not _is_placeholder(d.get("text", "")):
            return False

    return True


# ---------------------------------------------------------------------------
# Document context extraction
# ---------------------------------------------------------------------------

def _extract_document_context(structure: dict) -> dict:
    """Pull a concise context object from the document structure with cumulative word-count page estimation."""
    dom_children = structure.get("dom", {}).get("children", [])

    headings: list[dict] = []
    body_samples: list[str] = []
    cumulative_words = 0

    for el in dom_children:
        text = el.get("text", "").strip()
        word_count = len(text.split()) if text else 0
        if el.get("type") == "table":
            for row in el.get("rows", []):
                for cell in row:
                    word_count += len(str(cell).split())

        role = el.get("role", "body")
        if el.get("type") == "paragraph" and role == "heading" and text:
            est_page = max(1, (cumulative_words // 400) + 1)
            headings.append({
                "text": text,
                "level": el.get("heading_level", 1),
                "estimated_page": est_page,
            })
        elif role == "body" and text and len(body_samples) < 6:
            body_samples.append(text[:300])

        cumulative_words += word_count

    return {
        "headings": headings,
        "body_samples": body_samples,
    }


def _build_doc_summary(ctx: dict) -> str:
    """Build a short human-readable document summary for the LLM."""
    parts = []
    if ctx["headings"]:
        heading_list = ", ".join(h["text"] for h in ctx["headings"][:10])
        parts.append(f"Document sections: {heading_list}.")
    if ctx["body_samples"]:
        parts.append("Excerpt: " + " ".join(ctx["body_samples"][:2])[:400])
    return " ".join(parts) if parts else "General business document."


# ---------------------------------------------------------------------------
# Native Word ToC generation
# ---------------------------------------------------------------------------

def _sanitize_bookmark_name(name: str) -> str:
    import re
    cleaned = re.sub(r'[^a-zA-Z0-9_]', '_', str(name))
    if not cleaned or not (cleaned[0].isalpha() or cleaned[0] == '_'):
        cleaned = '_' + cleaned
    return cleaned[:40]


def _make_visible_toc(insert_toc_op: dict, ctx: dict) -> dict:
    """Convert an insert_toc op into an insert_block op with a 2-column table and PAGEREF fields."""
    before_id = insert_toc_op.get("parameters", {}).get("before_id")
    after_id = insert_toc_op.get("parameters", {}).get("after_id")

    headings = ctx.get("headings", [])
    if not headings:
        headings = [{"text": "No headings found", "level": 1, "estimated_page": 1}]

    rows = []
    for i, h in enumerate(headings):
        indent = "\u00a0\u00a0\u00a0\u00a0" * (h.get("level", 1) - 1)
        h_id = h.get("heading_id") or f"heading_{i+1}"
        bmk_name = f"_Ref_{_sanitize_bookmark_name(h_id)}"
        page_val = str(h.get("estimated_page", 1))
        heading_text = f"{indent}{h['text']}"
        rows.append([
            heading_text,
            {
                "pageref": bmk_name,
                "page": page_val,
                "heading_id": h.get("heading_id"),
            }
        ])

    data = [
        {"role": "heading", "text": "Table of Contents", "heading_level": 1},
        {
            "role": "table",
            "headers": ["Section", "Page"],
            "rows": rows,
            "style": "toc",
        },
    ]

    params: dict = {"action": "insert_block", "data": data}
    if before_id:
        params["before_id"] = before_id
    elif after_id:
        params["after_id"] = after_id

    return {"op_type": "layout_op", "target_id": None, "parameters": params}


# ---------------------------------------------------------------------------
# Batched content generation via LLM
# ---------------------------------------------------------------------------

def _enrich_sections_with_llm(
    ops_needing_enrichment: list[tuple[int, dict]],
    doc_summary: str,
    original_request: str,
    llm,
) -> dict[int, list[dict]]:
    """Call LLM once to generate content for all insert_block ops that need it."""

    sections = []
    for idx, op in ops_needing_enrichment:
        data = op.get("parameters", {}).get("data", [])
        headings_in_op = [d["text"] for d in data if d.get("role") == "heading" and d.get("text")]
        section_title = headings_in_op[0] if headings_in_op else "New Section"
        sections.append({"index": idx, "title": section_title})

    response = llm.complete(LLMRequest(
        system_prompt=(
            "You are a professional document writer. Given a document's context and a list of new sections "
            "that need to be added, generate substantive, professional content for each section.\n\n"
            "RULES:\n"
            "1. For each section, produce 2-3 paragraphs of real, relevant content based on the document context.\n"
            "2. Content must be specific to the document's domain — NOT generic filler or placeholders.\n"
            "3. Each paragraph should be 2-5 sentences.\n"
            "4. For a 'Conclusion' section: summarize the key points from the document and suggest next steps.\n"
            "5. For a 'Risks and Challenges' section: identify specific risks relevant to the document's domain.\n"
            "6. For any other section: write content appropriate to the section title and document context.\n"
            "7. Return ONLY a valid JSON object:\n"
            "   {\n"
            '     "sections": [\n'
            "       {\n"
            '         "index": <number>,\n'
            '         "title": "<section title>",\n'
            '         "data": [\n'
            '           {"role": "heading", "text": "<title>", "heading_level": 2},\n'
            '           {"role": "body", "text": "<paragraph 1>"},\n'
            '           {"role": "body", "text": "<paragraph 2>"},\n'
            '           {"role": "body", "text": "<paragraph 3>"}\n'
            "         ]\n"
            "       }\n"
            "     ]\n"
            "   }\n"
            "8. Do NOT include markdown fences, commentary, or any text outside the JSON object."
        ),
        user_prompt=(
            f"Document context:\n{doc_summary}\n\n"
            f"User's original request: {original_request}\n\n"
            f"Sections needing content:\n{json.dumps(sections, indent=2)}\n\n"
            f"Current date: {datetime.now().strftime('%B %d, %Y')}\n\n"
            "Generate professional, relevant content for each section above."
        ),
        temperature=0.4,
        max_tokens=4096,
        json_mode=True,
    ))

    parsed = response.json or {}
    result: dict[int, list[dict]] = {}
    for section in parsed.get("sections", []):
        idx = section.get("index")
        data = section.get("data", [])
        if isinstance(idx, int) and isinstance(data, list):
            result[idx] = data

    return result


def _format_kb_chunks_for_enricher(chunks: list[dict]) -> str:
    lines = []
    for i, c in enumerate(chunks[:15]):
        meta = c.get("metadata", {})
        source = meta.get("source") or meta.get("doc_id") or "KB Document"
        page = meta.get("page", "")
        loc = f"[{source}"
        if page:
            loc += f", p.{page}"
        loc += "]"
        lines.append(f"[chunk:{i+1}] {loc}\n{c.get('text', '')[:1200]}")
    return "\n\n".join(lines)


def _enrich_sections_with_kb_grounding(
    ops_needing_enrichment: list[tuple[int, dict]],
    doc_summary: str,
    original_request: str,
    kb_evidence: list[dict],
    llm,
) -> dict[int, list[dict]]:
    """Generate content for insert_block ops grounded strictly in provided KB evidence."""

    sections = []
    for idx, op in ops_needing_enrichment:
        data = op.get("parameters", {}).get("data", [])
        headings_in_op = [d["text"] for d in data if d.get("role") == "heading" and d.get("text")]
        section_title = headings_in_op[0] if headings_in_op else "New Section"
        sections.append({"index": idx, "title": section_title})

    formatted_kb = _format_kb_chunks_for_enricher(kb_evidence)

    system_prompt = (
        "You are an expert content writer for professional corporate and audit documents.\n"
        "Your task: Generate substantive content for new document section(s) based STRICTLY on the provided Knowledge Base evidence.\n\n"
        "RULES:\n"
        "1. STRICT GROUNDING: Every factual claim (numbers, percentages, dates, names, statistics, initiatives) MUST come directly from the provided KB context.\n"
        "2. If data for a specific topic is NOT in the KB evidence, do NOT fabricate or estimate details. Just write about what you DO have evidence for.\n"
        "3. NEVER write 'information not available', 'N/A', 'data not provided', or any disclaimers. Just omit topics without data.\n"
        "4. INLINE CITATION: Include citation tags [chunk:N] where N matches the chunk number in the KB CONTEXT when stating factual claims from evidence.\n"
        "5. PROFESSIONAL TONE: Formal, concise, third-person language.\n"
        "6. Return ONLY a valid JSON object:\n"
        "   {\n"
        '     "sections": [\n'
        "       {\n"
        '         "index": <number>,\n'
        '         "title": "<section title>",\n'
        '         "data": [\n'
        '           {"role": "heading", "text": "<title>", "heading_level": 2},\n'
        '           {"role": "body", "text": "<paragraph 1> [chunk:1]"},\n'
        '           {"role": "body", "text": "<paragraph 2> [chunk:2]"}\n'
        "         ]\n"
        "       }\n"
        "     ]\n"
        "   }\n"
        "7. Do NOT include markdown fences, commentary, or any text outside the JSON object."
    )

    user_prompt = (
        f"Document Summary:\n{doc_summary}\n\n"
        f"User Request: {original_request}\n\n"
        f"KB CONTEXT (Ground your content ONLY in these excerpts):\n{formatted_kb}\n\n"
        f"Sections needing content:\n{json.dumps(sections, indent=2)}\n\n"
        "Generate professional, grounded content JSON."
    )

    response = llm.complete(LLMRequest(
        system_prompt=system_prompt,
        user_prompt=user_prompt,
        temperature=0.2,
        max_tokens=4096,
        json_mode=True,
    ))

    parsed = response.json or {}
    result: dict[int, list[dict]] = {}
    for section in parsed.get("sections", []):
        idx = section.get("index")
        data = section.get("data", [])
        if isinstance(idx, int) and isinstance(data, list):
            clean_data = []
            for item in data:
                if isinstance(item, dict) and "text" in item:
                    item_copy = dict(item)
                    item_copy["text"] = re.sub(r'\s*\[chunk:\d+\]', '', str(item["text"]))
                    clean_data.append(item_copy)
                else:
                    clean_data.append(item)
            result[idx] = clean_data

    return result


# ---------------------------------------------------------------------------
# Main enricher class
# ---------------------------------------------------------------------------

class ContentEnricher:
    """Post-processes operations to fill in substantive content and fix ToC.

    Constructor no longer takes api_key/base_url/llm_model directly.
    Pass a pre-built LLMClient instance (or leave None to have one built lazily).
    """

    def __init__(self, llm=None) -> None:
        self._llm = llm

    def enrich(
        self,
        operations: list[dict],
        structure: dict,
        document_type: str,
        original_request: str,
        task: dict | None = None,
    ) -> list[dict]:
        """Enrich operations in-place and return the updated list.

        Only does work for DOCX documents. PPTX is a pass-through.
        """
        if document_type != "docx":
            return operations

        ctx = _extract_document_context(structure)
        doc_summary = _build_doc_summary(ctx)

        enriched = list(operations)

        # ---- Step 1: Convert insert_toc → visible table -----------------
        for i, op in enumerate(enriched):
            if op.get("op_type") == "layout_op":
                params = op.get("parameters", {})
                if params.get("action") == "insert_toc":
                    enriched[i] = _make_visible_toc(op, ctx)
                    log.info(
                        "ContentEnricher: converted insert_toc → visible table at op index %d", i
                    )

        # ---- Step 2: Fill in empty/thin insert_block sections -----------
        ops_needing = [
            (i, op) for i, op in enumerate(enriched)
            if _block_needs_enrichment(op)
        ]

        if not ops_needing:
            return enriched

        log.info(
            "ContentEnricher: %d insert_block op(s) need content enrichment",
            len(ops_needing),
        )

        llm = self._llm or LLMClient()

        try:
            kb_evidence = task.get("kb_evidence") if task else None
            if kb_evidence:
                enriched_data_map = _enrich_sections_with_kb_grounding(
                    ops_needing_enrichment=ops_needing,
                    doc_summary=doc_summary,
                    original_request=original_request,
                    kb_evidence=kb_evidence,
                    llm=llm,
                )
            else:
                enriched_data_map = _enrich_sections_with_llm(
                    ops_needing_enrichment=ops_needing,
                    doc_summary=doc_summary,
                    original_request=original_request,
                    llm=llm,
                )
            for idx, new_data in enriched_data_map.items():
                if new_data:
                    enriched[idx]["parameters"]["data"] = new_data
                    log.info(
                        "ContentEnricher: enriched op %d with %d content items",
                        idx, len(new_data),
                    )
        except Exception as exc:
            log.warning(
                "ContentEnricher: LLM call failed — %s. Using original data.", exc
            )

        return enriched

# ===== END services/content_enricher.py =====

# ===== BEGIN services/operation_generator.py =====
"""Operation Generator — generates structured document operations.

Supports two pipelines:
1. Legacy Pipeline: single-shot operations generation for the entire request.
2. New staged Pipeline: generates operations for a single focused task at a time.
"""

import json
import logging
from datetime import datetime
from typing import Any


log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Per-task JSON schemas for Phase 4
# ---------------------------------------------------------------------------

_TEXT_EDIT_SCHEMA = """Return a JSON array of text_edit operations:
[
  {
    "op_type": "text_edit",
    "target_id": "element_id_to_edit" or ["id1", "id2"],
    "parameters": {
      "new_text": "The complete new rewritten text content. NEVER use HTML tags (e.g. <b>, <span>) or Markdown here. NEVER use this tool to change fonts, colors, alignment, or spacing - use text_format for that!"
    }
  }
]

RULES:
1. For expanding/shortening an existing section, rewrite existing body paragraph(s) in place.
2. Do NOT add a new heading or duplicate the section title.
3. Do NOT target heading elements unless the user asked to change heading text.
4. If Target Element IDs include a heading and body paragraphs, use only the body paragraph IDs for content expansion.
5. For a target word count, write natural prose close to that count across the edited paragraph(s).
6. If Knowledge Base evidence is provided, factual additions MUST be grounded in that evidence. Do not invent or assume numbers, dates, initiatives, risks, opportunities, or outcomes.
7. If the evidence is thin, expand using cautious synthesis of the existing paragraph and evidence without adding unsupported specifics.
"""

_TEXT_FORMAT_SCHEMA = """Return a JSON array of text_format operations:
[
  {
    "op_type": "text_format",
    "target_id": "element_id_to_format" or "all" or ["id1", "id2"],
    "parameters": {
      "match_text": "exact substring to format (optional, if omitted formats the entire element)",
      "match_role": "heading" or "bullet_point" or "body" (optional),
      "bold": true/false/null,
      "italic": true/false/null,
      "underline": true/false/null,
      "font_family": "Arial"/"Calibri"/null,
      "font_size_pt": 11.5 or null,
      "color_hex": "0000FF" (MUST be 6-char hex like "FF0000" for red, never use words like "red") or null,
      "highlight_hex": "FFFF00" (6-char hex) or null,
      "alignment": "left"/"center"/"right"/"justify" or null,
      "line_spacing": 1.5 or null,
      "char_spacing": 1.0 or null,
      "space_before_pt": 12.0 or null,
      "space_after_pt": 12.0 or null
    }
  }
]
"""

_TABLE_OP_SCHEMA = """Return a JSON array of table_op operations:
[
  {
    "op_type": "table_op",
    "target_id": "table_element_id_or_null" or ["id1", "id2"],
    "parameters": {
      "action": "create"|"delete"|"add_row"|"remove_row"|"add_col"|"remove_col"|"merge_cells"|"set_cell_bg"|"set_borders"|"alternate_rows"|"populate"|"sort_data"|"apply_theme"|"set_header_format"|"set_cell_alignment"|"set_alignment"|"set_width_pct",
      "action_notes": "set_header_format makes the first row dark blue with white text. set_cell_bg sets the background of a specific row/col.",
      "rows": number_of_rows_or_null,
      "cols": number_of_cols_or_null,
      "header_row": true/false/null,
      "alternate_row_colors": ["FFFFFF", "F0F0F0"] or null,
      "data": [["cell1", "cell2"], ["cell3", "cell4"]] or null,
      "row_index": "0-based index or null",
      "col_index": "0-based index or null",
      "sort_by_column": "Name of the column to sort by (for sort_data)",
      "cell_alignment": "left|center|right|justify (for set_cell_alignment)",
      "alignment": "left|center|right (for set_alignment of the entire table)",
      "width_pct": 1.0,
      "before_id": "id_of_anchor_to_insert_before",
      "after_id": "id_of_anchor_to_insert_after",
      "cell_bg_hex": "HEX" or null,
      "border_color_hex": "HEX" or null,
      "theme_color_hex": "HEX" or null,
      "border_width_pt": number_or_null
    }
  }
]

RULES:
1. For creating a summary/findings/metrics table, use action "create" and populate the "data" field with a header row plus evidence-grounded rows.
2. If Knowledge Base evidence is provided, all metric names, values, dates, ESG facts, and findings must come from that evidence. Do not invent missing metric values.
3. For "under/after Conclusion section", use the provided Section insertion anchor as after_id.
4. Do not create separate sections for ESG metrics or other metrics when the user asked for one final table under Conclusion.
"""

_IMAGE_OP_SCHEMA = """Return a JSON array of image_op operations:
[
  {
    "op_type": "image_op",
    "target_id": "image_element_id_or_null" or ["id1", "id2"],
    "parameters": {
      "action": "insert"|"replace"|"remove"|"resize"|"reposition"|"move"|"place_inline"|"add_caption"|"reposition_caption"|"format_caption"|"rotate"|"rounded_corners"|"shadow",
      "image_path": "path_to_image_or_placeholder_or_null",
      "after_id": "element_id_to_insert_after_or_null",
      "before_id": "element_id_to_insert_before_or_null",
      "anchor_id": "paragraph_id_for_inline_title_placement_or_null",
      "placement": "top_right_page|top_page|null",
      "width_page_pct": 0.4,
      "height_page_pct": null,
      "alignment": "left|center|right|null",
      "float_position": "left|right|null",
      "caption_text": "caption text or null",
      "caption_style": "Caption",
      "font_family": "Arial or null",
      "font_size_pt": 12,
      "bold": true/false/null,
      "italic": true/false/null,
      "underline": true/false/null,
      "color_hex": "HEX or null",
      "alt_text": "short identity/description or null",
      "position": {
        "left_pct": 0.1,
        "top_pct": 0.2,
        "width_pct": 0.4,
        "height_pct": 0.5
      },
      "maintain_aspect_ratio": true/false/null,
      "rotation_degrees": number_or_null,
      "border_color_hex": "HEX" or null,
      "border_width_pt": number_or_null,
      "rounded_corners": true/false/null,
      "shadow": true/false/null
    }
  }
]

DOCX IMAGE RULES:
1. For "insert below/after X", set target_id to X and set parameters.after_id to X.
2. For "insert above/before X", set target_id to X and set parameters.before_id to X.
3. For "below the title", use the title/heading element ID from Target Element ID(s) as after_id.
4. For "replace the placeholder image", target the placeholder image ID and use action "replace".
5. For "resize to 40% of page width", use action "resize" and width_page_pct 0.4.
6. For "center/left/right align the image", use action "reposition" and alignment "center"|"left"|"right". This also aligns the adjacent caption when present.
7. For "move the image to the right/left side", use action "reposition" and alignment "right"|"left". Do not remove or replace the image for alignment requests.
8. For captions, use action "add_caption" with caption_text exactly as requested.
9. For "move/center/left/right align just the caption", use action "reposition_caption" and the requested alignment.
10. For moving an existing image to another document location, use action "move" with before_id or after_id. The backend moves the adjacent caption with the image.
11. For formatting an existing caption below an image, use action "format_caption" and include formatting fields such as font_size_pt, font_family, bold, italic, color_hex, or alignment.
12. For "top right corner of the page", use action "move", placement "top_right_page", and alignment "right" for an existing image. For inserting a new uploaded image there, use action "insert", before_id as the first content ID, and alignment "right".
13. For "right side/corner of the title", use action "place_inline" with anchor_id set to the title/heading paragraph ID.
14. Keep target_id as the resolved image ID for replace/resize/reposition/move/place_inline/add_caption/reposition_caption/format_caption/remove. Do not use null for existing-image edits.
"""

_LAYOUT_OP_SCHEMA = """Return a JSON array of layout_op operations.
Each operation has "op_type": "layout_op", "target_id": null, and a "parameters" object.

ACTIONS AND THEIR PARAMETERS:

=== move_block ===
Move a section (heading + content) to a new location. One section relocates; the other stays.
{
  "action": "move_block",
  "start_id": "ID of the first element of the section to MOVE (from Section Range start_id)",
  "end_id": "ID of the last element of the section to MOVE (from Section Range end_id)",
  "before_id": "ID of the element to insert BEFORE (use 'Move anchor' before_id if provided)",
  "after_id": "ID of the element to insert AFTER (use 'Move anchor' after_id if provided)"
}
NOTE: Provide EITHER before_id OR after_id, not both. Use the 'Move anchor' values from the context.

=== swap_sections ===
Exchange two sections — both sections trade positions.
{
  "action": "swap_sections",
  "section_a_start_id": "first element of section A (from Section A range)",
  "section_a_end_id": "last element of section A",
  "section_b_start_id": "first element of section B (from Section B range)",
  "section_b_end_id": "last element of section B"
}

=== insert_page_break ===
Insert a hard page break immediately before a specific element.
{
  "action": "insert_page_break",
  "before_id": "ID of the element to insert the page break before (the section heading ID)"
}
NOTE: Use the heading ID of the section you want to start on a new page as before_id.

=== insert_block ===
Insert new content (paragraphs, bullets, tables, headings) at a specific location.
{
  "action": "insert_block",
  "after_id": "ID of the element to insert AFTER (prefer the Section insertion anchor after_id)",
  "data": [
    {"role": "heading", "text": "Heading text", "heading_level": 2},
    {"role": "body", "text": "Full paragraph prose (2-4 sentences, NOT bullet points)"},
    {"role": "bullet_point", "text": "Single bullet item text"},
    {"role": "table", "headers": ["Col 1", "Col 2"], "rows": [["A", "B"], ["C", "D"]]}
  ]
}

=== set_columns ===
Set a multi-column page layout (e.g., two columns side by side).
{
  "action": "set_columns",
  "num_columns": 2,
  "column_gap_inches": 0.5
}

=== remove_block ===
Delete a range of elements from the document.
{
  "action": "remove_block",
  "start_id": "ID of first element to remove",
  "end_id": "ID of last element to remove"
}

=== duplicate_block ===
Duplicate a range of elements to another location.
{
  "action": "duplicate_block",
  "start_id": "ID of first element to duplicate",
  "end_id": "ID of last element to duplicate",
  "after_id": "ID of element to insert after"
}

=== insert_toc ===
Insert a Table of Contents.
{
  "action": "insert_toc",
  "before_id": "ID of element to insert BEFORE (or null)",
  "after_id": "ID of element to insert AFTER (or null)"
}

CRITICAL RULES:
1. move_block != swap_sections. Use move_block when the user says 'move X above/below Y' (only X relocates).
   Use swap_sections when the user says 'swap X and Y' (both sections change positions).
2. For move_block: set start_id/end_id from the 'Section Range' in context, and before_id/after_id
   from the 'Move anchor' values provided. DO NOT set both before_id and after_id.
3. For insert_page_break: set before_id to the heading element ID of the section that should start
   on the new page. Use the 'Target Element ID(s)' from context if it contains a heading ID.
4. For insert_block: GROUP ALL items into ONE operation. If the user asks for 3 paragraphs,
   emit ONE insert_block with ALL 3 items in data[]. NEVER emit multiple insert_block ops.
5. For insert_toc: When the task asks to add a Table of Contents (TOC), ALWAYS use action 'insert_toc'. NEVER use 'insert_block' or generate paragraph prose for a Table of Contents. The TOC is formatted as a 2-column table (Section | Page) with live Word PAGEREF fields.
6. For set_columns: this applies a document-level column layout — use when user asks for
   'two-column layout', 'multi-column', or 'side-by-side columns'.
7. For TOC dot leaders & formatting: Requests to adjust TOC dot leaders, extend dotted lines to the right, or format TOC entry alignment must NEVER emit 'set_alignment', 'text_format', or generic paragraph formatting operations. TOC tab stops and dot leaders are managed natively in Word via w:tab w:leader='dot' w:val='right'.
8. Output ONLY the raw JSON array — no markdown, no commentary.
"""


_LIST_OP_SCHEMA = """Return a JSON array of list_op operations:
[
  {
    "op_type": "list_op",
    "target_id": "null or 'all' or list of IDs",
    "parameters": {
      "action": "convert_type"|"add_items"|"sort_items"|"set_bullet_char",
      "start_id": "first_item_element_id_or_null",
      "end_id": "last_item_element_id_or_null",
      "after_id": "insert_after_this_element_id_or_null",
      "list_type": "bullet"|"numbered"|"checklist",
      "items": ["list item 1 text", "list item 2 text"],
      "bullet_char": "char_or_null"
    }
  }
]
"""

_FIND_REPLACE_SCHEMA = """Return a JSON array of find_replace operations:
[
  {
    "op_type": "find_replace",
    "target_id": "all",
    "parameters": {
      "find_text": "text_to_find",
      "replace_text": "text_to_replace_with",
      "is_regex": false,
      "match_case": false
    }
  }
]
"""

_THEME_OP_SCHEMA = """Return a JSON array of theme_op operations:
[
  {
    "op_type": "theme_op",
    "target_id": "null or 'all' or list of IDs",
    "parameters": {
      "action": "set_bg_color"|"set_margins"|"add_page_numbers"|"apply_theme_colors",
      "bg_color_hex": "HEX_or_null",
      "margin_inches": number_or_null,
      "accent_colors": ["HEX1", "HEX2"] or null
    }
  }
]
"""

_SLIDE_OP_SCHEMA = """Return a JSON array of slide_op operations (PPTX only):
[
  {
    "op_type": "slide_op",
    "target_id": "slide_id_or_null" or ["id1", "id2"],
    "parameters": {
      "action": "add"|"delete"|"duplicate"|"reorder"|"hide"|"unhide"|"rename_title"|"apply_layout",
      "after_index": 1_based_index_or_null,
      "from_index": 1_based_index_or_null,
      "to_index": 1_based_index_or_null,
      "layout_name": "layout_name_or_null",
      "title": "new_title_or_null"
    }
  }
]
"""

_AI_DESIGN_OP_SCHEMA = """Return a JSON array of ai_design_op operations:
[
  {
    "op_type": "ai_design_op",
    "target_id": null,
    "parameters": {
      "action": "normalize_fonts"|"normalize_spacing"|"improve_hierarchy"|"balance_whitespace"|"remove_overlaps"|"improve_readability",
      "scope": "all_slides"|"slide:1" or null,
      "target_font": "Calibri" or null,
      "base_font_size_pt": 11 or null
    }
  }
]
"""

_CONTENT_GENERATION_SCHEMA = """Return a JSON array of layout_op operations with action 'insert_block':
[
  {
    "op_type": "layout_op",
    "target_id": null,
    "parameters": {
      "action": "insert_block",
      "after_id": "ID of element to insert AFTER (or null)",
      "before_id": "ID of element to insert BEFORE (or null)",
      "data": [
        {"role": "heading", "text": "Section Title", "heading_level": 2},
        {"role": "body", "text": "[Content placeholder]"}
      ]
    }
  }
]

CRITICAL RULES FOR content_generation:
1. For the "body" item in data[], set text to EXACTLY "[Content placeholder]".
2. Do NOT write real paragraph prose at this stage — substantive content will be grounded in Knowledge Base evidence downstream by the enricher.
3. For the "heading" item in data[], set text to the section title requested by the user.
4. Set after_id (or before_id) to the anchor ID provided in context.
5. Output ONLY the raw JSON array — no markdown, no commentary.
"""


_SCHEMA_BY_TYPE = {
    "text_edit": _TEXT_EDIT_SCHEMA,
    "text_format": _TEXT_FORMAT_SCHEMA,
    "table_op": _TABLE_OP_SCHEMA,
    "image_op": _IMAGE_OP_SCHEMA,
    "layout_op": _LAYOUT_OP_SCHEMA,
    "content_generation": _CONTENT_GENERATION_SCHEMA,
    "list_op": _LIST_OP_SCHEMA,
    "find_replace": _FIND_REPLACE_SCHEMA,
    "theme_op": _THEME_OP_SCHEMA,
    "slide_op": _SLIDE_OP_SCHEMA,
    "ai_design_op": _AI_DESIGN_OP_SCHEMA,
}


# ---------------------------------------------------------------------------
# Legacy prompt (retained for backward compatibility of single-shot operations node)
# ---------------------------------------------------------------------------

_SYSTEM_PROMPT = """You are a precise document editing operations generator.
You convert document editing instructions into a list of structured operations.

Available Operation Types:
1. text_edit: Rewrite text content of a targeted paragraph. NEVER use this tool to change fonts, colors, or alignment - use text_format for that!
   - target_id: paragraph element ID or list of IDs
   - parameters: { new_text: str }

2. text_format: Apply formatting to a paragraph.
   - target_id: paragraph element ID or list of IDs
   - parameters: { bold: bool, italic: bool, underline: bool, strikethrough: bool, font_family: str, font_size_pt: float, color_hex: str, highlight_hex: str, alignment: 'left'|'center'|'right'|'justify', line_spacing: float, char_spacing: float, space_before_pt: float, space_after_pt: float, match_role: str }

3. table_op: Create, modify or style tables.
   - target_id: table element ID (or null for insert) or list of IDs
   - parameters: { action: 'create'|'delete'|'add_row'|'remove_row'|'add_col'|'remove_col'|'merge_cells'|'set_cell_bg'|'set_borders'|'alternate_rows'|'populate'|'sort_data'|'apply_theme'|'set_header_format'|'set_alignment'|'set_width_pct', alignment: 'left'|'center'|'right', width_pct: float, rows: int, cols: int, before_id: str, after_id: str, alternate_row_colors: list[str], data: list[list[str]], cell_bg_hex: str, border_color_hex: str, theme_color_hex: str, border_width_pt: float }

4. image_op: Insert or style images.
   - target_id: image element ID (or null for insert) or list of IDs
   - parameters: { action: 'insert'|'replace'|'remove'|'resize'|'reposition'|'rounded_corners'|'shadow', image_path: str, position: { left_pct: float, top_pct: float, width_pct: float, height_pct: float } }

5. layout_op: Manipulate document pages / section order.
   - target_id: null or 'all' or list of IDs
   - parameters: { action: 'move_block'|'insert_page_break'|'remove_block'|'duplicate_block'|'insert_block'|'insert_toc', start_id: str, end_id: str, before_id: str, after_id: str, data: list[dict] }
     (data array for insert_block supports {"role": "heading", "text": "..."}, {"role": "body", "text": "..."}, {"role": "table", "headers": [...], "rows": [...]})

6. list_op: List manipulations.
   - target_id: null or 'all' or list of IDs
   - parameters: { action: 'convert_type'|'add_items'|'sort_items'|'set_bullet_char', start_id: str, end_id: str, after_id: str, list_type: 'bullet'|'numbered'|'checklist', items: list[str] }

7. find_replace: Global search and replace.
   - target_id: 'all'
   - parameters: { find_text: str, replace_text: str, is_regex: bool }

8. theme_op: Theme setting operations.
   - target_id: null or 'all' or list of IDs
   - parameters: { action: 'set_bg_color'|'set_margins'|'add_page_numbers'|'apply_theme_colors', bg_color_hex: str, margin_inches: float }

Return ONLY a JSON array of operations.
"""


class OperationGenerator:
    """Generates structured operations from editing tasks."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    # ------------------------------------------------------------------
    # Stage-focused generate_for_task (Phase 4)
    # ------------------------------------------------------------------

    def generate_for_task(
        self,
        task: dict,
        resolved_ids: dict,
        element_context: dict[str, dict],
        outline: dict,
        attached_image_path: str | None = None,
        previous_ops: list[dict] | None = None,
        verifier_feedback: str | None = None,
        relevant_blocks: dict | None = None,
    ) -> list[dict]:
        """Generate operations for a SINGLE task in the planner's sequence."""
        import re
        task_type = task["task_type"]
        desc_lower_for_route = str(task.get("description", "")).lower()
        target_lower_for_route = str(task.get("target_hint", "")).lower()
        existing_rewrite_intent = (
            any(word in desc_lower_for_route for word in (
                "increase", "expand", "lengthen", "elaborate", "make longer",
                "shorten", "reduce", "condense", "summarize", "summarise",
            ))
            or bool(re.search(r"\b\d+\s*(?:words?|word)\b", desc_lower_for_route))
        )
        existing_section_target = (
            "section" in target_lower_for_route
            or "paragraph" in target_lower_for_route
            or any(
                str(s.get("heading", "")).lower() in target_lower_for_route
                for s in outline.get("sections", [])
                if s.get("heading")
            )
        )
        if task_type == "content_generation" and existing_rewrite_intent and existing_section_target:
            task_type = "text_edit"
            task = {**task, "task_type": "text_edit"}
        schema = _SCHEMA_BY_TYPE.get(task_type)
        if not schema:
            log.warning("No operational schema found for task type: %s", task_type)
            return []

        # If it's an image insertion task but no image is attached, return needs_image
        if task_type == "image_op" and not attached_image_path:
            desc_lower = task["description"].lower()
            caption_intent = "caption" in desc_lower
            requires_new_image = (
                "replace" in desc_lower
                or "insert" in desc_lower
                or (
                    "add" in desc_lower
                    and any(word in desc_lower for word in ("image", "picture", "photo", "logo", "chart", "graph"))
                    and not caption_intent
                )
            )
            if requires_new_image:
                return [needs_image_response(f"To satisfy: '{task['description']}', please upload an image.")]

        llm = self._llm or LLMClient()

        # Build task-specific context
        import json
        ids_list = resolved_ids.get("ids", []) if isinstance(resolved_ids, dict) else resolved_ids
        after_anchor_id = resolved_ids.get("after_anchor_id") if isinstance(resolved_ids, dict) else None
        section_range = resolved_ids.get("section_range") if isinstance(resolved_ids, dict) else None
        
        resolved_str = json.dumps(ids_list)
        context_str = json.dumps(element_context, indent=2)
        
        # Include full sections list in outline summary for layout_op and content_generation so the model 
        # can pick correct section boundaries for insertion/anchors
        if task_type in ("layout_op", "content_generation"):
            outline_summary = json.dumps({
                "document_type": outline.get("document_type"),
                "title": outline.get("title"),
                "indices": outline.get("indices"),
                "sections": [
                    {
                        "heading": s.get("heading"),
                        "heading_id": s.get("heading_id"),
                        "section_start_id": s.get("section_start_id"),
                        "section_end_id": s.get("section_end_id"),
                    }
                    for s in outline.get("sections", [])
                    if s.get("heading_id") != "start"
                ],
            }, indent=2)
        else:
            outline_summary = json.dumps({
                "document_type": outline.get("document_type"),
                "title": outline.get("title"),
                "indices": outline.get("indices"),
            }, indent=2)

        repair_str = ""
        if verifier_feedback and previous_ops:
            repair_str = (
                f"\n=== REPAIR FEEDBACK ===\n"
                f"Your previous operations for this task: {json.dumps(previous_ops)}\n"
                f"Verifier Feedback: {verifier_feedback}\n"
                f"Please fix the operations to satisfy this feedback.\n"
            )

        # Build anchor hint for section-end insertions, swaps, and moves
        before_anchor_id = resolved_ids.get("before_anchor_id") if isinstance(resolved_ids, dict) else None
        anchor_hint = ""
        if before_anchor_id:
            anchor_hint = f"\nMove anchor - insert the section BEFORE this element (use as before_id in move_block): {before_anchor_id}\n"
        elif after_anchor_id:
            anchor_hint = f"\nSection insertion anchor (use this as after_id): {after_anchor_id}\n"
        if section_range:
            anchor_hint += f"Section range (the content to move): start_id={section_range.get('start_id')}, end_id={section_range.get('end_id')}\n"
            
        section_a_range = resolved_ids.get("section_a_range") if isinstance(resolved_ids, dict) else None
        section_b_range = resolved_ids.get("section_b_range") if isinstance(resolved_ids, dict) else None
        
        if section_a_range and section_b_range:
            anchor_hint += f"Section A range (for swap): start_id={section_a_range.get('start_id')}, end_id={section_a_range.get('end_id')}\n"
            anchor_hint += f"Section B range (for swap): start_id={section_b_range.get('start_id')}, end_id={section_b_range.get('end_id')}\n"

        if task_type == "text_edit" and existing_rewrite_intent and existing_section_target:
            body_ids = [
                eid for eid in ids_list
                if element_context.get(eid, {}).get("role") != "heading"
                and element_context.get(eid, {}).get("type") == "paragraph"
            ]
            heading_ids = [
                eid for eid in ids_list
                if element_context.get(eid, {}).get("role") == "heading"
            ]
            if body_ids:
                anchor_hint += (
                    f"Existing-section rewrite body IDs (use these, not heading IDs): {json.dumps(body_ids)}\n"
                )
            if heading_ids:
                anchor_hint += (
                    f"Existing-section heading IDs to preserve unchanged: {json.dumps(heading_ids)}\n"
                )

        system_prompt = (
            f"You are a document operations generator specializing in '{task_type}' changes.\n"
            f"Given a target task, the target element IDs, the content of those elements, "
            f"and the document outline, generate the structured operations required to perform the change.\n\n"
            f"FORMAT SCHEMA FOR '{task_type}':\n{schema}\n"
            "RULES:\n"
            "1. Output ONLY the raw JSON array. Do not include markdown wraps (like ```json) or commentary.\n"
            "2. Use the resolved target IDs exactly as provided. Do NOT invent IDs.\n"
            "3. Make sure all parameters strictly match the schema fields.\n"
            "4. If you receive Repair Feedback, you MUST alter your previous operations to address the error. If the verifier repeatedly complains that a style (like paragraph spacing) was not applied, it may be unsupported by the extraction engine—in that case, do not emit the exact same JSON again; skip it or try an alternative.\n"
            "5. CRITICAL FOR list_op and layout_op: 'target_id' MUST be null! Place the provided Target Element IDs into 'start_id' and 'end_id' instead."
        )

        user_prompt = (
            f"Task: {task['description']}\n"
            f"Target Element ID(s): {resolved_str}\n"
            f"{anchor_hint}"
            f"Element Context: {context_str}\n"
            f"Outline: {outline_summary}\n"
            f"Attached Image: {attached_image_path or 'None'}\n"
        )
        task_kb_evidence = task.get("kb_evidence") if isinstance(task, dict) else None
        if task_kb_evidence:
            user_prompt += (
                "Relevant Knowledge Base Evidence (ground factual claims ONLY in these excerpts; "
                "do not invent numbers, dates, initiatives, risks, metrics, or conclusions not supported here):\n"
                f"{json.dumps(task_kb_evidence, indent=2)}\n"
            )
        elif relevant_blocks:
            user_prompt += f"Relevant Full-Text Blocks (from Knowledge Base):\n{json.dumps(relevant_blocks, indent=2)}\n"
        user_prompt += f"{repair_str}"

        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.0,
            max_tokens=2048,
            json_mode=True,
        ))

        parsed = response.json
        ops_raw = []
        if isinstance(parsed, list):
            ops_raw = parsed
        elif isinstance(parsed, dict):
            ops_raw = parsed.get("operations") or parsed.get("ops") or [parsed]

        # Normalize op_type and populate image_path if needed
        for op in ops_raw:
            if not isinstance(op, dict):
                continue
            if (
                op.get("op_type") == "text_format"
                and "highlight" in desc_lower_for_route
                and any(word in desc_lower_for_route for word in ("important", "key", "critical", "main"))
                and op.get("target_id") in ("all", None)
            ):
                op["_skip_broad_important_highlight"] = True
            if op.get("op_type") == "table_op":
                params = op.setdefault("parameters", {})
                action = str(params.get("action") or "").lower().strip()
                if action == "create" and after_anchor_id and not params.get("after_id") and not params.get("before_id"):
                    params["after_id"] = after_anchor_id
            if op.get("op_type") == "image_op" and attached_image_path:
                op.setdefault("parameters", {})
                if not op["parameters"].get("image_path"):
                    op["parameters"]["image_path"] = attached_image_path
            if op.get("op_type") == "image_op":
                params = op.setdefault("parameters", {})
                action = str(params.get("action") or "").lower().strip()
                desc_lower = task.get("description", "").lower()
                import re as _re
                caption_only = (
                    "caption" in desc_lower
                    and any(word in desc_lower for word in ("align", "move", "center", "centre", "left", "right"))
                    and not any(word in desc_lower for word in ("add", "insert", "create"))
                )
                caption_format = (
                    "caption" in desc_lower
                    and any(word in desc_lower for word in ("font", "size", "bold", "italic", "underline", "color", "colour"))
                    and not any(word in desc_lower for word in ("add", "insert", "create"))
                )

                def _title_anchor_id() -> str | None:
                    for section in outline.get("sections", []):
                        heading_id = section.get("heading_id")
                        if heading_id and heading_id != "start":
                            return heading_id
                    return outline.get("indices", {}).get("first_content_id")

                if caption_format:
                    params["action"] = "format_caption"
                elif caption_only:
                    params["action"] = "reposition_caption"
                elif action in ("caption", "add caption"):
                    params["action"] = "add_caption"
                elif not action and "caption" in desc_lower:
                    params["action"] = "add_caption"
                elif action == "move" and (params.get("before_id") or params.get("after_id")):
                    params["action"] = "move"
                elif action in ("move", "align", "alignment"):
                    params["action"] = "reposition"
                elif not action and any(word in desc_lower for word in ("align", "center", "centre", "left", "right", "move")):
                    params["action"] = "reposition"

                title_side_placement = (
                    "title" in desc_lower
                    and ("right" in desc_lower or "corner" in desc_lower or "side" in desc_lower)
                )
                page_top_placement = (
                    "top" in desc_lower
                    and ("right" in desc_lower or "corner" in desc_lower)
                    and ("page" in desc_lower or "top right" in desc_lower)
                )

                if title_side_placement:
                    anchor_id = params.get("anchor_id") or _title_anchor_id()
                    if anchor_id:
                        params["anchor_id"] = anchor_id
                        params["alignment"] = params.get("alignment") or "right"
                        if attached_image_path and params.get("action") in (None, "", "insert"):
                            params["action"] = "insert_into_paragraph"
                            op["target_id"] = anchor_id
                        else:
                            params["action"] = "place_inline"

                if page_top_placement and not title_side_placement:
                    params["alignment"] = params.get("alignment") or "right"
                    params["placement"] = params.get("placement") or "top_right_page"
                    if params.get("action") == "insert":
                        first_content_id = outline.get("indices", {}).get("first_content_id")
                        if first_content_id and not params.get("before_id") and not params.get("after_id"):
                            params["before_id"] = first_content_id
                    elif params.get("action") in (None, "", "reposition"):
                        params["action"] = "move"

                if params.get("action") == "add_caption" and not params.get("caption_text"):
                    desc = task.get("description", "")
                    quote_match = _re.search(r'["“](.+?)["”]', desc)
                    saying_match = _re.search(r"\bsaying\s+(.+)$", desc, flags=_re.IGNORECASE)
                    if quote_match:
                        params["caption_text"] = quote_match.group(1).strip()
                    elif saying_match:
                        params["caption_text"] = saying_match.group(1).strip().strip("'\"")

                if params.get("action") == "format_caption":
                    if params.get("font_size_pt") is None:
                        size_match = _re.search(r"(?:font\s+size|size)\s*(?:to|=|:)?\s*(\d+(?:\.\d+)?)", desc_lower)
                        if not size_match:
                            size_match = _re.search(r"\b(\d+(?:\.\d+)?)\s*(?:pt|points?)\b", desc_lower)
                        if size_match:
                            params["font_size_pt"] = float(size_match.group(1))
                    if params.get("bold") is None and "bold" in desc_lower:
                        params["bold"] = True
                    if params.get("italic") is None and "italic" in desc_lower:
                        params["italic"] = True
                    if params.get("underline") is None and "underline" in desc_lower:
                        params["underline"] = True

                position = params.get("position") if isinstance(params.get("position"), dict) else {}
                if params.get("width_page_pct") is None and position.get("width_pct") is not None:
                    params["width_page_pct"] = position.get("width_pct")

                if params.get("width_page_pct") is None:
                    import re as _re
                    pct_match = _re.search(r"(\d+(?:\.\d+)?)\s*%", desc_lower)
                    if pct_match:
                        params["width_page_pct"] = float(pct_match.group(1)) / 100.0

                if not params.get("alignment"):
                    if "center" in desc_lower or "centre" in desc_lower:
                        params["alignment"] = "center"
                    elif "right" in desc_lower:
                        params["alignment"] = "right"
                    elif "left" in desc_lower:
                        params["alignment"] = "left"

                if not params.get("float_position") and "right side" in desc_lower:
                    params["float_position"] = "right"

                if params.get("action") == "insert" and not params.get("after_id") and not params.get("before_id"):
                    if after_anchor_id:
                        params["after_id"] = after_anchor_id
                    elif ids_list:
                        params["after_id"] = ids_list[0]
                    if not op.get("target_id") and params.get("after_id"):
                        op["target_id"] = params["after_id"]

            # Normalize op_type for layout operations or content_generation
            raw_op_type = op.get("op_type")
            if task_type == "content_generation" or raw_op_type in ("insert_block", "move_block", "insert_toc", "insert_page_break", "duplicate_block"):
                if raw_op_type != "layout_op":
                    params = op.setdefault("parameters", {})
                    if "action" not in params and raw_op_type and raw_op_type != "content_generation":
                        params["action"] = raw_op_type
                    elif "action" not in params:
                        params["action"] = "insert_block"
                    op["op_type"] = "layout_op"

            if (
                existing_rewrite_intent
                and existing_section_target
                and op.get("op_type") == "layout_op"
                and op.get("parameters", {}).get("action") == "insert_block"
            ):
                op["_skip_existing_section_insert"] = True

        # Validate
        validated = []
        for op in ops_raw:
            try:
                if isinstance(op, dict) and op.get("_skip_broad_important_highlight"):
                    continue
                if isinstance(op, dict) and op.get("_skip_existing_section_insert"):
                    continue
                validated.append(validate_operation(op))
            except ValueError as e:
                log.warning("Task op validation failed: %s — %s", op, e)

        return validated

    # ------------------------------------------------------------------
    # Legacy generate (retained for backward compatibility)
    # ------------------------------------------------------------------

    def generate(
        self,
        request: str,
        structure: dict,
        document_type: str,
        chat_history: list[dict],
        intent: dict,
        attached_image_path: str | None = None,
        previous_ops: list[dict] | None = None,
        reviewer_feedback: str | None = None,
        missed_tasks: list[str] | None = None,
    ) -> list[dict]:
        """Legacy single-shot operation generator."""
        if settings.gemini_api_key or settings.openai_api_key:
            try:
                return self._generate_with_llm(
                    request, structure, document_type, chat_history,
                    intent, attached_image_path, previous_ops, reviewer_feedback,
                    missed_tasks,
                )
            except Exception as exc:
                log.exception("LLM legacy operations generation failed: %s", exc)
        
        return self._generate_fallback(request, intent, attached_image_path)

    def _generate_with_llm(
        self,
        request: str,
        structure: dict,
        document_type: str,
        chat_history: list[dict],
        intent: dict,
        attached_image_path: str | None,
        previous_ops: list[dict] | None,
        reviewer_feedback: str | None,
        missed_tasks: list[str] | None = None,
    ) -> list[dict]:
        llm = self._llm or LLMClient()
        sys_prompt = _SYSTEM_PROMPT.replace("{CURRENT_DATE}", datetime.now().strftime("%B %d, %Y"))
        
        outline = OutlineBuilder.build(structure, document_type)
        structure_summary = json.dumps(outline, indent=2)
        
        history_str = ""
        if chat_history:
            history_str = "Recent conversation:\n"
            for msg in chat_history[-4:]:
                role = "User" if msg["role"] == "user" else "Agent"
                history_str += f"{role}: {msg['content'][:200]}\n"
            history_str += "\n"

        image_str = ""
        if attached_image_path:
            image_str = f"Attached image path: {attached_image_path}\n"

        refinement_str = ""
        if reviewer_feedback and previous_ops:
            missed_str = ""
            if missed_tasks:
                missed_str = f"\nMissing tasks: {missed_tasks}\n"
            refinement_str = f"\nPrevious attempt: {json.dumps(previous_ops)}\nFeedback: {reviewer_feedback}\n{missed_str}"

        user_prompt = (
            f"{history_str}"
            f"{image_str}"
            f"Document outline:\n{structure_summary}\n\n"
            f"{refinement_str}"
            f"User instruction: {request}"
        )

        import logging
        log = logging.getLogger(__name__)

        response = llm.complete(LLMRequest(
            system_prompt=sys_prompt,
            user_prompt=user_prompt,
            temperature=0,
            max_tokens=4096,
            json_mode=True,
        ))

        log.warning(f"OPERATION GENERATOR RAW RESPONSE: {response.text}")
        print(f"OPERATION GENERATOR RAW RESPONSE: {response.text}")

        parsed = response.json or {}
        if isinstance(parsed, list):
            ops_raw = parsed
        elif isinstance(parsed, dict):
            ops_raw = parsed.get("operations") or parsed.get("ops") or [parsed]
        else:
            ops_raw = []

        validated = []
        for op in ops_raw:
            try:
                validated.append(validate_operation(op))
            except Exception as e:
                log.warning("Legacy op validation failed: %s", e)

        return validated

    def _generate_fallback(
        self,
        request: str,
        intent: dict,
        attached_image_path: str | None,
    ) -> list[dict]:
        category = intent.get("op_category", "")
        if category == "image_op" and attached_image_path:
            slide = intent.get("slide", 1) or 1
            return [validate_operation({
                "op_type": "image_op",
                "target": {"slide": slide, "shape_index": None},
                "parameters": {
                    "action": "insert",
                    "image_path": attached_image_path,
                    "position": {"left_pct": 0.1, "top_pct": 0.2, "width_pct": 0.4, "height_pct": 0.5},
                    "maintain_aspect_ratio": True,
                },
            })]

        if category == "slide_op":
            lowered = request.lower()
            if "add" in lowered or "new" in lowered:
                return [validate_operation({
                    "op_type": "slide_op",
                    "target": {"slide": intent.get("slide", 1) or 1},
                    "parameters": {"action": "add", "after_index": intent.get("slide", 1) or 1},
                })]
            if "delete" in lowered or "remove" in lowered:
                return [validate_operation({
                    "op_type": "slide_op",
                    "target": {"slide": intent.get("slide", 1) or 1},
                    "parameters": {"action": "delete"},
                })]

        return [validate_operation({
            "op_type": "ai_design_op",
            "target": {},
            "parameters": {
                "action": "improve_readability",
                "scope": "all_slides",
            },
        })]

# ===== END services/operation_generator.py =====

# ===== BEGIN services/editor.py =====
"""Content editor — rewrites document text blocks to fulfil a user request.

Uses LLMClient for provider-agnostic LLM calls (Gemini by default).
Falls back to deterministic heuristics when no LLM is available.
"""

import re


class ContentEditor:
    """Rewrites document text blocks to fulfil a user request.

    When an LLM is available the rewrite is delegated to it.
    Falls back to deterministic heuristics so the platform remains
    usable without credentials.
    """

    def __init__(self, llm=None) -> None:
        # Accept optional injected LLMClient; constructed lazily if not provided
        self._llm = llm

    def rewrite(
        self,
        request: str,
        text: str,
        metadata: dict | None = None,
        chat_history: list[dict] | None = None,
    ) -> str:
        if settings.gemini_api_key or settings.openai_api_key:
            try:
                return self._rewrite_with_llm(request, text, metadata or {}, chat_history or [])
            except Exception:
                None
        return self._rewrite_local(request, text)

    # ------------------------------------------------------------------
    # LLM path
    # ------------------------------------------------------------------

    def _rewrite_with_llm(
        self,
        request: str,
        text: str,
        metadata: dict,
        chat_history: list[dict],
    ) -> str:
        import json

        llm = self._llm or LLMClient()

        system_prompt = (
            "You are a precise document editing assistant evaluating ONE specific text block at a time.\n"
            "You will be given an editing instruction, the text block itself, and its metadata.\n\n"
            "IMPORTANT RULES:\n"
            "1. You must decide if the provided text block is the INTENDED TARGET of the instruction.\n"
            "2. If the user asks to edit a 'title', 'heading', or 'topic', you MUST ASSUME the provided "
            "text block IS the target and rewrite it, UNLESS the text block is obviously a footer, "
            "slide number (like `‹#›` or a plain digit), or a specific field label (like `Theme Name:`).\n"
            "3. If you decide the text block IS the target, return ONLY the new rewritten text.\n"
            "4. If you decide the text block is NOT the target, you MUST return the ORIGINAL TEXT EXACTLY AS-IS. "
            "Do not return any other text.\n"
            "5. Provide no commentary, markdown, or quotes.\n"
            "6. For formatting-related instructions (e.g. 'make it bold', 'center align', 'change font size'), "
            "you cannot change formatting directly — only change the TEXT content. Return the text as-is if the "
            "instruction is purely about formatting.\n"
            "7. When writing content, ensure it is professional, well-structured, and compelling. "
            "Use bullet points (•) for lists when appropriate."
        )

        meta_str = json.dumps(metadata) if metadata else "None"

        history_str = ""
        if chat_history:
            history_str = "Previous conversation context:\n"
            for msg in chat_history[-5:]:
                role = "User" if msg["role"] == "user" else "Agent"
                history_str += f"{role}: {msg['content']}\n"
            history_str += "\n"

        user_prompt = (
            f"{history_str}"
            f"Editing instruction: {request}\n\n"
            f"Text block to consider:\n{text}\n\n"
            f"Block metadata: {meta_str}"
        )

        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.3,
            max_tokens=1024,
            json_mode=False,  # Raw text response
        ))
        return response.text or text

    # ------------------------------------------------------------------
    # Local fallback heuristics
    # ------------------------------------------------------------------

    def _rewrite_local(self, request: str, text: str) -> str:
        lowered = request.lower()
        if "remove" in lowered:
            return self._remove_mentions(request, text)
        if "shorter" in lowered or "summarize" in lowered or "shorten" in lowered:
            return self._shorten(text)
        if "professional" in lowered or "professionally" in lowered:
            return self._professionalize(text)
        if "expand" in lowered:
            return self._expand(text)
        return self._professionalize(text)

    def _shorten(self, text: str) -> str:
        sentences = re.split(r"(?<=[.!?])\s+", text.strip())
        if len(sentences) > 1:
            return " ".join(sentences[: max(1, len(sentences) // 2)])
        words = text.split()
        return " ".join(words[: max(8, int(len(words) * 0.6))])

    def _professionalize(self, text: str) -> str:
        cleaned = re.sub(r"\s+", " ", text).strip()
        replacements = {
            "we got": "we achieved",
            "big": "significant",
            "lots of": "substantial",
            "things": "initiatives",
        }
        for source, target in replacements.items():
            cleaned = re.sub(source, target, cleaned, flags=re.IGNORECASE)
        return cleaned[:1].upper() + cleaned[1:] if cleaned else cleaned

    def _expand(self, text: str) -> str:
        cleaned = self._professionalize(text)
        return f"{cleaned} This reflects a focused, measurable improvement aligned with the broader strategy."

    def _remove_mentions(self, request: str, text: str) -> str:
        match = re.search(r"remove all mentions of\s+(.+?)(?:[.!?]|$)", request, flags=re.IGNORECASE)
        phrase = match.group(1).strip().strip('"') if match else ""
        if not phrase:
            return text
        updated = re.sub(re.escape(phrase), "", text, flags=re.IGNORECASE)
        updated = re.sub(r"\s{2,}", " ", updated)
        updated = re.sub(r"\s+([,.!?])", r"\1", updated)
        return updated.strip()

# ===== END services/editor.py =====

# ===== BEGIN services/verifier.py =====
"""Verifier — validates document operations structurally and semantically.

Replaces the legacy reviewer.py.
Compares document outlines before and after execution to generate a semantic
diff, then asks Gemini to verify if all planned tasks were satisfied.
"""

import json
import logging
from typing import Any


log = logging.getLogger(__name__)


class Verifier:
    """Performs structural and semantic post-execution verification."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    # ------------------------------------------------------------------
    # Stage-focused verify_semantic (Phase 5)
    # ------------------------------------------------------------------

    def verify_semantic(
        self,
        request: str,
        tasks: list[dict],
        before_outline: dict,
        after_outline: dict,
    ) -> dict:
        """Compute structural/textual diff and ask LLM if tasks were satisfied."""
        llm = self._llm or LLMClient()

        # Compute the outline diff
        diff = self._compute_outline_diff(before_outline, after_outline)

        system_prompt = (
            "You are a strict document editing verifier.\n"
            "Your job is to look at a user's original request, the list of planned tasks, "
            "and a structural/textual diff of what actually changed in the document.\n"
            "Determine if each planned task has been successfully satisfied.\n\n"
            "RULES:\n"
            "1. Be strict for structural changes. If a task says 'move section X before Y' but the order of headings in the outline "
            "shows X is still after Y, mark it as satisfied=false.\n"
            "2. Ensure you check the 'style_summary' field in the outline diff to verify formatting changes like spacing, font sizes, colors, and page breaks.\n"
            "3. If a task asks for complex aesthetic changes (like 'modernize tables' or 'make beautiful'), verify that at least some related style properties (like table_style, colors, or spacing) were updated in the diff.\n"
            "4. If all tasks are successfully satisfied, set 'all_satisfied' to true.\n"
            "5. Provide constructive feedback for any unsatisfied tasks outlining what is wrong.\n\n"
            "Return JSON matching this schema:\n"
            "{\n"
            '  "tasks": [\n'
            '    {\n'
            '      "index": 0,\n'
            '      "description": "...",\n'
            '      "satisfied": true/false,\n'
            '      "feedback": "Why it failed, or empty string"\n'
            "    }\n"
            "  ],\n"
            '  "all_satisfied": true/false\n'
            "}"
        )

        user_prompt = (
            f"Original Request: {request}\n\n"
            f"Decomposed Tasks:\n{json.dumps(tasks, indent=2)}\n\n"
            f"Outline Diff of Document Changes:\n{json.dumps(diff, indent=2)}"
        )

        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0,
            max_tokens=1024,
            json_mode=True,
        ))

        parsed = response.json or {}
        
        # Parse output
        validated_tasks = []
        all_satisfied = parsed.get("all_satisfied", True)
        
        diff_has_changes = bool(
            diff.get("added_elements") or
            diff.get("modified_elements") or
            diff.get("deleted_elements")
        )

        for idx, task in enumerate(tasks):
            match = next((t for t in parsed.get("tasks", []) if t.get("index") == idx), None)
            
            if task.get("insufficient_kb_evidence"):
                sat = True
                fb = f"Insufficient Knowledge Base evidence found for task '{task.get('description')}'. Task exited early with user notification."
            elif not diff_has_changes:
                # Deterministic Guard: Zero document changes occurred, task cannot be satisfied
                sat = False
                fb = f"No document operations or changes were applied for task '{task.get('description')}'."
                all_satisfied = False
            elif match:
                sat = bool(match.get("satisfied", False))
                fb = str(match.get("feedback", ""))
                if not sat:
                    all_satisfied = False
            else:
                sat = True
                fb = ""

            validated_tasks.append({
                "index": idx,
                "description": task.get("description", ""),
                "satisfied": sat,
                "feedback": fb,
            })

        return {
            "all_satisfied": all_satisfied,
            "tasks": validated_tasks,
            "diff": diff,
        }

    def review_plan(self, request: str, slide_plan: dict, intent: dict) -> dict:
        """Review the initial presentation structure generated by SlidePlanner."""
        llm = self._llm or LLMClient()
        
        system_prompt = (
            "You are a presentation plan reviewer.\n"
            "Review the slide outline against the user's request.\n"
            "If the plan fulfills the core request, return {'satisfied': true}.\n"
            "If it is missing crucial topics or has severe logical flaws, return {'satisfied': false, 'feedback': 'explain the flaws'}.\n"
            "Do NOT reject purely for stylistic reasons or minor omissions. Only reject if the core intent is missed.\n"
        )
        
        user_prompt = (
            f"User request: {request}\n"
            f"Intent: {json.dumps(intent, indent=2)}\n"
            f"Generated Plan: {json.dumps(slide_plan, indent=2)}"
        )
        
        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0,
            max_tokens=512,
            json_mode=True,
        ))
        
        parsed = response.json or {}
        return {
            "satisfied": parsed.get("satisfied", True),
            "feedback": parsed.get("feedback", "")
        }

    def _compute_outline_diff(self, before: dict, after: dict) -> dict:
        """Helper to compute a simplified diff of before and after outlines."""
        diff: dict[str, list] = {
            "heading_order_before": [],
            "heading_order_after": [],
            "modified_elements": [],
            "added_elements": [],
            "deleted_elements": [],
        }

        # Heading orders
        before_headings = [s.get("heading", "") for s in before.get("sections", []) if s.get("heading_id") != "start"]
        after_headings = [s.get("heading", "") for s in after.get("sections", []) if s.get("heading_id") != "start"]
        diff["heading_order_before"] = before_headings
        diff["heading_order_after"] = after_headings

        # Map elements by ID
        before_map = {}
        for s in before.get("sections", []):
            for el in s.get("elements", []):
                before_map[el["id"]] = el
        
        after_map = {}
        for s in after.get("sections", []):
            for el in s.get("elements", []):
                after_map[el["id"]] = el

        # Detect additions and modifications
        for el_id, after_el in after_map.items():
            if el_id not in before_map:
                diff["added_elements"].append({
                    "id": el_id,
                    "type": after_el.get("type"),
                    "preview": after_el.get("text_preview", ""),
                })
            else:
                before_el = before_map[el_id]
                # Compare text preview or other values
                changed = []
                if before_el.get("text_preview") != after_el.get("text_preview"):
                    changed.append("text")
                if before_el.get("rows") != after_el.get("rows") or before_el.get("cols") != after_el.get("cols"):
                    changed.append("table_dimensions")
                if before_el.get("size") != after_el.get("size"):
                    changed.append("image_size")
                if before_el.get("style_summary") != after_el.get("style_summary"):
                    changed.append("style")
                if before_el.get("alignment") != after_el.get("alignment"):
                    changed.append("alignment")
                
                if changed:
                    diff_entry = {
                        "id": el_id,
                        "type": after_el.get("type"),
                        "changes": changed,
                        "before_preview": before_el.get("text_preview", ""),
                        "after_preview": after_el.get("text_preview", ""),
                    }
                    if "style" in changed:
                        diff_entry["before_style"] = before_el.get("style_summary", [])
                        diff_entry["after_style"] = after_el.get("style_summary", [])
                    if "alignment" in changed:
                        diff_entry["before_alignment"] = before_el.get("alignment")
                        diff_entry["after_alignment"] = after_el.get("alignment")
                        
                    diff["modified_elements"].append(diff_entry)

        # Detect deletions
        for el_id, before_el in before_map.items():
            if el_id not in after_map:
                diff["deleted_elements"].append({
                    "id": el_id,
                    "type": before_el.get("type"),
                    "preview": before_el.get("text_preview", ""),
                })

        return diff

# ===== END services/verifier.py =====

# ===== BEGIN services/slide_planner.py =====
"""Slide Planner — generates a structured multi-slide content plan via LLM.

Given a template's rich structure and a user request, the planner produces
a complete slide plan in a single LLM call. The plan specifies:
- Which template slides to use as layout sources.
- Whether to populate, keep, or delete each slide.
- Text content for every text frame, with optional formatting overrides.

Uses LLMClient for provider-agnostic LLM calls (Gemini by default).
"""

import concurrent.futures
import json
import logging


log = logging.getLogger(__name__)

MAX_SLIDES = 20  # Hard cap on generated slides


class SlidePlanner:
    """Generates a structured slide plan from a template + user request."""

    def __init__(self, llm=None) -> None:
        self._llm = llm

    def plan(
        self,
        request: str,
        template_structure: dict,
        intent: dict,
        chat_history: list[dict] | None = None,
    ) -> dict:
        """Produce a slide plan dict."""
        if not (settings.gemini_api_key or settings.openai_api_key):
            return self._plan_local(request, template_structure, intent)

        # Generate outline
        outline = self._generate_outline(
            request=request,
            template_structure=template_structure,
            intent=intent,
            chat_history=chat_history or [],
        )

        if not outline:
            return self._plan_local(request, template_structure, intent)

        # Process each slide's content (populate, delete, keep)
        def process_slide(slide_outline: dict) -> dict:
            action = slide_outline.get("action", "populate")
            src_idx = slide_outline.get("source_slide_index", 1)

            if action == "delete":
                return {
                    "source_slide_index": src_idx,
                    "action": "delete",
                    "shapes": [],
                }
            if action == "keep":
                return {
                    "source_slide_index": src_idx,
                    "action": "keep",
                    "shapes": [],
                }

            return self._generate_slide_content(
                request=request,
                template_structure=template_structure,
                slide_outline=slide_outline,
                intent=intent
            )

        # Run concurrent generation (max 10 workers for speed)
        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            results = list(executor.map(process_slide, outline))

        # Build final plan
        plan = {"slides": results}
        
        # Validate and sanitise
        return self._sanitise_plan(plan, template_structure)

    def _generate_outline(
        self, request: str, template_structure: dict, intent: dict, chat_history: list[dict]
    ) -> list[dict]:

        llm = self._llm or LLMClient()

        template_desc = self._describe_template(template_structure)
        topic = intent.get("topic", "")

        slide_count = intent.get("slide_count")
        delete_slides = intent.get("delete_slides", [])
        add_count = intent.get("add_slides_count")
        existing_slide_count = template_structure.get("slide_count", 1)

        count_guidance = ""
        if delete_slides:
            count_guidance = f"The user wants to DELETE slides: {delete_slides}. Remove them from the output."
        elif add_count:
            count_guidance = f"Keep existing {existing_slide_count} slides and add {add_count} new ones."
        elif slide_count:
            count_guidance = f"Target exactly {min(slide_count, MAX_SLIDES)} slides."
        else:
            count_guidance = f"Template has {existing_slide_count} slides. Use 5–12 slides depending on topic depth."

        history_str = ""
        if chat_history:
            history_str = "Previous conversation:\n"
            for msg in chat_history[-5:]:
                role = "User" if msg["role"] == "user" else "Agent"
                history_str += f"{role}: {msg['content']}\n"
            history_str += "\n"

        system_prompt = (
            "You are an expert presentation outliner. Your job is to define the structure of a presentation.\n\n"
            "Given a user request, a topic, and a list of available template slide layouts, generate a slide-by-slide outline.\n\n"
            "RULES:\n"
            "1. Choose the best template layout for each slide (source_slide_index).\n"
            "2. Define the exact action ('populate', 'keep', 'delete').\n"
            "3. For 'populate' slides, write a detailed 'outline' containing the main talking points, bullet points, and data to be included on this specific slide.\n\n"
            "OUTPUT FORMAT: Return a JSON object with a key 'outline' containing an array of slides:\n"
            "Each slide: { source_slide_index: int, action: string, outline: string (detailed bullet points of what goes on the slide) }\n\n"
            "Return ONLY valid JSON."
        )

        user_prompt = (
            f"{history_str}User request: {request}\nTopic: {topic}\nSlide count guidance: {count_guidance}\n\n"
            f"Available Layouts:\n{template_desc}"
        )

        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.4,
            max_tokens=4096,
            json_mode=True,
        ))
        
        parsed = response.json or {}
        return parsed.get("outline", [])

    def _generate_slide_content(
        self, request: str, template_structure: dict, slide_outline: dict, intent: dict
    ) -> dict:

        llm = self._llm or LLMClient()

        src_idx = slide_outline.get("source_slide_index", 1)
        action = slide_outline.get("action", "populate")
        outline_text = slide_outline.get("outline", "")
        topic = intent.get("topic", "")

        slides = template_structure.get("slides", [])
        slide_info = next((s for s in slides if s["slide_index"] == src_idx), slides[0] if slides else {})
        layout_desc = self._describe_template({"slides": [slide_info]})

        system_prompt = (
            "You are an expert copywriter and presentation designer. Your task is to write deep, substantive content "
            "for a SINGLE slide, matching its exact layout structure.\n\n"
            "CRITICAL DESIGN RULES:\n"
            "1. NO WALLS OF TEXT. Never write long paragraphs. For body text placeholders, you MUST use short, scannable bullet points.\n"
            "2. PREFIX BULLETS. Prefix each bullet point with the bullet character '• '. Keep bullets under 2 sentences.\n"
            "3. PREVENT OVERLAP (FONT SIZE). If your generated text is longer than the original placeholder text, you MUST "
            "drastically decrease the `font_size_pt` (e.g., from 44 down to 24, or 18 to 12). If you do not reduce the font size, your text will overlap other shapes and ruin the slide!\n"
            "4. PREVENT WEIRD SPACING (ALIGNMENT). Template titles often use 'distributed' alignment which puts massive spaces between letters of long words. "
            "You MUST override the `alignment` property to `'left'` or `'center'` for any title or heading you modify.\n"
            "5. To populate a table, include a 'table_rows' array where each element is an array of cells, and each cell has a 'paragraphs' array.\n\n"
            "FORMATTING OVERRIDES (null = keep template default):\n"
            "- font_size_pt: number or null\n"
            "- bold: true/false or null\n"
            "- italic: true/false or null\n"
            "- color_hex: 6-char hex string (e.g. '2B579A') or null\n"
            "- alignment: 'left'|'center'|'right'|'justify' or null\n\n"
            "OUTPUT FORMAT: Return a JSON object with a key 'shapes' containing an array:\n"
            "[{shape_index, paragraphs: [{para_index, text, formatting}], table_rows: [[{paragraphs: ...}]]}]\n\n"
            "Return ONLY valid JSON."
        )

        user_prompt = (
            f"Topic: {topic}\nSlide Outline to write content for:\n{outline_text}\n\n"
            f"Template Layout Shapes for this slide:\n{layout_desc}"
        )

        response = llm.complete(LLMRequest(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.5,
            max_tokens=4096,
            json_mode=True,
        ))
        
        parsed = response.json or {}
        shapes = parsed.get("shapes", [])
        
        return {
            "source_slide_index": src_idx,
            "action": action,
            "shapes": shapes
        }

    # ------------------------------------------------------------------
    # Local fallback (no LLM)
    # ------------------------------------------------------------------

    def _plan_local(self, request: str, template_structure: dict, intent: dict) -> dict:
        """Simple rule-based plan: keep all slides, populate with basic text."""
        slides = template_structure.get("slides", [])
        delete_slides = set(intent.get("delete_slides", []))
        topic = intent.get("topic", "") or request

        plan_slides: list[dict] = []
        for slide_data in slides:
            slide_idx = slide_data["slide_index"]

            if slide_idx in delete_slides:
                plan_slides.append({
                    "source_slide_index": slide_idx,
                    "action": "delete",
                    "shapes": [],
                })
                continue

            shapes: list[dict] = []
            for shape in slide_data.get("shapes", []):
                if not shape.get("has_text_frame"):
                    continue
                paras = []
                for para in shape.get("paragraphs", []):
                    text = para.get("text", "")
                    if not text:
                        text = topic[:80]
                    paras.append({
                        "para_index": para["para_index"],
                        "text": text,
                        "formatting": {},
                    })
                if paras:
                    shapes.append({
                        "shape_index": shape["shape_index"],
                        "paragraphs": paras,
                    })

            plan_slides.append({
                "source_slide_index": slide_idx,
                "action": "populate",
                "shapes": shapes,
            })

        return {"slides": plan_slides}

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _describe_template(self, structure: dict) -> str:
        """Create a concise text description of the template for the LLM prompt."""
        lines: list[str] = []
        slides = structure.get("slides", [])
        lines.append(f"Total template slides: {len(slides)}")

        for slide in slides:
            slide_idx = slide["slide_index"]
            layout = slide.get("layout_name", "unknown")
            lines.append(f"\n--- Slide {slide_idx} (layout: {layout}) ---")

            for shape in slide.get("shapes", []):
                has_tf = shape.get("has_text_frame")
                has_t = shape.get("has_table")
                if not (has_tf or has_t):
                    continue
                shape_idx = shape["shape_index"]
                name = shape.get("shape_name", "?")
                is_ph = shape.get("is_placeholder", False)
                ph_type = shape.get("placeholder_type", "")

                ph_str = f", placeholder: {ph_type}" if is_ph else ""
                lines.append(f"  Shape {shape_idx} ({name}{ph_str}):")

                if has_tf:
                    for para in shape.get("paragraphs", []):
                        pi = para["para_index"]
                        text = para.get("text", "")
                        fmt = para.get("formatting", {})
                        fmt_parts = []
                        if fmt.get("font_size_pt"):
                            fmt_parts.append(f"{fmt['font_size_pt']}pt")
                        if fmt.get("bold"):
                            fmt_parts.append("bold")
                        if fmt.get("alignment"):
                            fmt_parts.append(fmt["alignment"])
                        fmt_str = f" [{', '.join(fmt_parts)}]" if fmt_parts else ""
                        display_text = text if text else "(empty)"
                        if len(display_text) > 80:
                            display_text = display_text[:77] + "..."
                        lines.append(f"    Para {pi}: \"{display_text}\"{fmt_str}")
                
                if has_t:
                    table_rows = shape.get("table_rows", [])
                    lines.append(f"    [TABLE with {len(table_rows)} rows]")
                    for r_idx, row in enumerate(table_rows):
                        row_texts = []
                        for cell in row:
                            cell_text = " ".join(p.get("text", "") for p in cell.get("paragraphs", []))
                            if len(cell_text) > 20:
                                  cell_text = cell_text[:17] + "..."
                            row_texts.append(cell_text if cell_text else "(empty)")
                        lines.append(f"      Row {r_idx}: {row_texts}")

        return "\n".join(lines)

    def _sanitise_plan(self, plan: dict, template_structure: dict) -> dict:
        """Validate and sanitise an LLM-generated plan."""
        slides = plan.get("slides", [])
        template_slide_count = template_structure.get("slide_count", 1)

        sanitised: list[dict] = []
        for slide in slides[:MAX_SLIDES]:
            src = slide.get("source_slide_index", 1)
            if not isinstance(src, int) or src < 1 or src > template_slide_count:
                src = 1

            action = slide.get("action", "populate")
            if action not in ("populate", "keep", "delete"):
                action = "populate"

            shapes: list[dict] = []
            for shape in slide.get("shapes", []):
                si = shape.get("shape_index", 0)
                if not isinstance(si, int) or si < 0:
                    continue
                paras: list[dict] = []
                for para in (shape.get("paragraphs") or []):
                    pi = para.get("para_index", 0)
                    text = str(para.get("text", ""))
                    fmt = para.get("formatting") or {}
                    if not isinstance(fmt, dict):
                        fmt = {}
                    paras.append({
                        "para_index": pi,
                        "text": text,
                        "formatting": {
                            "font_size_pt": fmt.get("font_size_pt"),
                            "bold": fmt.get("bold"),
                            "italic": fmt.get("italic"),
                            "color_hex": fmt.get("color_hex"),
                            "alignment": fmt.get("alignment"),
                        },
                    })
                table_rows = []
                for row in (shape.get("table_rows") or []):
                    if not isinstance(row, list):
                        continue
                    clean_row = []
                    for cell in row:
                        if not isinstance(cell, dict):
                            continue
                        cell_paras = []
                        for para in (cell.get("paragraphs") or []):
                            pi = para.get("para_index", 0)
                            text = str(para.get("text", ""))
                            fmt = para.get("formatting") or {}
                            if not isinstance(fmt, dict):
                                fmt = {}
                            cell_paras.append({
                                "para_index": pi,
                                "text": text,
                                "formatting": {
                                    "font_size_pt": fmt.get("font_size_pt"),
                                    "bold": fmt.get("bold"),
                                    "italic": fmt.get("italic"),
                                    "color_hex": fmt.get("color_hex"),
                                    "alignment": fmt.get("alignment"),
                                },
                            })
                        clean_row.append({"paragraphs": cell_paras})
                    table_rows.append(clean_row)

                clean_shape = {
                    "shape_index": si,
                    "paragraphs": paras,
                }
                if table_rows:
                    clean_shape["table_rows"] = table_rows
                
                shapes.append(clean_shape)

            sanitised.append({
                "source_slide_index": src,
                "action": action,
                "shapes": shapes,
            })

        return {"slides": sanitised}

# ===== END services/slide_planner.py =====

# ===== BEGIN services/template_analyzer.py =====
"""Template Analyzer — deep extraction of DOCX template structure and styling.

Extracts:
  1. Structural skeleton: ordered sections with heading levels and names
  2. Style catalog: fonts, sizes, colors, spacing for every heading/body style
  3. Table styles: border, header formatting, alternating colors
  4. Page layout: margins, orientation, page size
  5. Numbering conventions: table/figure prefixes, list styles
  6. Header/footer text
"""

import logging
import re
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)


def _hex_from_color(color_obj) -> str | None:
    """Safely extract hex color from a python-docx color object."""
    try:
        if color_obj and color_obj.rgb:
            return str(color_obj.rgb).upper()
    except Exception:
        None
    return None


def _pt_from_emu(emu_val) -> float | None:
    """Convert EMU to points (1 pt = 12700 EMU)."""
    if emu_val is None:
        return None
    try:
        return round(emu_val / 12700, 1)
    except Exception:
        return None


def _pt_from_twips(twips) -> float | None:
    """Convert twips to points (1 pt = 20 twips)."""
    if twips is None:
        return None
    try:
        return round(twips / 20, 1)
    except Exception:
        return None


class TemplateAnalyzer:
    """Deep analysis of a DOCX template for structure and style extraction."""

    def analyze(self, docx_path: Path) -> dict:
        """Full template analysis.

        Returns a dict with:
          sections, style_catalog, table_style, page_layout, numbering,
          headers_footers, has_cover_page
        """
        try:
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(str(docx_path))
        except Exception as exc:
            log.error("TemplateAnalyzer: could not open %s: %s", docx_path, exc)
            return self._empty_analysis()

        sections = self._extract_sections(doc)
        style_catalog = self._extract_style_catalog(doc)
        table_style = self._extract_table_style(doc)
        page_layout = self._extract_page_layout(doc)
        numbering = self._infer_numbering_conventions(doc)
        headers_footers = self._extract_headers_footers(doc)

        return {
            "sections": sections,
            "style_catalog": style_catalog,
            "table_style": table_style,
            "page_layout": page_layout,
            "numbering": numbering,
            "headers_footers": headers_footers,
            "has_cover_page": self._detect_cover_page(doc),
        }

    # ------------------------------------------------------------------
    # Section Structure
    # ------------------------------------------------------------------

    def _extract_sections(self, doc) -> list[dict]:
        """Extract ordered section headings with hierarchy."""
        sections: list[dict] = []
        current_h1: dict | None = None

        for para in doc.paragraphs:
            text = para.text.strip()
            style_name = para.style.name if para.style else ""

            if not style_name.startswith("Heading"):
                continue
            if not text:
                continue

            try:
                level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                level = 1

            section = {
                "heading_text": text,
                "heading_level": level,
                "style_name": style_name,
                "subsections": [],
            }

            if level == 1:
                sections.append(section)
                current_h1 = section
            elif level == 2 and current_h1:
                current_h1["subsections"].append(section)
            else:
                sections.append(section)

        return sections

    # ------------------------------------------------------------------
    # Style Catalog
    # ------------------------------------------------------------------

    def _extract_style_catalog(self, doc) -> dict:
        """Extract formatting for each named style in the document."""
        catalog: dict[str, dict] = {}

        for style in doc.styles:
            style_name = style.name
            # Only capture heading and body-relevant styles
            relevant = (
                style_name.startswith("Heading")
                or style_name in ("Normal", "Body Text", "Caption", "List Bullet",
                                  "List Number", "Table Grid", "Title", "Subtitle",
                                  "Intense Quote", "Quote", "No Spacing")
            )
            if not relevant:
                continue

            info: dict[str, Any] = {"style_name": style_name}

            try:
                font = style.font
                if font:
                    if font.name:
                        info["font_family"] = font.name
                    if font.size:
                        info["font_size_pt"] = round(font.size.pt, 1)
                    if font.bold is not None:
                        info["bold"] = font.bold
                    if font.italic is not None:
                        info["italic"] = font.italic
                    if font.underline is not None:
                        info["underline"] = font.underline
                    color = _hex_from_color(font.color)
                    if color:
                        info["color_hex"] = color

                pf = style.paragraph_format
                if pf:
                    if pf.alignment is not None:
                        info["alignment"] = str(pf.alignment).split(".")[-1].lower()
                    if pf.space_before is not None:
                        info["space_before_pt"] = _pt_from_emu(pf.space_before)
                    if pf.space_after is not None:
                        info["space_after_pt"] = _pt_from_emu(pf.space_after)
                    if pf.line_spacing is not None:
                        # line_spacing can be a Length (pt) or a float multiplier
                        ls = pf.line_spacing
                        try:
                            info["line_spacing"] = round(float(ls), 2)
                        except Exception:
                            None

            except Exception as exc:
                log.debug("Style extraction failed for %s: %s", style_name, exc)

            catalog[style_name] = info

        return catalog

    # ------------------------------------------------------------------
    # Table Style
    # ------------------------------------------------------------------

    def _extract_table_style(self, doc) -> dict:
        """Infer table formatting from the first table in the document."""
        table_info: dict[str, Any] = {
            "style_name": "Table Grid",
            "has_header_row": True,
            "header_bg_hex": None,
            "header_font_color_hex": None,
            "header_bold": True,
            "alt_row_colors": [],
            "border_color_hex": None,
            "cell_padding_pt": 3.0,
        }

        if not doc.tables:
            return table_info

        try:
            table = doc.tables[0]
            if table.style:
                table_info["style_name"] = table.style.name

            # Inspect header row (first row) formatting
            if table.rows:
                first_row = table.rows[0]
                for cell in first_row.cells:
                    # Try to get fill color from XML
                    try:
                        from docx.oxml.ns import qn
                        tc = cell._tc
                        tcPr = tc.find(qn("w:tcPr"))
                        if tcPr is not None:
                            shd = tcPr.find(qn("w:shd"))
                            if shd is not None:
                                fill = shd.get(qn("w:fill"))
                                if fill and fill != "auto" and len(fill) == 6:
                                    table_info["header_bg_hex"] = fill.upper()
                    except Exception:
                        None

                    # Font color
                    for para in cell.paragraphs:
                        for run in para.runs:
                            color = _hex_from_color(run.font.color)
                            if color:
                                table_info["header_font_color_hex"] = color
                            if run.font.bold:
                                table_info["header_bold"] = True
                    break

            # Inspect second row for alternating color
            if len(table.rows) > 1:
                second_row = table.rows[1]
                for cell in second_row.cells:
                    try:
                        from docx.oxml.ns import qn
                        tc = cell._tc
                        tcPr = tc.find(qn("w:tcPr"))
                        if tcPr is not None:
                            shd = tcPr.find(qn("w:shd"))
                            if shd is not None:
                                fill = shd.get(qn("w:fill"))
                                if fill and fill != "auto" and len(fill) == 6:
                                    table_info["alt_row_colors"].append(fill.upper())
                    except Exception:
                        None
                    break

        except Exception as exc:
            log.debug("Table style extraction failed: %s", exc)

        return table_info

    # ------------------------------------------------------------------
    # Page Layout
    # ------------------------------------------------------------------

    def _extract_page_layout(self, doc) -> dict:
        """Extract page dimensions and margins from the first section."""
        layout: dict[str, Any] = {
            "orientation": "portrait",
            "page_width_pt": 612.0,
            "page_height_pt": 792.0,
            "margins": {"top": 72.0, "bottom": 72.0, "left": 72.0, "right": 72.0},
        }

        try:
            if not doc.sections:
                return layout

            sect = doc.sections[0]
            w_pt = _pt_from_emu(sect.page_width)
            h_pt = _pt_from_emu(sect.page_height)
            if w_pt and h_pt:
                layout["page_width_pt"] = w_pt
                layout["page_height_pt"] = h_pt
                layout["orientation"] = "landscape" if w_pt > h_pt else "portrait"

            margins: dict[str, float | None] = {
                "top": _pt_from_emu(sect.top_margin),
                "bottom": _pt_from_emu(sect.bottom_margin),
                "left": _pt_from_emu(sect.left_margin),
                "right": _pt_from_emu(sect.right_margin),
            }
            layout["margins"] = {k: v for k, v in margins.items() if v is not None}

        except Exception as exc:
            log.debug("Page layout extraction failed: %s", exc)

        return layout

    # ------------------------------------------------------------------
    # Numbering / Conventions
    # ------------------------------------------------------------------

    def _infer_numbering_conventions(self, doc) -> dict:
        """Detect numbering prefixes for tables and figures from paragraph text."""
        conventions: dict[str, Any] = {
            "table_prefix": "Table",
            "figure_prefix": "Figure",
            "section_numbered": False,
        }

        table_pattern = re.compile(r"^(Table|Tabel|Tab\.?)\s*\d+", re.IGNORECASE)
        figure_pattern = re.compile(r"^(Figure|Fig\.?)\s*\d+", re.IGNORECASE)

        for para in doc.paragraphs:
            text = para.text.strip()
            if table_pattern.match(text):
                prefix = table_pattern.match(text).group(1)
                conventions["table_prefix"] = prefix.rstrip(".")
            if figure_pattern.match(text):
                prefix = figure_pattern.match(text).group(1)
                conventions["figure_prefix"] = prefix.rstrip(".")

        # Check if Heading 1 text starts with a number (e.g. "1. Introduction")
        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""
            if style_name == "Heading 1":
                text = para.text.strip()
                if re.match(r"^\d+[\.\)]\s", text):
                    conventions["section_numbered"] = True
                break

        return conventions

    # ------------------------------------------------------------------
    # Headers / Footers
    # ------------------------------------------------------------------

    def _extract_headers_footers(self, doc) -> dict:
        """Extract text from page headers and footers."""
        result: dict[str, str | None] = {"header": None, "footer": None}

        try:
            if doc.sections:
                sect = doc.sections[0]
                header = sect.header
                if header and not header.is_linked_to_previous:
                    header_text = " ".join(p.text for p in header.paragraphs if p.text.strip())
                    if header_text:
                        result["header"] = header_text.strip()

                footer = sect.footer
                if footer and not footer.is_linked_to_previous:
                    footer_text = " ".join(p.text for p in footer.paragraphs if p.text.strip())
                    if footer_text:
                        result["footer"] = footer_text.strip()

        except Exception as exc:
            log.debug("Header/footer extraction failed: %s", exc)

        return result

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _detect_cover_page(self, doc) -> bool:
        """Heuristic: first paragraph has a 'Title' style → likely a cover page."""
        for para in doc.paragraphs:
            if para.text.strip():
                style_name = para.style.name if para.style else ""
                return "title" in style_name.lower()
        return False

    def _empty_analysis(self) -> dict:
        return {
            "sections": [],
            "style_catalog": {},
            "table_style": {},
            "page_layout": {},
            "numbering": {"table_prefix": "Table", "figure_prefix": "Figure"},
            "headers_footers": {},
            "has_cover_page": False,
        }

# ===== END services/template_analyzer.py =====

# ===== BEGIN services/document_planner.py =====
"""Document Planner — plans the full document structure before generation.

Given:
  - User's generation request
  - Template analysis (section skeleton, styles)
  - KB context chunks (retrieved from the knowledge base)

The planner produces a section-by-section plan specifying:
  - Which headings to use (template-derived by default)
  - Content type per section (narrative, bullet list, table, toc, skip)
  - What KB chunks are relevant to each section
  - Table specifications (columns, auto-numbered caption)
  - Whether to include a table of contents
  - Any cross-references or special elements

Rules:
  1. Template headings are STRICTLY preserved by default.
  2. The LLM may add sub-sections or add a TOC/Conclusion only if
     the user explicitly requests it.
  3. If no KB data is relevant to a section, LLM uses general knowledge.
"""

import json
import logging
import re
from typing import Any


log = logging.getLogger(__name__)

_PLANNER_SYSTEM_PROMPT = """You are an expert document planning assistant for professional, audit-grade business documents.

Your job: Given a user's request, the document template's section structure, and retrieved knowledge base context, produce a comprehensive section-by-section document plan.

CRITICAL RULES:
1. TEMPLATE HEADINGS: You MUST preserve ALL headings defined in the template. Do not rename, skip, or reorder template headings unless the user explicitly says to.
2. NEW SECTIONS: You may ADD new sections (e.g., Table of Contents, Conclusion, additional sub-sections) ONLY when the user explicitly requests them.
3. GROUNDING: For each section, identify which KB chunks are most relevant by examining the chunk contents. Assign chunks generously — it's better to assign too many than too few. If no KB data is relevant, use "kb_chunk_ids": [] and instruct the generator to write a brief structural paragraph without specific claims.
4. TABLE NUMBERING: If a section contains a table, assign it a sequential number (Table 1, Table 2, ...) with a descriptive caption.
5. PROFESSIONAL FORMAT: Plan content that looks like a professional internal document: concise narrative paragraphs, well-structured tables with headers, numbered bullet points for key findings.
6. TOC: Include a "toc" section at the very beginning ONLY if the user requests it or the template already has one.
7. DATA OWNERSHIP: Each piece of KB data should appear in exactly ONE section. Assign financial metrics to "Financial Highlights", customer data to "Customer Insights", etc. The Executive Summary may reference key figures but should not contain detailed tables.
8. FACT vs. GOAL DISTINCTION: The KB contains historical data (what happened). Do NOT reframe historical results as future targets or objectives unless the KB explicitly states them as goals. If a section is called "Business Objectives", note that the KB data is historical and instruct the generator accordingly.
9. TABLE BUDGET: Assign at most ONE table per section. If the same data (e.g., revenue figures) appears relevant to multiple sections, assign the detailed table to the most specific section and reference it briefly in others.

Return ONLY valid JSON in this exact structure:
{
  "document_title": "...",
  "sections": [
    {
      "heading": "Section heading text (from template)",
      "heading_level": 1,
      "style_name": "Heading 1",
      "content_type": "narrative" | "bullet_list" | "table" | "mixed" | "toc" | "skip",
      "instructions": "What content to write here (specific, grounded in the data available)",
      "kb_chunk_ids": ["chunk_id_1", "chunk_id_2"],
      "table_spec": null | {
        "caption": "Table N: Descriptive title",
        "columns": ["Col1", "Col2", "Col3"],
        "data_source": "kb",
        "notes": "what data to extract from KB for table rows"
      },
      "bullet_count": null | 5,
      "subsections": []
    }
  ],
  "generation_notes": "Any high-level notes for the generator about tone, focus, or constraints"
}
"""


class DocumentPlanner:
    """Plans the full document structure for generation."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    def plan(
        self,
        user_request: str,
        template_analysis: dict,
        kb_context: list[dict],
        chat_history: list[dict] | None = None,
    ) -> dict:
        """Produce a comprehensive document plan.

        Args:
            user_request: What the user wants to generate.
            template_analysis: Output of TemplateAnalyzer.analyze().
            kb_context: Retrieved KB chunks (may be empty if no KB).
            chat_history: Recent conversation for context.
        """
        llm = self._llm or LLMClient()

        # Build context strings
        sections_summary = self._format_template_sections(template_analysis)
        kb_summary = self._format_kb_context(kb_context)
        history_str = self._format_history(chat_history)

        user_prompt = (
            f"{history_str}"
            f"USER REQUEST:\n{user_request}\n\n"
            f"TEMPLATE SECTION STRUCTURE (these headings MUST be preserved):\n{sections_summary}\n\n"
            f"KNOWLEDGE BASE CONTEXT (use this data to ground content):\n{kb_summary}\n\n"
            "Now produce the complete document plan following all rules."
        )

        response = llm.complete(LLMRequest(
            system_prompt=_PLANNER_SYSTEM_PROMPT,
            user_prompt=user_prompt,
            temperature=0.1,
            max_tokens=4096,
            json_mode=True,
        ))

        plan = response.json or {}
        if not plan.get("sections"):
            # Fallback: generate a minimal plan from template sections
            plan = self._fallback_plan(user_request, template_analysis)

        # Enrich plan with KB chunk references
        plan = self._enrich_with_kb_refs(plan, kb_context)
        
        # Deduplicate KB references
        plan = self._deduplicate_plan(plan)

        return plan

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _format_template_sections(self, analysis: dict) -> str:
        sections = analysis.get("sections", [])
        if not sections:
            return "(No heading structure detected — generate logical sections based on the request)"

        lines: list[str] = []
        for s in sections:
            indent = "  " * (s.get("heading_level", 1) - 1)
            lines.append(f"{indent}- [{s['style_name']}] {s['heading_text']}")
            for sub in s.get("subsections", []):
                sub_indent = "  " * (sub.get("heading_level", 2) - 1)
                lines.append(f"{sub_indent}  - [{sub['style_name']}] {sub['heading_text']}")

        return "\n".join(lines)

    def _format_kb_context(self, chunks: list[dict]) -> str:
        if not chunks:
            return "(No knowledge base documents uploaded)"
        
        by_source: dict[str, list[dict]] = {}
        for i, chunk in enumerate(chunks):
            source = chunk.get("metadata", {}).get("source", "Unknown")
            by_source.setdefault(source, []).append({"idx": i, "text": chunk.get("text", "")})
        
        lines = []
        for source, items in by_source.items():
            lines.append(f"\n📄 {source} ({len(items)} chunks):")
            for item in items[:10]:
                preview = item["text"][:300].replace("\n", " ")
                lines.append(f"  [Chunk {item['idx']}] {preview}")
        return "\n".join(lines)

    def _format_history(self, history: list[dict] | None) -> str:
        if not history:
            return ""
        lines: list[str] = ["Recent conversation:\n"]
        for msg in history[-3:]:
            role = "User" if msg["role"] == "user" else "Agent"
            lines.append(f"{role}: {msg['content'][:200]}")
        return "\n".join(lines) + "\n\n"

    def _enrich_with_kb_refs(self, plan: dict, kb_context: list[dict]) -> dict:
        """Tag each section with relevant KB chunk indices based on keyword overlap."""
        if not kb_context:
            return plan

        for section in plan.get("sections", []):
            if section.get("kb_chunk_ids"):
                continue  # LLM already assigned
            # Simple keyword match to assign chunks
            heading = section.get("heading", "").lower()
            instructions = section.get("instructions", "").lower()
            relevant: list[str] = []
            for i, chunk in enumerate(kb_context):
                text = chunk.get("text", "").lower()
                # Score based on heading and instruction words
                words = set(re.findall(r"[a-z]+", heading + " " + instructions))
                score = sum(1 for w in words if len(w) > 3 and w in text)
                if score >= 2:
                    relevant.append(str(i))
            section["kb_chunk_ids"] = relevant[:5]  # Top 5 relevant chunks

        return plan

    def _deduplicate_plan(self, plan: dict) -> dict:
        """Ensure each KB chunk is assigned to at most 2 sections
        (one primary, one brief reference in Executive Summary)."""
        chunk_assignments: dict[str, list[str]] = {}
        
        for section in plan.get("sections", []):
            for cid in section.get("kb_chunk_ids", []):
                chunk_assignments.setdefault(cid, []).append(section.get("heading", ""))
        
        for section in plan.get("sections", []):
            new_cids = []
            heading = section.get("heading", "")
            for cid in section.get("kb_chunk_ids", []):
                assigned = chunk_assignments.get(cid, [])
                if len(assigned) > 2:
                    # Keep if it's the executive summary or if this is the most specific section.
                    # Since we don't have a perfect "most specific" heuristic, we just keep the first two
                    # assignments or Executive Summary.
                    if "summary" in heading.lower() or assigned.index(heading) < 2:
                        new_cids.append(cid)
                else:
                    new_cids.append(cid)
            section["kb_chunk_ids"] = new_cids
            
        return plan

    def _fallback_plan(self, request: str, analysis: dict) -> dict:
        """Minimal fallback plan when LLM fails."""
        sections_raw = analysis.get("sections", [])
        sections = []
        for s in sections_raw:
            sections.append({
                "heading": s.get("heading_text", "Section"),
                "heading_level": s.get("heading_level", 1),
                "style_name": s.get("style_name", "Heading 1"),
                "content_type": "narrative",
                "instructions": f"Write professional content for: {s.get('heading_text', 'this section')}",
                "kb_chunk_ids": [],
                "table_spec": None,
                "subsections": [],
            })
        if not sections:
            sections = [{
                "heading": "Content",
                "heading_level": 1,
                "style_name": "Heading 1",
                "content_type": "narrative",
                "instructions": request,
                "kb_chunk_ids": [],
                "table_spec": None,
                "subsections": [],
            }]
        return {"document_title": request[:80], "sections": sections, "generation_notes": ""}

# ===== END services/document_planner.py =====

# ===== BEGIN services/section_generator.py =====
"""Section Generator — generates structured content for each planned section.

For each section in the document plan:
  1. Retrieves the relevant KB chunks specified in the plan
  2. Calls LLM to generate grounded, professional content
  3. Returns structured output (paragraphs, bullets, tables) that maps
     directly to python-docx operations in DocumentAssembler

The generator runs sections concurrently to reduce latency.
"""

import concurrent.futures
import json
import logging
import re
from typing import Any


log = logging.getLogger(__name__)

_SECTION_SYSTEM_PROMPT = """You are an expert content writer for professional corporate and audit documents.

Your task: Generate the content for a single section of a formal business document.

RULES:
1. STRICT GROUNDING: Every factual claim (numbers, percentages, currency values, dates, names, statistics) MUST come from the provided KB context.
   - If data for a topic is NOT in the KB context, simply SKIP that topic entirely. Write about what you DO have evidence for.
   - NEVER write "information not available", "N/A", "data not provided", or any similar disclaimer. Just omit topics without data.
   - You MAY use general knowledge ONLY for: transitions, explanations of concepts, and structural sentences.
   - For tables: only include rows/columns you have actual data for. A smaller table with real data is far better than a large table with gaps.
2. PROFESSIONAL TONE: Formal, concise, third-person language.
3. STRUCTURE: Return structured JSON.
4. INLINE CITATION: Every factual claim from the KB MUST include a citation tag [chunk:N] where N matches the chunk number in the KB CONTEXT above.
   Example: "Revenue grew by 11.2% year-over-year [chunk:3]."
   Claims without [chunk:N] tags must NOT contain specific numbers, dates, or statistics.
5. TABLE CITATION: Each table must include a source annotation as the last element: {"type": "table_source", "chunks": [3, 7]}
   listing which chunks the table data was drawn from.
6. QUANTITATIVE ACCURACY: Numbers must appear verbatim in the cited chunk. Do not round, estimate, or "improve" numbers.
7. LISTS: Concise bullet points, one sentence each.
8. NO PLACEHOLDERS: Never "[placeholder]" or "Lorem ipsum".

Return JSON in this exact structure:
{
  "elements": [
    {"type": "paragraph", "text": "Full paragraph text [chunk:1].", "style": "Normal"},
    {"type": "bullet_list", "items": ["Item 1 [chunk:2]"], "style": "List Bullet"},
    {
      "type": "table_caption",
      "text": "Table N: Descriptive Caption",
      "style": "Caption"
    },
    {
      "type": "table",
      "headers": ["Column 1", "Column 2"],
      "rows": [["val", "val"], ...],
      "style": "Table Grid"
    },
    {"type": "table_source", "chunks": [1, 2]}
  ]
}
"""


class SectionGenerator:
    """Generates structured content for individual document sections."""

    def __init__(self, llm: LLMClient | None = None) -> None:
        self._llm = llm

    def generate_all_sections(
        self,
        document_plan: dict,
        kb_context: list[dict],
        template_analysis: dict,
        workspace_id: str,
        kb_retrieval: Any = None,
        fact_verifier: Any = None,
        max_workers: int = 5,
    ) -> list[dict]:
        """Generate all sections sequentially.

        Returns a list of generated section dicts with 'heading' and 'elements'.
        """
        sections = document_plan.get("sections", [])
        generation_notes = document_plan.get("generation_notes", "")
        doc_title = document_plan.get("document_title", "")
        table_counter = {"n": 1}

        results: list[dict] = []
        covered_themes: list[str] = []

        for idx, section_plan in enumerate(sections):
            # Per-section retrieval via Qdrant
            section_query = f"{section_plan.get('heading', '')}: {section_plan.get('instructions', '')}"
            used_semantic = True
            
            if kb_retrieval:
                section_chunks, used_semantic = kb_retrieval.retrieve_for_section(
                    workspace_id=workspace_id,
                    section_query=section_query,
                    fallback_chunks=kb_context,
                    limit=15,
                )
            else:
                section_chunks = kb_context[:15]
            
            if not used_semantic:
                log.warning("Section '%s': used keyword fallback (semantic unavailable)", section_plan.get('heading'))

            ctx = {
                "section_plan": section_plan,
                "section_chunks": section_chunks,
                "doc_title": doc_title,
                "generation_notes": generation_notes,
                "style_catalog": template_analysis.get("style_catalog", {}),
                "table_style": template_analysis.get("table_style", {}),
                "section_idx": idx,
            }

            try:
                result = self._generate_section(ctx, table_counter, covered_themes)
                
                # Inline verification
                if fact_verifier and section_chunks:
                    result = self._verify_inline(result, section_chunks, fact_verifier)
                
                results.append(result)
                
                # Extract themes
                theme = self._extract_theme(result)
                covered_themes.append(
                    f"Section '{result.get('heading', '')}': covered {theme}"
                )
            except Exception as exc:
                log.error("Section %d generation failed: %s", idx, exc)
                results.append({
                    "heading": section_plan.get("heading", "Section"),
                    "heading_level": section_plan.get("heading_level", 1),
                    "style_name": section_plan.get("style_name", "Heading 1"),
                    "elements": [
                        {
                            "type": "paragraph",
                            "text": f"[Content for {section_plan.get('heading', 'this section')} could not be generated. Please review and complete manually.]",
                            "style": "Normal",
                        }
                    ],
                })

        return results

    def _verify_inline(self, result: dict, section_chunks: list[dict], verifier) -> dict:
        """Verify a section's elements against the chunks it was generated from."""
        verified = verifier.verify_section(
            section_elements=result.get("elements", []),
            kb_chunks=section_chunks,
            section_heading=result.get("heading", ""),
        )
        return {**result, "elements": verified.verified_elements}

    def _extract_theme(self, result: dict) -> str:
        """Extract topic-level summary of what a section covered (no specific numbers)."""
        themes = []
        for el in result.get("elements", []):
            if el.get("type") == "paragraph":
                text = el.get("text", "")[:80]
                # Strip numbers to keep it theme-level
                text = re.sub(r'[\d,]+\.?\d*[%$₹€£]?', '', text).strip()
                if text:
                    themes.append(text)
            elif el.get("type") == "table":
                themes.append("data table")
        return "; ".join(themes[:3]) if themes else "general narrative"

    def _generate_section(self, ctx: dict, table_counter: dict, covered_themes: list[str]) -> dict:
        """Generate content for a single section."""
        section_plan = ctx["section_plan"]
        content_type = section_plan.get("content_type", "narrative")

        # TOC sections are handled by the assembler
        if content_type == "toc":
            return {
                "heading": section_plan.get("heading", "Table of Contents"),
                "heading_level": section_plan.get("heading_level", 1),
                "style_name": section_plan.get("style_name", "Heading 1"),
                "elements": [{"type": "toc_placeholder"}],
            }

        # Skip sections
        if content_type == "skip":
            return {
                "heading": section_plan.get("heading", ""),
                "heading_level": section_plan.get("heading_level", 1),
                "style_name": section_plan.get("style_name", "Heading 1"),
                "elements": [],
            }

        llm = self._llm or LLMClient()

        # Assign table number atomically
        table_spec = section_plan.get("table_spec")
        current_table_n = None
        if table_spec:
            current_table_n = table_counter["n"]
            table_counter["n"] += 1
            # Update caption with correct number
            caption = table_spec.get("caption", f"Table {current_table_n}")
            if "Table N:" in caption or "Table N " in caption:
                caption = caption.replace("Table N", f"Table {current_table_n}")
            elif not re.match(r"Table \d+", caption):
                caption = f"Table {current_table_n}: {caption}"
            table_spec = {**table_spec, "caption": caption}

        section_chunks = ctx.get("section_chunks", [])
        kb_text = self._format_kb_chunks(section_chunks)

        user_prompt = (
            f"DOCUMENT TITLE: {ctx['doc_title']}\n\n"
            f"SECTION HEADING: {section_plan.get('heading')}\n"
            f"CONTENT TYPE: {content_type}\n"
            f"INSTRUCTIONS: {section_plan.get('instructions', 'Write professional content for this section.')}\n\n"
            f"KB CONTEXT (ground your content in this data):\n{kb_text}\n\n"
        )

        if table_spec:
            user_prompt += (
                f"TABLE SPECIFICATION:\n"
                f"  Caption: {table_spec.get('caption')}\n"
                f"  Columns: {table_spec.get('columns')}\n"
                f"  Notes: {table_spec.get('notes', 'Extract data from KB context')}\n"
                f"  Table Number: Must be exactly Table {current_table_n}\n\n"
            )
        if ctx.get("generation_notes"):
            user_prompt += f"OVERALL DOCUMENT NOTES: {ctx['generation_notes']}\n\n"

        if covered_themes:
            user_prompt += "TOPICS COVERED IN EARLIER SECTIONS (do not reproduce the detailed analysis, but you may reference these topics briefly if relevant to this section):\n"
            for theme in covered_themes:
                user_prompt += f"- {theme}\n"
            user_prompt += "\n"

        user_prompt += "Generate the section content as JSON following the required structure."

        response = llm.complete(LLMRequest(
            system_prompt=_SECTION_SYSTEM_PROMPT,
            user_prompt=user_prompt,
            temperature=0.2,
            max_tokens=3000,
            json_mode=True,
        ))

        elements = []
        parsed = response.json or {}
        if isinstance(parsed, dict) and "elements" in parsed:
            elements = parsed["elements"]
        elif isinstance(parsed, list):
            elements = parsed

        # Validate and clean elements
        elements = self._validate_elements(elements, content_type)

        return {
            "heading": section_plan.get("heading", ""),
            "heading_level": section_plan.get("heading_level", 1),
            "style_name": section_plan.get("style_name", "Heading 1"),
            "elements": elements,
        }

    def _format_kb_chunks(self, chunks: list[dict]) -> str:
        if not chunks:
            return "(No KB data is directly relevant to this section. Write a brief, general structural paragraph appropriate for this section heading. Do NOT include any specific numbers, statistics, or claims. Keep it to 2-3 sentences maximum.)"
        lines: list[str] = []
        for i, chunk in enumerate(chunks[:12]):
            meta = chunk.get("metadata", {})
            source = meta.get("source", "Unknown")
            page = meta.get("page", "")
            section = meta.get("section", "")
            loc = f"[{source}"
            if page:
                loc += f", p.{page}"
            if section:
                loc += f", §{section}"
            loc += "]"
            lines.append(f"[{i+1}] {loc}\n{chunk.get('text', '')[:1500]}")
        return "\n\n".join(lines)

    def _validate_elements(self, elements: list, content_type: str) -> list[dict]:
        """Ensure elements have the required structure and strip raw citation tags."""
        valid = []
        for el in elements:
            if not isinstance(el, dict):
                continue
            el_type = el.get("type", "paragraph")
            if el_type == "paragraph":
                if el.get("text"):
                    clean_text = re.sub(r'\s*\[chunk:\d+\]', '', str(el["text"]))
                    valid.append({
                        "type": "paragraph",
                        "text": clean_text,
                        "style": el.get("style", "Normal"),
                    })
            elif el_type == "bullet_list":
                items = el.get("items", [])
                if items and isinstance(items, list):
                    clean_items = [re.sub(r'\s*\[chunk:\d+\]', '', str(i)) for i in items if i]
                    valid.append({
                        "type": "bullet_list",
                        "items": clean_items,
                        "style": el.get("style", "List Bullet"),
                    })
            elif el_type == "numbered_list":
                items = el.get("items", [])
                if items and isinstance(items, list):
                    clean_items = [re.sub(r'\s*\[chunk:\d+\]', '', str(i)) for i in items if i]
                    valid.append({
                        "type": "numbered_list",
                        "items": clean_items,
                        "style": el.get("style", "List Number"),
                    })
            elif el_type == "table_caption":
                if el.get("text"):
                    valid.append({
                        "type": "table_caption",
                        "text": str(el["text"]),
                        "style": el.get("style", "Caption"),
                    })
            elif el_type == "table":
                headers = el.get("headers", [])
                rows = el.get("rows", [])
                if headers:
                    valid.append({
                        "type": "table",
                        "headers": [str(h) for h in headers],
                        "rows": [[str(c) for c in row] for row in rows if row],
                        "style": el.get("style", "Table Grid"),
                    })
            elif el_type == "table_source":
                chunks = el.get("chunks", [])
                if chunks:
                    valid.append({
                        "type": "table_source",
                        "chunks": [int(c) for c in chunks],
                    })
            elif el_type == "toc_placeholder":
                valid.append({"type": "toc_placeholder"})
        return valid

# ===== END services/section_generator.py =====

# ===== BEGIN services/fact_verifier.py =====
"""Fact Verifier — validates generated quantitative claims against the KB context."""

import logging
import re
from typing import NamedTuple

log = logging.getLogger(__name__)

class VerificationResult(NamedTuple):
    verified_elements: list[dict]
    violations: list[dict]
    was_modified: bool

class FactVerifier:
    """Verifies that quantitative claims in generated text are grounded in KB data."""

    def __init__(self, llm=None) -> None:
        self._llm = llm  # Kept for future use, not used in verification

    def verify_section(
        self,
        section_elements: list[dict],
        kb_chunks: list[dict],
        section_heading: str,
    ) -> VerificationResult:
        """Verify citations and scan for uncited numbers. Returns clean output."""
        violations = []
        verified_elements = []
        was_modified = False

        for el in section_elements:
            verified_el = self._verify_element(el, kb_chunks, violations)
            if verified_el is not None:
                clean_el = self._strip_citations(verified_el)
                verified_elements.append(clean_el)
                if clean_el != el:
                    was_modified = True

        return VerificationResult(verified_elements, violations, was_modified)

    def _verify_element(self, el: dict, chunks: list[dict], violations: list[dict]) -> dict | None:
        el_type = el.get("type", "paragraph")
        
        if el_type == "table_source":
            # Table source annotation — verify chunk indices are valid
            for idx in el.get("chunks", []):
                if not (0 <= idx - 1 < len(chunks)):
                    violations.append({"type": "invalid_table_source", "chunk": idx})
            return None  # Don't include the source annotation in final output
        
        if el_type == "table_caption":
            return el  # Captions are structural, never touch them
        
        if el_type == "table":
            return self._verify_table(el, chunks, violations)
        
        if el_type == "paragraph":
            text = el.get("text", "")
            
            # 1. Verify cited claims
            citations = re.findall(r'\[chunk:(\d+)\]', text)
            for ref in citations:
                chunk_idx = int(ref) - 1  # Prompt uses 1-indexed
                if not (0 <= chunk_idx < len(chunks)):
                    violations.append({"claim": f"[chunk:{ref}]", "reason": "invalid index"})
                    text = self._remove_sentence_at(text, f"[chunk:{ref}]")
                    continue
                
                # ACTUAL VERIFICATION: check that numbers near this citation
                # appear in the cited chunk
                chunk_text = chunks[chunk_idx].get("text", "").lower()
                sentence = self._get_sentence_containing(text, f"[chunk:{ref}]")
                numbers_in_sentence = re.findall(r'[\d,]+\.?\d*', sentence)
                
                for num in numbers_in_sentence:
                    clean_num = num.replace(",", "")
                    if clean_num and len(clean_num) > 1:  # Skip single digits
                        if clean_num not in chunk_text.replace(",", ""):
                            violations.append({
                                "claim": sentence.strip(),
                                "value": num,
                                "cited_chunk": ref,
                                "reason": "number not found in cited chunk"
                            })
                            text = self._remove_sentence_at(text, f"[chunk:{ref}]")
                            break
            
            # 2. Scan for UNCITED multi-digit numbers — backstop for claims that bypass citations
            remaining_text = re.sub(r'\[chunk:\d+\]', '', text)
            # Match any 2+ digit number not inside Table/Section/Chapter/Q references
            uncited_numbers = re.findall(r'(?<!Table\s)(?<!Section\s)(?<!Chapter\s)(?<!Q)(\d{2,}[\d,.]*)', remaining_text)
            if uncited_numbers:
                for num_match in uncited_numbers:
                    # Skip if it looks like a year in context (e.g., "Q3 2024")
                    if re.search(r'Q[1-4]\s*' + re.escape(num_match), remaining_text):
                        continue
                    sentence = self._get_sentence_containing(text, num_match)
                    if sentence:
                        violations.append({
                            "claim": sentence.strip(),
                            "value": num_match,
                            "reason": "uncited quantitative claim"
                        })
                        text = self._remove_sentence_at(text, num_match)
            
            if not text.strip():
                return None  # Entire paragraph was removed
            
            return {**el, "text": text}
        
        return el  # bullet_list, numbered_list — pass through

    def _verify_table(self, el: dict, chunks: list[dict], violations: list[dict]) -> dict:
        """Verify table cell values against the cited source chunks."""
        # Find the table_source annotation that should accompany this table
        # (passed via _verify_section which tracks cited_table_chunks)
        # For now, verify against ALL provided chunks
        combined_chunk_text = "\n".join(c.get("text", "") for c in chunks).lower()
        
        new_rows = []
        for row in el.get("rows", []):
            new_row = []
            for cell in row:
                # Extract numbers from cell
                numbers = re.findall(r'[\d,]+\.?\d*', str(cell))
                cell_ok = True
                for num in numbers:
                    clean_num = num.replace(",", "")
                    if clean_num and len(clean_num) > 1:
                        if clean_num not in combined_chunk_text.replace(",", ""):
                            violations.append({
                                "claim": f"Table cell: {cell}",
                                "value": num,
                                "reason": "table cell number not found in source chunks"
                            })
                            cell_ok = False
                            break
                # Keep the cell if valid, drop the entire row if any cell fails
                new_row.append(cell if cell_ok else "—")
            new_rows.append(new_row)
        
        return {**el, "rows": new_rows}

    def _strip_citations(self, el: dict) -> dict:
        new_el = dict(el)
        if new_el.get("type") == "paragraph":
            new_el["text"] = re.sub(r'\s*\[chunk:\d+\]', '', new_el["text"]).strip()
        elif new_el.get("type") in ("bullet_list", "numbered_list"):
            new_el["items"] = [re.sub(r'\s*\[chunk:\d+\]', '', item).strip() 
                               for item in new_el.get("items", [])]
        return new_el

    def _remove_sentence_at(self, text: str, needle: str) -> str:
        """Remove ONLY the first sentence containing `needle` from text.
        Uses index-based removal to avoid collateral damage when the same
        substring appears in multiple sentences."""
        sentences = re.split(r'(?<=[.!?])\s+', text)
        for i, s in enumerate(sentences):
            if needle in s:
                return ' '.join(sentences[:i] + sentences[i+1:])
        return text

    def _get_sentence_containing(self, text: str, needle: str) -> str:
        """Return the first sentence containing `needle`."""
        sentences = re.split(r'(?<=[.!?])\s+', text)
        for s in sentences:
            if needle in s:
                return s
        return ""

# ===== END services/fact_verifier.py =====

# ===== BEGIN services/document_assembler.py =====
"""Document Assembler — builds the final DOCX from generated section content.

Strategy:
  1. Clone the template DOCX.
  2. Clear all body content below the cover page (if any), preserving styles.
  3. For each generated section, insert:
     - Heading paragraph with the correct heading style
     - Content elements: narrative paragraphs, bullet lists, tables with captions
  4. Apply template-matched formatting to all inserted elements:
     - Font family, size, color from style catalog
     - Table borders, header colors, alternating row colors
  5. Build a visible Table of Contents from headings (if requested).
  6. Preserve headers, footers, and page layout from the template.
"""

import copy
import logging
import shutil
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)


class DocumentAssembler:
    """Assembles a professional DOCX from generated section content."""

    def assemble(
        self,
        template_path: Path,
        target_path: Path,
        generated_sections: list[dict],
        template_analysis: dict,
    ) -> tuple[bool, list[str]]:
        """Assemble the final document.

        Args:
            template_path: Source DOCX template.
            target_path: Where to write the assembled document.
            generated_sections: Output of SectionGenerator.generate_all_sections().
            template_analysis: Output of TemplateAnalyzer.analyze().

        Returns:
            (success, list_of_section_summaries)
        """
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.oxml.ns import qn
        except ImportError:
            raise RuntimeError("python-docx is required.")

        # Clone template
        shutil.copy2(template_path, target_path)

        doc = Document(str(target_path))

        # Clear existing body content (keep styles intact)
        self._clear_body(doc, template_analysis)

        style_catalog = template_analysis.get("style_catalog", {})
        table_style = template_analysis.get("table_style", {})

        summaries: list[str] = []
        toc_headings: list[dict] = []

        for section in generated_sections:
            heading_text = section.get("heading", "")
            heading_level = section.get("heading_level", 1)
            style_name = section.get("style_name", f"Heading {heading_level}")
            elements = section.get("elements", [])

            if not heading_text:
                continue

            # Insert heading
            heading_para = self._add_heading(doc, heading_text, style_name)
            toc_headings.append({"text": heading_text, "level": heading_level})

            # Insert elements
            has_content = False
            for element in elements:
                el_type = element.get("type", "paragraph")

                if el_type == "toc_placeholder":
                    self._add_toc(doc, toc_headings)
                    has_content = True

                elif el_type == "paragraph":
                    text = element.get("text", "")
                    if text:
                        self._add_paragraph(doc, text, element.get("style", "Normal"), style_catalog)
                        has_content = True

                elif el_type in ("bullet_list", "numbered_list"):
                    items = element.get("items", [])
                    list_style = "List Bullet" if el_type == "bullet_list" else "List Number"
                    effective_style = element.get("style", list_style)
                    for item in items:
                        if item:
                            self._add_paragraph(doc, item, effective_style, style_catalog)
                    if items:
                        has_content = True

                elif el_type == "table_caption":
                    text = element.get("text", "")
                    if text:
                        self._add_paragraph(doc, text, element.get("style", "Caption"), style_catalog)

                elif el_type == "table":
                    headers = element.get("headers", [])
                    rows = element.get("rows", [])
                    if headers:
                        self._add_table(doc, headers, rows, table_style)
                        has_content = True

            if has_content:
                summaries.append(f"Generated section: {heading_text}")

        # Add page break between major sections for readability
        # (already done by heading styles in most templates)

        doc.save(str(target_path))
        log.info("DocumentAssembler: saved %s with %d sections", target_path, len(summaries))
        return True, summaries

    # ------------------------------------------------------------------
    # Body Clearing
    # ------------------------------------------------------------------

    def _clear_body(self, doc, template_analysis: dict) -> None:
        """Remove all body content while preserving styles.

        If the template has a cover page (detected via 'has_cover_page'),
        we keep the cover page content and clear only after the first
        section break / page break.
        """
        from docx.oxml.ns import qn

        body = doc.element.body
        has_cover = template_analysis.get("has_cover_page", False)

        # Collect all body children (paragraphs, tables, sectPr)
        children = list(body)

        # Find sectPr (section properties) — must be kept
        sectPr = body.find(qn("w:sectPr"))

        # Determine where to start clearing
        # If cover page detected, keep first paragraph cluster before first heading
        skip_until = 0
        if has_cover:
            for i, child in enumerate(children):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == "p":
                    style_elem = child.find(f".//{qn('w:pStyle')}")
                    if style_elem is not None:
                        style_val = style_elem.get(qn("w:val"), "")
                        if "Heading" in style_val:
                            skip_until = i
                            break

        # Remove all non-cover-page children
        for child in children[skip_until:]:
            if child is sectPr:
                continue
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("p", "tbl", "sdt"):
                body.remove(child)

        # Ensure sectPr stays at end
        if sectPr is not None:
            body.remove(sectPr)
            body.append(sectPr)

    # ------------------------------------------------------------------
    # Element Insertion
    # ------------------------------------------------------------------

    def _add_heading(self, doc, text: str, style_name: str):
        """Add a heading paragraph, falling back to a numbered heading if style missing."""
        from docx.shared import Pt
        try:
            # Try exact style name first
            para = doc.add_paragraph(style=style_name)
        except Exception:
            # Fallback: try "Heading 1" etc.
            level = 1
            try:
                level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                None
            try:
                para = doc.add_paragraph(style=f"Heading {level}")
            except Exception:
                para = doc.add_paragraph()

        run = para.add_run(text)
        return para

    def _add_paragraph(self, doc, text: str, style_name: str, style_catalog: dict):
        """Add a body paragraph with template-matched formatting."""
        try:
            para = doc.add_paragraph(style=style_name)
        except Exception:
            para = doc.add_paragraph()

        run = para.add_run(text)

        # Apply additional formatting from style catalog if style has explicit settings
        # (python-docx paragraph styles already carry most formatting, so we only
        # override where the catalog explicitly specifies something not captured by style)
        style_info = style_catalog.get(style_name, {})
        self._apply_run_formatting(run, style_info)

        return para

    def _apply_run_formatting(self, run, style_info: dict) -> None:
        """Apply explicit run-level formatting from style catalog."""
        from docx.shared import Pt, RGBColor
        try:
            if style_info.get("font_family"):
                run.font.name = style_info["font_family"]
            if style_info.get("font_size_pt"):
                run.font.size = Pt(style_info["font_size_pt"])
            if style_info.get("bold") is not None:
                run.font.bold = style_info["bold"]
            if style_info.get("italic") is not None:
                run.font.italic = style_info["italic"]
            color_hex = style_info.get("color_hex")
            if color_hex and len(color_hex) == 6:
                r = int(color_hex[0:2], 16)
                g = int(color_hex[2:4], 16)
                b = int(color_hex[4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)
        except Exception as exc:
            log.debug("Run formatting failed: %s", exc)

    def _add_table(
        self,
        doc,
        headers: list[str],
        rows: list[list[str]],
        table_style: dict,
    ):
        """Add a professionally formatted table matching the template style."""
        from docx.shared import Pt, RGBColor
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        # Try to use the template's table style
        style_name = table_style.get("style_name", "Table Grid")
        try:
            table = doc.add_table(rows=num_rows, cols=num_cols, style=style_name)
        except Exception:
            try:
                table = doc.add_table(rows=num_rows, cols=num_cols, style="Table Grid")
            except Exception:
                table = doc.add_table(rows=num_rows, cols=num_cols)

        table.autofit = True

        # Set header row
        header_row = table.rows[0]
        header_bg = table_style.get("header_bg_hex")
        header_fc = table_style.get("header_font_color_hex")
        header_bold = table_style.get("header_bold", True)

        for col_idx, header_text in enumerate(headers):
            cell = header_row.cells[col_idx]
            cell.text = ""
            para = cell.paragraphs[0]
            run = para.add_run(header_text)
            run.font.bold = header_bold

            if header_fc and len(header_fc) == 6:
                try:
                    r = int(header_fc[0:2], 16)
                    g = int(header_fc[2:4], 16)
                    b = int(header_fc[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except Exception:
                    None

            if header_bg and len(header_bg) == 6:
                self._set_cell_bg(cell, header_bg)

        # Data rows
        alt_colors = table_style.get("alt_row_colors", [])
        for row_idx, row_data in enumerate(rows):
            table_row = table.rows[row_idx + 1]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    break
                cell = table_row.cells[col_idx]
                cell.text = str(cell_text)

            # Alternate row coloring
            if alt_colors and len(alt_colors) >= 1:
                bg = alt_colors[row_idx % len(alt_colors)]
                if bg and len(bg) == 6:
                    for cell in table_row.cells:
                        self._set_cell_bg(cell, bg)

        return table

    def _set_cell_bg(self, cell, hex_color: str) -> None:
        """Set the background fill color of a table cell."""
        try:
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), hex_color.upper())
            # Remove existing shd if any
            existing = tcPr.find(qn("w:shd"))
            if existing is not None:
                tcPr.remove(existing)
            tcPr.append(shd)
        except Exception as exc:
            log.debug("Cell background set failed: %s", exc)

    def _add_toc(self, doc, headings: list[dict]) -> None:
        """Insert a native Word Table of Contents field."""
        from docx.shared import Pt
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls

        # Add TOC heading paragraph
        try:
            toc_para = doc.add_paragraph("Table of Contents", style="TOC Heading")
        except Exception:
            toc_para = doc.add_paragraph("Table of Contents")
            for run in toc_para.runs:
                run.font.bold = True
                run.font.size = Pt(14)

        fld_xml = (
            r'<w:p %s>'
            r'  <w:r>'
            r'    <w:fldChar w:fldCharType="begin"/>'
            r'  </w:r>'
            r'  <w:r>'
            r'    <w:instrText xml:space="preserve"> TOC \o "1-3" \h \z \u </w:instrText>'
            r'  </w:r>'
            r'  <w:r>'
            r'    <w:fldChar w:fldCharType="separate"/>'
            r'  </w:r>'
            r'  <w:r>'
            r'    <w:fldChar w:fldCharType="end"/>'
            r'  </w:r>'
            r'</w:p>' % nsdecls('w')
        )
        p_toc = parse_xml(fld_xml)
        doc.element.body.append(p_toc)

        _enable_update_fields(doc)

# ===== END services/document_assembler.py =====

# ===== BEGIN services/run_store.py =====
"""In-memory run store for the fire-and-forget pipeline.

The /chat endpoint creates a run entry and spawns the pipeline as a
background task.  The /polling endpoint reads from this store to return
live status, events (thoughts, version tiles), and the final workspace
snapshot.

Entries are kept for 30 minutes after completion, then lazily evicted.
"""

import time
from dataclasses import dataclass, field
from typing import Any


@dataclass
class RunEntry:
    run_id: str
    workspace_id: str
    status: str = "running"           # "running" | "completed" | "error"
    events: list[dict] = field(default_factory=list)
    workspace: dict | None = None     # serialised workspace, set on completion
    error_message: str | None = None
    created_at: float = field(default_factory=time.time)
    completed_at: float | None = None

    TTL_SECONDS: int = 1800           # 30 min


class RunStore:
    """Thread-safe (asyncio-safe) in-memory store for agent runs."""

    def __init__(self) -> None:
        self._runs: dict[str, RunEntry] = {}

    # ------------------------------------------------------------------
    # Write API (called by the background task / graph)
    # ------------------------------------------------------------------

    def create(self, run_id: str, workspace_id: str) -> RunEntry:
        entry = RunEntry(run_id=run_id, workspace_id=workspace_id)
        self._runs[run_id] = entry
        return entry

    def push_event(self, run_id: str, event: dict) -> None:
        entry = self._runs.get(run_id)
        if entry:
            entry.events.append(event)

    def complete(self, run_id: str, workspace: dict) -> None:
        entry = self._runs.get(run_id)
        if entry:
            entry.status = "completed"
            entry.workspace = workspace
            entry.completed_at = time.time()

    def fail(self, run_id: str, error_message: str) -> None:
        entry = self._runs.get(run_id)
        if entry:
            entry.status = "error"
            entry.error_message = error_message
            entry.completed_at = time.time()

    # ------------------------------------------------------------------
    # Read API (called by /polling endpoint)
    # ------------------------------------------------------------------

    def get(self, run_id: str) -> RunEntry | None:
        self._evict()
        return self._runs.get(run_id)

    def snapshot(self, run_id: str) -> dict[str, Any]:
        """Return a JSON-serialisable poll response."""
        entry = self.get(run_id)
        if entry is None:
            return {"status": "not_found", "events": [], "workspace": None}
        return {
            "status": entry.status,
            "events": entry.events,
            "workspace": entry.workspace,
            "error": entry.error_message,
        }

    # ------------------------------------------------------------------
    # Eviction
    # ------------------------------------------------------------------

    def _evict(self) -> None:
        now = time.time()
        stale = [
            rid for rid, e in self._runs.items()
            if e.completed_at and (now - e.completed_at) > RunEntry.TTL_SECONDS
        ]
        for rid in stale:
            del self._runs[rid]


# Singleton — shared across the process
run_store = RunStore()

# ===== END services/run_store.py =====

# ===== BEGIN services/serializers.py =====
import json



def file_url(version: DocumentVersion, path_value: str) -> str:
    return f"/api/files/{version.workspace_id}/{path_value.split('/')[-1]}"


def _parse_content(content: str) -> dict | None:
    """Try to parse message content as structured JSON. Returns None for plain text."""
    try:
        parsed = json.loads(content)
        if isinstance(parsed, dict) and parsed.get("type") == "agent_response":
            return parsed
    except Exception:
        None
    return None


def serialize_workspace(workspace: Workspace, repo: WorkspaceRepository) -> WorkspaceOut:
    versions = repo.versions(workspace.id)
    messages = repo.messages(workspace.id)
    knowledge_docs = repo.list_knowledge_documents(workspace.id)
    return WorkspaceOut(
        id=workspace.id,
        document_type=workspace.document_type,
        original_filename=workspace.original_filename,
        current_version=workspace.current_version,
        created_at=workspace.created_at,
        updated_at=workspace.updated_at,
        messages=[
            MessageOut(
                id=m.id,
                role=m.role,
                content=m.content,
                content_parsed=_parse_content(m.content),
                created_at=m.created_at,
            )
            for m in messages
        ],
        versions=[
            VersionOut(
                id=v.id,
                version_number=v.version_number,
                document_url=file_url(v, v.document_path),
                pdf_url=file_url(v, v.pdf_path),
                created_at=v.created_at,
            )
            for v in versions
        ],
        knowledge_documents=[
            KnowledgeDocumentOut(
                id=d.id,
                filename=d.filename,
                file_type=d.file_type,
                file_size_bytes=d.file_size_bytes,
                chunk_count=d.chunk_count,
                status=d.status,
                error_message=d.error_message,
                created_at=d.created_at,
            )
            for d in knowledge_docs
        ],
    )

# ===== END services/serializers.py =====

# ===== BEGIN services/uploads.py =====
from pathlib import Path

from fastapi import UploadFile
from sqlalchemy.orm import Session



class UploadService:
    def __init__(self, db: Session):
        self.db = db
        self.storage = StorageService()
        self.processor = DocumentProcessor()
        self.preview = PreviewService()
        self.retrieval = RetrievalService()

    async def create_workspace(self, file: UploadFile) -> Workspace:
        document_type = Path(file.filename or "").suffix.lower().replace(".", "")
        workspace = Workspace(document_type=document_type, original_filename=file.filename or "document")
        self.db.add(workspace)
        self.db.flush()

        document_path = self.storage.version_document_path(workspace.id, 1, document_type)
        with document_path.open("wb") as handle:
            while chunk := await file.read(1024 * 1024):
                handle.write(chunk)

        pdf_path = self.storage.version_pdf_path(workspace.id, 1)
        await self.preview.convert_to_pdf(document_path, pdf_path)
        structure = self.processor.extract(document_path, document_type)

        import json
        json_path = document_path.with_suffix(".json")
        with json_path.open("w", encoding="utf-8") as f:
            json.dump(structure, f, indent=2)

        self.db.add(
            DocumentVersion(
                workspace_id=workspace.id,
                version_number=1,
                document_path=str(document_path),
                pdf_path=str(pdf_path),
            )
        )
        self.db.add(
            DocumentStructure(
                workspace_id=workspace.id,
                version_number=1,
                structure_json=structure,
            )
        )
        self.db.add(Message(workspace_id=workspace.id, role="assistant", content="Document uploaded and indexed."))
        self.db.commit()
        self.db.refresh(workspace)
        # Index blocks into Qdrant for semantic retrieval (no-op if not configured).
        self.retrieval.index_workspace(workspace.id, structure)
        return workspace

# ===== END services/uploads.py =====
